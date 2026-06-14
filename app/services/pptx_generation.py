from __future__ import annotations

import re
import tempfile
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from app.schemas.presentation_plan import PresentationPlan, TopicSection
from app.services.gemini_planner import GeminiPresentationPlanner
from app.services.presentation_templates import PresentationTemplateRegistry


class PptxGenerationService:
    def __init__(
        self,
        output_dir: str | None = None,
        gemini_planner: GeminiPresentationPlanner | None = None,
    ) -> None:
        base_dir = Path(output_dir) if output_dir else Path(tempfile.gettempdir()) / 'slide_generator_outputs'
        base_dir.mkdir(parents=True, exist_ok=True)
        self.output_dir = base_dir
        self.gemini_planner = gemini_planner
        self.template_registry = PresentationTemplateRegistry()

    def generate(self, payload: dict) -> str:
        plan = self.build_plan(payload)
        return self.render(payload, plan)

    def build_plan(self, payload: dict) -> PresentationPlan:
        topic, presenter_name, slide_count, language_code, pack = self._extract_core(payload)
        if self.gemini_planner and self.gemini_planner.enabled:
            return self.gemini_planner.build_plan(
                topic=topic,
                presenter_name=presenter_name,
                slide_count=slide_count,
                language_code=language_code,
                template=pack.get('_template'),
            )
        return self._build_fallback_plan(topic=topic, slide_count=slide_count, pack=pack)

    def render(self, payload: dict, plan: PresentationPlan) -> str:
        topic, presenter_name, _, language_code, pack = self._extract_core(payload)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        chapters = self._chapter_groups(len(plan.sections))
        # title + agenda + chapter dividers + sections + summary + closing
        total_slides = 4 + len(plan.sections) + len(chapters)

        self._add_title_slide(
            prs,
            presentation_title=plan.presentation_title,
            subtitle=plan.title_subtitle,
            presenter_name=presenter_name,
            agenda_preview=plan.agenda_items[:4],
            page_number=1,
            total_slides=total_slides,
            pack=pack,
        )
        self._add_agenda_slide(
            prs,
            agenda_items=plan.agenda_items,
            agenda_notes=[],
            presenter_name=presenter_name,
            page_number=2,
            total_slides=total_slides,
            pack=pack,
        )

        current_page = 3
        for chapter_index, section_indices in enumerate(chapters):
            self._add_divider_slide(
                prs,
                chapter_number=chapter_index + 1,
                total_chapters=len(chapters),
                section_titles=[plan.sections[idx].title for idx in section_indices],
                presenter_name=presenter_name,
                page_number=current_page,
                total_slides=total_slides,
                pack=pack,
                accent_rgb=self._section_accent_tuple(pack, chapter_index),
            )
            current_page += 1

            for section_index in section_indices:
                section = plan.sections[section_index]
                if section.content_type == 'table' and section.table is not None:
                    self._add_table_slide(
                        prs,
                        section=section,
                        presenter_name=presenter_name,
                        page_number=current_page,
                        total_slides=total_slides,
                        pack=pack,
                        section_index=section_index,
                    )
                elif section.content_type == 'process':
                    self._add_process_slide(
                        prs,
                        section=section,
                        presenter_name=presenter_name,
                        page_number=current_page,
                        total_slides=total_slides,
                        pack=pack,
                        section_index=section_index,
                    )
                else:
                    self._add_facts_slide(
                        prs,
                        section=section,
                        presenter_name=presenter_name,
                        page_number=current_page,
                        total_slides=total_slides,
                        pack=pack,
                        section_index=section_index,
                    )
                current_page += 1

        self._add_summary_slide(
            prs,
            title=str(pack['summary']),
            summary_points=plan.summary_points,
            presenter_name=presenter_name,
            page_number=total_slides - 1,
            total_slides=total_slides,
            pack=pack,
        )

        self._add_closing_slide(
            prs,
            presentation_title=plan.presentation_title,
            presenter_name=presenter_name,
            page_number=total_slides,
            total_slides=total_slides,
            pack=pack,
        )

        template_id = str((pack.get('_template') or {}).get('id') or 'template')
        safe_stem = self._safe_filename(f"{topic[:60]}_{presenter_name[:30]}_{language_code}_{template_id}")
        file_path = self.output_dir / f'{safe_stem}.pptx'
        prs.save(file_path)
        return str(file_path)

    @staticmethod
    def _safe_filename(value: str) -> str:
        cleaned = re.sub(r'[^A-Za-z0-9А-Яа-яЎўҚқҒғҲҳ_\- ]+', '', value).strip()
        cleaned = re.sub(r'\s+', '_', cleaned)
        return cleaned or 'presentation'

    def _extract_core(self, payload: dict) -> tuple[str, str, int, str, dict]:
        language_code = str(payload.get('language_code', 'uz') or 'uz')
        pack = self._language_pack(language_code)
        pack['_template'] = self.template_registry.get(str(payload.get('template_id') or '').strip())
        topic = str(payload.get('topic', 'Untitled presentation')).strip()
        presenter_name = str(payload.get('presenter_name', 'Unknown')).strip()
        slide_count = int(payload.get('slide_count', 6) or 6)
        slide_count = min(max(slide_count, 6), 12)
        return topic, presenter_name, slide_count, language_code, pack


    @staticmethod
    def _language_pack(language_code: str) -> dict[str, str | list[str]]:
        if language_code == 'ru':
            return {
                'language_code': 'ru',
                'prepared_by': 'Подготовил',
                'agenda': 'Основные разделы',
                'agenda_note_title': 'Краткий обзор',
                'summary': 'Итоговые выводы',
                'key_focus': 'Ключевые темы',
                'cover_points': ['Истоки', 'Развитие', 'Ключевые факты', 'Значение'],
                'thanks_title': 'Спасибо за внимание!',
                'thanks_subtitle': 'Есть вопросы? С удовольствием отвечу и обсужу детали.',
                'thanks_tag': 'Завершение',
                'chapter_word': 'ЧАСТЬ',
            }
        if language_code == 'en':
            return {
                'language_code': 'en',
                'prepared_by': 'Prepared by',
                'agenda': 'Main sections',
                'agenda_note_title': 'Quick overview',
                'summary': 'Final conclusions',
                'key_focus': 'Key themes',
                'cover_points': ['Origins', 'Development', 'Key facts', 'Significance'],
                'thanks_title': 'Thank you for your attention!',
                'thanks_subtitle': 'Questions are welcome — happy to discuss the details.',
                'thanks_tag': 'The end',
                'chapter_word': 'PART',
            }
        return {
            'language_code': 'uz',
            'prepared_by': 'Tayyorlagan',
            'agenda': 'Asosiy bo‘limlar',
            'agenda_note_title': 'Qisqacha reja',
            'summary': 'Yakuniy xulosalar',
            'key_focus': 'Asosiy yo‘nalishlar',
            'cover_points': ['Boshlanish', 'Rivojlanish', 'Muhim faktlar', 'Ahamiyati'],
            'thanks_title': 'E’tiboringiz uchun rahmat!',
            'thanks_subtitle': 'Savollaringiz bo‘lsa, javob berishdan mamnun bo‘laman.',
            'thanks_tag': 'Yakun',
            'chapter_word': 'QISM',
        }

    @staticmethod
    def _hex_to_tuple(value: str, default: tuple[int, int, int]) -> tuple[int, int, int]:
        text = str(value or '').strip().lstrip('#')
        if len(text) != 6:
            return default
        try:
            return tuple(int(text[index:index + 2], 16) for index in (0, 2, 4))  # type: ignore[return-value]
        except ValueError:
            return default

    def _theme_value(self, pack: dict, key: str, default: str) -> str:
        template = pack.get('_template') or {}
        theme = template.get('theme') if isinstance(template, dict) else {}
        if isinstance(theme, dict):
            return str(theme.get(key) or default)
        return default

    def _rgb_tuple(self, pack: dict, key: str, default: tuple[int, int, int]) -> tuple[int, int, int]:
        return self._hex_to_tuple(self._theme_value(pack, key, ''), default)

    def _rgb(self, pack: dict, key: str, default: tuple[int, int, int]) -> RGBColor:
        return RGBColor(*self._rgb_tuple(pack, key, default))

    def _font_family(self, pack: dict) -> str:
        return self._theme_value(pack, 'font_family', 'Calibri')

    def _heading_font(self, pack: dict) -> str:
        # Falls back to the body font when a template defines no dedicated heading face.
        return self._theme_value(pack, 'heading_font', self._font_family(pack))

    def _layout_value(self, pack: dict, key: str, default: str) -> str:
        template = pack.get('_template') or {}
        layout = template.get('layout') if isinstance(template, dict) else {}
        if isinstance(layout, dict):
            return str(layout.get(key) or default)
        return default

    def _process_color_tuples(self, pack: dict) -> list[tuple[int, int, int]]:
        template = pack.get('_template') or {}
        theme = template.get('theme') if isinstance(template, dict) else {}
        raw_colors = theme.get('process_colors') if isinstance(theme, dict) else None
        if not isinstance(raw_colors, list) or not raw_colors:
            raw_colors = ['#dbeafe', '#e0e7ff', '#dcfce7', '#fef3c7', '#fee2e2']
        defaults = [(219, 234, 254), (224, 231, 255), (220, 252, 231), (254, 243, 199), (254, 226, 226)]
        return [
            self._hex_to_tuple(str(color), defaults[index % len(defaults)])
            for index, color in enumerate(raw_colors[:6])
        ]

    def _process_color_list(self, pack: dict) -> list[RGBColor]:
        return [RGBColor(*color) for color in self._process_color_tuples(pack)]

    def _section_accent_tuple(self, pack: dict, index: int) -> tuple[int, int, int]:
        colors = self._process_color_tuples(pack)
        return colors[index % len(colors)] if colors else (191, 219, 254)

    @staticmethod
    def _chapter_groups(section_count: int) -> list[list[int]]:
        # Group the flat section list into a few chapters so we add 2-4 divider
        # slides regardless of deck length, instead of one before every section.
        if section_count <= 0:
            return []
        size = 1 if section_count <= 4 else 2 if section_count <= 8 else 3
        groups: list[list[int]] = []
        index = 0
        while index < section_count:
            groups.append(list(range(index, min(index + size, section_count))))
            index += size
        return groups

    @staticmethod
    def _relative_luminance(rgb: tuple[int, int, int]) -> float:
        def channel(value: int) -> float:
            c = value / 255.0
            return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
        r, g, b = rgb
        return 0.2126 * channel(r) + 0.7152 * channel(g) + 0.0722 * channel(b)

    def _contrast_ratio(self, a: tuple[int, int, int], b: tuple[int, int, int]) -> float:
        la, lb = self._relative_luminance(a), self._relative_luminance(b)
        high, low = max(la, lb), min(la, lb)
        return (high + 0.05) / (low + 0.05)

    def _visible_accent(self, pack: dict, accent_rgb: tuple[int, int, int]) -> tuple[int, int, int]:
        # For thin accent bars the colour must read against the background. On dark
        # templates the pastel section colours wash out, so fall back to the brand tones.
        background = self._rgb_tuple(pack, 'background', (248, 250, 252))
        candidates = (
            accent_rgb,
            self._rgb_tuple(pack, 'secondary', (96, 165, 250)),
            self._rgb_tuple(pack, 'primary', (30, 64, 175)),
        )
        for candidate in candidates:
            if self._contrast_ratio(candidate, background) >= 1.6:
                return candidate
        return accent_rgb

    # Icons are drawn from PowerPoint preset geometries (not an icon font) so they
    # render identically in PowerPoint, Google Slides, mobile and the PDF export.
    ICON_PRESETS = {
        'idea': MSO_AUTO_SHAPE_TYPE.SUN,
        'gear': MSO_AUTO_SHAPE_TYPE.GEAR_6,
        'flow': MSO_AUTO_SHAPE_TYPE.CIRCULAR_ARROW,
        'star': MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
        'doc': MSO_AUTO_SHAPE_TYPE.FLOWCHART_DOCUMENT,
        'pie': MSO_AUTO_SHAPE_TYPE.PIE,
        'donut': MSO_AUTO_SHAPE_TYPE.DONUT,
        'cube': MSO_AUTO_SHAPE_TYPE.CUBE,
        'hex': MSO_AUTO_SHAPE_TYPE.HEXAGON,
        'pentagon': MSO_AUTO_SHAPE_TYPE.PENTAGON,
        'diamond': MSO_AUTO_SHAPE_TYPE.DIAMOND,
        'plaque': MSO_AUTO_SHAPE_TYPE.PLAQUE,
        'bolt': MSO_AUTO_SHAPE_TYPE.LIGHTNING_BOLT,
        'chevron': MSO_AUTO_SHAPE_TYPE.CHEVRON,
        'arrow': MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        'plus': MSO_AUTO_SHAPE_TYPE.MATH_PLUS,
        'heart': MSO_AUTO_SHAPE_TYPE.HEART,
        'cloud': MSO_AUTO_SHAPE_TYPE.CLOUD,
    }

    _ICON_RULES = (
        (('tarix', 'history', 'davr', 'asr', 'dynasty', 'otmish', 'qadim', 'meros', 'heritage'), 'plaque'),
        (('iqtisod', 'economy', 'savdo', 'moliya', 'bozor', 'business', 'market', 'daromad', 'narx', 'budjet'), 'pie'),
        (('texnologiya', 'technology', 'innovatsiya', 'innovation', 'raqamli', 'digital', 'dastur', 'tizim', 'system', 'kompyuter'), 'gear'),
        (('talim', 'education', 'oqit', 'maktab', 'oquv', 'pedagog', 'learning', 'teaching', 'dars'), 'idea'),
        (('fan', 'science', 'biolog', 'fizika', 'kimyo', 'tadqiqot', 'research', 'tajriba', 'kashfiyot'), 'bolt'),
        (('jarayon', 'process', 'bosqich', 'rivojlanish', 'development', 'evolution', 'ketma', 'sequence'), 'flow'),
        (('natija', 'xulosa', 'result', 'outcome', 'yutuq', 'ahamiyat', 'significance', 'impact', 'samara'), 'star'),
        (('omil', 'factor', 'sabab', 'cause', 'asos', 'principle', 'tushuncha', 'concept', 'tamoyil'), 'diamond'),
        (('kelajak', 'future', 'reja', 'plan', 'strategiya', 'strategy', 'maqsad', 'goal', 'istiqbol'), 'arrow'),
        (('madaniyat', 'culture', 'sanat', 'art', 'adabiyot', 'literature', 'meʼmor', 'memor'), 'heart'),
        (('tabiat', 'nature', 'ekolog', 'environment', 'iqlim', 'climate', 'suv', 'energiya', 'energy'), 'cloud'),
    )

    @classmethod
    def _icon_for_text(cls, text: str, *, fallback_index: int = 0) -> str:
        lowered = re.sub(r"['ʻʼ`’]", '', str(text or '').lower())
        for keys, icon in cls._ICON_RULES:
            if any(key in lowered for key in keys):
                return icon
        cycle = ['diamond', 'hex', 'star', 'pentagon', 'cube', 'donut']
        return cycle[fallback_index % len(cycle)]

    def _draw_icon(self, slide, name: str, *, left: float, top: float, size: float, color: tuple[int, int, int]) -> None:
        shape_type = self.ICON_PRESETS.get(name, MSO_AUTO_SHAPE_TYPE.DIAMOND)
        shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(size), Inches(size))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        shape.line.fill.background()

    def _apply_alpha(self, shape, alpha_pct: float) -> None:
        # python-pptx has no public transparency API; add an <a:alpha> child to the
        # fill colour so big decorative shapes can sit softly behind content.
        try:
            sp_pr = shape._element.spPr
            solid_fill = sp_pr.find(qn('a:solidFill'))
            if solid_fill is None:
                return
            srgb = solid_fill.find(qn('a:srgbClr'))
            if srgb is None:
                return
            for existing in srgb.findall(qn('a:alpha')):
                srgb.remove(existing)
            value = int(max(0.0, min(100.0, alpha_pct)) * 1000)
            srgb.append(srgb.makeelement(qn('a:alpha'), {'val': str(value)}))
        except Exception:
            pass

    def _soft_oval(self, slide, pack: dict, *, left: float, top: float, size: float, color_key: str, default: tuple[int, int, int], alpha: float) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(pack, color_key, default)
        shape.line.fill.background()
        self._apply_alpha(shape, alpha)

    def _soft_rect(self, slide, pack: dict, *, left: float, top: float, width: float, height: float, color_key: str, default: tuple[int, int, int], alpha: float, rotation: float = 0.0) -> None:
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(pack, color_key, default)
        shape.line.fill.background()
        if rotation:
            shape.rotation = rotation
        self._apply_alpha(shape, alpha)

    def _draw_background_motif(self, slide, pack: dict, *, bold: bool = False) -> None:
        # Soft, corner-anchored decoration. Kept in the margins and behind content
        # via low alpha so it enriches the page without hurting readability.
        style = self._layout_value(pack, 'accent_style', 'academic_grid')
        k = 1.5 if bold else 1.0

        def a(value: float) -> float:
            return min(value * k, 100.0)

        if style == 'neon_lines':
            # Diagonal glow strokes (a wide faint stroke behind a thin bright one).
            for (lx, ly), rot in (((9.7, 0.1), -32), ((9.2, 5.6), -32)):
                self._soft_rect(slide, pack, left=lx, top=ly, width=4.2, height=0.16, color_key='primary', default=(6, 182, 212), alpha=a(16), rotation=rot)
                self._soft_rect(slide, pack, left=lx, top=ly + 0.06, width=4.2, height=0.045, color_key='secondary', default=(34, 211, 238), alpha=a(70), rotation=rot)
            self._soft_oval(slide, pack, left=11.9, top=5.45, size=0.95, color_key='secondary', default=(34, 211, 238), alpha=a(16))
            self._soft_oval(slide, pack, left=12.18, top=5.73, size=0.40, color_key='secondary', default=(34, 211, 238), alpha=a(85))
            return

        if style == 'analytic_grid':
            # A neat measurement grid tucked into the top-right corner + a soft bloom.
            for gx in (10.95, 11.55, 12.15, 12.75):
                self._soft_rect(slide, pack, left=gx, top=0.0, width=0.014, height=2.25, color_key='border', default=(153, 246, 228), alpha=a(38))
            for gy in (0.50, 1.12, 1.74):
                self._soft_rect(slide, pack, left=10.75, top=gy, width=2.30, height=0.014, color_key='border', default=(153, 246, 228), alpha=a(38))
            self._soft_rect(slide, pack, left=10.78, top=0.30, width=1.72, height=0.09, color_key='secondary', default=(20, 184, 166), alpha=a(80))
            self._soft_oval(slide, pack, left=-1.0, top=5.7, size=2.5, color_key='accent', default=(153, 246, 228), alpha=a(13))
            return

        if style == 'editorial_corner':
            # Bold magazine blocks bleeding off opposite corners.
            self._soft_rect(slide, pack, left=-1.2, top=5.25, width=3.9, height=2.7, color_key='accent', default=(254, 215, 170), alpha=a(24))
            self._soft_rect(slide, pack, left=0.0, top=6.66, width=2.7, height=0.13, color_key='secondary', default=(13, 148, 136), alpha=a(85))
            self._soft_rect(slide, pack, left=11.52, top=-0.4, width=2.3, height=3.5, color_key='surface_alt', default=(255, 237, 213), alpha=a(34))
            self._soft_rect(slide, pack, left=11.42, top=-0.4, width=0.10, height=3.5, color_key='accent', default=(254, 215, 170), alpha=a(85))
            return

        # academic_grid (default): layered soft blooms.
        self._soft_oval(slide, pack, left=10.3, top=-1.55, size=4.9, color_key='accent', default=(186, 230, 253), alpha=a(14))
        self._soft_oval(slide, pack, left=11.3, top=-0.6, size=3.0, color_key='secondary', default=(96, 165, 250), alpha=a(10))
        self._soft_oval(slide, pack, left=-1.15, top=5.65, size=2.7, color_key='accent', default=(186, 230, 253), alpha=a(12))

    def _style_run(
        self,
        run,
        pack: dict,
        *,
        size: float,
        color_key: str = 'text',
        default_color: tuple[int, int, int] = (15, 23, 42),
        bold: bool = False,
        font_name: str | None = None,
        color_rgb: tuple[int, int, int] | None = None,
    ) -> None:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font_name or self._font_family(pack)
        run.font.color.rgb = RGBColor(*color_rgb) if color_rgb is not None else self._rgb(pack, color_key, default_color)

    def _build_fallback_plan(self, *, topic: str, slide_count: int, pack: dict) -> PresentationPlan:
        language_code = str(pack.get('language_code', 'uz'))
        if language_code == 'en':
            subtitle = f'{topic}: the main ideas, key facts, and practical significance.'
            sections = [
                TopicSection(content_type='facts', title='Core concepts', focus=f'This section explains the main ideas behind {topic}.', facts=[
                    f'{topic} can be understood through several connected dimensions.',
                    'The key terms and examples help define the topic more clearly.',
                    'Causes, features, and outcomes make the topic easier to compare and explain.',
                    'A structured overview helps show why the topic matters in practice or theory.',
                ]),
                TopicSection(content_type='process', title='Development or sequence', focus=f'This section follows the main stages or sequence related to {topic}.', facts=[
                    'The initial context or starting conditions are identified first.',
                    'The main stages or changes are then described in logical order.',
                    'Important factors that influence the development are highlighted.',
                    'The final state or outcome is summarized clearly.',
                ]),
                TopicSection(content_type='facts', title='Key implications', focus=f'This section highlights why {topic} is important.', facts=[
                    'Examples and comparisons help clarify the meaning of the topic.',
                    'The topic often becomes clearer when viewed from both theory and practice.',
                    'Well-chosen facts make the content easier to remember and explain.',
                    'A concise summary strengthens the overall understanding of the subject.',
                ]),
            ]
            summary = [
                f'{topic} is best understood through its core ideas, examples, and implications.',
                'A clear sequence of facts helps connect the topic’s causes, features, and outcomes.',
                'Carefully selected examples make the subject more memorable and easier to explain.',
            ]
        elif language_code == 'ru':
            subtitle = f'{topic}: основные идеи, ключевые факты и практическое значение темы.'
            sections = [
                TopicSection(content_type='facts', title='Основные понятия', focus=f'Раздел раскрывает ключевые идеи темы {topic}.', facts=[
                    f'{topic} раскрывается через несколько взаимосвязанных аспектов.',
                    'Ключевые термины и примеры помогают точнее понять содержание темы.',
                    'Причины, признаки и последствия делают тему более понятной для сравнения.',
                    'Структурированный обзор помогает показать значимость темы на практике и в теории.',
                ]),
                TopicSection(content_type='process', title='Последовательность развития', focus=f'Раздел показывает основные этапы или последовательность, связанные с {topic}.', facts=[
                    'Сначала обозначается исходный контекст и стартовые условия.',
                    'Затем в логическом порядке раскрываются основные этапы или изменения.',
                    'Подчеркиваются важные факторы, влияющие на развитие темы.',
                    'В завершение кратко показывается итоговое состояние или результат.',
                ]),
                TopicSection(content_type='facts', title='Ключевое значение', focus=f'Раздел показывает, почему {topic} имеет значение.', facts=[
                    'Примеры и сопоставления помогают точнее раскрыть содержание темы.',
                    'Тема становится понятнее при рассмотрении с теоретической и практической сторон.',
                    'Хорошо подобранные факты делают содержание более запоминающимся.',
                    'Краткие выводы помогают закрепить целостное понимание темы.',
                ]),
            ]
            summary = [
                f'{topic} лучше всего раскрывается через основные идеи, примеры и значение.',
                'Последовательное изложение фактов помогает связать причины, признаки и результаты темы.',
                'Удачно подобранные примеры делают содержание более понятным и запоминающимся.',
            ]
        else:
            subtitle = f'{topic} mavzusining asosiy g‘oyalari, muhim faktlari va amaliy ahamiyati.'
            sections = [
                TopicSection(content_type='facts', title='Asosiy tushunchalar', focus=f'Bu bo‘lim {topic} mavzusining tayanch g‘oyalarini yoritadi.', facts=[
                    f'{topic} bir nechta o‘zaro bog‘liq jihatlar orqali tushuniladi.',
                    'Asosiy atamalar va misollar mavzu mazmunini aniqroq ochib beradi.',
                    'Sabab, belgi va natijalarni ajratish mavzuni solishtirishni yengillashtiradi.',
                    'Tartibli sharh mavzuning nazariy va amaliy ahamiyatini ko‘rsatadi.',
                ]),
                TopicSection(content_type='process', title='Rivojlanish ketma-ketligi', focus=f'Bu bo‘lim {topic} bilan bog‘liq asosiy bosqich yoki ketma-ketlikni ko‘rsatadi.', facts=[
                    'Avval boshlang‘ich sharoit va kontekst aniqlanadi.',
                    'So‘ng asosiy bosqichlar yoki o‘zgarishlar mantiqiy tartibda ko‘rsatiladi.',
                    'Rivojlanishga ta’sir qiluvchi muhim omillar ajratib ko‘rsatiladi.',
                    'Yakunida umumiy holat yoki natija qisqacha jamlanadi.',
                ]),
                TopicSection(content_type='facts', title='Muhim xulosalar', focus=f'Bu bo‘lim {topic} mavzusining ahamiyatini ko‘rsatadi.', facts=[
                    'Misollar va taqqoslashlar mavzuning mazmunini ravshanlashtiradi.',
                    'Mavzu nazariy va amaliy nuqtai nazardan ko‘rilganda yaxshiroq tushuniladi.',
                    'To‘g‘ri tanlangan faktlar mavzuni esda qolarli qiladi.',
                    'Qisqa xulosalar mavzu bo‘yicha yaxlit tasavvur beradi.',
                ]),
            ]
            summary = [
                f'{topic} mavzusi asosiy g‘oyalar, misollar va ahamiyat orqali yaxshiroq tushuniladi.',
                'Faktlarni ketma-ket ko‘rish mavzuning sabab, belgi va natijalarini bog‘laydi.',
                'To‘g‘ri tanlangan misollar mazmunni aniqroq va esda qolarli qiladi.',
            ]

        section_count = max(3, slide_count - 3)
        while len(sections) < section_count:
            sections.append(sections[len(sections) % 3])
        agenda_items = [section.title for section in sections[: min(8, len(sections))]]
        for item in list(pack.get('cover_points') or []):
            if len(agenda_items) >= 4:
                break
            if item not in agenda_items:
                agenda_items.append(str(item))
        while len(agenda_items) < 4:
            agenda_items.append(topic)
        return PresentationPlan(
            presentation_title=topic,
            title_subtitle=subtitle,
            agenda_items=agenda_items,
            sections=sections[:section_count],
            summary_points=summary[:5],
        )

    def _topic_profile(self, topic: str) -> dict:
        lower = topic.lower()
        if any(marker in lower for marker in ('oʻzbekiston tarixi', "o'zbekiston tarixi", 'uzbekiston tarixi', 'history of uzbekistan', 'ozbekiston tarixi')):
            return {
                'subtitle': 'Qadimgi davrlardan mustaqillik davrigacha bo‘lgan siyosiy, madaniy va ijtimoiy taraqqiyot yo‘li.',
                'sections': [
                    TopicSection(content_type='facts', title='Qadimgi va ilk davlatlar', focus='Hududdagi ilk sivilizatsiyalar va davlat birlashmalarining shakllanishi.', facts=[
                        'Amudaryo va Sirdaryo oralig‘i qadimdan dehqonchilik, hunarmandchilik va savdo uchun qulay hudud bo‘lgan.',
                        'Baxtariya, So‘g‘d va Xorazm kabi qadimgi birliklar Markaziy Osiyo tarixida muhim siyosiy markazlar bo‘lib xizmat qilgan.',
                        'Hudud Buyuk Ipak yo‘li orqali Sharq va G‘arb madaniyatlari tutashgan maydonga aylangan.',
                        'Arxeologik topilmalar shaharlashuv va yozma madaniyatning juda erta shakllanganini ko‘rsatadi.',
                    ]),
                    TopicSection(content_type='process', title='Islom davri va uyg‘onish bosqichi', focus='VIII–XII asrlarda ilm-fan, shaharlar va madaniyatning yuksalish jarayoni.', facts=[
                        'Arablar kirib kelishi bilan islom dini hududda keng tarqala boshladi.',
                        'Mahalliy sulolalar shakllanib, Buxoro va Samarqand ilmiy markaz sifatida kuchaydi.',
                        'Madrasalar, kutubxonalar va ilmiy muhit rivojlanib, ko‘plab olimlar yetishib chiqdi.',
                        'Hudud musulmon Sharqidagi yirik ilm-fan va madaniyat markazlaridan biriga aylandi.',
                    ]),
                    TopicSection(content_type='facts', title='Temuriylar davri', focus='Amir Temur va temuriylar zamonida davlat boshqaruvi, bunyodkorlik va ilmning yuksalishi.', facts=[
                        'Amir Temur markazlashgan davlat tuzib, Movarounnahr siyosiy qudratini tikladi.',
                        'Samarqand yirik siyosiy va me’moriy markazga aylantirildi.',
                        'Ulug‘bek davrida astronomiya, matematika va madaniy hayot yangi bosqichga ko‘tarildi.',
                        'Registon, observatoriya va ko‘plab me’moriy obidalar bu davr merosini ifodalaydi.',
                    ]),
                    TopicSection(content_type='table', title='Tarixiy davrlar taqqoslanishi', focus='Asosiy davrlar bo‘yicha siyosiy, madaniy va ilmiy xususiyatlarni solishtirish.', facts=[
                        'Har bir davrda davlat boshqaruvi va madaniy rivojlanish darajasi turlicha bo‘lgan.',
                        'Ilm-fan va me’morchilikning yuksalishi ayniqsa uyg‘onish va temuriylar davrida kuchli ko‘ringan.',
                    ], table={
                        'columns': ['Davr', 'Asosiy markaz', 'Ajralib turgan jihat'],
                        'rows': [
                            ['Qadimgi davr', 'Xorazm, So‘g‘d', 'Ilk davlatchilik va savdo yo‘llari'],
                            ['IX–XII asrlar', 'Buxoro, Samarqand', 'Ilm-fan va madaniy uyg‘onish'],
                            ['Temuriylar', 'Samarqand, Hirot', 'Bunyodkorlik va ilmiy taraqqiyot'],
                            ['Mustaqillik davri', 'Toshkent', 'Milliy tiklanish va merosni asrash'],
                        ],
                    }),
                    TopicSection(content_type='facts', title='Xonliklar va yangi davr', focus='Buxoro, Xiva va Qo‘qon xonliklari davrida siyosiy raqobat va mahalliy boshqaruv.', facts=[
                        'Hudud bir necha xonliklarga bo‘linib, markazlashuv zaiflashdi.',
                        'Buxoro, Xiva va Qo‘qon o‘ziga xos siyosiy va iqtisodiy tizimlarga ega bo‘lgan.',
                        'Savdo, hunarmandchilik va qishloq xo‘jaligi ichki iqtisodiyotning asosini tashkil etgan.',
                        'Ichki nizolar tashqi bosim oldida hududning zaiflashishiga sabab bo‘lgan.',
                    ]),
                    TopicSection(content_type='facts', title='Mustamlaka va jadidchilik', focus='Rossiya imperiyasi davri, islohot g‘oyalari va milliy uyg‘onishning kuchayishi.', facts=[
                        'XIX asr oxiriga kelib hudud Rossiya imperiyasi ta’siri ostiga tushdi.',
                        'Mustamlaka boshqaruvi iqtisodiy va siyosiy hayotda keskin o‘zgarishlar olib keldi.',
                        'Jadidlar yangi usul maktablari, matbuot va ma’rifat orqali jamiyatni isloh qilishga intildi.',
                        'Milliy o‘zlik, ta’lim va zamonaviy fikr jadidchilik harakatining asosiy mavzulari bo‘ldi.',
                    ]),
                    TopicSection(content_type='facts', title='Mustaqillik va madaniy meros', focus='1991-yildan keyingi davrda tarixiy xotira, milliy tiklanish va merosni asrash ishlari.', facts=[
                        '1991-yilda O‘zbekiston mustaqillikka erishib, yangi davlat taraqqiyoti yo‘lini boshladi.',
                        'Tarixiy obidalarni restavratsiya qilish va milliy merosni tiklash davlat siyosatining muhim qismiga aylandi.',
                        'Samarqand, Buxoro va Xiva kabi shaharlarga xalqaro e’tibor kuchaydi.',
                        'Tarix va madaniyat milliy o‘zlikni mustahkamlovchi asosiy omillardan biri bo‘lib qoldi.',
                    ]),
                ],
                'summary_points': [
                    'O‘zbekiston tarixi qadimgi sivilizatsiyalar, yirik davlatlar va boy madaniy meros uzviyligiga tayangan.',
                    'Hudud ilm-fan, savdo va me’morchilik markazi sifatida Markaziy Osiyoda alohida o‘rin tutgan.',
                    'Temuriylar davri va jadidchilik harakati tarixiy rivojlanishda burilish nuqtalari bo‘lib xizmat qilgan.',
                    'Mustaqillik davrida tarixiy xotira va madaniy merosni asrash yangi bosqichga ko‘tarildi.',
                ],
            }

        base_facts = [
            f'{topic} mavzusi bir nechta muhim yo‘nalishlar orqali tushuntiriladi.',
            'Asosiy tushunchalar va faktlarni tarixiy yoki amaliy kontekst bilan ko‘rib chiqish mazmunni boyitadi.',
            'Mavzudagi sabab va natijalarni ajratish tushunishni ancha yengillashtiradi.',
            'Misollar va taqqoslashlar mavzuning ahamiyatini yanada ravshanlashtiradi.',
        ]
        return {
            'subtitle': f'{topic} mavzusining asosiy yo‘nalishlari, muhim faktlari va amaliy ahamiyati.',
            'sections': [
                TopicSection(content_type='facts', title='Asosiy mazmun', focus=f'{topic} mavzusining markaziy g‘oyalari va tayanch faktlari.', facts=base_facts),
                TopicSection(content_type='facts', title='Muhim omillar', focus=f'{topic}ga ta’sir qiluvchi sabablar va sharoitlar.', facts=base_facts),
                TopicSection(content_type='process', title='Rivojlanish ketma-ketligi', focus=f'{topic} bo‘yicha asosiy bosqichlar yoki jarayonlar.', facts=[
                    'Boshlang‘ich sharoit va kontekstni aniqlash.',
                    'Asosiy o‘zgarishlar yoki rivojlanish nuqtalarini ko‘rsatish.',
                    'Natijaga ta’sir qilgan muhim omillarni ajratish.',
                    'Bugungi holat yoki yakuniy natijani baholash.',
                ]),
                TopicSection(content_type='table', title='Qisqa taqqoslash', focus=f'{topic} bo‘yicha asosiy jihatlarni jadval orqali jamlash.', facts=[
                    'Jadval mavzudagi asosiy jihatlarni bir joyga to‘playdi.',
                    'Taqqoslash orqali farq va o‘xshashliklar tezroq ko‘rinadi.',
                ], table={
                    'columns': ['Jihat', 'Tavsif', 'Ahamiyati'],
                    'rows': [
                        ['Mazmun', topic, 'Asosiy yo‘nalish'],
                        ['Asosiy omil', 'Tayanch faktlar', 'Tushunishni kuchaytiradi'],
                        ['Natija', 'Xulosa va ta’sir', 'Amaliy ahamiyat'],
                    ],
                }),
            ],
            'summary_points': [
                f'{topic} mavzusi mazmunan bir nechta o‘zaro bog‘liq yo‘nalishlardan tashkil topadi.',
                'Asosiy faktlar va jarayonlarni ketma-ket ko‘rish umumiy tasavvurni kuchaytiradi.',
                'Taqqoslash va misollar mavzuni esda qolarli va aniq qiladi.',
            ],
        }

    @staticmethod
    def _normalize_text(value: Any, max_chars: int | None = None) -> str:
        text = '' if value is None else str(value)
        text = text.replace(' ', ' ')
        text = text.replace('\"', '"').replace("\'", "'")
        text = text.replace('“', '"').replace('”', '"').replace('‘', "'").replace('’', "'")
        text = re.sub(r'\s+', ' ', text).strip()
        if (text.startswith('{') and text.endswith('}')) or (text.startswith('[') and text.endswith(']')):
            text = text[1:-1].strip()
        text = text.replace('{', '').replace('}', '').replace('[', '').replace(']', '')
        text = re.sub(r'\s*([,:;])\s*', r'\1 ', text)
        text = re.sub(r'\s+', ' ', text).strip(" \"'\n\t-–—")
        if max_chars and len(text) > max_chars:
            text = text[:max_chars].rstrip(' ,;:-.')
        return text


    @staticmethod
    def _emu_to_inches(value: Any) -> float | None:
        try:
            return float(value) / 914400.0
        except (TypeError, ValueError):
            return None

    def _frame_inner_size(
        self,
        text_frame,
        avail_width: float | None,
        avail_height: float | None,
    ) -> tuple[float, float]:
        width = avail_width
        height = avail_height
        shape = getattr(text_frame, '_parent', None)
        if width is None and shape is not None:
            width = self._emu_to_inches(getattr(shape, 'width', None))
        if height is None and shape is not None:
            height = self._emu_to_inches(getattr(shape, 'height', None))
        margin_l = self._emu_to_inches(text_frame.margin_left) or 0.0
        margin_r = self._emu_to_inches(text_frame.margin_right) or 0.0
        margin_t = self._emu_to_inches(text_frame.margin_top) or 0.0
        margin_b = self._emu_to_inches(text_frame.margin_bottom) or 0.0
        inner_w = (width - margin_l - margin_r) if width else 0.0
        inner_h = (height - margin_t - margin_b) if height else 0.0
        return max(inner_w, 0.0), max(inner_h, 0.0)

    # _estimate_lines is tuned optimistically (it assumes more characters per line
    # than a real proportional font such as Aptos renders), so we inflate the
    # estimated block height by this factor when deciding whether text fits.
    _FIT_SAFETY = 1.08

    def _text_block_height(
        self,
        paragraphs: list[tuple[str, float, float, float]],
        inner_w: float,
        *,
        scale: float,
        min_size: float,
    ) -> float:
        total = 0.0
        for text, design_size, space_after, line_spacing in paragraphs:
            font_size = max(min_size, design_size * scale)
            lines = self._estimate_lines(text, inner_w, font_size)
            line_height = (font_size * 1.2 * max(line_spacing, 1.0)) / 72.0
            total += (lines * line_height) + (space_after / 72.0)
        return total * self._FIT_SAFETY

    def _fit_frame(
        self,
        text_frame,
        *,
        max_size: int | float,
        min_size: int | float = 10,
        bold: bool = False,
        font_family: str = 'Calibri',
        avail_width: float | None = None,
        avail_height: float | None = None,
    ) -> None:
        # Deterministic, font-independent shrink-to-fit. python-pptx's fit_text
        # relies on locally installed font metrics (e.g. Aptos), which are absent
        # both in the dev font registry and on the production server, so it would
        # silently fail and let text overflow its shape. We size text ourselves
        # and scale every run by a shared factor so size hierarchy is preserved.
        text_frame.word_wrap = True
        max_size = float(max_size)
        min_size = min(float(min_size), max_size)

        inner_w, inner_h = self._frame_inner_size(text_frame, avail_width, avail_height)

        def design_size_of(paragraph) -> float:
            run_sizes = [run.font.size.pt for run in paragraph.runs if run.font.size is not None]
            base = max(run_sizes) if run_sizes else max_size
            return min(base, max_size)

        paragraphs: list[tuple[str, float, float, float]] = []
        for paragraph in text_frame.paragraphs:
            text = ''.join(run.text for run in paragraph.runs) or (paragraph.text or '')
            space_after = paragraph.space_after.pt if paragraph.space_after is not None else 0.0
            raw_spacing = paragraph.line_spacing
            line_spacing = float(raw_spacing) if isinstance(raw_spacing, (int, float)) else 1.0
            paragraphs.append((text, design_size_of(paragraph), float(space_after), line_spacing))

        scale = 1.0
        has_text = any(text.strip() for text, _, _, _ in paragraphs)
        if inner_w > 0 and inner_h > 0 and has_text:
            if self._text_block_height(paragraphs, inner_w, scale=1.0, min_size=min_size) > inner_h:
                scale = 0.0
                candidate = 1.0
                while candidate > 0.05:
                    if self._text_block_height(paragraphs, inner_w, scale=candidate, min_size=min_size) <= inner_h:
                        scale = candidate
                        break
                    candidate -= 0.02

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                original = run.font.size.pt if run.font.size is not None else max_size
                original = min(original, max_size)
                if scale <= 0:
                    new_size = min_size
                else:
                    new_size = max(min_size, original * scale)
                run.font.name = font_family
                run.font.size = Pt(round(new_size, 1))

    @staticmethod
    def _content_bounds(has_subtitle: bool) -> tuple[float, float]:
        return (2.08 if has_subtitle else 1.86, 6.58)

    @staticmethod
    def _estimate_lines(text: str, width_inches: float, font_size: float, *, bullet: bool = False) -> int:
        cleaned = re.sub(r'\s+', ' ', str(text or '')).strip()
        if not cleaned:
            return 1
        words = cleaned.split()
        capacity = max(
            12,
            int(
                width_inches
                * (
                    21 if font_size <= 11 else 19 if font_size <= 12 else 17 if font_size <= 13.5 else 15 if font_size <= 15 else 13
                )
            ),
        )
        current = 2 if bullet else 0
        lines = 1
        for word in words:
            token_len = len(word) + 1
            if current + token_len > capacity:
                lines += 1
                current = len(word)
            else:
                current += token_len
        return max(1, lines)

    def _estimate_bullet_block_height(
        self,
        items: list[str],
        *,
        width_inches: float,
        font_size: float,
        space_after_pt: float,
    ) -> float:
        line_height_inches = (font_size * 1.18) / 72
        gap_inches = space_after_pt / 72
        total = 0.0
        for item in items:
            lines = self._estimate_lines(item, width_inches, font_size, bullet=True)
            total += (lines * line_height_inches) + gap_inches
        return total

    def _balance_items_for_columns(self, items: list[str], *, width_inches: float, font_size: float) -> list[list[str]]:
        if len(items) <= 3:
            return [items, []]
        best_split: list[list[str]] | None = None
        best_score: float | None = None
        for split in range(1, len(items)):
            left = items[:split]
            right = items[split:]
            left_h = self._estimate_bullet_block_height(left, width_inches=width_inches, font_size=font_size, space_after_pt=8)
            right_h = self._estimate_bullet_block_height(right, width_inches=width_inches, font_size=font_size, space_after_pt=8)
            score = abs(left_h - right_h) + (abs(len(left) - len(right)) * 0.12)
            if best_score is None or score < best_score:
                best_score = score
                best_split = [left, right]
        return best_split or [items, []]

    def _select_facts_layout(self, items: list[str], *, panel_height: float) -> dict:
        presets = [
            {'columns': 1, 'font_size': 18.0, 'space_after': 12, 'width': 10.35},
            {'columns': 1, 'font_size': 16.8, 'space_after': 11, 'width': 10.35},
            {'columns': 1, 'font_size': 15.5, 'space_after': 10, 'width': 10.35},
            {'columns': 2, 'font_size': 15.0, 'space_after': 10, 'width': 4.92},
            {'columns': 2, 'font_size': 14.0, 'space_after': 9, 'width': 4.92},
            {'columns': 2, 'font_size': 13.0, 'space_after': 8, 'width': 4.92},
            {'columns': 2, 'font_size': 12.0, 'space_after': 7, 'width': 4.92},
        ]
        best: dict | None = None
        best_score: float | None = None
        for preset in presets:
            if preset['columns'] == 1:
                columns = [items]
                required_h = self._estimate_bullet_block_height(items, width_inches=preset['width'], font_size=preset['font_size'], space_after_pt=preset['space_after'])
            else:
                columns = self._balance_items_for_columns(items, width_inches=preset['width'], font_size=preset['font_size'])
                non_empty = [col for col in columns if col]
                required_h = max(
                    self._estimate_bullet_block_height(col, width_inches=preset['width'], font_size=preset['font_size'], space_after_pt=preset['space_after'])
                    for col in non_empty
                )
            fill_ratio = required_h / max(panel_height, 0.1)
            if required_h <= panel_height:
                score = abs(0.84 - fill_ratio)
                if preset['columns'] == 2 and len(items) <= 3:
                    score += 0.15
                if best_score is None or score < best_score:
                    best_score = score
                    best = {**preset, 'columns_data': columns, 'required_h': required_h}
        if best:
            return best
        fallback = presets[-1]
        columns = self._balance_items_for_columns(items, width_inches=fallback['width'], font_size=fallback['font_size'])
        return {**fallback, 'columns_data': columns, 'required_h': panel_height}

    def _write_bullet_block(
        self,
        frame,
        items: list[str],
        *,
        font_size: float,
        color: RGBColor,
        space_after_pt: float,
        align=PP_ALIGN.LEFT,
        min_size: int = 9,
        font_family: str = 'Calibri',
    ) -> None:
        frame.clear()
        frame.word_wrap = True
        frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        frame.margin_left = Inches(0.05)
        frame.margin_right = Inches(0.05)
        frame.margin_top = Inches(0.05)
        frame.margin_bottom = Inches(0.05)
        
        for idx, item in enumerate(items):
            p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            p.alignment = align
            p.text = f'• {self._normalize_text(item)}'
            p.level = 0
            p.space_after = Pt(space_after_pt)
            p.line_spacing = 1.1  # Better readability for long text
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.name = font_family
                run.font.color.rgb = color
        
        self._fit_frame(frame, max_size=font_size, min_size=min_size, font_family=font_family)

    def _set_cell_text(
        self,
        cell,
        text: str,
        *,
        font_size: float,
        bold: bool = False,
        color: tuple[int, int, int] = (31, 41, 55),
        min_size: float = 9.5,
        font_family: str = 'Calibri',
        avail_width: float | None = None,
        avail_height: float | None = None,
    ) -> None:
        cell.text = self._normalize_text(text)
        frame = cell.text_frame
        frame.word_wrap = True
        frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        frame.margin_left = Inches(0.06)
        frame.margin_right = Inches(0.06)
        frame.margin_top = Inches(0.03)
        frame.margin_bottom = Inches(0.03)
        for p in frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.name = font_family
                run.font.color.rgb = RGBColor(*color)
        self._fit_frame(
            frame,
            max_size=font_size,
            min_size=min_size,
            bold=bold,
            font_family=font_family,
            avail_width=avail_width,
            avail_height=avail_height,
        )


    def _facts_variant(self, items: list[str], *, rotation_index: int = 0) -> str:
        # Rotate the layout per content slide so consecutive fact slides never look
        # identical, while only picking a style the content actually fits into.
        if not items:
            return 'bullets'
        count = len(items)
        max_len = max(len(item) for item in items)
        avg_len = sum(len(item) for item in items) / max(1, count)
        feasible = {'bullets'}
        if 2 <= count <= 4 and max_len <= 95 and avg_len <= 82:
            feasible.add('cards')
        if 3 <= count <= 4 and avg_len <= 120:
            feasible.add('spotlight')
        order = ['cards', 'spotlight', 'bullets']
        start = rotation_index % len(order)
        for offset in range(len(order)):
            candidate = order[(start + offset) % len(order)]
            if candidate in feasible:
                return candidate
        return 'bullets'

    @staticmethod
    def _fact_card_coords(item_count: int) -> list[tuple[float, float, float, float]]:
        if item_count == 3:
            return [
                (0.92, 2.42, 3.72, 2.58),
                (4.81, 2.42, 3.72, 2.58),
                (8.70, 2.42, 3.72, 2.58),
            ]
        if item_count == 4:
            return [
                (0.92, 2.20, 5.12, 1.82),
                (7.30, 2.20, 5.12, 1.82),
                (0.92, 4.18, 5.12, 1.82),
                (7.30, 4.18, 5.12, 1.82),
            ]
        if item_count == 2:
            return [
                (1.00, 2.58, 5.40, 2.24),
                (6.93, 2.58, 5.40, 2.24),
            ]
        return []

    def _render_fact_cards(self, slide, *, items: list[str], pack: dict, accent: tuple[int, int, int] | None = None) -> None:
        coords = self._fact_card_coords(len(items))
        colors = self._process_color_list(pack)
        stripe_color = RGBColor(*accent) if accent is not None else None
        font_size = 14.4 if max((len(item) for item in items), default=0) <= 72 else 13.2
        for index, item in enumerate(items):
            x, y, w, h = coords[index]
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(pack, 'surface', (255, 255, 255))
            shape.line.color.rgb = self._rgb(pack, 'border', (191, 219, 254))

            stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.18))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = stripe_color if stripe_color is not None else colors[index % len(colors)]
            stripe.line.fill.background()

            badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.18), Inches(y + 0.26), Inches(0.46), Inches(0.46))
            badge.fill.solid()
            badge.fill.fore_color.rgb = self._rgb(pack, 'primary', (30, 64, 175))
            badge.line.fill.background()
            badge_frame = badge.text_frame
            badge_frame.clear()
            badge_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            badge_p = badge_frame.paragraphs[0]
            badge_p.alignment = PP_ALIGN.CENTER
            badge_run = badge_p.add_run()
            badge_run.text = str(index + 1)
            self._style_run(badge_run, pack, size=12, color_key='on_primary', default_color=(255, 255, 255), bold=True)

            box = slide.shapes.add_textbox(Inches(x + 0.38), Inches(y + 0.86), Inches(w - 0.60), Inches(h - 1.00))
            frame = box.text_frame
            frame.clear()
            frame.word_wrap = True
            frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            frame.margin_left = Inches(0.04)
            frame.margin_right = Inches(0.04)
            p = frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = item
            self._style_run(run, pack, size=font_size, color_key='text', default_color=(30, 41, 59))
            self._fit_frame(frame, max_size=font_size, min_size=10.8, font_family=self._font_family(pack))

    def _render_focus_spotlight(self, slide, *, items: list[str], pack: dict, accent: tuple[int, int, int] | None = None) -> None:
        highlight = items[0] if items else ''
        remaining_items = items[1:] if len(items) > 1 else items
        left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(2.18), Inches(4.12), Inches(3.98))
        left.fill.solid()
        left.fill.fore_color.rgb = RGBColor(*accent) if accent is not None else self._rgb(pack, 'surface_alt', (239, 246, 255))
        left.line.color.rgb = self._rgb(pack, 'border', (191, 219, 254))

        frame = left.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Inches(0.18)
        frame.margin_right = Inches(0.18)
        frame.margin_top = Inches(0.14)
        header = frame.paragraphs[0]
        header.alignment = PP_ALIGN.LEFT
        run = header.add_run()
        run.text = '01'
        self._style_run(run, pack, size=13.2, color_key='primary', default_color=(30, 64, 175), bold=True, font_name=self._heading_font(pack))

        focus_p = frame.add_paragraph()
        focus_p.alignment = PP_ALIGN.LEFT
        focus_p.space_before = Pt(8)
        focus_p.space_after = Pt(10)
        focus_run = focus_p.add_run()
        focus_run.text = self._normalize_text(highlight)
        self._style_run(focus_run, pack, size=15, color_key='text', default_color=(30, 41, 59))
        self._fit_frame(frame, max_size=15, min_size=11.2, font_family=self._font_family(pack))

        right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.36), Inches(2.18), Inches(7.00), Inches(3.98))
        right.fill.solid()
        right.fill.fore_color.rgb = self._rgb(pack, 'surface', (255, 255, 255))
        right.line.color.rgb = self._rgb(pack, 'border', (203, 213, 225))

        inner = slide.shapes.add_textbox(Inches(5.68), Inches(2.46), Inches(6.34), Inches(3.42))
        self._write_bullet_block(
            inner.text_frame,
            remaining_items,
            font_size=13.6,
            color=self._rgb(pack, 'text', (31, 41, 55)),
            space_after_pt=9,
            min_size=10.2,
            font_family=self._font_family(pack),
        )

    @staticmethod
    def _process_coords(item_count: int) -> list[tuple[float, float, float, float]]:
        if item_count <= 3:
            width = 3.55 if item_count == 3 else 4.45
            gap = 0.38
            total_width = (width * item_count) + (gap * max(0, item_count - 1))
            start_x = max(0.78, (13.333 - total_width) / 2)
            return [(start_x + i * (width + gap), 2.42, width, 1.92) for i in range(item_count)]
        if item_count == 4:
            return [
                (0.92, 2.12, 5.15, 1.78),
                (7.26, 2.12, 5.15, 1.78),
                (0.92, 4.18, 5.15, 1.78),
                (7.26, 4.18, 5.15, 1.78),
            ]
        if item_count == 5:
            return [
                (0.58, 2.02, 4.02, 1.62),
                (4.66, 2.02, 4.02, 1.62),
                (8.74, 2.02, 4.02, 1.62),
                (2.62, 4.12, 4.02, 1.62),
                (6.70, 4.12, 4.02, 1.62),
            ]
        if item_count == 6:
            return [
                (0.78, 2.12, 3.75, 1.68),
                (4.78, 2.12, 3.75, 1.68),
                (8.78, 2.12, 3.75, 1.68),
                (0.78, 4.12, 3.75, 1.68),
                (4.78, 4.12, 3.75, 1.68),
                (8.78, 4.12, 3.75, 1.68),
            ]
        # For 7-8 items
        width = 3.00
        height = 1.60
        gap_x = 0.15
        gap_y = 0.20
        coords = []
        rows = 2
        cols = (item_count + 1) // 2
        start_x = (13.333 - (cols * width + (cols - 1) * gap_x)) / 2
        start_y = 2.15
        for i in range(item_count):
            r = i // cols
            c = i % cols
            coords.append((start_x + c * (width + gap_x), start_y + r * (height + gap_y), width, height))
        return coords

    @staticmethod
    def _agenda_layout(item_count: int) -> dict[str, Any]:
        if item_count <= 5:
            return {'columns': 1, 'left': 0.90, 'top': 2.04, 'col_width': 5.72, 'item_height': 0.58, 'gap': 0.12}
        if item_count <= 7:
            return {'columns': 1, 'left': 0.88, 'top': 2.00, 'col_width': 5.74, 'item_height': 0.50, 'gap': 0.09}
        return {'columns': 2, 'left': 0.82, 'top': 2.00, 'col_width': 2.76, 'item_height': 0.52, 'gap': 0.10, 'col_gap': 0.16}

    def _base_slide(
        self,
        prs: Presentation,
        *,
        title: str,
        presenter_name: str,
        page_number: int,
        total_slides: int,
        pack: dict,
        subtitle: str | None = None,
        background_rgb: tuple[int, int, int] | None = None,
        accent_rgb: tuple[int, int, int] | None = None,
        section_number: int | None = None,
        motif_bold: bool = False,
    ):
        header_style = self._layout_value(pack, 'header_style', 'top_band')
        heading_font = self._heading_font(pack)
        accent_color = RGBColor(*accent_rgb) if accent_rgb is not None else None
        bar_color = RGBColor(*self._visible_accent(pack, accent_rgb)) if accent_rgb is not None else None
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(*(background_rgb or self._rgb_tuple(pack, 'background', (248, 250, 252))))
        self._draw_background_motif(slide, pack, bold=motif_bold)

        title_x = 0.96
        title_y = 0.80
        title_w = 9.40 if section_number is not None else 11.10
        subtitle_y = 1.48
        footer_y = 6.84

        # Large, faint chapter number in the top-right — gives each content slide
        # its own identity and a sense of progression through the deck.
        if section_number is not None:
            watermark = slide.shapes.add_textbox(Inches(10.55), Inches(0.50), Inches(2.0), Inches(1.34))
            wm_frame = watermark.text_frame
            wm_frame.clear()
            wm_frame.word_wrap = False
            wm_p = wm_frame.paragraphs[0]
            wm_p.alignment = PP_ALIGN.RIGHT
            wm_run = wm_p.add_run()
            wm_run.text = f'{section_number:02d}'
            wm_run.font.size = Pt(66)
            wm_run.font.bold = True
            wm_run.font.name = heading_font
            wm_run.font.color.rgb = accent_color if accent_color is not None else self._rgb(pack, 'accent', (186, 230, 253))

        if header_style == 'left_rail':
            rail = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.52), Inches(7.5))
            rail.fill.solid()
            rail.fill.fore_color.rgb = self._rgb(pack, 'primary', (4, 120, 87))
            rail.line.fill.background()
            marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.52), Inches(0.76), Inches(0.08), Inches(1.42 if subtitle else 0.92))
            marker.fill.solid()
            marker.fill.fore_color.rgb = bar_color if bar_color is not None else self._rgb(pack, 'secondary', (20, 184, 166))
            marker.line.fill.background()
            title_x = 0.92
            title_y = 0.74
            subtitle_y = 1.44
        elif header_style == 'corner_band':
            corner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(10.18), Inches(0), Inches(3.16), Inches(1.16))
            corner.fill.solid()
            corner.fill.fore_color.rgb = self._rgb(pack, 'primary', (194, 65, 12))
            corner.line.fill.background()
            accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.76), Inches(0.72), Inches(1.24), Inches(0.08))
            accent.fill.solid()
            accent.fill.fore_color.rgb = bar_color if bar_color is not None else self._rgb(pack, 'secondary', (15, 118, 110))
            accent.line.fill.background()
            title_y = 0.84
            subtitle_y = 1.52
        else:
            band_h = 0.68 if header_style == 'dark_bar' else 0.62
            top_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(band_h))
            top_band.fill.solid()
            top_band.fill.fore_color.rgb = self._rgb(pack, 'primary', (30, 64, 175))
            top_band.line.fill.background()

            accent = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(0.72),
                Inches(0.82),
                Inches(0.10),
                Inches(1.22 if subtitle else 0.82),
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = bar_color if bar_color is not None else self._rgb(pack, 'secondary', (96, 165, 250))
            accent.line.fill.background()

        if title:
            title_box = slide.shapes.add_textbox(Inches(title_x), Inches(title_y), Inches(title_w), Inches(0.92))
            title_frame = title_box.text_frame
            title_frame.clear()
            title_frame.word_wrap = True
            title_frame.margin_left = 0
            title_frame.margin_right = 0
            p = title_frame.paragraphs[0]
            run = p.add_run()
            run.text = self._normalize_text(title)
            self._style_run(run, pack, size=25, color_key='text', default_color=(15, 23, 42), bold=True, font_name=heading_font)
            self._fit_frame(title_frame, max_size=25, min_size=18, bold=True, font_family=heading_font)

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(Inches(title_x), Inches(subtitle_y), Inches(11.25), Inches(0.68))
            subtitle_frame = subtitle_box.text_frame
            subtitle_frame.clear()
            subtitle_frame.word_wrap = True
            subtitle_frame.margin_left = 0
            subtitle_frame.margin_right = 0
            p = subtitle_frame.paragraphs[0]
            run = p.add_run()
            run.text = self._normalize_text(subtitle)
            self._style_run(run, pack, size=12.5, color_key='muted', default_color=(71, 85, 105))
            self._fit_frame(subtitle_frame, max_size=12, min_size=10, font_family=self._font_family(pack))

        footer_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(footer_y), Inches(11.86), Inches(0.02))
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = self._rgb(pack, 'border', (203, 213, 225))
        footer_line.line.fill.background()

        author_box = slide.shapes.add_textbox(Inches(0.76), Inches(6.90), Inches(5.8), Inches(0.26))
        author_frame = author_box.text_frame
        author_frame.clear()
        author_run = author_frame.paragraphs[0].add_run()
        author_run.text = f"{pack['prepared_by']}: {self._normalize_text(presenter_name, max_chars=40)}"
        self._style_run(author_run, pack, size=10, color_key='footer', default_color=(100, 116, 139))

        page_box = slide.shapes.add_textbox(Inches(11.45), Inches(6.90), Inches(1.05), Inches(0.26))
        page_frame = page_box.text_frame
        page_frame.clear()
        page_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        page_run = page_frame.paragraphs[0].add_run()
        page_run.text = f'{page_number}/{total_slides}'
        self._style_run(page_run, pack, size=10, color_key='footer', default_color=(100, 116, 139))
        return slide

    def _add_title_slide(
        self,
        prs: Presentation,
        *,
        presentation_title: str,
        subtitle: str,
        presenter_name: str,
        agenda_preview: list[str],
        page_number: int,
        total_slides: int,
        pack: dict,
    ) -> None:
        slide = self._base_slide(
            prs,
            title='',
            presenter_name=presenter_name,
            page_number=page_number,
            total_slides=total_slides,
            pack=pack,
            background_rgb=self._rgb_tuple(pack, 'cover_background', (239, 246, 255)),
            motif_bold=True,
        )

        cover_style = self._layout_value(pack, 'cover_style', 'split_card')
        if cover_style == 'center_focus':
            title_coords = (1.18, 1.02, 10.90, 1.36)
            subtitle_coords = (1.74, 2.36, 9.78, 0.80)
            card_coords = (1.28, 3.36, 10.74, 2.28)
        elif cover_style == 'magazine':
            block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(7.38), Inches(0.86), Inches(4.68), Inches(5.28))
            block.fill.solid()
            block.fill.fore_color.rgb = self._rgb(pack, 'accent', (254, 215, 170))
            block.line.fill.background()
            title_coords = (0.88, 1.02, 6.00, 1.82)
            subtitle_coords = (0.92, 3.02, 5.80, 1.06)
            card_coords = (7.06, 1.58, 4.94, 3.74)
        elif cover_style == 'dashboard':
            title_coords = (0.88, 1.08, 7.00, 1.72)
            subtitle_coords = (0.92, 2.96, 6.42, 1.06)
            card_coords = (8.12, 1.36, 4.16, 4.02)
        else:
            title_coords = (0.92, 1.28, 6.6, 1.68)
            subtitle_coords = (0.95, 2.92, 6.2, 1.18)
            card_coords = (7.7, 1.56, 4.58, 3.78)

        title_box = slide.shapes.add_textbox(Inches(title_coords[0]), Inches(title_coords[1]), Inches(title_coords[2]), Inches(title_coords[3]))
        frame = title_box.text_frame
        frame.clear()
        frame.word_wrap = True
        p = frame.paragraphs[0]
        if cover_style == 'center_focus':
            p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = self._normalize_text(presentation_title)
        self._style_run(run, pack, size=30, color_key='text', default_color=(15, 23, 42), bold=True, font_name=self._heading_font(pack))
        self._fit_frame(frame, max_size=30, min_size=20, bold=True, font_family=self._heading_font(pack))

        # Accent rule between the title and the subtitle.
        rule_width = 1.8
        rule_left = (title_coords[0] + title_coords[2] / 2 - rule_width / 2) if cover_style == 'center_focus' else title_coords[0] + 0.04
        rule_top = subtitle_coords[1] - 0.24
        accent_rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(rule_left), Inches(rule_top), Inches(rule_width), Inches(0.08))
        accent_rule.fill.solid()
        accent_rule.fill.fore_color.rgb = RGBColor(*self._visible_accent(pack, self._rgb_tuple(pack, 'secondary', (14, 165, 233))))
        accent_rule.line.fill.background()

        subtitle_box = slide.shapes.add_textbox(Inches(subtitle_coords[0]), Inches(subtitle_coords[1]), Inches(subtitle_coords[2]), Inches(subtitle_coords[3]))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_frame.word_wrap = True
        if cover_style == 'center_focus':
            subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_run = subtitle_frame.paragraphs[0].add_run()
        subtitle_run.text = self._normalize_text(subtitle)
        self._style_run(subtitle_run, pack, size=15, color_key='muted', default_color=(51, 65, 85))
        self._fit_frame(subtitle_frame, max_size=15, min_size=11, font_family=self._font_family(pack))

        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(card_coords[0]), Inches(card_coords[1]), Inches(card_coords[2]), Inches(card_coords[3]))
        card.fill.solid()
        card.fill.fore_color.rgb = self._rgb(pack, 'surface', (255, 255, 255))
        card.line.color.rgb = self._rgb(pack, 'border', (191, 219, 254))
        card_frame = card.text_frame
        card_frame.clear()
        card_frame.word_wrap = True
        card_frame.margin_left = Inches(0.18)
        card_frame.margin_right = Inches(0.18)
        card_frame.margin_top = Inches(0.12)
        card_frame.margin_bottom = Inches(0.12)

        header = card_frame.paragraphs[0]
        header.alignment = PP_ALIGN.CENTER
        run = header.add_run()
        run.text = str(pack['key_focus'])
        self._style_run(run, pack, size=17, color_key='primary', default_color=(30, 64, 175), bold=True)

        for item in (agenda_preview or list(pack['cover_points'])[:4])[:4]:
            paragraph = card_frame.add_paragraph()
            paragraph.text = f'• {self._normalize_text(item)}'
            paragraph.space_after = Pt(9)
            for run in paragraph.runs:
                self._style_run(run, pack, size=15, color_key='text', default_color=(30, 41, 59))
        self._fit_frame(card_frame, max_size=17, min_size=11, font_family=self._font_family(pack))

    def _add_agenda_slide(self, prs: Presentation, *, agenda_items: list[str], agenda_notes: list[str], presenter_name: str, page_number: int, total_slides: int, pack: dict) -> None:
        slide = self._base_slide(
            prs,
            title=str(pack['agenda']),
            presenter_name=presenter_name,
            page_number=page_number,
            total_slides=total_slides,
            pack=pack,
        )
        visible_items = [self._normalize_text(item) for item in agenda_items[:8] if self._normalize_text(item)]
        agenda_style = self._layout_value(pack, 'agenda_style', 'list_note')
        if agenda_style in {'grid', 'tiles'}:
            coords: list[tuple[float, float, float, float]] = []
            cols = 4 if len(visible_items) > 4 else 2
            rows = 2 if len(visible_items) > 4 else max(1, (len(visible_items) + 1) // 2)
            card_w = 2.84 if cols == 4 else 5.56
            card_h = 1.28 if rows == 2 else 1.44
            start_x = 0.82
            start_y = 2.28
            gap_x = 0.30
            gap_y = 0.34
            for index in range(len(visible_items)):
                row = index // cols
                col = index % cols
                coords.append((start_x + col * (card_w + gap_x), start_y + row * (card_h + gap_y), card_w, card_h))

            for index, item in enumerate(visible_items, start=1):
                x, y, w, h = coords[index - 1]
                shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._rgb(pack, 'surface', (255, 255, 255))
                shape.line.color.rgb = self._rgb(pack, 'border', (191, 219, 254))

                badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.16), Inches(y + 0.16), Inches(0.42), Inches(0.42))
                badge.fill.solid()
                badge.fill.fore_color.rgb = self._rgb(pack, 'primary', (30, 64, 175))
                badge.line.fill.background()
                self._draw_icon(
                    slide,
                    self._icon_for_text(item, fallback_index=index - 1),
                    left=x + 0.16 + 0.115,
                    top=y + 0.16 + 0.115,
                    size=0.19,
                    color=self._rgb_tuple(pack, 'on_primary', (255, 255, 255)),
                )

                text_box = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 0.56), Inches(w - 0.48), Inches(h - 0.68))
                frame = text_box.text_frame
                frame.clear()
                frame.word_wrap = True
                frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                p = frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = item
                self._style_run(run, pack, size=12.6 if cols == 4 else 14.0, color_key='text', default_color=(15, 23, 42), bold=agenda_style == 'tiles')
                self._fit_frame(frame, max_size=14, min_size=9.4, bold=agenda_style == 'tiles', font_family=self._font_family(pack))
            return

        layout = self._agenda_layout(len(visible_items))
        note_left = 7.08
        note_width = 5.1
        body_font = 13.8 if layout['columns'] == 1 and len(visible_items) <= 5 else 12.4 if layout['columns'] == 1 else 11.7

        for index, item in enumerate(visible_items, start=1):
            if layout['columns'] == 1:
                x = layout['left']
                y = layout['top'] + (layout['item_height'] + layout['gap']) * (index - 1)
            else:
                left_count = (len(visible_items) + 1) // 2
                col = 0 if index <= left_count else 1
                row = (index - 1) if col == 0 else (index - 1 - left_count)
                x = layout['left'] + col * (layout['col_width'] + layout['col_gap'])
                y = layout['top'] + row * (layout['item_height'] + layout['gap'])
                note_left = 6.82
                note_width = 5.34

            item_h = layout['item_height']
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(layout['col_width']), Inches(item_h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(pack, 'surface', (255, 255, 255))
            shape.line.color.rgb = self._rgb(pack, 'border', (191, 219, 254))

            icon_size = min(0.30, item_h - 0.18)
            self._draw_icon(
                slide,
                self._icon_for_text(item, fallback_index=index - 1),
                left=x + 0.16,
                top=y + (item_h - icon_size) / 2,
                size=icon_size,
                color=self._rgb_tuple(pack, 'primary', (30, 64, 175)),
            )

            frame = shape.text_frame
            frame.clear()
            frame.word_wrap = True
            frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            frame.margin_left = Inches(0.16 + icon_size + 0.14)
            frame.margin_right = Inches(0.10)
            run = frame.paragraphs[0].add_run()
            run.text = f'{index}. {item}'
            self._style_run(run, pack, size=body_font, color_key='text', default_color=(15, 23, 42))
            self._fit_frame(frame, max_size=body_font, min_size=10.8, font_family=self._font_family(pack))

        note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(note_left), Inches(2.10), Inches(note_width), Inches(3.62))
        note.fill.solid()
        note.fill.fore_color.rgb = self._rgb(pack, 'surface_alt', (219, 234, 254))
        note.line.fill.background()
        note_frame = note.text_frame
        note_frame.clear()
        note_frame.word_wrap = True
        note_frame.margin_left = Inches(0.16)
        note_frame.margin_right = Inches(0.16)
        note_frame.margin_top = Inches(0.10)
        note_frame.margin_bottom = Inches(0.10)
        title_paragraph = note_frame.paragraphs[0]
        title_paragraph.alignment = PP_ALIGN.CENTER
        run = title_paragraph.add_run()
        run.text = str(pack['agenda_note_title'])
        self._style_run(run, pack, size=15.5, color_key='primary', default_color=(30, 64, 175), bold=True)
        for item in (agenda_notes or visible_items)[:4]:
            p = note_frame.add_paragraph()
            p.text = f'• {self._normalize_text(item)}'
            p.space_after = Pt(7)
            for run in p.runs:
                self._style_run(run, pack, size=12.0, color_key='text', default_color=(30, 41, 59))
        self._fit_frame(note_frame, max_size=15, min_size=10.5, font_family=self._font_family(pack))


    def _add_facts_slide(self, prs: Presentation, *, section: TopicSection, presenter_name: str, page_number: int, total_slides: int, pack: dict, section_index: int = 0) -> None:
        accent = self._section_accent_tuple(pack, section_index)
        slide = self._base_slide(
            prs,
            title=section.title,
            presenter_name=presenter_name,
            page_number=page_number,
            total_slides=total_slides,
            pack=pack,
            accent_rgb=accent,
            section_number=section_index + 1,
        )

        facts = [self._normalize_text(item) for item in section.facts[:10] if self._normalize_text(item)]
        if not facts:
            return

        # Rotate layouts across the deck so fact slides feel distinct from one another.
        variant = self._facts_variant(facts, rotation_index=section_index)
        if variant == 'cards':
            self._render_fact_cards(slide, items=facts[:4], pack=pack, accent=accent)
            return
        if variant == 'spotlight':
            self._render_focus_spotlight(slide, items=facts[:4], pack=pack, accent=accent)
            return

        content_top, content_bottom = self._content_bounds(False)
        panel_left = 0.78
        panel_top = content_top
        panel_width = 11.78
        panel_height = content_bottom - content_top

        panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(panel_left), Inches(panel_top), Inches(panel_width), Inches(panel_height))
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._rgb(pack, 'surface', (255, 255, 255))
        panel.line.color.rgb = self._rgb(pack, 'border', (203, 213, 225))

        accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(panel_left + 0.18), Inches(panel_top + 0.22), Inches(0.10), Inches(panel_height - 0.44))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor(*self._visible_accent(pack, accent))
        accent_bar.line.fill.background()

        layout = self._select_facts_layout(facts, panel_height=panel_height - 0.48)
        if layout['columns'] == 1:
            text_box = slide.shapes.add_textbox(Inches(panel_left + 0.42), Inches(panel_top + 0.22), Inches(10.92), Inches(panel_height - 0.38))
            self._write_bullet_block(
                text_box.text_frame,
                layout['columns_data'][0],
                font_size=layout['font_size'],
                color=self._rgb(pack, 'text', (31, 41, 55)),
                space_after_pt=layout['space_after'],
                min_size=10.2,
                font_family=self._font_family(pack),
            )
            return

        left_box = slide.shapes.add_textbox(Inches(panel_left + 0.40), Inches(panel_top + 0.22), Inches(5.02), Inches(panel_height - 0.36))
        right_box = slide.shapes.add_textbox(Inches(panel_left + 5.98), Inches(panel_top + 0.22), Inches(5.02), Inches(panel_height - 0.36))

        self._write_bullet_block(left_box.text_frame, layout['columns_data'][0], font_size=layout['font_size'], color=self._rgb(pack, 'text', (31, 41, 55)), space_after_pt=layout['space_after'], min_size=10.2, font_family=self._font_family(pack))
        self._write_bullet_block(right_box.text_frame, layout['columns_data'][1], font_size=layout['font_size'], color=self._rgb(pack, 'text', (31, 41, 55)), space_after_pt=layout['space_after'], min_size=10.2, font_family=self._font_family(pack))

        divider = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(panel_left + 5.56), Inches(panel_top + 0.26), Inches(0.02), Inches(panel_height - 0.52))
        divider.fill.solid()
        divider.fill.fore_color.rgb = self._rgb(pack, 'border', (226, 232, 240))
        divider.line.fill.background()

    def _add_process_slide(self, prs: Presentation, *, section: TopicSection, presenter_name: str, page_number: int, total_slides: int, pack: dict, section_index: int = 0) -> None:
        slide = self._base_slide(
            prs,
            title=section.title,
            presenter_name=presenter_name,
            page_number=page_number,
            total_slides=total_slides,
            pack=pack,
            accent_rgb=self._section_accent_tuple(pack, section_index),
            section_number=section_index + 1,
        )

        items = [self._normalize_text(item) for item in section.facts[:5] if self._normalize_text(item)]
        if not items:
            return
        coords = self._process_coords(len(items))
        colors = self._process_color_list(pack)
        max_len = max((len(item) for item in items), default=0)
        body_size = 12.5 if max_len <= 70 else 11.5 if max_len <= 100 else 10.5
        process_style = self._layout_value(pack, 'process_style', 'cards')

        if process_style == 'timeline':
            line_y = 3.42
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.92), Inches(line_y), Inches(11.48), Inches(0.04))
            line.fill.solid()
            line.fill.fore_color.rgb = self._rgb(pack, 'accent', (191, 219, 254))
            line.line.fill.background()
            item_w = min(2.32, 10.9 / max(1, len(items)))
            gap = (11.20 - (item_w * len(items))) / max(1, len(items) - 1) if len(items) > 1 else 0
            start_x = 1.06
            for idx, item in enumerate(items):
                x = start_x + idx * (item_w + gap)
                node = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + item_w / 2 - 0.20), Inches(line_y - 0.18), Inches(0.40), Inches(0.40))
                node.fill.solid()
                node.fill.fore_color.rgb = self._rgb(pack, 'primary', (30, 64, 175))
                node.line.fill.background()
                box_y = 2.08 if idx % 2 == 0 else 3.86
                shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(box_y), Inches(item_w), Inches(1.10))
                shape.fill.solid()
                shape.fill.fore_color.rgb = colors[idx % len(colors)]
                shape.line.color.rgb = self._rgb(pack, 'border', (191, 219, 254))
                frame = shape.text_frame
                frame.clear()
                frame.word_wrap = True
                frame.margin_left = Inches(0.08)
                frame.margin_right = Inches(0.08)
                frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
                p = frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = item
                self._style_run(run, pack, size=10.6, color_key='text', default_color=(51, 65, 85))
                self._fit_frame(frame, max_size=10.6, min_size=8.2, font_family=self._font_family(pack))
            return

        for idx, item in enumerate(items):
            x, y, w, h = coords[idx]
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = colors[idx % len(colors)]
            shape.line.color.rgb = self._rgb(pack, 'border', (191, 219, 254))

            badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.16), Inches(y + 0.16), Inches(0.54), Inches(0.54))
            badge.fill.solid()
            badge.fill.fore_color.rgb = self._rgb(pack, 'primary', (30, 64, 175))
            badge.line.fill.background()
            badge_frame = badge.text_frame
            badge_frame.clear()
            badge_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            badge_p = badge_frame.paragraphs[0]
            badge_p.alignment = PP_ALIGN.CENTER
            badge_run = badge_p.add_run()
            badge_run.text = str(idx + 1)
            self._style_run(badge_run, pack, size=14, color_key='on_primary', default_color=(255, 255, 255), bold=True)

            text_box = slide.shapes.add_textbox(Inches(x + 0.28), Inches(y + 0.82), Inches(w - 0.56), Inches(h - 0.98))
            frame = text_box.text_frame
            frame.clear()
            frame.word_wrap = True
            frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            p = frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = item
            self._style_run(run, pack, size=body_size, color_key='text', default_color=(51, 65, 85))
            self._fit_frame(frame, max_size=int(round(body_size)), min_size=9, font_family=self._font_family(pack))

        if len(items) == 5:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(5.64), Inches(3.34), Inches(2.05), Inches(0.48))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = self._rgb(pack, 'accent', (191, 219, 254))
            arrow.line.fill.background()


    @staticmethod
    def _parse_number(text: str) -> float | None:
        cleaned = str(text or '').strip()
        if not cleaned:
            return None
        match = re.search(r'[-+]?\d[\d\s.,]*\d|\d', cleaned)
        if not match:
            return None
        token = match.group(0).replace(' ', '').replace(' ', '')
        if ',' in token and '.' in token:
            token = token.replace(',', '')
        elif ',' in token:
            token = token.replace('.', '').replace(',', '.') if re.search(r',\d{1,2}$', token) else token.replace(',', '')
        try:
            return float(token)
        except ValueError:
            return None

    def _build_chart_spec(self, table) -> dict | None:
        columns = [self._normalize_text(col) for col in table.columns]
        rows = table.rows
        if len(columns) < 2 or not (2 <= len(rows) <= 6):
            return None
        categories = [self._normalize_text(row[0]) for row in rows if row]
        if len(categories) != len(rows) or any(not cat for cat in categories):
            return None
        series: list[tuple[str, list[float]]] = []
        for col_idx in range(1, len(columns)):
            values = [self._parse_number(row[col_idx]) if col_idx < len(row) else None for row in rows]
            if all(value is not None for value in values):
                series.append((columns[col_idx], [float(value) for value in values]))  # type: ignore[arg-type]
            if len(series) >= 3:
                break
        if not series:
            return None
        return {
            'categories': categories,
            'series': series,
            'horizontal': any(len(cat) > 14 for cat in categories),
        }

    @staticmethod
    def _apply_solid_fill(sp_pr, rgb: tuple[int, int, int]) -> None:
        hex_value = '%02X%02X%02X' % rgb
        for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
            for element in sp_pr.findall(qn(tag)):
                sp_pr.remove(element)
        solid_fill = sp_pr.makeelement(qn('a:solidFill'), {})
        srgb = solid_fill.makeelement(qn('a:srgbClr'), {'val': hex_value})
        solid_fill.append(srgb)
        sp_pr.insert(0, solid_fill)

    def _style_chart_surface(self, chart, pack: dict) -> None:
        # Charts sit directly on the slide, so paint a surface + readable text colour
        # explicitly; otherwise dark templates get black axis labels on a dark slide.
        surface = self._rgb_tuple(pack, 'surface', (255, 255, 255))
        border = self._rgb_tuple(pack, 'border', (203, 213, 225))
        try:
            chart_space = chart._chartSpace
            sp_pr = chart_space.find(qn('c:spPr'))
            if sp_pr is None:
                sp_pr = chart_space.makeelement(qn('c:spPr'), {})
                tx_pr = chart_space.find(qn('c:txPr'))
                if tx_pr is not None:
                    tx_pr.addprevious(sp_pr)
                else:
                    chart_space.append(sp_pr)
            self._apply_solid_fill(sp_pr, surface)
            ln = sp_pr.makeelement(qn('a:ln'), {})
            ln_fill = ln.makeelement(qn('a:solidFill'), {})
            ln_fill.append(ln.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % border}))
            ln.append(ln_fill)
            sp_pr.append(ln)
        except Exception:
            pass

    def _render_chart(self, slide, *, spec: dict, section: TopicSection, pack: dict) -> None:
        content_top, content_bottom = self._content_bounds(False)
        area_height = content_bottom - content_top
        body_font = self._font_family(pack)
        side_facts = [self._normalize_text(fact) for fact in section.facts[:3] if self._normalize_text(fact)]
        has_note = bool(side_facts)

        chart_width = 9.05 if has_note else 11.70
        note_left = 10.06
        note_width = 2.34

        chart_data = CategoryChartData()
        chart_data.categories = spec['categories']
        for name, values in spec['series']:
            chart_data.add_series(name, values)

        chart_type = XL_CHART_TYPE.BAR_CLUSTERED if spec['horizontal'] else XL_CHART_TYPE.COLUMN_CLUSTERED
        graphic_frame = slide.shapes.add_chart(
            chart_type, Inches(0.82), Inches(content_top), Inches(chart_width), Inches(area_height), chart_data,
        )
        chart = graphic_frame.chart
        chart.font.name = body_font
        chart.font.size = Pt(11)
        chart.font.color.rgb = self._rgb(pack, 'text', (15, 23, 42))
        chart.has_title = False
        self._style_chart_surface(chart, pack)

        multi = len(spec['series']) > 1
        if multi:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
        else:
            chart.has_legend = False

        series_palette = [
            self._rgb_tuple(pack, 'primary', (30, 64, 175)),
            self._rgb_tuple(pack, 'secondary', (96, 165, 250)),
            self._rgb_tuple(pack, 'accent', (191, 219, 254)),
        ]
        plot = chart.plots[0]
        plot.gap_width = 70
        for index, series in enumerate(plot.series):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(*series_palette[index % len(series_palette)])
            series.format.line.fill.background()

        if not multi:
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.number_format = '0.#'
            data_labels.number_format_is_linked = False
            data_labels.font.size = Pt(10)
            data_labels.font.name = body_font
            data_labels.font.bold = True
            data_labels.font.color.rgb = self._rgb(pack, 'text', (15, 23, 42))
            try:
                data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            except Exception:
                pass

        for axis in (chart.category_axis, chart.value_axis):
            try:
                axis.tick_labels.font.size = Pt(10)
                axis.tick_labels.font.name = body_font
                axis.tick_labels.font.color.rgb = self._rgb(pack, 'text', (15, 23, 42))
                axis.format.line.color.rgb = self._rgb(pack, 'border', (203, 213, 225))
            except Exception:
                pass
        try:
            chart.value_axis.has_major_gridlines = True
            chart.value_axis.major_gridlines.format.line.color.rgb = self._rgb(pack, 'border', (226, 232, 240))
            chart.category_axis.has_major_gridlines = False
        except Exception:
            pass

        if not has_note:
            return

        note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(note_left), Inches(content_top), Inches(note_width), Inches(area_height))
        note.fill.solid()
        note.fill.fore_color.rgb = self._rgb(pack, 'surface_alt', (239, 246, 255))
        note.line.color.rgb = self._rgb(pack, 'border', (191, 219, 254))
        note_body = slide.shapes.add_textbox(Inches(note_left + 0.16), Inches(content_top + 0.18), Inches(note_width - 0.32), Inches(area_height - 0.34))
        note_frame = note_body.text_frame
        note_frame.clear()
        note_frame.word_wrap = True
        for index, fact in enumerate(side_facts):
            p = note_frame.paragraphs[0] if index == 0 else note_frame.add_paragraph()
            p.text = f'• {fact}'
            p.space_after = Pt(6)
            for run in p.runs:
                self._style_run(run, pack, size=11, color_key='text', default_color=(31, 41, 55))
        self._fit_frame(note_frame, max_size=11, min_size=9.0, font_family=body_font, avail_width=note_width - 0.32, avail_height=area_height - 0.34)

    def _add_table_slide(self, prs: Presentation, *, section: TopicSection, presenter_name: str, page_number: int, total_slides: int, pack: dict, section_index: int = 0) -> None:
        assert section.table is not None
        accent = self._section_accent_tuple(pack, section_index)

        # When the table is really numeric comparison data, a chart communicates it
        # far better than a grid of numbers.
        chart_spec = self._build_chart_spec(section.table)
        if chart_spec is not None:
            slide = self._base_slide(
                prs,
                title=section.title,
                presenter_name=presenter_name,
                page_number=page_number,
                total_slides=total_slides,
                pack=pack,
                accent_rgb=accent,
                section_number=section_index + 1,
            )
            self._render_chart(slide, spec=chart_spec, section=section, pack=pack)
            return

        slide = self._base_slide(
            prs,
            title=section.title,
            presenter_name=presenter_name,
            page_number=page_number,
            total_slides=total_slides,
            pack=pack,
            accent_rgb=accent,
            section_number=section_index + 1,
        )

        content_top, content_bottom = self._content_bounds(False)
        area_height = content_bottom - content_top

        rows = len(section.table.rows) + 1
        cols = len(section.table.columns)
        max_cell_len = max((len(self._normalize_text(cell)) for row in section.table.rows for cell in row), default=20)

        use_side_note = cols >= 3
        if use_side_note:
            table_left = 0.76
            table_top = content_top
            table_width = 9.35
            note_left = 10.32
            note_width = 2.06
        else:
            table_left = 0.78
            table_top = content_top + 0.06
            table_width = 11.78
            note_left = note_width = 0.0

        table_shape = slide.shapes.add_table(rows, cols, Inches(table_left), Inches(table_top), Inches(table_width), Inches(area_height - (0.06 if not use_side_note else 0.0)))
        table = table_shape.table

        def col_weight(index: int) -> float:
            header = self._normalize_text(section.table.columns[index])
            header_weight = max(8, len(header) * 1.9)
            row_weight = max((len(self._normalize_text(row[index])) for row in section.table.rows if index < len(row)), default=12)
            return max(header_weight, row_weight)

        weights = [col_weight(i) for i in range(cols)]
        total_weight = sum(weights) or cols
        min_width = 2.15 if cols == 2 else 1.35
        col_widths = [max(min_width, round(table_width * (weights[idx] / total_weight), 2)) for idx in range(cols)]
        # Rescale so the columns never extend past the allotted table width.
        width_sum = sum(col_widths)
        if width_sum > table_width:
            scale = table_width / width_sum
            col_widths = [round(width * scale, 2) for width in col_widths]
        for idx in range(cols):
            table.columns[idx].width = Inches(col_widths[idx])

        header_height = 0.56 if cols <= 3 else 0.54
        body_height = max(0.62 if rows <= 4 else 0.54, (area_height - header_height) / max(1, rows - 1))
        table.rows[0].height = Inches(header_height)
        for row_idx in range(1, rows):
            table.rows[row_idx].height = Inches(body_height)

        if cols == 2:
            header_font = 12.4
            body_font = 12.0 if max_cell_len <= 54 else 11.2
            body_min = 10.4
        else:
            header_font = 11.5 if cols <= 3 else 10.8
            body_font = 11.1 if rows <= 4 and max_cell_len <= 48 else 10.2 if max_cell_len <= 70 else 9.6
            body_min = 8.8

        for col_idx, column_name in enumerate(section.table.columns):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._rgb(pack, 'table_header_bg', (219, 234, 254))
            self._set_cell_text(
                cell,
                column_name,
                font_size=header_font,
                bold=True,
                color=self._rgb_tuple(pack, 'table_header_text', (30, 64, 175)),
                min_size=max(10.0, header_font - 1.2),
                font_family=self._font_family(pack),
                avail_width=col_widths[col_idx],
                avail_height=header_height,
            )

        for row_idx, row in enumerate(section.table.rows, start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._rgb(pack, 'table_row_odd', (255, 255, 255)) if row_idx % 2 else self._rgb(pack, 'table_row_even', (248, 250, 252))
                self._set_cell_text(
                    cell,
                    value,
                    font_size=body_font,
                    bold=False,
                    color=self._rgb_tuple(pack, 'text', (31, 41, 55)),
                    min_size=body_min,
                    font_family=self._font_family(pack),
                    avail_width=col_widths[col_idx],
                    avail_height=body_height,
                )

        if not use_side_note:
            return

        note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(note_left), Inches(content_top), Inches(note_width), Inches(area_height))
        note.fill.solid()
        note.fill.fore_color.rgb = self._rgb(pack, 'surface_alt', (239, 246, 255))
        note.line.color.rgb = self._rgb(pack, 'border', (191, 219, 254))

        note_body_box = slide.shapes.add_textbox(Inches(note_left + 0.12), Inches(content_top + 0.18), Inches(note_width - 0.24), Inches(area_height - 0.34))
        note_frame = note_body_box.text_frame
        note_frame.clear()
        note_frame.word_wrap = True
        side_facts = section.facts[:2] if rows >= 5 else section.facts[:3]
        for index, fact in enumerate(side_facts):
            p = note_frame.paragraphs[0] if index == 0 else note_frame.add_paragraph()
            p.text = f'• {self._normalize_text(fact)}'
            p.space_after = Pt(4)
            for run in p.runs:
                self._style_run(run, pack, size=9.8, color_key='text', default_color=(31, 41, 55))
        self._fit_frame(note_frame, max_size=11.2, min_size=9.0, font_family=self._font_family(pack))

    def _add_summary_slide(self, prs: Presentation, *, title: str, summary_points: list[str], presenter_name: str, page_number: int, total_slides: int, pack: dict) -> None:
        slide = self._base_slide(
            prs,
            title=title,
            presenter_name=presenter_name,
            page_number=page_number,
            total_slides=total_slides,
            pack=pack,
        )

        items = [self._normalize_text(item) for item in summary_points[:5] if self._normalize_text(item)]
        if not items:
            return

        if len(items) == 4:
            coords = [
                (0.90, 2.18, 5.40, 1.92),
                (7.03, 2.18, 5.40, 1.92),
                (0.90, 4.28, 5.40, 1.92),
                (7.03, 4.28, 5.40, 1.92),
            ]
        elif len(items) == 5:
            coords = [
                (0.82, 2.10, 3.88, 1.82),
                (4.73, 2.10, 3.88, 1.82),
                (8.64, 2.10, 3.88, 1.82),
                (2.78, 4.12, 3.88, 1.82),
                (6.69, 4.12, 3.88, 1.82),
            ]
        elif len(items) == 3:
            coords = [
                (0.92, 2.55, 3.72, 2.48),
                (4.80, 2.55, 3.72, 2.48),
                (8.68, 2.55, 3.72, 2.48),
            ]
        else:
            coords = [
                (1.18, 2.60, 5.20, 2.34),
                (6.95, 2.60, 5.20, 2.34),
            ]

        max_len = max((len(item) for item in items), default=0)
        font_size = 15.2 if max_len <= 75 else 14.0 if max_len <= 110 else 12.8
        summary_style = self._layout_value(pack, 'summary_style', 'cards')

        icon_color = self._rgb_tuple(pack, 'primary', (30, 64, 175))
        for idx, item in enumerate(items):
            x, y, w, h = coords[idx]
            shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb(pack, 'surface', (255, 255, 255))
            shape.line.color.rgb = self._rgb(pack, 'border', (191, 219, 254))

            stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.16))
            stripe.fill.solid()
            stripe.fill.fore_color.rgb = self._rgb(pack, 'primary' if summary_style == 'grid' else 'accent', (219, 234, 254))
            stripe.line.fill.background()

            self._draw_icon(slide, self._icon_for_text(item, fallback_index=idx), left=x + 0.24, top=y + 0.34, size=0.50, color=icon_color)

            text_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 1.02), Inches(w - 0.44), Inches(h - 1.18))
            frame = text_box.text_frame
            frame.clear()
            frame.word_wrap = True
            frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
            frame.margin_left = Inches(0.04)
            frame.margin_right = Inches(0.04)
            p = frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.08
            run = p.add_run()
            run.text = item
            self._style_run(run, pack, size=font_size, color_key='text', default_color=(30, 41, 59))
            self._fit_frame(frame, max_size=font_size, min_size=11.0, font_family=self._font_family(pack), avail_width=w - 0.44, avail_height=h - 1.18)

    def _add_divider_slide(
        self,
        prs: Presentation,
        *,
        chapter_number: int,
        total_chapters: int,
        section_titles: list[str],
        presenter_name: str,
        page_number: int,
        total_slides: int,
        pack: dict,
        accent_rgb: tuple[int, int, int],
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(*self._rgb_tuple(pack, 'cover_background', (239, 246, 255)))
        self._draw_background_motif(slide, pack, bold=True)
        heading_font = self._heading_font(pack)
        body_font = self._font_family(pack)
        primary = self._rgb(pack, 'primary', (30, 64, 175))
        bar_accent = self._visible_accent(pack, accent_rgb)

        # Oversized chapter number on the left.
        number_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.30), Inches(4.30), Inches(3.70))
        number_frame = number_box.text_frame
        number_frame.clear()
        number_frame.word_wrap = False
        number_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        number_p = number_frame.paragraphs[0]
        number_p.alignment = PP_ALIGN.LEFT
        number_run = number_p.add_run()
        number_run.text = f'{chapter_number:02d}'
        number_run.font.size = Pt(170)
        number_run.font.bold = True
        number_run.font.name = heading_font
        number_run.font.color.rgb = primary

        rule = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(5.32), Inches(1.66), Inches(0.05), Inches(3.30))
        rule.fill.solid()
        rule.fill.fore_color.rgb = RGBColor(*bar_accent)
        rule.line.fill.background()

        kicker_box = slide.shapes.add_textbox(Inches(5.66), Inches(1.74), Inches(6.84), Inches(0.50))
        kicker_frame = kicker_box.text_frame
        kicker_frame.clear()
        kicker_run = kicker_frame.paragraphs[0].add_run()
        kicker_run.text = f"{pack['chapter_word']} {chapter_number} / {total_chapters}"
        self._style_run(kicker_run, pack, size=15, color_key='primary', default_color=(30, 64, 175), bold=True, font_name=heading_font)
        self._fit_frame(kicker_frame, max_size=15, min_size=10, bold=True, font_family=heading_font)

        titles = [self._normalize_text(title) for title in section_titles if self._normalize_text(title)]
        title_size = 32.0 if len(titles) <= 1 else 25.0 if len(titles) == 2 else 21.0
        titles_box = slide.shapes.add_textbox(Inches(5.62), Inches(2.36), Inches(6.92), Inches(3.36))
        titles_frame = titles_box.text_frame
        titles_frame.clear()
        titles_frame.word_wrap = True
        titles_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        for index, title in enumerate(titles):
            paragraph = titles_frame.paragraphs[0] if index == 0 else titles_frame.add_paragraph()
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = Pt(10)
            paragraph.line_spacing = 1.05
            run = paragraph.add_run()
            run.text = title
            self._style_run(run, pack, size=title_size, color_key='text', default_color=(15, 23, 42), bold=True, font_name=heading_font)
        self._fit_frame(titles_frame, max_size=title_size, min_size=15, bold=True, font_family=heading_font, avail_width=6.92, avail_height=3.36)

        # Progress indicator: one segment per chapter, current one highlighted.
        seg_total_w = 6.84
        seg_gap = 0.12
        seg_w = (seg_total_w - seg_gap * (total_chapters - 1)) / max(1, total_chapters)
        for index in range(total_chapters):
            segment = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(5.66 + index * (seg_w + seg_gap)), Inches(6.06), Inches(seg_w), Inches(0.12),
            )
            segment.fill.solid()
            if index == chapter_number - 1:
                segment.fill.fore_color.rgb = primary
            else:
                segment.fill.fore_color.rgb = self._rgb(pack, 'border', (203, 213, 225))
            segment.line.fill.background()

        self._add_frame_footer(slide, presenter_name=presenter_name, page_number=page_number, total_slides=total_slides, pack=pack)

    def _add_frame_footer(self, slide, *, presenter_name: str, page_number: int, total_slides: int, pack: dict) -> None:
        footer_line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(6.84), Inches(11.86), Inches(0.02))
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = self._rgb(pack, 'border', (203, 213, 225))
        footer_line.line.fill.background()

        author_box = slide.shapes.add_textbox(Inches(0.76), Inches(6.90), Inches(5.8), Inches(0.26))
        author_frame = author_box.text_frame
        author_frame.clear()
        author_run = author_frame.paragraphs[0].add_run()
        author_run.text = f"{pack['prepared_by']}: {self._normalize_text(presenter_name, max_chars=40)}"
        self._style_run(author_run, pack, size=10, color_key='footer', default_color=(100, 116, 139))

        page_box = slide.shapes.add_textbox(Inches(11.45), Inches(6.90), Inches(1.05), Inches(0.26))
        page_frame = page_box.text_frame
        page_frame.clear()
        page_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        page_run = page_frame.paragraphs[0].add_run()
        page_run.text = f'{page_number}/{total_slides}'
        self._style_run(page_run, pack, size=10, color_key='footer', default_color=(100, 116, 139))

    def _add_closing_slide(
        self,
        prs: Presentation,
        *,
        presentation_title: str,
        presenter_name: str,
        page_number: int,
        total_slides: int,
        pack: dict,
    ) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(*self._rgb_tuple(pack, 'cover_background', (239, 246, 255)))
        self._draw_background_motif(slide, pack, bold=True)
        font_family = self._font_family(pack)

        panel_w, panel_h = 10.6, 3.18
        panel_x = (13.333 - panel_w) / 2
        panel_y = 2.18
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(panel_x), Inches(panel_y), Inches(panel_w), Inches(panel_h),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._rgb(pack, 'primary', (30, 64, 175))
        panel.line.fill.background()

        stripe = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(panel_x), Inches(panel_y), Inches(panel_w), Inches(0.16),
        )
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = self._rgb(pack, 'secondary', (96, 165, 250))
        stripe.line.fill.background()

        # Decorative accents flanking the banner.
        for offset, color_key, default in ((-0.52, 'secondary', (96, 165, 250)), (panel_w + 0.18, 'accent', (191, 219, 254))):
            dot = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(panel_x + offset), Inches(panel_y + panel_h / 2 - 0.17), Inches(0.34), Inches(0.34),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = self._rgb(pack, color_key, default)
            dot.line.fill.background()

        tag_box = slide.shapes.add_textbox(Inches(panel_x), Inches(panel_y + 0.40), Inches(panel_w), Inches(0.40))
        tag_frame = tag_box.text_frame
        tag_frame.clear()
        tag_frame.word_wrap = True
        tag_p = tag_frame.paragraphs[0]
        tag_p.alignment = PP_ALIGN.CENTER
        tag_run = tag_p.add_run()
        tag_run.text = str(pack['thanks_tag']).upper()
        self._style_run(tag_run, pack, size=13, color_key='accent', default_color=(191, 219, 254), bold=True)
        self._fit_frame(tag_frame, max_size=13, min_size=9, bold=True, font_family=font_family)

        title_box = slide.shapes.add_textbox(Inches(panel_x + 0.4), Inches(panel_y + 0.92), Inches(panel_w - 0.8), Inches(1.36))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        title_p = title_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_run = title_p.add_run()
        title_run.text = self._normalize_text(str(pack['thanks_title']))
        self._style_run(title_run, pack, size=40, color_key='on_primary', default_color=(255, 255, 255), bold=True)
        self._fit_frame(title_frame, max_size=40, min_size=24, bold=True, font_family=font_family)

        underline = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(panel_x + panel_w / 2 - 0.85), Inches(panel_y + panel_h - 0.52), Inches(1.70), Inches(0.06),
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = self._rgb(pack, 'secondary', (96, 165, 250))
        underline.line.fill.background()

        subtitle_box = slide.shapes.add_textbox(Inches(panel_x + 0.4), Inches(panel_y + panel_h + 0.26), Inches(panel_w - 0.8), Inches(0.80))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_frame.word_wrap = True
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.alignment = PP_ALIGN.CENTER
        subtitle_run = subtitle_p.add_run()
        subtitle_run.text = self._normalize_text(str(pack['thanks_subtitle']))
        self._style_run(subtitle_run, pack, size=15, color_key='muted', default_color=(71, 85, 105))
        self._fit_frame(subtitle_frame, max_size=15, min_size=11, font_family=font_family)

        # Footer (kept consistent with the rest of the deck).
        self._add_frame_footer(slide, presenter_name=presenter_name, page_number=page_number, total_slides=total_slides, pack=pack)
