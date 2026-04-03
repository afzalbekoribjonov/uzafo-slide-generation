from __future__ import annotations

import json
import logging
import re
import tempfile
from collections import deque
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Sequence

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.text.text import _Paragraph
from pptx.util import Pt

logger = logging.getLogger(__name__)

try:
    from google import genai
    from google.genai import types
except Exception:  # pragma: no cover
    genai = None
    types = None


STATIC_TEXT_EXACT = {
    'here is where your presentation begins',
    'table of contents',
    'contents of this template',
    'content of this template',
    'for more info, click here',
    'slidesgo',
    'template:',
    'faqs',
}
STATIC_TEXT_CONTAINS = (
    'here’s what you’ll find in this',
    'here抯 what you抣l find in this',
    'download and install the fonts',
    'an assortment of graphic resources',
    'an assortment of illustrations',
    'you can delete this slide',
    'credits for our design',
    'download the backgrounds you want from this template',
    'select the slide you want to save',
    'click file > download',
    'png image/jpeg image',
    'download the pc version of zoom',
    'click on the virtual fund option',
    'add an image or video will appear',
)
DESC_PLACEHOLDER_TEXTS = {
    'you can describe the topic of the section here',
    'here you could describe the topic of the section',
    'you could enter a subtitle here',
}
SOURCE_PLACEHOLDER_EXACT = {
    'awesome words',
    'key events',
    'critical analysis',
    'context',
    'problems',
    'relevant theories',
    'our framework',
    'theoretical framework',
    'key terms',
    'literature review',
    'schedule',
    'theory 1',
    'theory 2',
    'phase 01',
    'phase 02',
    'discussion 1',
    'discussion 2',
    'hypothesis 1',
    'hypothesis 2',
    'hypothesis 3',
    'method',
    'type of data',
    'motives',
    'data collection',
    'specific sampling',
    'response',
    'monitor',
    'develop',
    'current situation',
    'source',
    'transmission',
    'control',
    'outbreak',
    'prevalence',
    'incidence',
    'case control',
    'intervention',
    'social shifts',
    'influences',
    'cultural impact',
    'financial statements',
    'director’s report',
    "director's report",
    'management discussion',
    'capital structure',
    'earnings calls',
    'corporate report',
    'prospectus',
    'auditors report',
    'mockups',
    'english',
    'art',
    'french',
    'math',
    'markers',
    'your name',
    'thanks!',
    'for your attention',
}
SOURCE_PLACEHOLDER_CONTAINS = (
    'mercury is the closest planet to the sun',
    'mercury is the smallest planet',
    'mercury is very small',
    'mercury is quite a small planet',
    'mercury is small',
    'venus has a beautiful name',
    'venus has extremely high temperatures',
    'venus has a very beautiful name',
    "jupiter's rotation period",
    'jupiter is a gas giant',
    'jupiter is the biggest planet',
    'saturn is a gas giant',
    'saturn has rings',
    'neptune is the farthest planet from the sun',
    'neptune is an ice giant',
    'neptune is far away from earth',
    'neptune is far away from us',
    'neptune is very far',
    'despite being red, mars',
    'mars is full of iron oxide dust',
    'mars is made of basalt',
    'earth is the planet where we live',
    'earth is the third planet from the sun',
    'pluto is considered a dwarf planet',
    'pluto is now considered a dwarf planet',
    'ceres is located in the main asteroid belt',
    'the moon is earth’s natural satellite',
    "the moon is earth's natural satellite",
    'the sun’s mass compared to earth’s',
    "the sun's mass compared to earth's",
    'distance between earth and the moon',
    'follow the link in the graph',
    'for more info, click here',
    'you can replace the image on the screen with your own',
    'you can replace the images on the screen with your own',
    'author. (year). title of the publication',
    'surname, a. (year). name of the source',
    'venus is the second planet from the sun',
    'saturn was named after a roman god',
    'saturn is the ringed one and a gas giant',
    'earth is the planet with life',
    'venus is a very hot planet',
    'neptune is the fourth-largest planet by diameter in the solar system',
    'download the backgrounds you want from this template',
    'select the slide you want to save',
    'modify it and add text if needed',
    'click file > download > png image/jpeg image',
    'download the pc version of zoom',
    'click on the virtual fund option',
    'if you want to use a background downloaded to your computer',
    'you can describe the topic of the section here',
    'big numbers catch your audience',
    'images reveal large amounts of data',
    'an alphabetical type',
    'associated the letters to the phoneme',
    'they adapted to the phoenician alphabet',
    'with the paper appeared',
    'the invention of the printing press',
)
LISTISH_SLOT_NAMES = {
    'steps',
    'items',
    'materials',
    'entries',
    'ok_items',
    'not_ok_items',
    'rows',
    'columns',
    'table_headers',
    'data_points',
    'key_terms_items',
    'phase_1_points',
    'phase_2_points',
    'supporting_points',
    'summary_points',
    'bullet_points',
}
TRAILING_CONNECTOR_WORDS = {
    'va',
    'yoki',
    'bilan',
    'uchun',
    'orqali',
    'bo‘yicha',
    "bo'yicha",
    'ga',
    'da',
    'ni',
    'ning',
    'hamda',
}


@dataclass(slots=True)
class MagicTemplateSpec:
    template_id: str
    template_name: str
    category: str
    source_file: str
    template_path: Path
    json_path: Path
    template_profile: dict[str, Any]
    global_rules: dict[str, Any]
    slides: list[dict[str, Any]]
    gemini_prompt_package: dict[str, Any]

    @property
    def generated_slides(self) -> list[dict[str, Any]]:
        return [slide for slide in self.slides if slide.get('generate')]

    @property
    def uses_token_contract(self) -> bool:
        render_engine = str(self.template_profile.get('render_engine') or '').strip().lower()
        if render_engine == 'token_replace_v1':
            return True
        generated = self.generated_slides
        return bool(generated) and all(
            str(slide.get('render_strategy') or '').strip().lower() == 'token_replace'
            for slide in generated
        )


@dataclass(slots=True)
class ShapeInfo:
    shape: Any
    text: str
    top: int
    left: int
    paragraph_count: int


class MagicTemplateRegistry:
    def __init__(self, *, templates_root: str | Path = 'templates', json_dir: str | Path = 'templates/json') -> None:
        self.templates_root = Path(templates_root)
        self.json_dir = Path(json_dir)
        self._templates = self._load_templates()

    def _load_templates(self) -> dict[str, MagicTemplateSpec]:
        templates: dict[str, MagicTemplateSpec] = {}
        for json_path in sorted(self.json_dir.glob('*.json')):
            payload = json.loads(json_path.read_text(encoding='utf-8'))
            category = str(payload.get('category') or '').strip()
            source_file = str(payload.get('source_file') or '').strip()
            if not category or not source_file:
                continue

            template_path = self.templates_root / category / source_file
            spec = MagicTemplateSpec(
                template_id=str(payload['template_id']),
                template_name=str(payload.get('template_name') or payload['template_id']),
                category=category,
                source_file=source_file,
                template_path=template_path,
                json_path=json_path,
                template_profile=dict(payload.get('template_profile') or {}),
                global_rules=dict(payload.get('global_rules') or {}),
                slides=list(payload.get('slides') or []),
                gemini_prompt_package=dict(payload.get('gemini_prompt_package') or {}),
            )
            self._validate_template(spec)
            templates[spec.template_id] = spec
        return templates

    @staticmethod
    def _validate_template(template: MagicTemplateSpec) -> None:
        if not template.template_path.exists():
            raise FileNotFoundError(f"Template fayli topilmadi: {template.template_path}")
        if not template.uses_token_contract:
            raise ValueError(f"Template token kontraktga o'tmagan: {template.template_id}")
        if '-tokenized.' not in template.source_file:
            raise ValueError(f"Template source_file tokenized fayl bo'lishi kerak: {template.template_id}")
        recommended_slides = int(template.template_profile.get('recommended_output_slides') or 0)
        if recommended_slides < 13:
            raise ValueError(f"Premium template kamida 13 slayd bo'lishi kerak: {template.template_id}")
        if len(template.generated_slides) < 13:
            raise ValueError(f"Premium template generated slide'lari kamida 13 ta bo'lishi kerak: {template.template_id}")
        for slide in template.generated_slides:
            if not isinstance(slide.get('content_slots'), dict):
                raise ValueError(
                    f"Template slide content_slots object bo'lishi kerak: {template.template_id} slide={slide.get('id')}"
                )

    def get(self, template_id: str) -> MagicTemplateSpec:
        try:
            return self._templates[template_id]
        except KeyError as exc:
            raise ValueError(f'Template topilmadi: {template_id}') from exc

    def has(self, template_id: str) -> bool:
        return template_id in self._templates


class MagicContentGenerator:
    def __init__(self, *, api_key: str | None, model_name: str = 'gemini-2.5-flash-lite') -> None:
        self.api_key = (api_key or '').strip()
        self.model_name = (model_name or 'gemini-2.5-flash-lite').strip()

    @property
    def enabled(self) -> bool:
        return bool(self.api_key and genai is not None and types is not None)

    def build_content(self, template: MagicTemplateSpec, payload: dict[str, Any]) -> list[dict[str, Any]]:
        if self.enabled:
            try:
                return self._call_gemini(template, payload)
            except Exception as exc:  # pragma: no cover
                logger.warning('Magic slide Gemini generation failed for %s, falling back: %s', template.template_id, exc)
        return self._build_fallback_content(template, payload)

    def _call_gemini(self, template: MagicTemplateSpec, payload: dict[str, Any]) -> list[dict[str, Any]]:
        client = genai.Client(api_key=self.api_key)
        response = client.models.generate_content(
            model=self.model_name,
            contents=self._build_prompt(template, payload),
            config=types.GenerateContentConfig(
                temperature=0.35,
                top_p=0.9,
                response_mime_type='application/json',
            ),
        )
        raw_text = (getattr(response, 'text', None) or '').strip()
        if not raw_text:
            raise RuntimeError('Gemini premium generatsiyasi bo‘sh javob qaytardi.')
        data = json.loads(self._extract_json(raw_text))
        if not isinstance(data, dict):
            raise RuntimeError('Gemini premium generatsiyasi JSON object qaytarmadi.')
        return self._sanitize_result(template, payload, data)

    def _build_prompt(self, template: MagicTemplateSpec, payload: dict[str, Any]) -> str:
        package = template.gemini_prompt_package
        variables = dict(payload.get('variables') or {})
        extra_context = dict(payload.get('extra_context') or {})

        prompt_lines = [str(line).strip() for line in package.get('prompt_lines') or [] if str(line).strip()]
        if not prompt_lines:
            prompt_lines = [
                f"You are generating slide content for the template '{template.template_name}'.",
                'Fill only the requested slots and return JSON only.',
            ]

        slide_contract_lines = []
        constraint_lines = []
        for slide in template.generated_slides:
            slot_names = ', '.join((slide.get('content_slots') or {}).keys())
            slide_contract_lines.append(
                f"- Slide {slide['index']} | id={slide['id']} | role={slide.get('role')} | slots: {slot_names}"
            )
            slot_constraints = dict(slide.get('slot_constraints') or {})
            for slot_name, raw_constraints in slot_constraints.items():
                if not isinstance(raw_constraints, dict):
                    continue
                parts = [f"{name}={value}" for name, value in raw_constraints.items() if value not in (None, '', False)]
                if parts:
                    constraint_lines.append(
                        f"- Slide {slide['index']} | {slot_name}: " + ', '.join(parts)
                    )

        group_slot_lines = []
        for slide in template.generated_slides:
            for slot_name in (slide.get('content_slots') or {}).keys():
                if self._is_group_slot(slot_name):
                    group_slot_lines.append(f"- {slot_name}: array of strings is allowed when useful")

        variable_lines = [f"- {key}: {value}" for key, value in variables.items() if str(value).strip()]
        extra_lines = [f"- {key}: {value}" for key, value in extra_context.items() if str(value).strip()]

        expected_shape = {
            'template_id': template.template_id,
            'topic': variables.get('topic', ''),
            'language': variables.get('language', 'uz'),
            'slides_content': [
                {
                    'slide_index': slide['index'],
                    'slide_id': slide['id'],
                    'role': slide.get('role'),
                    'content': {
                        slot_name: '...'
                        for slot_name in (slide.get('content_slots') or {}).keys()
                    },
                }
                for slide in template.generated_slides[:2]
            ],
        }

        sections = [
            '\n'.join(prompt_lines),
            'User variables:\n' + ('\n'.join(variable_lines) if variable_lines else '- topic: not provided'),
            'Extra context:\n' + ('\n'.join(extra_lines) if extra_lines else '- none'),
            'Slide contract:\n' + '\n'.join(slide_contract_lines),
            (
                'Response rules:\n'
                '- Return valid JSON only.\n'
                '- Keep slides_content in the same order as the slide contract.\n'
                '- content must only include slots declared for that slide.\n'
                '- All generated text must be in the requested language only.\n'
                '- Keep headings concise and slide-friendly.\n'
                '- Keep body text presentation-ready, not essay-like.\n'
                '- Use the minimum text needed for compact side panels, cards and small text boxes.\n'
                '- Do not repeat the same title, explanation or label across multiple slides unless the contract explicitly requires it.\n'
                '- Section dividers must introduce different section focuses, not repeat the same heading.\n'
                '- Prefer short, concrete wording over abstract or decorative phrasing.'
            ),
        ]
        if group_slot_lines:
            sections.append('Slots that may be arrays:\n' + '\n'.join(group_slot_lines))
        if constraint_lines:
            sections.append('Slot limits:\n' + '\n'.join(constraint_lines))
        sections.append('Example JSON shape:\n' + json.dumps(expected_shape, ensure_ascii=False, indent=2))
        return '\n\n'.join(sections)

    def _sanitize_result(
        self,
        template: MagicTemplateSpec,
        payload: dict[str, Any],
        data: dict[str, Any],
    ) -> list[dict[str, Any]]:
        raw_items = data.get('slides_content')
        if not isinstance(raw_items, list):
            raw_items = []

        sanitized: list[dict[str, Any]] = []
        for position, slide in enumerate(template.generated_slides, start=1):
            item = self._find_matching_slide(raw_items, slide, position)
            source_content = item.get('content') if isinstance(item.get('content'), dict) else item
            if not isinstance(source_content, dict):
                source_content = {}

            content: dict[str, Any] = {}
            slot_constraints = dict(slide.get('slot_constraints') or {})
            for slot_name in (slide.get('content_slots') or {}).keys():
                if str(slide.get('role') or '') == 'subject_task_grid' and slot_name == 'items':
                    content[slot_name] = self._normalize_task_subject_pairs(
                        source_content.get(slot_name),
                        payload=payload,
                    )
                    continue
                content[slot_name] = self._normalize_slot_value(
                    slot_name=slot_name,
                    raw_value=source_content.get(slot_name),
                    payload=payload,
                    fallback_index=position,
                    constraints=slot_constraints.get(slot_name) if isinstance(slot_constraints.get(slot_name), dict) else None,
                )
            content = self._postprocess_slide_content(
                slide=slide,
                content=content,
                payload=payload,
                fallback_index=position,
            )
            sanitized.append(
                {
                    'slide_index': int(slide['index']),
                    'slide_id': str(slide['id']),
                    'role': str(slide.get('role') or ''),
                    'content': content,
                }
            )
        return sanitized

    def _postprocess_slide_content(
        self,
        *,
        slide: dict[str, Any],
        content: dict[str, Any],
        payload: dict[str, Any],
        fallback_index: int,
    ) -> dict[str, Any]:
        role = str(slide.get('role') or '')

        if role == 'step_by_step_instructions':
            steps = self._normalize_string_list(content.get('steps'))
            content['steps'] = steps[:6] if steps else self._fallback_group_items('steps', payload, fallback_index)[:4]
        elif role == 'process_guidance':
            steps = self._normalize_string_list(content.get('steps'))
            content['steps'] = steps[:5] if steps else self._fallback_group_items('steps', payload, fallback_index)[:4]
        elif role == 'platform_guide':
            content['part_1_steps'] = self._ensure_string_list_length(
                content.get('part_1_steps'),
                expected_count=3,
                payload=payload,
                fallback_index=fallback_index,
            )[:3]
            content['part_2_steps'] = self._ensure_string_list_length(
                content.get('part_2_steps'),
                expected_count=4,
                payload=payload,
                fallback_index=fallback_index,
            )[:4]
        elif role == 'dos_and_donts':
            content['ok_items'] = self._ensure_string_list_length(
                content.get('ok_items'),
                expected_count=3,
                payload=payload,
                fallback_index=fallback_index,
            )[:3]
            content['not_ok_items'] = self._ensure_string_list_length(
                content.get('not_ok_items'),
                expected_count=3,
                payload=payload,
                fallback_index=fallback_index,
            )[:3]
        elif role == 'review_checklist':
            content['items'] = self._ensure_string_list_length(
                content.get('items'),
                expected_count=3,
                payload=payload,
                fallback_index=fallback_index,
            )[:3]
        elif role == 'materials_checklist':
            materials = self._normalize_string_list(content.get('materials'))
            content['materials'] = (
                materials[:6] if materials else self._fallback_group_items('materials', payload, fallback_index)[:5]
            )
        elif role == 'subject_task_grid':
            task_pairs = content.get('items')
            if not isinstance(task_pairs, list) or not task_pairs:
                task_pairs = self._fallback_task_subject_pairs(payload)
            content['items'] = task_pairs[:4]
        elif role == 'section_divider':
            number = re.sub(r'\D+', '', str(content.get('section_number') or ''))
            if number:
                content['section_number'] = number.zfill(2)[:2]

        return content

    def _ensure_string_list_length(
        self,
        value: Any,
        *,
        expected_count: int,
        payload: dict[str, Any],
        fallback_index: int,
    ) -> list[str]:
        items = self._normalize_string_list(value)
        if len(items) >= expected_count:
            return items[:expected_count]

        fallback_items = self._fallback_group_items('items', payload, fallback_index)
        for item in fallback_items:
            if len(items) >= expected_count:
                break
            if item not in items:
                items.append(item)
        return items[:expected_count]

    def _normalize_task_subject_pairs(self, raw_value: Any, *, payload: dict[str, Any]) -> list[dict[str, str]]:
        pairs: list[dict[str, str]] = []
        if isinstance(raw_value, list):
            for item in raw_value:
                if isinstance(item, dict):
                    task = self._normalize_text(
                        item.get('task') or item.get('assignment') or item.get('item') or item.get('title')
                    )
                    subject = self._normalize_text(
                        item.get('subject') or item.get('focus') or item.get('category') or item.get('label')
                    )
                    if task or subject:
                        pairs.append({'task': task, 'subject': subject})
                        continue

                text = self._normalize_text(item)
                if not text:
                    continue
                task, subject = self._split_task_subject_text(text)
                pairs.append({'task': task, 'subject': subject})

        if pairs:
            return pairs[:4]
        return self._fallback_task_subject_pairs(payload)

    def _split_task_subject_text(self, value: str) -> tuple[str, str]:
        normalized = str(value or '').strip()
        if not normalized:
            return '', ''

        task_match = re.search(r'task\s*:\s*([^\n|]+)', normalized, flags=re.IGNORECASE)
        subject_match = re.search(r'subject\s*:\s*([^\n|]+)', normalized, flags=re.IGNORECASE)
        if task_match or subject_match:
            return (
                self._normalize_text(task_match.group(1) if task_match else ''),
                self._normalize_text(subject_match.group(1) if subject_match else ''),
            )

        if '|' in normalized:
            left, _, right = normalized.partition('|')
            return self._normalize_text(left), self._normalize_text(right)

        return self._normalize_text(normalized), ''

    def _fallback_task_subject_pairs(self, payload: dict[str, Any]) -> list[dict[str, str]]:
        language = str((payload.get('variables') or {}).get('language') or 'uz').strip()
        topic = str((payload.get('variables') or {}).get('topic') or 'Magic Slayd').strip()
        return [
            {
                'task': self._pick_phrase(
                    language,
                    uz=f'{topic} bo‘yicha qoida',
                    ru=f'Правило по теме {topic}',
                    en=f'Rule from {topic}',
                ),
                'subject': self._pick_phrase(language, uz='Asos', ru='Основа', en='Core'),
            },
            {
                'task': self._pick_phrase(
                    language,
                    uz=f'{topic} bo‘yicha misol',
                    ru=f'Пример по теме {topic}',
                    en=f'Example from {topic}',
                ),
                'subject': self._pick_phrase(language, uz='Misol', ru='Пример', en='Example'),
            },
            {
                'task': self._pick_phrase(
                    language,
                    uz=f'{topic} bo‘yicha taqqoslash',
                    ru=f'Сравнение по теме {topic}',
                    en=f'Comparison for {topic}',
                ),
                'subject': self._pick_phrase(language, uz='Qiyos', ru='Сравнение', en='Compare'),
            },
            {
                'task': self._pick_phrase(
                    language,
                    uz=f'{topic} bo‘yicha mashq',
                    ru=f'Практика по теме {topic}',
                    en=f'Practice for {topic}',
                ),
                'subject': self._pick_phrase(language, uz='Mashq', ru='Практика', en='Practice'),
            },
        ]

    def _find_matching_slide(self, raw_items: list[Any], slide: dict[str, Any], position: int) -> dict[str, Any]:
        for item in raw_items:
            if not isinstance(item, dict):
                continue
            if str(item.get('slide_id') or '') == str(slide['id']):
                return item
            if str(item.get('role') or '') == str(slide.get('role') or ''):
                return item
            if int(item.get('slide_index') or 0) == int(slide['index']):
                return item
        if 0 <= position - 1 < len(raw_items) and isinstance(raw_items[position - 1], dict):
            return raw_items[position - 1]
        return {}

    def _normalize_slot_value(
        self,
        *,
        slot_name: str,
        raw_value: Any,
        payload: dict[str, Any],
        fallback_index: int,
        constraints: dict[str, Any] | None = None,
    ) -> Any:
        if slot_name in {'rows'}:
            rows = self._normalize_rows(raw_value)
            if rows:
                return rows
            return self._fallback_rows(payload, fallback_index)

        if slot_name in {'columns', 'table_headers'}:
            values = self._normalize_string_list(raw_value)
            if values:
                return values
            return self._fallback_columns(slot_name)

        if self._is_group_slot(slot_name):
            values = self._normalize_string_list(raw_value)
            if values:
                return self._apply_list_constraints(values, constraints)
            return self._apply_list_constraints(self._fallback_group_items(slot_name, payload, fallback_index), constraints)

        text_value = self._normalize_text(raw_value)
        if text_value:
            return self._apply_text_constraints(text_value, constraints, slot_name=slot_name)
        return self._apply_text_constraints(
            self._fallback_text(slot_name, payload, fallback_index),
            constraints,
            slot_name=slot_name,
        )

    @staticmethod
    def _normalize_text(value: Any) -> str:
        if value is None:
            return ''
        if isinstance(value, (int, float)):
            return str(value)
        if isinstance(value, list):
            flattened = [str(item).strip() for item in value if str(item).strip()]
            return '\n'.join(flattened)
        if isinstance(value, dict):
            parts = [f'{key}: {item}' for key, item in value.items() if str(item).strip()]
            return '\n'.join(parts)
        return str(value).strip()

    @classmethod
    def _normalize_string_list(cls, value: Any) -> list[str]:
        if value is None:
            return []
        if isinstance(value, str):
            parts = [part.strip(' -•\t') for part in re.split(r'[\r\n]+', value) if part.strip(' -•\t')]
            return parts
        if isinstance(value, list):
            result: list[str] = []
            for item in value:
                text = cls._normalize_text(item)
                if text:
                    result.append(text)
            return result
        text = cls._normalize_text(value)
        return [text] if text else []

    def _normalize_rows(self, value: Any) -> list[list[str]]:
        rows: list[list[str]] = []
        if isinstance(value, list):
            for item in value:
                if isinstance(item, dict):
                    row = [self._normalize_text(cell) for cell in item.values() if self._normalize_text(cell)]
                elif isinstance(item, list):
                    row = [self._normalize_text(cell) for cell in item if self._normalize_text(cell)]
                else:
                    row = [part.strip() for part in re.split(r'\s*[|;]\s*', self._normalize_text(item)) if part.strip()]
                if row:
                    rows.append(row)
        elif isinstance(value, str):
            for line in value.splitlines():
                row = [part.strip() for part in re.split(r'\s*[|;]\s*', line) if part.strip()]
                if row:
                    rows.append(row)
        return rows

    def _apply_list_constraints(self, values: list[str], constraints: dict[str, Any] | None) -> list[str]:
        if not constraints:
            return values
        constrained: list[str] = []
        for value in values:
            constrained.append(self._apply_text_constraints(value, constraints))
        max_items = constraints.get('max_items')
        if isinstance(max_items, int) and max_items > 0:
            return constrained[:max_items]
        return constrained

    def _apply_text_constraints(
        self,
        value: str,
        constraints: dict[str, Any] | None,
        *,
        slot_name: str = '',
    ) -> str:
        text = self._clean_inline_text(value)
        if not text:
            return text

        text = self._tighten_text_for_slot(text, slot_name=slot_name)
        effective_constraints = dict(constraints or {})

        default_char_cap = self._default_slot_char_cap(slot_name)
        max_chars = effective_constraints.get('max_chars')
        if default_char_cap and (not isinstance(max_chars, int) or max_chars <= 0 or max_chars > default_char_cap):
            effective_constraints['max_chars'] = default_char_cap

        default_word_cap = self._default_slot_word_cap(slot_name)
        max_words = effective_constraints.get('max_words')
        if default_word_cap and (not isinstance(max_words, int) or max_words <= 0 or max_words > default_word_cap):
            effective_constraints['max_words'] = default_word_cap

        if not effective_constraints:
            return text

        max_words = effective_constraints.get('max_words')
        if isinstance(max_words, int) and max_words > 0:
            words = text.split()
            if len(words) > max_words:
                text = ' '.join(words[:max_words]).strip(' ,;:-')
                text = self._trim_trailing_connector(text)

        max_chars = effective_constraints.get('max_chars')
        if isinstance(max_chars, int) and max_chars > 0 and len(text) > max_chars:
            candidate = text[:max_chars].rstrip()
            sentence_break = max(candidate.rfind('.'), candidate.rfind('!'), candidate.rfind('?'))
            if sentence_break >= int(max_chars * 0.55):
                text = candidate[: sentence_break + 1].strip()
            else:
                shortened = candidate.rsplit(' ', 1)[0].strip(' ,;:-')
                text = shortened or candidate.strip()
            text = self._trim_trailing_connector(text)

        return text

    @staticmethod
    def _clean_inline_text(value: Any) -> str:
        text = str(value or '').replace('\u200b', ' ').replace('\xa0', ' ')
        text = re.sub(r'[\t\r]+', ' ', text)
        text = re.sub(r' +', ' ', text)
        text = re.sub(r' ?\n ?', '\n', text)
        return text.strip()

    def _tighten_text_for_slot(self, text: str, *, slot_name: str) -> str:
        slot = str(slot_name or '').strip().lower()
        if not text:
            return ''

        if slot.endswith(('_heading', '_label', '_name')):
            compact = re.split(r'\s*[:;]\s*|\s+[—–-]\s+', text, maxsplit=1)[0].strip()
            if compact:
                text = compact
        elif slot in {'intro_text', 'caption'} or slot.endswith(('_desc', '_subtitle')):
            if len(text) > 72 and ',' in text:
                compact = text.split(',', 1)[0].strip()
                if compact:
                    text = compact
        elif slot.endswith(('_text', '_body')):
            if len(text) > 70:
                compact = re.split(r'(?<=[.!?])\s+', text, maxsplit=1)[0].strip()
                if compact and len(compact) >= 28:
                    text = compact
            if len(text) > 62 and ',' in text:
                compact = text.split(',', 1)[0].strip()
                if compact and len(compact) >= 24:
                    text = compact
        elif slot.startswith('bullet_'):
            if len(text) > 58 and ',' in text:
                compact = text.split(',', 1)[0].strip()
                if compact:
                    text = compact

        return self._clean_inline_text(text)

    @staticmethod
    def _default_slot_char_cap(slot_name: str) -> int:
        slot = str(slot_name or '').strip().lower()
        if not slot:
            return 0
        if slot == 'title':
            return 72
        if slot.endswith('_title'):
            return 42
        if slot.endswith(('_heading', '_label', '_name')):
            return 28
        if slot.endswith(('_subtitle', '_desc', '_caption')):
            return 52
        if slot.startswith('bullet_'):
            return 42
        if slot == 'intro_text':
            return 72
        if slot.endswith(('_text', '_body')):
            return 60
        return 0

    @staticmethod
    def _default_slot_word_cap(slot_name: str) -> int:
        slot = str(slot_name or '').strip().lower()
        if not slot:
            return 0
        if slot == 'title':
            return 10
        if slot.endswith('_title'):
            return 6
        if slot.endswith(('_heading', '_label', '_name')):
            return 4
        if slot.endswith(('_subtitle', '_desc', '_caption')):
            return 8
        if slot.startswith('bullet_'):
            return 8
        if slot == 'intro_text':
            return 12
        if slot.endswith(('_text', '_body')):
            return 10
        return 0

    @staticmethod
    def _trim_trailing_connector(text: str) -> str:
        cleaned = str(text or '').strip(' ,;:-')
        if not cleaned:
            return ''
        parts = cleaned.split()
        while parts and parts[-1].strip(".,:;!?").lower() in TRAILING_CONNECTOR_WORDS:
            parts.pop()
        return ' '.join(parts).strip(' ,;:-') or cleaned

    def _fallback_text(self, slot_name: str, payload: dict[str, Any], fallback_index: int) -> str:
        variables = dict(payload.get('variables') or {})
        topic = str(variables.get('topic') or 'Premium presentation').strip()
        language = str(variables.get('language') or 'uz').strip()
        category = str(payload.get('category') or '').strip().lower()
        template_id = str(payload.get('template_id') or '').strip()
        base_context = self._build_context_tail(variables)
        ordinal = fallback_index
        slot_root = self._slot_root(slot_name)
        slot_index = self._slot_index(slot_name) or ordinal

        if slot_name == 'title':
            override = self._fallback_title_override(
                template_id=template_id,
                fallback_index=fallback_index,
                slot_name=slot_name,
                language=language,
                topic=topic,
                variables=variables,
            )
            return override or topic
        if slot_name == 'subtitle':
            if template_id == 'business_1_annual_report' and fallback_index == 1:
                company = str(variables.get('company_name') or '').strip() or topic
                period = str(variables.get('reporting_period') or '').strip() or self._pick_phrase(
                    language,
                    uz='joriy davr',
                    ru='tekushchiy period',
                    en='the current period',
                )
                return self._pick_phrase(
                    language,
                    uz=f"{company} uchun {period} yakunlari va strategik ustuvorliklar",
                    ru=f"Itogi {period} i strategicheskie prioritety dlya {company}",
                    en=f"{period} results and strategic priorities for {company}",
                )
            if template_id == 'science_2_epidemiology_thesis_defense' and fallback_index == 1:
                focus = str(variables.get('optional_focus_area') or '').strip() or base_context or topic
                return self._pick_phrase(
                    language,
                    uz=f"{focus} bo'yicha ilmiy himoya taqdimoti",
                    ru=f"Nauchnaya zashchita po teme {focus}",
                    en=f"Scientific defense built around {focus}",
                )
            if template_id == 'history_2_victorian_era_thesis_defense' and fallback_index == 1:
                period = str(variables.get('optional_focus_period') or '').strip() or base_context or topic
                return self._pick_phrase(
                    language,
                    uz=f"{period} kesimidagi tarixiy himoya taqdimoti",
                    ru=f"Istoricheskaya zashchita v razreze {period}",
                    en=f"History defense framed through {period}",
                )
            return self._pick_phrase(
                language,
                uz=f"{base_context or topic} bo‘yicha premium taqdimot",
                ru=f"Премиальная презентация по теме: {base_context or topic}",
                en=f"Premium presentation on {base_context or topic}",
            )
        if slot_name == 'author':
            return self._pick_phrase(
                language,
                uz='Ilmiy taqdimot',
                ru='Nauchnaya prezentatsiya',
                en='Research presentation',
            ) if category == 'science' else ''
        if slot_name.endswith('_title'):
            return self._fallback_section_title(
                slot_root=slot_root,
                language=language,
                ordinal=slot_index,
                topic=topic,
                category=category,
            )
        if slot_name.endswith('_subtitle') or slot_name.endswith('_desc'):
            return self._fallback_short_description(
                slot_root=slot_root,
                language=language,
                ordinal=slot_index,
                topic=topic,
                category=category,
            )
        if slot_name.startswith('question_'):
            return self._fallback_question(language=language, ordinal=slot_index, topic=topic)
        if slot_name.endswith('_heading') or slot_name.endswith('_label') or slot_name.endswith('_name'):
            return self._fallback_short_label(
                slot_root=slot_root,
                language=language,
                ordinal=slot_index,
                category=category,
            )
        if slot_name.endswith('_ref'):
            return self._pick_phrase(
                language,
                uz=f'Manba {slot_index}',
                ru=f'Источник {slot_index}',
                en=f'Source {slot_index}',
            )
        if slot_name.endswith('_year'):
            start_year = 2021 + max(0, ordinal - 1)
            return str(start_year)
        if slot_name.endswith('_percent') or slot_name == 'highlight_percent':
            return f'{18 + ((slot_index - 1) % 5) * 9}%'
        if slot_name.endswith('_value') or slot_name == 'main_value':
            return self._fallback_compact_value(slot_root=slot_root, language=language, ordinal=slot_index)
        if slot_name == 'section_number':
            return f'{ordinal:02d}'
        if slot_name == 'quote':
            return self._pick_phrase(
                language,
                uz=f'“{topic} bo‘yicha asosiy xulosa: tizimli yondashuv natijani kuchaytiradi.”',
                ru=f'“Ключевой вывод по теме {topic}: системный подход усиливает результат.”',
                en=f'"Key insight on {topic}: a structured approach strengthens the outcome."',
            )
        if slot_name in {'caption', 'chart_caption'}:
            return self._pick_phrase(
                language,
                uz=f'{topic} bo‘yicha asosiy ko‘rsatkich va kuzatuvlar jamlandi.',
                ru=f'Сводка ключевых показателей и наблюдений по теме {topic}.',
                en=f'A concise summary of the main indicators and observations for {topic}.',
            )
        if slot_name == 'headline':
            return self._fallback_section_title(
                slot_root='headline',
                language=language,
                ordinal=slot_index,
                topic=topic,
                category=category,
            )
        if slot_name in {'body', 'intro', 'intro_text', 'intro_line', 'closing_line', 'supporting_context', 'framework_text'}:
            return self._fallback_explanation(
                slot_root=slot_root or slot_name,
                language=language,
                topic=topic,
                category=category,
                context_tail=base_context,
            )
        if slot_name.endswith('_text') or slot_name.endswith('_note'):
            return self._fallback_explanation(
                slot_root=slot_root,
                language=language,
                topic=topic,
                category=category,
                context_tail=base_context,
            )
        return self._pick_phrase(
            language,
            uz=f'{base_context or topic} bo‘yicha asosiy fikr ixcham va tushunarli bayon qilinadi.',
            ru=f'Краткий и уместный для слайда текст по теме {topic}.',
            en=f'A concise slide-ready statement about {topic}.',
        )

    def _fallback_group_items(self, slot_name: str, payload: dict[str, Any], fallback_index: int) -> list[str]:
        language = str((payload.get('variables') or {}).get('language') or 'uz').strip()
        topic = str((payload.get('variables') or {}).get('topic') or 'Premium presentation').strip()
        category = str(payload.get('category') or '').strip().lower()
        template_id = str(payload.get('template_id') or '').strip()
        count = 4
        if 'materials' in slot_name:
            count = 5
        if 'five' in slot_name:
            count = 5
        if 'six' in slot_name:
            count = 6
        if 'entries' in slot_name:
            count = 5
        if 'phase' in slot_name:
            if template_id == 'science_2_epidemiology_thesis_defense':
                phase_sets = {
                    'uz': ['Signal topildi', 'Xatar ajratildi', 'Javob solishtirildi', 'Natija talqin qilindi'],
                    'ru': ['Signal opredelen', 'Risk vydelen', 'Otvet sravnen', 'Rezultat istolkovan'],
                    'en': ['Signal identified', 'Risk isolated', 'Response compared', 'Results interpreted'],
                }
                options = phase_sets.get(language) or phase_sets['uz']
                return options[:count]
            if template_id == 'history_2_victorian_era_thesis_defense':
                phase_sets = {
                    'uz': ['Davr bosimi ajratildi', 'Aktorlar bogʻlandi', 'Hududiy farq solishtirildi', 'Meros talqin qilindi'],
                    'ru': ['Davlenie perioda vydeleno', 'Aktery svyazany', 'Regionalnye razlichiya sravneny', 'Nasledie istolkovano'],
                    'en': ['Period pressures mapped', 'Actors connected', 'Regional differences compared', 'Legacy interpreted'],
                }
                options = phase_sets.get(language) or phase_sets['uz']
                return options[:count]
            return [
                self._pick_phrase(
                    language,
                    uz=f'{index + 1}. {topic} bo‘yicha muhim bosqich',
                    ru=f'{index + 1}. Аналитический вывод по теме {topic}',
                    en=f'{index + 1}. Analytical point about {topic}',
                )
                for index in range(count)
            ]
        if language == 'uz':
            category_points = {
                'business': [
                    'Bozor holati',
                    'Daromad drayveri',
                    'Mijoz ehtiyoji',
                    'Operatsion ustuvorlik',
                    'Xatar nuqtasi',
                    'Keyingi qadam',
                ],
                'education': [
                    'Asosiy qoida',
                    'Farqli jihat',
                    'Amaliy misol',
                    'Mustahkamlash mashqi',
                    'Ko‘p uchraydigan xato',
                    'Yakuniy eslatma',
                ],
                'history': [
                    'Davr manzarasi',
                    'Muhim hodisa',
                    'Asosiy omil',
                    'Tarixiy dalil',
                    'Ta’sir doirasi',
                    'Meros',
                ],
                'science': [
                    'Kuzatuv nuqtasi',
                    'Tajriba bosqichi',
                    'O‘lchov natijasi',
                    'Tahlil xulosasi',
                    'Amaliy tavsiya',
                    'Keyingi kuzatuv',
                ],
            }
            options = category_points.get(category)
            if options:
                return [options[index % len(options)] for index in range(count)]
        if 'key_terms' in slot_name:
            if category == 'science':
                terms = {
                    'uz': ['xatar zanjiri', 'ekspozitsiya', 'surveillance', 'klaster', 'profilaktika', 'time window'],
                    'ru': ['tsep riska', 'ekspozitsiya', 'surveillance', 'klaster', 'profilaktika', 'time window'],
                    'en': ['risk chain', 'exposure', 'surveillance', 'cluster', 'prevention', 'time window'],
                }
                options = terms.get(language) or terms['uz']
                return options[:count]
            return [
                self._pick_phrase(
                    language,
                    uz=f'Asosiy atama {index + 1}',
                    ru=f'Ключевой термин {index + 1}',
                    en=f'Key term {index + 1}',
                )
                for index in range(count)
            ]
        if 'supporting' in slot_name or 'summary' in slot_name or 'bullet' in slot_name:
            if category == 'history':
                options = {
                    'uz': ['davr bosimi kuchaydi', 'matbuot tilni o‘zgartirdi', 'institutlar tayanch bo‘ldi', 'meros keyingi davrga o‘tdi'],
                    'ru': ['davlenie perioda usililos', 'pressa izmenila yazyk', 'instituty stali oporoy', 'nasledie pereshlo v sleduyushchiy period'],
                    'en': ['period pressure intensified', 'the press reshaped language', 'institutions became anchors', 'legacy moved forward'],
                }
                values = options.get(language) or options['uz']
                return values[:count]
            if category == 'science':
                options = {
                    'uz': ['erta signal foyda berdi', 'zich aloqa xavfni oshirdi', 'profilaktika barqarorlik yaratdi', 'moslashgan javob natijani kuchaytirdi'],
                    'ru': ['ranniy signal dal effekt', 'plotnyy kontakt povysil risk', 'profilaktika ukrepila stabilnost', 'adaptivnyy otvet uluchshil rezultat'],
                    'en': ['early warning helped', 'dense contact raised risk', 'prevention improved stability', 'adaptive response strengthened the result'],
                }
                values = options.get(language) or options['uz']
                return values[:count]
            return [
                self._pick_phrase(
                    language,
                    uz=f'{index + 1}. Asosiy xulosa',
                    ru=f'{index + 1}. Ключевой вывод',
                    en=f'{index + 1}. Key takeaway',
                )
                for index in range(count)
            ]
        return [
            self._pick_phrase(
                language,
                uz=f'{index + 1}. {topic} bo‘yicha muhim nuqta',
                ru=f'{index + 1}. Важный пункт по теме {topic}',
                en=f'{index + 1}. Key point about {topic}',
            )
            for index in range(count)
        ]

    @staticmethod
    def _fallback_columns(slot_name: str) -> list[str]:
        if slot_name == 'columns':
            return ['Ko‘rsatkich', 'Qiymat', 'Izoh']
        if slot_name == 'table_headers':
            return ['Ko‘rsatkich', 'Tavsif', 'Natija']
        return ['Ustun 1', 'Ustun 2', 'Ustun 3']

    @staticmethod
    def _fallback_rows(payload: dict[str, Any], fallback_index: int) -> list[list[str]]:
        topic = str((payload.get('variables') or {}).get('topic') or 'Premium presentation').strip()
        category = str(payload.get('category') or '').strip().lower()
        template_id = str(payload.get('template_id') or '').strip()
        if template_id == 'history_2_victorian_era_thesis_defense':
            return [
                ['Manbalarni yig‘ish', 'birlamchi va ikkilamchi manbalarni saralash', '1-oy', 'yakunlandi'],
                ['Kontekst xaritasi', 'davr va hudud bo‘yicha tahlil karkasi', '2-oy', 'yakunlandi'],
                ['Farazlarni tekshirish', 'asosiy aktorlar va kuchlarni qiyoslash', '3-oy', 'davom etmoqda'],
                ['Talqin yozuvi', 'boblararo mantiq va dalillarni bog‘lash', '4-oy', 'rejalashtirilgan'],
                ['Himoya tayyorgarligi', 'yakuniy xulosa va ko‘rgazmali materiallar', '5-oy', 'navbatda'],
            ]
        if category == 'business':
            return [
                ['Bozor yo‘nalishi', topic, 'O‘sish sur’ati kuzatildi'],
                ['Mijoz segmenti', topic, 'Asosiy ehtiyoj aniqlandi'],
                ['Keyingi qadam', topic, 'Ustuvor reja belgilandi'],
            ]
        if category == 'education':
            return [
                ['Asosiy qoida', topic, 'Qisqa qiyosiy izoh'],
                ['Amaliy misol', topic, 'Qo‘llash usuli ko‘rsatildi'],
                ['Mustahkamlash', topic, 'Mashq orqali tekshirildi'],
            ]
        if category == 'history':
            return [
                ['Davr', topic, 'Asosiy tarixiy fon'],
                ['Jarayon', topic, 'Muhim o‘zgarish yo‘nalishi'],
                ['Meros', topic, 'Uzoq muddatli ta’sir'],
            ]
        if category == 'science':
            return [
                ['Ko‘rsatkich', topic, 'O‘lchov natijasi berildi'],
                ['Usul', topic, 'Kuzatuv tartibi qisqa yozildi'],
                ['Xulosa', topic, 'Amaliy talqin berildi'],
            ]
        return [
            ['Asosiy jihat', topic, 'Qisqa izoh'],
            ['Muhim nuqta', topic, 'Amaliy ma’no'],
            ['Yakuniy xulosa', topic, 'Qisqa tavsiya'],
        ]

    @staticmethod
    def _build_context_tail(variables: dict[str, Any]) -> str:
        priority = (
            'company_name',
            'reporting_period',
            'optional_focus_period',
            'optional_focus_area',
            'optional_region',
            'optional_subject',
            'optional_event_name',
            'optional_company_context',
        )
        for key in priority:
            value = str(variables.get(key) or '').strip()
            if value:
                return value
        return ''

    @staticmethod
    def _humanize_slot(slot_name: str, ordinal: int) -> str:
        cleaned = re.sub(r'_(\d+)_', ' ', slot_name.replace('_', ' ')).strip()
        cleaned = re.sub(r'\s+', ' ', cleaned)
        cleaned = cleaned.title() or f'Item {ordinal}'
        return cleaned

    @staticmethod
    def _slot_index(slot_name: str) -> int | None:
        match = re.search(r'_(\d+)(?:_|$)', slot_name)
        if not match:
            return None
        return int(match.group(1))

    @staticmethod
    def _slot_root(slot_name: str) -> str:
        cleaned = re.sub(r'_(\d+)(?:_|$)', '_', slot_name).strip('_')
        for suffix in (
            '_heading',
            '_label',
            '_name',
            '_text',
            '_note',
            '_desc',
            '_title',
            '_caption',
            '_subtitle',
            '_value',
            '_percent',
            '_year',
            '_points',
            '_items',
            '_ref',
        ):
            if cleaned.endswith(suffix):
                cleaned = cleaned[: -len(suffix)]
                break
        for prefix in ('left_', 'right_'):
            if cleaned.startswith(prefix):
                cleaned = cleaned[len(prefix) :]
                break
        return cleaned.strip('_')

    def _fallback_title_override(
        self,
        *,
        template_id: str,
        fallback_index: int,
        slot_name: str,
        language: str,
        topic: str,
        variables: dict[str, Any],
    ) -> str | None:
        if slot_name != 'title':
            return None

        company = str(variables.get('company_name') or '').strip() or topic
        business_titles = [
            company,
            '',
            self._pick_phrase(language, uz='Hisobot mazmuni', ru='Obzor perioda', en='Report overview'),
            self._pick_phrase(language, uz='Bozor va ichki fon', ru='Rynok i vnutrenniy fon', en='Market and internal context'),
            self._pick_phrase(language, uz='Yilning asosiy natijalari', ru='Klyuchevye itogi goda', en='Year in highlights'),
            self._pick_phrase(language, uz='Ikki asosiy drayver', ru='Dva klyuchevykh drayvera', en='Two primary drivers'),
            self._pick_phrase(language, uz='Strategik mavzular', ru='Strategicheskie temy', en='Strategic themes'),
            self._pick_phrase(language, uz="To'rtta tayanch", ru='Chetyre opory', en='Four pillars'),
            self._pick_phrase(language, uz='Natijani izohlovchi omillar', ru='Faktory rezultata', en='Factors behind the result'),
            '',
            self._pick_phrase(language, uz="Boshqaruv paneli ko'rinishi", ru='Vid paneli upravleniya', en='Dashboard view'),
            self._pick_phrase(language, uz="Yil davomida bajarilgan yo'l", ru='Khod raboty v techenie goda', en='Execution through the year'),
            self._pick_phrase(language, uz='Qiymat zanjiri qanday ishladi', ru='Kak srabotal tsepochka tsennosti', en='How the value chain worked'),
            self._pick_phrase(language, uz='Rahbariyat kuzatadigan hisobotlar', ru='Otchety dlya rukovodstva', en='Management reporting stack'),
            self._pick_phrase(language, uz='Tashabbuslar va tayanch nuqtalar', ru='Initsiativy i opornye vekhi', en='Initiatives and milestone points'),
            self._pick_phrase(language, uz='Daromad tarkibi', ru='Struktura vyruchki', en='Revenue mix'),
            self._pick_phrase(language, uz="Asosiy KPI ko'rsatkichlar", ru='Klyuchevye KPI', en='Key KPI highlights'),
            self._pick_phrase(language, uz='Operatsion marja talqini', ru='Interpretatsiya operatsionnoy marzhi', en='Operating margin interpretation'),
        ]
        science_titles = [
            topic,
            '',
            '',
            self._pick_phrase(language, uz='Kirish', ru='Vvedenie', en='Introduction'),
            self._pick_phrase(language, uz="Tadqiqotni yo'naltiruvchi savollar", ru='Napravlyayushchie voprosy', en='Guiding questions'),
            '',
            self._pick_phrase(language, uz='Hozirgi holat va muammolar', ru='Tekushchaya situatsiya i problemy', en='Current situation and problems'),
            '',
            self._pick_phrase(language, uz='Asosiy farazlar', ru='Osnovnye gipotezy', en='Core hypotheses'),
            self._pick_phrase(language, uz='Tadqiqot maqsadlari', ru='Tseli issledovaniya', en='Study objectives'),
            self._pick_phrase(language, uz='Tayanch epidemiologik tushunchalar', ru='Opornye epidemiologicheskie ponyatiya', en='Core epidemiological concepts'),
            self._pick_phrase(language, uz='Adabiyotlar sharhi', ru='Obzor literatury', en='Literature review'),
            self._pick_phrase(language, uz='Nazariy ramka', ru='Teoreticheskaya ramka', en='Theoretical framework'),
            self._pick_phrase(language, uz='Metodologiya', ru='Metodologiya', en='Methodology'),
            self._pick_phrase(language, uz="Tahlilning ikki bosqichi", ru='Dva etapa analiza', en='Two phases of analysis'),
            self._pick_phrase(language, uz='Kuzatuv natijalari dinamikasi', ru='Dinamika nablyudeniy', en='Observed trend'),
            self._pick_phrase(language, uz='Natijalar tahlili', ru='Analiz rezultatov', en='Results analysis'),
            self._pick_phrase(language, uz='Muhokama', ru='Obsuzhdenie', en='Discussion'),
            self._pick_phrase(language, uz='Tadqiqot jarayonining besh qadami', ru='Pyat shagov issledovaniya', en='Five study steps'),
            self._pick_phrase(language, uz='Asosiy xulosalar', ru='Osnovnye vyvody', en='Conclusions'),
            self._pick_phrase(language, uz="Amaliy qo'llash ko'rinishi", ru='Vid prakticheskogo primeneniya', en='Applied view'),
        ]
        history_titles = [
            topic,
            '',
            '',
            self._pick_phrase(language, uz="Tadqiqotning tayanch savollari", ru='Opornye voprosy issledovaniya', en='Core research questions'),
            self._pick_phrase(language, uz='Tarixiy manzara', ru='Istoricheskiy plan', en='Historical backdrop'),
            '',
            self._pick_phrase(language, uz='Asosiy farazlar', ru='Osnovnye gipotezy', en='Working hypotheses'),
            '',
            '',
            self._pick_phrase(language, uz='Tahlil bosqichi', ru='Etap analiza', en='Analytical phase'),
            self._pick_phrase(language, uz='Mavjud talqin va muammolar', ru='Sostoyanie interpretatsii i problemy', en='Current framing and tensions'),
            self._pick_phrase(language, uz='Tadqiqot vazifalari', ru='Zadachi issledovaniya', en='Study objectives'),
            self._pick_phrase(language, uz="Tadqiqotni tutib turgan olti g'oya", ru='Shest opornykh idey', en='Six analytical themes'),
            self._pick_phrase(language, uz='Adabiyotlar sharhi', ru='Obzor literatury', en='Literature review'),
            self._pick_phrase(language, uz='Nazariy ramka', ru='Teoreticheskaya ramka', en='Theoretical framework'),
            self._pick_phrase(language, uz='Tadqiqot jadvali', ru='Grafik issledovaniya', en='Research schedule'),
            self._pick_phrase(language, uz='Metodologiya', ru='Metodologiya', en='Methodology'),
            self._pick_phrase(language, uz="Tahlilning ikki bosqichi", ru='Dva etapa analiza', en='Two phases of analysis'),
            self._pick_phrase(language, uz='Tahliliy taqsimot', ru='Analiticheskoe raspredelenie', en='Analytical distribution'),
            self._pick_phrase(language, uz='Natijalar tahlili', ru='Analiz rezultatov', en='Results analysis'),
            self._pick_phrase(language, uz='Hududiy tayanch nuqtalar', ru='Regionalnye opornye tochki', en='Regional anchors'),
            self._pick_phrase(language, uz='Muhokama', ru='Obsuzhdenie', en='Discussion'),
            self._pick_phrase(language, uz='Asosiy xulosalar', ru='Osnovnye vyvody', en='Conclusions'),
            self._pick_phrase(language, uz='Qisqa bibliografiya', ru='Kratkaya bibliografiya', en='Selected bibliography'),
        ]
        overrides = {
            'business_1_annual_report': business_titles,
            'science_2_epidemiology_thesis_defense': science_titles,
            'history_2_victorian_era_thesis_defense': history_titles,
        }.get(template_id)
        if not overrides or fallback_index > len(overrides):
            return None
        value = overrides[fallback_index - 1]
        return value or None

    def _fallback_section_title(self, *, slot_root: str, language: str, ordinal: int, topic: str, category: str) -> str:
        if slot_root in {'section', 'part'}:
            category_titles = {
                'business': {
                    'uz': ['Kontekst', 'Natijalar', 'Ko‘rsatkichlar', 'Jarayon', 'Xatarlar', 'Xulosa'],
                    'ru': ['Контекст', 'Результаты', 'Показатели', 'Процессы', 'Риски', 'Вывод'],
                    'en': ['Context', 'Results', 'Metrics', 'Process', 'Risks', 'Conclusion'],
                },
                'education': {
                    'uz': ['Mavzu', 'Qoida', 'Qiyos', 'Misol', 'Mashq', 'Xulosa'],
                    'ru': ['Тема', 'Правило', 'Сравнение', 'Пример', 'Практика', 'Вывод'],
                    'en': ['Topic', 'Rule', 'Comparison', 'Example', 'Practice', 'Conclusion'],
                },
                'history': {
                    'uz': ['Davr', 'Manzara', 'Jarayon', 'Dalillar', 'Ta’sir', 'Xulosa'],
                    'ru': ['Период', 'Контекст', 'Процесс', 'Доводы', 'Влияние', 'Вывод'],
                    'en': ['Period', 'Context', 'Process', 'Evidence', 'Impact', 'Conclusion'],
                },
                'science': {
                    'uz': ['Muammo', 'Usul', 'Kuzatuv', 'Natija', 'Tahlil', 'Xulosa'],
                    'ru': ['Задача', 'Метод', 'Наблюдение', 'Результат', 'Анализ', 'Вывод'],
                    'en': ['Question', 'Method', 'Observation', 'Result', 'Analysis', 'Conclusion'],
                },
            }
            titles = {
                'uz': ['Kirish', 'Asosiy g‘oya', 'Tahlil', 'Natijalar', 'Muhokama', 'Xulosa'],
                'ru': ['Введение', 'Основная идея', 'Анализ', 'Результаты', 'Обсуждение', 'Вывод'],
                'en': ['Introduction', 'Core Idea', 'Analysis', 'Results', 'Discussion', 'Conclusion'],
            }
            options = (category_titles.get(category) or titles).get(language) or (category_titles.get(category) or titles)['uz']
            return options[(ordinal - 1) % len(options)]
        if slot_root == 'section_subtitle':
            return self._fallback_short_description(
                slot_root=slot_root,
                language=language,
                ordinal=ordinal,
                topic=topic,
                category=category,
            )
        return self._pick_phrase(
            language,
            uz=f'{topic} bo‘yicha bo‘lim',
            ru=f'Раздел по теме {topic}',
            en=f'Section on {topic}',
        )[:80]

    def _fallback_short_description(self, *, slot_root: str, language: str, ordinal: int, topic: str, category: str) -> str:
        if slot_root in {'section', 'part'}:
            category_descriptions = {
                'business': {
                    'uz': [
                        'Bozor va davr manzarasi',
                        'Asosiy natija va o‘sish',
                        'Raqamlar va ulushlar',
                        'Bajarilgan ishlar',
                        'Xatar va imkonlar',
                        'Yakuni va keyingi qadam',
                    ],
                },
                'education': {
                    'uz': [
                        'Mavzuning kirish qismi',
                        'Asosiy qoida va tushuncha',
                        'Farqli va o‘xshash jihat',
                        'Tushunarli amaliy misol',
                        'Mashq va mustahkamlash',
                        'Yakuniy xulosa',
                    ],
                },
                'history': {
                    'uz': [
                        'Davrning asosiy manzarasi',
                        'Jarayon va omillar',
                        'Muhim burilishlar',
                        'Dalil va izohlar',
                        'Ta’sir doirasi',
                        'Tarixiy xulosa',
                    ],
                },
                'science': {
                    'uz': [
                        'Masala va maqsad',
                        'Kuzatuv usuli',
                        'Topilgan natijalar',
                        'Tahlil va talqin',
                        'Amaliy ahamiyat',
                        'Yakuniy xulosa',
                    ],
                },
            }
            options = ((category_descriptions.get(category) or {}).get(language or 'uz')) or ((category_descriptions.get(category) or {}).get('uz'))
            if options:
                return options[(ordinal - 1) % len(options)]
        mapping = {
            'uz': {
                'section': 'Mavzuning asosiy yo‘nalishi',
                'part': 'Asosiy mazmun va talqin',
            },
            'ru': {
                'section': 'Краткое направление',
                'part': 'Основное содержание',
            },
            'en': {
                'section': 'Short direction',
                'part': 'Core content',
            },
        }
        value = (mapping.get(language) or mapping['uz']).get(slot_root)
        if value:
            return value
        return self._pick_phrase(
            language,
            uz=f'{topic} bo‘yicha qisqa sharh',
            ru=f'Краткое пояснение по теме {topic}',
            en=f'A short note on {topic}',
        )

    def _fallback_question(self, *, language: str, ordinal: int, topic: str) -> str:
        questions = {
            'uz': [
                f'{topic}ning asosiy sababi nimada?',
                f'{topic}ga qaysi omillar ta’sir qiladi?',
                f'{topic}dan qanday xulosa chiqarish mumkin?',
            ],
            'ru': [
                f'В чем основная причина темы {topic}?',
                f'Какие факторы влияют на тему {topic}?',
                f'Какой вывод можно сделать по теме {topic}?',
            ],
            'en': [
                f'What is the main driver behind {topic}?',
                f'Which factors shape {topic}?',
                f'What conclusion follows from {topic}?',
            ],
        }
        options = questions.get(language) or questions['uz']
        return options[(ordinal - 1) % len(options)]

    def _fallback_short_label(self, *, slot_root: str, language: str, ordinal: int, category: str) -> str:
        category_specific = {
            'business': {
                'card': ['Daromad sifati', 'Mijoz qadri', 'Bajarish ritmi', 'Marja drayveri'],
                'item': ['Segment tanlovi', 'Narx intizomi', 'Servis sifati', 'Hudud samarasi'],
                'block': ['Talab', 'Taklif', 'Yetkazish', 'Natija'],
                'segment': ['Takroriy paketlar', 'Kengayish savdosi', 'Yangi mijozlar'],
                'kpi': ['Takroriy tushum', 'Ushlab qolish', 'Time to value', 'Marjali miks', 'Prognoz aniqligi'],
                'period': ['Yil boshida', 'Yil oxirida', 'Keyingi davr'],
            },
            'science': {
                'problem': ['Kechikkan signal', 'Notekis nazorat', 'Sust javob'],
                'hypothesis': ['Aloqa zichligi', 'Erta kuzatuv', 'Qatlamli javob'],
                'objective': ['Omillarni ajratish', 'Signalni topish', 'Javobni baholash', 'Tavsiya berish'],
                'concept': ['Xatar', 'Ekspozitsiya', 'Surveillance', 'Klaster', 'Vaqt oynasi', 'Profilaktika'],
                'method': ['Yondashuv'],
                'data_type': ["Ma'lumot turi"],
                'motives': ['Tadqiqot sababi'],
                'collection': ["Yig'ish usuli"],
                'sampling': ['Tanlanma'],
                'result': ['Asosiy xatar', 'Himoya omili', 'Amaliy xulosa'],
                'discussion': ['Talqin', 'Cheklov', 'Qiyos', 'Tavsiya'],
                'step': ['Kuzatuv', 'Saralash', 'Qiyoslash', 'Sinov', 'Tavsiya'],
            },
            'history': {
                'hypothesis': ['Institutsional omil', "Ijtimoiy ko'prik", 'Uzoq oqibat'],
                'objective': ['Kontekstni tiklash', "Aktorlarni bog'lash", 'Merosni ochish'],
                'insight': ['Matbuot', 'Maktab', 'Ziyolilar', 'Hudud', 'Meros', 'Talqin'],
                'result': ['Kontekst', 'Aktorlar', 'Meros'],
                'location': ['Toshkent', 'Samarqand', 'Buxoro'],
                'discussion': ['Kuchli tomon', 'Bahsli nuqta', 'Qiyos', 'Meros'],
                'point': ['Kontekst', 'Aktorlar', 'Hudud', 'Meros'],
            },
        }
        if language == 'uz':
            options = (category_specific.get(category) or {}).get(slot_root)
            if options:
                return options[(ordinal - 1) % len(options)]
        premium_label_options = {
            'objective': ['Asosiy maqsad', 'Kutilgan natija', 'Amaliy vazifa', 'Tadqiqot savoli'],
            'step': ['Boshlang‘ich bosqich', 'Asosiy bosqich', 'Tahlil bosqichi', 'Yakuniy bosqich'],
            'phase': ['Boshlang‘ich bosqich', 'Asosiy bosqich', 'Tahlil bosqichi', 'Yakuniy bosqich'],
            'stage': ['Boshlang‘ich bosqich', 'Asosiy bosqich', 'Tahlil bosqichi', 'Yakuniy bosqich'],
            'method': ['Yondashuv', 'Tahlil usuli', 'Kuzatuv tartibi', 'Amaliy usul'],
            'data_type': ['Maʼlumot turi', 'Kuzatuv bazasi', 'Dalil turi', 'Oʻlchov turi'],
            'collection': ['Yigʻish usuli', 'Kuzatuv kanali', 'Qayd usuli', 'Manba yigʻimi'],
            'sampling': ['Tanlanma', 'Saralash mezoni', 'Tanlash usuli', 'Kuzatuv guruhi'],
            'motives': ['Tadqiqot sababi', 'Asosiy sabab', 'Tahlil motivi', 'Izlanish sababi'],
            'key_terms': ['Tayanch atamalar', 'Asosiy tushunchalar', 'Muhim atamalar', 'Terminlar bloki'],
            'item': ['Asosiy nuqta', 'Muhim jihat', 'Tahlil bloki', 'Amaliy xulosa'],
            'card': ['Asosiy nuqta', 'Muhim jihat', 'Tahlil bloki', 'Amaliy xulosa'],
            'block': ['Asosiy blok', 'Muhim blok', 'Tahlil bloki', 'Yakuniy blok'],
            'type': ['Asosiy yo‘nalish', 'Muhim yo‘nalish', 'Tahlil yo‘nalishi', 'Amaliy yo‘nalish'],
            'column': ['Asosiy yo‘nalish', 'Muhim ko‘rsatkich', 'Tahlil qismi', 'Yakuniy band'],
            'segment': ['Asosiy yo‘nalish', 'Muhim ko‘rsatkich', 'Tahlil qismi', 'Yakuniy band'],
            'bullet': ['Asosiy nuqta', 'Muhim jihat', 'Tahlil bloki', 'Amaliy xulosa'],
            'result': ['Asosiy natija', 'Muhim xulosa', 'Tahliliy kuzatuv', 'Amaliy xulosa'],
            'insight': ['Asosiy natija', 'Muhim xulosa', 'Tahliliy kuzatuv', 'Amaliy xulosa'],
            'event': ['Boshlanish', 'Muhim burilish', 'Kengayish', 'Natija', 'Davom'],
            'location': ['Asosiy hudud', 'Ta’sir hududi', 'Muhim nuqta'],
            'region': ['Asosiy hudud', 'Ta’sir hududi', 'Muhim nuqta'],
            'node': ['Asosiy bog‘lanish', 'Sabab aloqasi', 'Ta’sir zanjiri', 'Natija aloqasi'],
            'series': ['Asosiy ko‘rsatkich', 'Taqqoslash ko‘rsatkichi', 'Dinamik ko‘rsatkich'],
            'stat': ['Asosiy ko‘rsatkich', 'Taqqoslash ko‘rsatkichi', 'Dinamik ko‘rsatkich'],
            'kpi': ['Asosiy ko‘rsatkich', 'Taqqoslash ko‘rsatkichi', 'Dinamik ko‘rsatkich'],
            'point': ['Tayanch dalil', 'Muhim xulosa', 'Asosiy jihat', 'Qo‘shimcha izoh'],
            'review': ['Qisqa fikr', 'Muhim izoh', 'Asosiy mulohaza', 'Yakuniy qarash'],
            'person': ['Jamoa vakili', 'Mas’ul ishtirokchi', 'Tadqiqotchi', 'Kuzatuvchi'],
            'highlight': ['Asosiy urg‘u', 'Muhim nuqta', 'Tayanch dalil', 'Qisqa xulosa'],
            'task': ['Asosiy mashq', 'Qo‘llash vazifasi', 'Taqqoslash topshirig‘i', 'Mustahkamlash vazifasi'],
            'subject': ['Asosiy bo‘lim', 'Qiyos nuqtasi', 'Misol turi', 'Tahlil bloki'],
            'concept': ['Asosiy tushuncha', 'Muhim atama', 'Qo‘llanish ko‘rinishi', 'Tahlil jihati'],
            'factor': ['Asosiy omil', 'Ta’sir nuqtasi', 'Bog‘liq omil', 'Qo‘shimcha sabab'],
            'aspect': ['Asosiy jihat', 'Muhim farq', 'O‘xshash tomon', 'Izoh nuqtasi'],
            'detail': ['Asosiy jihat', 'Muhim farq', 'O‘xshash tomon', 'Izoh nuqtasi'],
            'part': ['Asosiy qism', 'Taqqoslash qismi', 'Amaliy qism', 'Yakuniy qism'],
            'topic': ['Asosiy mavzu', 'Muhim bo‘lim', 'Tahlil yo‘nalishi', 'Yakuniy jihat'],
            'discussion': ['Muhokama nuqtasi', 'Asosiy savol', 'Tahlil chizig‘i', 'Yakuniy mulohaza'],
        }
        if language == 'uz':
            options = premium_label_options.get(slot_root)
            if options:
                return options[(ordinal - 1) % len(options)]
        uz_map = {
            'objective': 'Maqsad',
            'question': 'Savol',
            'hypothesis': 'Faraz',
            'problem': 'Muammo',
            'step': 'Bosqich',
            'note': 'Eslatma',
            'task': 'Topshiriq',
            'subject': 'Bo‘lim',
            'factor': 'Omil',
            'aspect': 'Jihat',
            'concept': 'Tushuncha',
            'type': 'Yo‘nalish',
            'card': 'Yo‘nalish',
            'item': 'Band',
            'block': 'Blok',
            'column': 'Yo‘nalish',
            'topic': 'Mavzu',
            'event': 'Voqea',
            'node': 'Bog‘lanish',
            'detail': 'Jihat',
            'highlight': 'Muhim nuqta',
            'result': 'Natija',
            'insight': 'Xulosa',
            'phase': 'Bosqich',
            'stage': 'Bosqich',
            'theory': 'Nazariya',
            'person': 'Ishtirokchi',
            'location': 'Hudud',
            'region': 'Hudud',
            'series': 'Ko‘rsatkich',
            'segment': 'Yo‘nalish',
            'kpi': 'Ko‘rsatkich',
            'review': 'Fikr',
            'entry': 'Manba',
            'part': 'Qism',
            'section': 'Bo‘lim',
            'milestone': 'Bosqich',
            'period': 'Davr',
            'framework': 'Yondashuv',
            'method': 'Usul',
            'data_type': 'Maʼlumot',
            'collection': 'Yigʻish',
            'sampling': 'Tanlanma',
            'motives': 'Sabab',
            'key_terms': 'Atamalar',
            'key_term': 'Atama',
            'stat': 'Ko‘rsatkich',
            'point': 'Tayanch fikr',
            'discussion': 'Muhokama',
            'location_name': 'Hudud',
            'year': 'Yil',
            'month': 'Oy',
            'bullet': 'Yo‘nalish',
        }
        ru_map = {
            'objective': 'Цель',
            'question': 'Вопрос',
            'hypothesis': 'Гипотеза',
            'problem': 'Проблема',
            'step': 'Этап',
            'note': 'Примечание',
            'task': 'Задача',
            'subject': 'Раздел',
            'factor': 'Фактор',
            'aspect': 'Аспект',
            'concept': 'Понятие',
            'type': 'Направление',
            'card': 'Блок',
            'item': 'Пункт',
            'block': 'Блок',
            'column': 'Направление',
            'topic': 'Тема',
            'event': 'Событие',
            'node': 'Связь',
            'detail': 'Деталь',
            'highlight': 'Акцент',
            'result': 'Результат',
            'insight': 'Вывод',
            'phase': 'Этап',
            'stage': 'Этап',
            'theory': 'Теория',
            'person': 'Участник',
            'location': 'Регион',
            'region': 'Регион',
            'series': 'Показатель',
            'segment': 'Сегмент',
            'kpi': 'Показатель',
            'review': 'Мнение',
            'entry': 'Источник',
            'part': 'Часть',
            'section': 'Раздел',
            'milestone': 'Этап',
            'period': 'Период',
            'framework': 'Подход',
            'method': 'Метод',
            'data_type': 'Данные',
            'collection': 'Сбор',
            'sampling': 'Выборка',
            'motives': 'Мотив',
            'key_terms': 'Термины',
            'key_term': 'Термин',
            'stat': 'Показатель',
            'point': 'Тезис',
            'discussion': 'Обсуждение',
            'year': 'Год',
            'month': 'Месяц',
            'bullet': 'Направление',
        }
        en_map = {
            'objective': 'Objective',
            'question': 'Question',
            'hypothesis': 'Hypothesis',
            'problem': 'Problem',
            'step': 'Step',
            'note': 'Note',
            'task': 'Task',
            'subject': 'Section',
            'factor': 'Factor',
            'aspect': 'Aspect',
            'concept': 'Concept',
            'type': 'Theme',
            'card': 'Theme',
            'item': 'Point',
            'block': 'Block',
            'column': 'Stream',
            'topic': 'Theme',
            'event': 'Event',
            'node': 'Link',
            'detail': 'Detail',
            'highlight': 'Highlight',
            'result': 'Result',
            'insight': 'Insight',
            'phase': 'Phase',
            'stage': 'Stage',
            'theory': 'Theory',
            'person': 'Person',
            'location': 'Location',
            'region': 'Region',
            'series': 'Metric',
            'segment': 'Segment',
            'kpi': 'KPI',
            'review': 'View',
            'entry': 'Source',
            'part': 'Part',
            'section': 'Section',
            'milestone': 'Milestone',
            'period': 'Period',
            'framework': 'Framework',
            'method': 'Method',
            'data_type': 'Data',
            'collection': 'Collection',
            'sampling': 'Sampling',
            'motives': 'Rationale',
            'key_terms': 'Key Terms',
            'key_term': 'Key Term',
            'stat': 'Metric',
            'point': 'Key Point',
            'discussion': 'Discussion',
            'year': 'Year',
            'month': 'Month',
            'bullet': 'Focus',
        }
        mapping = {'uz': uz_map, 'ru': ru_map, 'en': en_map}
        base_label = (mapping.get(language) or mapping['uz']).get(slot_root, self._humanize_slot(slot_root, ordinal))
        return f'{base_label} {ordinal}'

    def _fallback_compact_value(self, *, slot_root: str, language: str, ordinal: int) -> str:
        if slot_root in {'stat', 'kpi', 'series', 'segment'}:
            return f'{12 + ordinal * 7}%'
        if slot_root in {'main', 'highlight'}:
            return f'20{20 + ordinal}'
        return self._pick_phrase(
            language,
            uz=f'{10 + ordinal}',
            ru=f'{10 + ordinal}',
            en=f'{10 + ordinal}',
        )

    def _fallback_explanation(self, *, slot_root: str, language: str, topic: str, category: str, context_tail: str) -> str:
        if language == 'uz':
            subject = context_tail or topic
            category_body_options = {
                'business': [
                    f'{subject} bo‘yicha daromad sifati, marja va bajarish ritmi ixcham yoritiladi.',
                    f'{subject} bo‘yicha bozor signali, mijoz talabi va keyingi qaror bir jumlada bog‘lanadi.',
                    f'{subject} bo‘yicha ustuvor yo‘nalish, xatar va amaliy natija qisqa tushuntiriladi.',
                ],
                'education': [
                    f'{topic} bo‘yicha qoida, qiyos yoki misol ixcham tushuntiriladi.',
                    f'{topic} bo‘yicha tushuncha va qo‘llanish sodda ifodalanadi.',
                    f'{topic} bo‘yicha farq va o‘xshashliklar qisqa yoritiladi.',
                ],
                'history': [
                    f'{topic} bo‘yicha davr bosimi, asosiy aktor va tarixiy ahamiyat qisqacha ochiladi.',
                    f'{topic} bo‘yicha jarayonning ichki mantiqi va uning oqibati ixcham yoritiladi.',
                    f'{topic} bo‘yicha tarixiy kontekst, hududiy farq va talqin qisqa bayon qilinadi.',
                ],
                'science': [
                    f'{subject} bo‘yicha usul, kuzatuv va epidemiologik xulosa qisqa bayon qilinadi.',
                    f'{subject} bo‘yicha o‘lchov natijasi, xatar signali va talqin ixcham beriladi.',
                    f'{subject} bo‘yicha kuzatilgan holat, javob chorasi va amaliy ahamiyat yoritiladi.',
                ],
            }
            category_objective = {
                'business': 'Asosiy biznes vazifa va kutilgan natija qisqa ko‘rsatiladi.',
                'education': 'Mavzu bo‘yicha o‘quv vazifa va kutilgan natija qisqa beriladi.',
                'history': 'Tadqiqot maqsadi va tarixiy savol qisqacha bayon qilinadi.',
                'science': 'Tajriba maqsadi va tekshiriladigan faraz qisqa ko‘rsatiladi.',
            }
            category_event = {
                'business': 'Ushbu bosqichning ta’siri va amaliy ahamiyati qisqa izohlanadi.',
                'education': 'Ushbu qadamning maqsadi va foydasi bir jumlada beriladi.',
                'history': 'Bu voqeaning mazmuni va tarixiy ahamiyati qisqa yoritiladi.',
                'science': 'Bu bosqichning natijaga ta’siri qisqacha tushuntiriladi.',
            }
            category_result = {
                'business': 'Natija va uning biznes talqini qisqa shaklda beriladi.',
                'education': 'Natija va undan olinadigan xulosa qisqa ko‘rsatiladi.',
                'history': 'Natija va uning tarixiy talqini ixcham yoritiladi.',
                'science': 'Natija va ilmiy talqin qisqa ifodalanadi.',
            }
        else:
            category_body_options = {}
            category_objective = {}
            category_event = {}
            category_result = {}
        body_options = category_body_options.get(category) or []
        if not body_options and language == 'uz':
            body_options = [
                f'{context_tail or topic} bo‘yicha asosiy mazmun ixcham yoritiladi.',
                f'{context_tail or topic} bo‘yicha muhim izoh qisqa beriladi.',
                f'{context_tail or topic} bo‘yicha tushunarli tahliliy sharh taqdim etiladi.',
            ]
        body_value = body_options[sum(ord(char) for char in slot_root) % len(body_options)] if body_options else f'{topic}ning asosiy mazmuni slaydga mos ravishda qisqa va ravon ifodalanadi.'
        explanations = {
            'objective': {
                'uz': category_objective.get(category) or 'Mavzu bo‘yicha aniq vazifa va kutilgan natija qisqacha bayon qilinadi.',
                'ru': 'Кратко раскрываются задача и ожидаемый результат по теме.',
                'en': 'The task and expected outcome are stated briefly.',
            },
            'supporting_context': {
                'uz': 'Asosiy fikrni asoslaydigan qisqa izoh va mantiq beriladi.',
                'ru': 'Дается краткое пояснение, поддерживающее основную мысль.',
                'en': 'A short supporting explanation is provided.',
            },
            'problem': {
                'uz': 'Dolzarb muammo va uning amaliy oqibati qisqacha tushuntiriladi.',
                'ru': 'Кратко объясняется проблема и ее практическое последствие.',
                'en': 'The problem and its practical consequence are explained briefly.',
            },
            'hypothesis': {
                'uz': 'Tekshiriladigan faraz va uning mantiqiy asosi qisqa bayon qilinadi.',
                'ru': 'Кратко формулируется гипотеза и ее логическое основание.',
                'en': 'The hypothesis and its rationale are stated briefly.',
            },
            'factor': {
                'uz': 'Ta’sir mexanizmi va kuzatilgan bog‘liqlik qisqacha yoritiladi.',
                'ru': 'Кратко описываются механизм влияния и наблюдаемая связь.',
                'en': 'The influence and observed connection are described briefly.',
            },
            'aspect': {
                'uz': 'Asosiy jihatning mazmuni bir-ikki qisqa ibora bilan ochiladi.',
                'ru': 'Смысл аспекта раскрывается в одной-двух коротких фразах.',
                'en': 'The aspect is explained in one or two short phrases.',
            },
            'event': {
                'uz': category_event.get(category) or 'Shu bosqichning mazmuni va ahamiyati bir jumlada beriladi.',
                'ru': 'Смысл и значение этого этапа раскрываются одной фразой.',
                'en': 'The meaning and significance of this step are stated in one sentence.',
            },
            'result': {
                'uz': category_result.get(category) or 'Natijaning mazmuni va talqini qisqa shaklda beriladi.',
                'ru': 'Содержание и интерпретация результата подаются кратко.',
                'en': 'The result and its interpretation are stated concisely.',
            },
            'discussion': {
                'uz': 'Muhokama uchun muhim tahliliy izoh qisqa bayon qilinadi.',
                'ru': 'Кратко формулируется аналитический комментарий для обсуждения.',
                'en': 'A concise analytical discussion point is provided.',
            },
            'framework': {
                'uz': 'Yondashuvning asosiy mantiqi va qo‘llanish doirasi tushuntiriladi.',
                'ru': 'Объясняются логика подхода и сфера его применения.',
                'en': 'The logic of the framework and its use are explained.',
            },
            'body': {
                'uz': body_value,
                'ru': f'Основное содержание темы {topic} передается кратко и ясно.',
                'en': f'The main idea of {topic} is expressed clearly and briefly.',
            },
        }
        language_map = explanations.get(slot_root) or explanations.get('body') or {}
        return language_map.get(language) or language_map.get('uz') or self._pick_phrase(
            language,
            uz=f'{context_tail or topic} bo‘yicha qisqa va aniq sharh beriladi.',
            ru=f'Даётся краткое и ясное пояснение по теме {topic}.',
            en=f'A short, clear explanation is provided for {topic}.',
        )

    @staticmethod
    def _pick_phrase(language: str, *, uz: str, ru: str, en: str) -> str:
        if language == 'ru':
            return ru
        if language == 'en':
            return en
        return uz

    @staticmethod
    def _extract_json(raw_text: str) -> str:
        fenced = re.search(r'```(?:json)?\s*(\{.*\})\s*```', raw_text, flags=re.DOTALL)
        if fenced:
            return fenced.group(1)
        start = raw_text.find('{')
        end = raw_text.rfind('}')
        if start == -1 or end == -1 or end <= start:
            raise RuntimeError('Gemini javobidan JSON topilmadi.')
        return raw_text[start:end + 1]

    @staticmethod
    def _is_group_slot(slot_name: str) -> bool:
        if slot_name in LISTISH_SLOT_NAMES:
            return True
        return slot_name.endswith('_items') or slot_name.endswith('_points')

    def _build_fallback_content(self, template: MagicTemplateSpec, payload: dict[str, Any]) -> list[dict[str, Any]]:
        if template.template_id == 'business_2_elegant_workplan':
            return self._build_business_workplan_fallback(template, payload)
        if template.template_id == 'business_1_annual_report':
            return self._build_business_annual_report_fallback(template, payload)
        if template.template_id == 'education_1_elegant_education_pack':
            return self._build_education_elegant_pack_fallback(template, payload)
        if template.template_id == 'education_2_analysis_conclusions_high_school':
            return self._build_education_analysis_fallback(template, payload)
        if template.template_id == 'history_1_writing_history_thesis':
            return self._build_history_writing_thesis_fallback(template, payload)
        if template.template_id == 'history_2_victorian_era_thesis_defense':
            return self._build_history_thesis_defense_fallback(template, payload)
        if template.template_id == 'science_1_science_fair_newsletter':
            return self._build_science_newsletter_fallback(template, payload)
        if template.template_id == 'science_2_epidemiology_thesis_defense':
            return self._build_science_epidemiology_fallback(template, payload)
        return self._sanitize_result(template, payload, {'slides_content': []})

    @staticmethod
    def _slide_item(slide_id: str, content: dict[str, Any]) -> dict[str, Any]:
        return {'slide_id': slide_id, 'content': content}

    def _build_education_elegant_pack_fallback(
        self,
        template: MagicTemplateSpec,
        payload: dict[str, Any],
    ) -> list[dict[str, Any]]:
        variables = dict(payload.get('variables') or {})
        language = str(variables.get('language') or 'uz').strip().lower()
        if language != 'uz':
            return self._sanitize_result(template, payload, {'slides_content': []})

        topic = str(variables.get('topic') or template.template_name).strip()
        subject = str(variables.get('optional_subject') or '').strip() or 'asosiy mavzu'
        mode = str(variables.get('optional_delivery_mode') or '').strip() or "dars jarayoni"

        words = topic.split()
        title_line_1 = ' '.join(words[:4]) or topic
        title_line_2 = ' '.join(words[4:8]) or 'o‘quv taqdimoti'

        raw = {
            'slides_content': [
                self._slide_item('cover', {
                    'title_line_1': title_line_1,
                    'title_line_2': title_line_2,
                    'subtitle': f"{subject} bo'yicha amaliy va tushunarli o'quv taqdimoti",
                }),
                self._slide_item('table_of_contents_primary', {
                    'toc_title': 'Mundarija',
                    'section_1_title': 'Kirish',
                    'section_1_desc': 'mavzu va asosiy ma’no',
                    'section_2_title': 'Tushunish',
                    'section_2_desc': 'bosqichma-bosqich izoh',
                    'section_3_title': 'Qiyos',
                    'section_3_desc': 'to‘g‘ri va xato holat',
                    'section_4_title': 'Mashq',
                    'section_4_desc': 'vazifa va eslatmalar',
                    'section_5_title': 'Reja',
                    'section_5_desc': 'ketma-ket bajarish',
                    'section_6_title': 'Xulosa',
                    'section_6_desc': 'asosiy tayanch fikrlar',
                }),
                self._slide_item('topic_intro', {
                    'title': 'Mavzu kirishi',
                    'body': f"{topic} {mode} davomida bosqichma-bosqich tushuntirilsa, o'quvchi asosiy qoidani, qiyosiy farqni va amaliy qo'llash usulini tezroq anglaydi.",
                }),
                self._slide_item('guided_instructions', {
                    'title': 'Mavzuni ochish',
                    'intro_note': f"{topic}ni tushunishda asosiy e'tiborni mazmun, misol va qo'llashga qarating.",
                    'step_1': 'Mavzuning asosiy maqsadini aniqlang.',
                    'step_2': 'Tayanch qoidani sodda ifoda bilan ayting.',
                    'step_3': 'Muhim atamalarni alohida ajrating.',
                    'step_4': 'Bir oddiy misol bilan izohlang.',
                    'step_5': 'Qiyos yoki farqni ko‘rsating.',
                    'step_6': 'Tez-tez uchraydigan xatoni belgilang.',
                    'step_7': 'Qisqa xulosa bilan yakunlang.',
                }),
                self._slide_item('during_process_guidance', {
                    'title': 'Amaliy ishlash',
                    'step_1': "Matnni yoki misolni diqqat bilan o'qing.",
                    'step_2': 'Qaysi qoida ishlayotganini toping.',
                    'step_3': "Asosiy farqni qisqa yozib oling.",
                    'step_4': "O'xshash misol bilan tekshirib ko'ring.",
                    'step_5': 'Xato variantni nega noto‘g‘ri ekanini ayting.',
                    'step_6': "To'g'ri variantni mustaqil tuzing.",
                    'tip': 'Har bir javobni qoidaga tayangan holda izohlang.',
                    'warning': "Faqat yodlash bilan cheklanib qolmang, qo'llashni ham mashq qiling.",
                }),
                self._slide_item('platform_or_tool_instructions', {
                    'title': 'Ikki yo‘l bilan o‘rganish',
                    'part_1_heading': 'Nazariy qism',
                    'part_1_step_1': 'Asosiy qoida va atamalarni ajrating.',
                    'part_1_step_2': "Mavzu ichidagi muhim farqlarni yozing.",
                    'part_1_step_3': "Qisqa ta'riflarni takrorlab chiqing.",
                    'part_1_step_4': 'Bir namunani to‘liq tahlil qiling.',
                    'part_2_heading': 'Amaliy qism',
                    'part_2_step_1': "Misolni mustaqil ishlab ko'ring.",
                    'part_2_step_2': "To'g'ri va xato holatni solishtiring.",
                    'part_2_step_3': 'Yakuniy javobni qisqa izoh bilan mustahkamlang.',
                }),
                self._slide_item('school_rules', {
                    'title': 'To‘g‘ri va xato holat',
                    'ok_label': "To'g'ri",
                    'not_ok_label': 'Xato',
                    'ok_1': 'Qoidani misol bilan bog‘lash',
                    'ok_2': 'Farqni aniq izohlash',
                    'ok_3': 'Xulosani qisqa aytish',
                    'not_ok_1': 'Faqat yodlangan javob aytish',
                    'not_ok_2': 'Misolsiz qoida berish',
                    'not_ok_3': 'Farqni chalkashtirib yuborish',
                }),
                self._slide_item('reminder_tasks', {
                    'title': 'Vazifalarni eslang',
                    'tasks_heading': 'Vazifa',
                    'subjects_heading': 'Bo‘lim',
                    'task_1': 'Qoida yozing',
                    'subject_1': 'Asos',
                    'task_2': 'Misol toping',
                    'subject_2': 'Amal',
                    'task_3': 'Farqni yozing',
                    'subject_3': 'Qiyos',
                    'task_4': 'Xulosa ayting',
                    'subject_4': 'Yakun',
                }),
                self._slide_item('class_materials', {
                    'title': 'Kerakli tayanchlar',
                    'note_label': 'Eslatma',
                    'material_1': 'qisqa qoida',
                    'material_2': 'namuna matn',
                    'material_3': 'mashq daftari',
                    'material_4': 'taqqoslash jadvali',
                    'material_5': 'rangli belgi',
                    'material_6': 'yakun savoli',
                }),
                self._slide_item('review_pending_items', {
                    'title': 'Qayta ko‘rish',
                    'item_1': 'Asosiy qoidani yana bir bor ayting',
                    'item_2': "To'g'ri va xato holatni solishtiring",
                    'item_3': 'Bitta mustaqil misol tuzing',
                }),
                self._slide_item('task_organization', {
                    'title': 'Bajarish rejasi',
                    'step_1_heading': 'Tushun',
                    'step_1_text': 'Qoidani va tayanch atamani aniqlang.',
                    'step_2_heading': 'Qiyosla',
                    'step_2_text': "Misol va qarshi misolni solishtiring.",
                    'step_3_heading': 'Qo‘lla',
                    'step_3_text': 'Yangi vazifada qoidadan foydalaning.',
                }),
                self._slide_item('infographic_cards', {
                    'title': 'Tayanch fikrlar',
                    'card_1_heading': 'Qoida',
                    'card_1_text': 'Mavzuning asosiy tayanchi shu yerda turadi.',
                    'card_2_heading': 'Misol',
                    'card_2_text': "Nazariyani amaliy holat bilan bog'laydi.",
                    'card_3_heading': 'Farq',
                    'card_3_text': "O'xshash tushunchalarni ajratishga yordam beradi.",
                    'card_4_heading': 'Yakun',
                    'card_4_text': "O'rgangan fikrni mustahkamlaydi.",
                }),
                self._slide_item('infographic_comparison', {
                    'title': 'Qisqa yakun',
                    'item_1_heading': 'Bilim',
                    'item_1_text': 'Qoidani aniq tushunish kerak.',
                    'item_2_heading': 'Mashq',
                    'item_2_text': "Har bir fikr misol bilan sinab ko'riladi.",
                    'item_3_heading': 'Natija',
                    'item_3_text': "Mavzu mustaqil qo'llashga tayyor bo'ladi.",
                }),
            ]
        }
        return self._sanitize_result(template, payload, raw)

    def _build_education_analysis_fallback(
        self,
        template: MagicTemplateSpec,
        payload: dict[str, Any],
    ) -> list[dict[str, Any]]:
        variables = dict(payload.get('variables') or {})
        language = str(variables.get('language') or 'uz').strip().lower()
        if language != 'uz':
            return self._sanitize_result(template, payload, {'slides_content': []})

        topic = str(variables.get('topic') or template.template_name).strip()
        grade = str(variables.get('optional_grade_level') or '').strip() or 'yuqori sinf'
        subject = str(variables.get('optional_subject') or '').strip() or "o'quv mavzusi"
        data_context = str(variables.get('optional_dataset_context') or '').strip() or "sinf kuzatuvlari va qisqa so'rov natijalari"
        goal = str(variables.get('presentation_goal') or '').strip() or 'yakuniy tahlil'

        raw = {
            'slides_content': [
                self._slide_item('cover', {
                    'title': topic,
                    'author': f'{grade} taqdimoti',
                }),
                self._slide_item('table_of_contents', {
                    'toc_title': 'Mundarija',
                    'section_1_title': 'Kirish',
                    'section_1_desc': 'mavzu va dolzarblik',
                    'section_2_title': 'Maqsad',
                    'section_2_desc': "asosiy va aniq vazifalar",
                    'section_3_title': "Ma'lumot",
                    'section_3_desc': 'namuna va dalillar',
                    'section_4_title': 'Tahlil',
                    'section_4_desc': 'bogʻlanish va dinamika',
                    'section_5_title': 'Natija',
                    'section_5_desc': 'topilmalar va talqin',
                    'section_6_title': 'Xulosa',
                    'section_6_desc': 'amaliy tavsiyalar',
                }),
                self._slide_item('introduction', {
                    'title': 'Kirish',
                    'body': f"{topic} mavzusi {subject} doirasida {grade} o'quvchilari uchun qanday foyda, imkoniyat va cheklovlar borligini ko'rsatadi. Taqdimot {data_context} asosida {goal} yo'nalishida tuzilgan.",
                }),
                self._slide_item('general_objective', {
                    'title': 'Umumiy maqsad',
                    'objective_text': f"{topic} bo'yicha asosiy kuzatuvlarni tartiblab, o'quv jarayoniga ta'sir qilayotgan omillarni aniq ko'rsatish.",
                    'supporting_context': f"Tahlil {grade} darajasidagi ehtiyojlar, {subject} bilan bog'liq tajriba va {data_context}dan olingan dalillarni birlashtiradi.",
                }),
                self._slide_item('specific_objectives', {
                    'title': 'Aniq vazifalar',
                    'intro_line': "Tahlil quyidagi vazifalar asosida olib borildi.",
                    'bullet_1': "Asosiy ehtiyojlarni ajratish",
                    'bullet_2': "Foydali amaliyotlarni taqqoslash",
                    'bullet_3': "Cheklovlarni ochiq ko'rsatish",
                    'bullet_4': "Yakuniy tavsiyalarni jamlash",
                }),
                self._slide_item('objective_limitations', {
                    'title': 'Qamrov va cheklov',
                    'card_1_heading': 'Qamrov',
                    'card_1_text': f"Tahlil {grade} darajasidagi o'quvchilar, dars jarayoni va {subject}ga oid amaliy misollarni qamrab oladi.",
                    'card_2_heading': 'Cheklov',
                    'card_2_text': "Natijalar kichik namuna va qisqa kuzatuv davri bilan cheklanadi, shu bois talqin ehtiyotkor qo'llanadi.",
                }),
                self._slide_item('objectives_relevance', {
                    'title': 'Mavzuning ahamiyati',
                    'type_1_heading': "Ta'limiy",
                    'type_1_text': "Mavzu o'quvchi tushunishini chuqurlashtiradi.",
                    'type_2_heading': 'Amaliy',
                    'type_2_text': "Darsdagi qo'llash usullarini aniq ko'rsatadi.",
                    'type_3_heading': 'Tahliliy',
                    'type_3_text': "Qaror qabul qilish uchun dalil beradi.",
                }),
                self._slide_item('study_population', {
                    'title': "Tadqiqot guruhi va namuna",
                    'factor_1_heading': 'Guruh',
                    'factor_1_text': f"{grade} o'quvchilari asosiy guruh sifatida olindi.",
                    'factor_2_heading': 'Tanlash',
                    'factor_2_text': "Faol qatnashgan va kuzatuvga mos sinflar tanlandi.",
                    'factor_3_heading': 'Hajm',
                    'factor_3_text': "Natijani ko'rish uchun ixcham, ammo yetarli namuna olindi.",
                    'factor_4_heading': 'Belgi',
                    'factor_4_text': "Faollik, tushunish va qo'llash ko'rsatkichlari kuzatildi.",
                }),
                self._slide_item('data_presentation', {
                    'title': "Dalillar qanday ko'rsatildi",
                    'aspect_1_heading': "Ko'rsatkich",
                    'aspect_1_text': "Asosiy o'lchovlar oldindan tanlandi.",
                    'aspect_2_heading': "So'rov",
                    'aspect_2_text': "Qisqa savollar bilan fikr yig'ildi.",
                    'aspect_3_heading': 'Topshiriq',
                    'aspect_3_text': "Amaliy ish orqali natija tekshirildi.",
                    'aspect_4_heading': 'Jadval',
                    'aspect_4_text': "Guruhlar orasidagi farq tartiblandi.",
                    'aspect_5_heading': 'Trend',
                    'aspect_5_text': "Davriy o'zgarishlar alohida kuzatildi.",
                    'aspect_6_heading': 'Talqin',
                    'aspect_6_text': "Har ko'rsatkich uchun qisqa xulosa berildi.",
                }),
                self._slide_item('quote_slide', {
                    'quote': "Tahlilning kuchi murakkab fikrni aniq va foydali qarorga aylantirishidadir.",
                    'author': "Tadqiqot xulosasi",
                }),
                self._slide_item('image_break_title', {
                    'title': "Mavzuning markazida o'quvchi turadi",
                }),
                self._slide_item('visual_explanation', {
                    'headline_line_1': 'Asosiy kuzatuv',
                    'headline_line_2': 'bir qarashda',
                    'body': f"{topic} bo'yicha eng muhim dalil shuki, to'g'ri yondashuv tanlansa, {grade} o'quvchilarida tushunish va ishtirok birga oshadi.",
                }),
                self._slide_item('computer_mockup', {
                    'title': "Kompyuter ko'rinishi",
                    'caption': "Kengroq ekran natijalar, jadval va taqqoslashlarni bitta ko'rinishda yig'ishga qulay bo'ldi.",
                }),
                self._slide_item('tablet_mockup', {
                    'title': "Planshet ko'rinishi",
                    'caption': "Harakatdagi dars yoki guruh ishlari uchun ixcham va tez ko'rinadigan variant shakllandi.",
                }),
                self._slide_item('phone_mockup', {
                    'title': "Telefon ko'rinishi",
                    'caption': "Qisqa eslatma, tezkor aloqa va yakuniy kuzatuvlarni mobil formatda ko'rsatish mumkin bo'ldi.",
                }),
                self._slide_item('global_impact_map', {
                    'title': "Amaliy qo'llash nuqtalari",
                    'location_1_name': 'Maktab',
                    'location_1_text': "Asosiy kuzatuv va dars sinovi shu yerda o'tdi.",
                    'location_2_name': 'Sinf',
                    'location_2_text': "O'quvchi faolligi bevosita shu muhitda baholandi.",
                    'location_3_name': 'Hamjamiyat',
                    'location_3_text': "Natijalar ota-ona va ustozlar bilan muhokama qilindi.",
                }),
                self._slide_item('key_events_timeline', {
                    'title': 'Jarayon bosqichlari',
                    'event_1_heading': 'Boshlanish',
                    'event_1_text': "Mavzu va asosiy savollar aniqlashtirildi.",
                    'event_2_heading': 'Kuzatuv',
                    'event_2_text': "Dars jarayonidagi birlamchi dalillar yig'ildi.",
                    'event_3_heading': 'Taqqoslash',
                    'event_3_text': "Guruhlar va natijalar o'zaro solishtirildi.",
                    'event_4_heading': 'Talqin',
                    'event_4_text': "Kuchli va zaif jihatlar yakunlandi.",
                    'milestone_heading': 'Marralar',
                    'milestone_text': "Amaliy tavsiyalar bitta xulosaga jamlandi.",
                }),
                self._slide_item('semantic_map', {
                    'title': 'Tahlil xaritasi',
                    'subtitle': "Bog'lanishlar",
                    'node_1_heading': 'Motivatsiya',
                    'node_1_text': "Qiziqish boshlang'ich harakatni kuchaytiradi.",
                    'node_2_heading': "Nazorat",
                    'node_2_text': "Kuzatuv xatolarni erta ko'rsatadi.",
                    'node_3_heading': "Ko'nikma",
                    'node_3_text': "Amaliyot orqali mustahkamlanadi.",
                    'node_4_heading': 'Dalil',
                    'node_4_text': "Jadval va misol bilan tasdiqlanadi.",
                    'node_5_heading': 'Natija',
                    'node_5_text': "Yakuniy sifatni belgilab beradi.",
                }),
                self._slide_item('key_findings_table', {
                    'title': 'Asosiy topilmalar',
                    'header_aspect': "Ko'rsatkich",
                    'header_subject_a': 'A guruh',
                    'header_subject_b': 'B guruh',
                    'header_subject_c': 'C guruh',
                    'row_1_aspect': 'Faollik',
                    'row_1_subject_a': '78%',
                    'row_1_subject_b': '69%',
                    'row_1_subject_c': '73%',
                    'row_2_aspect': 'Tushunish',
                    'row_2_subject_a': '81%',
                    'row_2_subject_b': '72%',
                    'row_2_subject_c': '76%',
                    'row_3_aspect': "Qo'llash",
                    'row_3_subject_a': '74%',
                    'row_3_subject_b': '63%',
                    'row_3_subject_c': '70%',
                    'row_4_aspect': 'Barqarorlik',
                    'row_4_subject_a': '77%',
                    'row_4_subject_b': '68%',
                    'row_4_subject_c': '71%',
                }),
                self._slide_item('periodic_analysis_chart', {
                    'title': 'Davriy tahlil',
                    'chart_caption': "To'rtta yo'nalish bo'yicha umumiy o'sish kuzatildi, ayniqsa muntazam nazorat va amaliy topshiriq qo'llangan bosqichlarda.",
                    'series_1_label': 'Faollik',
                    'series_1_text': "Muntazam ishtirok oyma-oy ko'tarildi.",
                    'series_2_label': 'Tushunish',
                    'series_2_text': "Izohli misollar bilan natija yaxshilandi.",
                    'series_3_label': "Qo'llash",
                    'series_3_text': "Amaliy vazifa samarasi aniq ko'rindi.",
                    'series_4_label': "Qiziqish",
                    'series_4_text': "Mavzuga ijobiy munosabat mustahkamlandi.",
                }),
                self._slide_item('organized_data_summary', {
                    'title': "Tartiblangan ma'lumot",
                    'header_data': "Ma'lumot",
                    'header_description': 'Tavsif',
                    'row_1_data': 'Turi',
                    'row_1_description': "So'rov, kuzatuv va topshiriq natijalari.",
                    'row_2_data': 'Manba',
                    'row_2_description': "Sinf jarayoni, ustoz izohi va o'quvchi javobi.",
                    'row_3_data': 'Birlik',
                    'row_3_description': "Foiz, son va qisqa sifat tavsifi.",
                    'row_4_data': 'Oraliq',
                    'row_4_description': "Bir necha dars va qisqa bosqichlar kesimi.",
                    'row_5_data': "O'rtacha",
                    'row_5_description': "Guruhlar bo'yicha umumiy yo'nalishni ko'rsatadi.",
                }),
            ]
        }
        return self._sanitize_result(template, payload, raw)

    def _build_history_writing_thesis_fallback(
        self,
        template: MagicTemplateSpec,
        payload: dict[str, Any],
    ) -> list[dict[str, Any]]:
        variables = dict(payload.get('variables') or {})
        language = str(variables.get('language') or 'uz').strip().lower()
        if language != 'uz':
            return self._sanitize_result(template, payload, {'slides_content': []})

        topic = str(variables.get('topic') or template.template_name).strip()
        focus_period = str(variables.get('optional_focus_period') or '').strip() or 'asosiy davr'
        region = str(variables.get('optional_region') or '').strip() or 'asosiy hudud'

        words = topic.split()
        title_line_1 = ' '.join(words[:4]) or topic
        title_line_2 = ' '.join(words[4:8]) or 'tarixiy tahlil'

        raw = {
            'slides_content': [
                self._slide_item('cover', {
                    'title_line_1': title_line_1,
                    'title_line_2': title_line_2,
                    'subtitle': f"{focus_period} va {region} kesimidagi tarixiy taqdimot",
                }),
                self._slide_item('table_of_contents', {
                    'toc_title': 'Mundarija',
                    'section_1_title': 'Kirish',
                    'section_1_desc': 'mavzu va tarixiy kontekst',
                    'section_2_title': 'Rivoj',
                    'section_2_desc': 'bosqich va burilishlar',
                    'section_3_title': 'Tahlil',
                    'section_3_desc': 'hudud va natijalar',
                    'section_4_title': 'Xulosa',
                    'section_4_desc': 'meros va yakun',
                }),
                self._slide_item('section_divider_intro', {
                    'section_title': 'Kirish',
                    'section_number': '01',
                    'section_subtitle_line_1': focus_period,
                    'section_subtitle_line_2': 'tarixiy manzara',
                }),
                self._slide_item('intro_narrative', {
                    'title': 'Tarixiy fon',
                    'body': f"{topic} {focus_period} davrida siyosiy, ijtimoiy va madaniy o'zgarishlar kesishgan maydonda shakllandi. {region} doirasidagi institutlar, jamoaviy kayfiyat va ma'rifiy oqimlar bu jarayonning ichki mantiqini belgiladi.",
                }),
                self._slide_item('quote_slide', {
                    'quote': f"{topic}ni anglash uchun voqea bilan birga uni ko'targan muhitni ham o'qish kerak.",
                    'author': 'Tadqiqot tezisi',
                }),
                self._slide_item('visual_interlude_1', {
                    'title': 'Tarixiy kayfiyat',
                }),
                self._slide_item('study_objectives', {
                    'title': 'Tadqiqot vazifalari',
                    'card_1_heading': 'Kontekst',
                    'card_1_text': 'Davr bosimlari va tarixiy fonni tiklash.',
                    'card_2_heading': 'Aktorlar',
                    'card_2_text': 'Guruhlar va institutlar rolini bog‘lash.',
                    'card_3_heading': 'Meros',
                    'card_3_text': 'Uzoq davom etgan ta’sirni ko‘rsatish.',
                }),
                self._slide_item('evolution_part_1', {
                    'title': 'Rivoj bosqichlari',
                    'step_1_heading': 'Boshlanish',
                    'step_1_text': f"{focus_period} boshida g'oyaviy tayanchlar shakllandi.",
                    'step_2_heading_line_1': 'Faollashuv',
                    'step_2_text': "Institutlar va jamoatchilik o'rtasida yangi aloqa kuchaydi.",
                    'step_3_heading_line_1': 'Tarqalish',
                    'step_3_text': f"{region} bo'ylab mazmun va shakllar kengaydi.",
                }),
                self._slide_item('evolution_part_2', {
                    'title': 'Rivoj bosqichlari',
                    'step_4_heading_line_1': 'Burilish',
                    'step_4_text': 'Bahslar va qarama-qarshi kuchlar yo‘nalishni o‘zgartirdi.',
                    'step_5_heading_line_1': 'Meros',
                    'step_5_text': "Keyingi davrlarga o'tgan ma'naviy va ijtimoiy iz qoldi.",
                }),
                self._slide_item('literature_review', {
                    'title': 'Adabiyotlar talqini',
                    'left_heading': 'Klassik talqin',
                    'left_text': f"{topic}ni markaziy voqea va yetakchi aktorlar atrofida o'qiydi, ammo ijtimoiy baza kamroq ochiladi.",
                    'right_heading': 'Yangi qarash',
                    'right_text': "Jarayonni hudud, institut, jamoa va uzoq meros birligida ko'rishni taklif qiladi.",
                }),
                self._slide_item('taxonomy_network', {
                    'title': 'Tayanch tushunchalar',
                    'label_1': 'davr',
                    'label_2': 'aktorlar',
                    'label_3': 'institut',
                    'label_4': 'hudud',
                    'label_5': 'matbuot',
                    'label_6': 'meros',
                }),
                self._slide_item('section_divider_2', {
                    'section_title': 'Tahlil',
                    'section_number': '02',
                    'section_subtitle_line_1': region,
                    'section_subtitle_line_2': 'manba va natija',
                }),
                self._slide_item('methodology_materials', {
                    'title': 'Manba va usullar',
                    'card_1_heading': 'Matbuot',
                    'card_1_text': "Davr tili va jamoaviy kayfiyatni ochadi.",
                    'card_2_heading': 'Arxiv',
                    'card_2_text': 'Rasmiy qaror va amaliy jarayonni ko‘rsatadi.',
                    'card_3_heading': 'Xotira',
                    'card_3_text': 'Keyingi talqin va merosni izohlaydi.',
                    'card_4_heading': 'Qiyos',
                    'card_4_text': 'Hududlar va bosqichlarni solishtiradi.',
                }),
                self._slide_item('summary_table', {
                    'title': 'Materiallar jamlanmasi',
                    'row_1_label': 'Matbuot',
                    'row_1_text': 'G‘oya va tilning ommaviy shaklini beradi.',
                    'row_2_label': 'Arxiv',
                    'row_2_text': 'Institutlar qarori va rasmiy yo‘nalishni ko‘rsatadi.',
                    'row_3_label': 'Xotira',
                    'row_3_text': 'Uzoq ta’sir va qayta talqinni yoritadi.',
                    'row_4_label': 'Qiyos',
                    'row_4_text': 'Hududiy farqlarni bir mantiqqa soladi.',
                }),
                self._slide_item('chart_slide', {
                    'title': 'Tahlil kesimi',
                    'point_1_label': 'davr bosimi',
                    'point_2_label': 'jamoa kayfiyati',
                    'point_3_label': 'institut roli',
                    'point_4_label': 'uzoq meros',
                }),
                self._slide_item('visual_interlude_2', {
                    'title': 'Tarixiy burilish nuqtasi',
                }),
                self._slide_item('infographic_focus', {
                    'title': 'Asosiy yo‘nalish',
                    'era_label': focus_period,
                    'intro_text': f"{topic} bir vaqtning o'zida g'oyaviy, ijtimoiy va amaliy qatlamlarda rivojlangan.",
                    'detail_1_heading': 'Tarkib',
                    'detail_1_text': "Ichki mazmun ta'lim, matbuot va jamoa bilan boyidi.",
                    'detail_2_heading': 'Natija',
                    'detail_2_text': 'Keyingi davr uchun tayanch ma’naviy iz qoldi.',
                }),
                self._slide_item('results_flow', {
                    'title': 'Natijalar oqimi',
                    'stage_1': 'g‘oya',
                    'stage_2': 'tarqalish',
                    'stage_3': 'meros',
                }),
                self._slide_item('milestone_timeline', {
                    'event_1_year': '1-bosqich',
                    'event_1_text': 'Dastlabki g‘oyaviy tayanchlar shakllandi.',
                    'event_2_year': '2-bosqich',
                    'event_2_text': 'Faol institut va jamoa aloqasi kuchaydi.',
                    'event_3_year': '3-bosqich',
                    'event_3_text': 'Uzoq davom etgan madaniy ta’sir belgilandi.',
                }),
                self._slide_item('map_slide', {
                    'title': 'Hududiy ko‘rinish',
                    'region_1_name': 'Markaz',
                    'region_1_text': "G'oya va tashabbus markazi sifatida ko'rinadi.",
                    'region_2_name': 'Chekka',
                    'region_2_text': 'Qabul va moslashuvning o‘ziga xos shaklini beradi.',
                }),
                self._slide_item('single_highlight', {
                    'highlight_value': focus_period,
                    'highlight_text': "Tarixiy burilish aynan shu davrda ko'rinadigan tus oldi.",
                }),
                self._slide_item('conclusions', {
                    'title': 'Xulosalar',
                    'point_1_heading': 'Fon',
                    'point_1_text': 'Jarayon murakkab tarixiy bosimlarda shakllandi.',
                    'point_2_heading': 'Aktor',
                    'point_2_text': 'Institutsiya va ziyoli tarmog‘i muhim rol o‘ynadi.',
                    'point_3_heading': 'Hudud',
                    'point_3_text': 'Hududiy farq umumiy talqinni chuqurlashtiradi.',
                    'point_4_heading': 'Til',
                    'point_4_text': 'Matbuot va muhokama mazmunni yoydi.',
                    'point_5_heading': 'Meros',
                    'point_5_text': 'Ta’sir keyingi davrlarda ham saqlanib qoldi.',
                    'point_6_heading': 'Talqin',
                    'point_6_text': 'Voqea va muhit birgalikda o‘qilishi kerak.',
                }),
                self._slide_item('chronology_early', {
                    'title': 'Xronologiya',
                    'event_1_year': 'Boshlanish',
                    'event_1_heading': 'Tayanch',
                    'event_1_text': 'Dastlabki g‘oyalar shakllandi.',
                    'event_2_year': 'Faollashuv',
                    'event_2_heading': 'Tarmoq',
                    'event_2_text': 'Aktorlar va institutlar bog‘landi.',
                    'event_3_year': 'Kengayish',
                    'event_3_heading': 'Tarqalish',
                    'event_3_text': 'Mazmun yangi maydonlarga yoyildi.',
                }),
                self._slide_item('chronology_late', {
                    'title': 'Xronologiya',
                    'event_4_year': 'Burilish',
                    'event_4_heading': 'Bahs',
                    'event_4_text': 'Qarama-qarshi kuchlar yo‘nalishga ta’sir qildi.',
                    'event_5_year': 'Davom',
                    'event_5_heading': 'Iz',
                    'event_5_text': 'Madaniy xotirada mavzu saqlanib qoldi.',
                    'event_6_year': 'Meros',
                    'event_6_heading': 'Ta’sir',
                    'event_6_text': 'Keyingi avlodlar uchun tayanch bo‘ldi.',
                }),
                self._slide_item('multimedia_or_legacy', {
                    'title': 'Uzoq meros',
                    'body': f"{topic} bugungi tarixiy xotira, ma'rifiy qarash va jamoaviy o'zlikni tushunishda ham muhim tayanch bo'lib qoladi.",
                }),
                self._slide_item('references', {
                    'title': 'Adabiyotlar',
                    'entry_1': "Tarixiy monografiyalar va umumiy tadqiqotlar",
                    'entry_2': "Davr matbuoti va sharhlar to'plami",
                    'entry_3': "Hududiy tarix va ijtimoiy jarayonlar ishlari",
                    'entry_4': "Madaniy xotira hamda merosga oid tahlillar",
                }),
            ]
        }
        return self._sanitize_result(template, payload, raw)

    def _build_business_workplan_fallback(
        self,
        template: MagicTemplateSpec,
        payload: dict[str, Any],
    ) -> list[dict[str, Any]]:
        variables = dict(payload.get('variables') or {})
        language = str(variables.get('language') or 'uz').strip().lower()
        if language != 'uz':
            return self._sanitize_result(template, payload, {'slides_content': []})

        topic = str(variables.get('topic') or template.template_name).strip()
        company = (
            str(variables.get('company_name') or '').strip()
            or str(variables.get('optional_company_context') or '').strip()
            or 'NovaPay'
        )
        industry = (
            str(variables.get('industry') or '').strip()
            or str(variables.get('optional_industry') or '').strip()
            or 'raqamli savdo'
        )
        timeframe = (
            str(variables.get('reporting_period') or '').strip()
            or str(variables.get('optional_timeframe') or '').strip()
            or '90 kun'
        )
        region = str(variables.get('region') or '').strip() or 'asosiy hududlar'

        raw = {
            'slides_content': [
                self._slide_item('cover', {
                    'title': topic,
                    'subtitle': f"{company} uchun {timeframe}lik ustuvor reja",
                }),
                self._slide_item('table_of_contents', {
                    'toc_title': 'Yoʻl xarita',
                    'section_1_title': 'Kontekst',
                    'section_1_desc': 'start nuqta va asosiy fokus',
                    'section_2_title': 'Ijro',
                    'section_2_desc': 'kanallar va ish oqimlari',
                    'section_3_title': 'Nazorat',
                    'section_3_desc': 'xatarlar va boshqaruv ritmi',
                    'section_4_title': 'Yakun',
                    'section_4_desc': 'tayanchlar va nazorat jadvali',
                }),
                self._slide_item('section_divider_context', {
                    'section_title': 'Kontekst',
                    'section_number': '01',
                    'section_subtitle': f"{industry} bozorida {timeframe}lik fokus",
                }),
                self._slide_item('priority_snapshot', {
                    'title': '90 kunlik maqsadlar',
                    'intro_text': 'Asosiy ustuvor yo‘nalishlar:',
                    'bullet_1': 'savdo hajmini 15% oshirish',
                    'bullet_2': 'konversiyani 5% ko‘tarish',
                }),
                self._slide_item('channel_comparison', {
                    'title': 'O‘sish va sifat',
                    'concept_1_heading': 'Kanal o‘sishi',
                    'concept_1_text': 'Toshkent va viloyatlarda talab yuqori segmentlar kuchaytiriladi.',
                    'concept_2_heading': 'Xizmat sifati',
                    'concept_2_text': 'Yetkazish tezligi va mijoz javobi haftalik tartibda nazorat qilinadi.',
                }),
                self._slide_item('workstreams_overview', {
                    'title': 'Uch ish oqimi',
                    'column_1_heading': 'Talab',
                    'column_1_text': 'Eng issiq kanal va segmentlar qayta ustuvorlashtiriladi.',
                    'column_2_heading': 'Jarayon',
                    'column_2_text': 'Leaddan savdogacha uzilishlar qisqartirilib, javob muddati standartlanadi.',
                    'column_3_heading': 'Nazorat',
                    'column_3_text': 'Haftalik ko‘rsatkichlar va og‘ishlar bitta panelda jamlanadi.',
                }),
                self._slide_item('section_divider_execution', {
                    'section_title': 'Ijro',
                    'section_number': '02',
                    'section_subtitle': f"{company} uchun tezlik va barqarorlik muvozanati",
                }),
                self._slide_item('execution_balances', {
                    'title': 'Tezlik va barqarorlik',
                    'concept_1_heading': 'Tezkor harakat',
                    'concept_1_text': 'Issiq segmentlarda taklif va reklama tezroq yangilanadi.',
                    'concept_2_heading': 'Barqaror xizmat',
                    'concept_2_text': 'Ombor, logistika va mijoz javobi bir maromda ushlab turiladi.',
                }),
                self._slide_item('risk_control_snapshot', {
                    'title': 'Asosiy xatarlar',
                    'intro_text': 'Ijro ritmiga taʼsir qiluvchi nuqtalar:',
                    'bullet_1': 'ombor va yetkazish uzilishi',
                    'bullet_2': 'kanal bo‘yicha konversiya pasayishi',
                }),
                self._slide_item('image_focus', {
                    'title': 'Bitta panel, bitta qaror ritmi',
                    'caption': 'Rahbar paneli oqim, konversiya va javob tezligini bir joyda ko‘rsatadi.',
                }),
                self._slide_item('section_divider_governance', {
                    'section_title': 'Nazorat',
                    'section_number': '03',
                    'section_subtitle': f"{region} bo‘yicha ritm va javob mexanizmi",
                }),
                self._slide_item('infographic', {
                    'title': 'Yetkazishni ushlab turuvchi toʻrtta tayanch',
                    'item_1_heading': 'Segment',
                    'item_1_text': 'yuqori qaytuvli mijozlar alohida boshqariladi',
                    'item_2_heading': 'Taklif',
                    'item_2_text': 'takliflar qisqa paketlar atrofida soddalashtiriladi',
                    'item_3_heading': 'Jamoa',
                    'item_3_text': 'vazifalar haftalik sprint va egaga bogʻlanadi',
                    'item_4_heading': 'Nazorat',
                    'item_4_text': 'ogʻishlar 48 soat ichida qayta ko‘rib chiqiladi',
                }),
                self._slide_item('table_summary', {
                    'title': '90 kunlik nazorat jadvali',
                    'table_stub_header': 'Yoʻnalish',
                    'table_column_1': '30 kun',
                    'table_column_2': '60 kun',
                    'table_column_3': '90 kun',
                    'row_1_label': 'Talab',
                    'row_1_value_1': 'segment tanlandi',
                    'row_1_value_2': 'kanal kuchaydi',
                    'row_1_value_3': 'barqaror oqim',
                    'row_2_label': 'Jarayon',
                    'row_2_value_1': 'handoff tozalandi',
                    'row_2_value_2': 'javob tezlashdi',
                    'row_2_value_3': 'standart yopildi',
                    'row_3_label': 'Nazorat',
                    'row_3_value_1': 'panel ishga tushdi',
                    'row_3_value_2': 'ogʻish ko‘rindi',
                    'row_3_value_3': 'qaror ritmi mustahkam',
                }),
            ]
        }
        return self._sanitize_result(template, payload, raw)

    def _build_business_annual_report_fallback(
        self,
        template: MagicTemplateSpec,
        payload: dict[str, Any],
    ) -> list[dict[str, Any]]:
        variables = dict(payload.get('variables') or {})
        language = str(variables.get('language') or 'uz').strip().lower()
        if language != 'uz':
            return self._sanitize_result(template, payload, {'slides_content': []})

        topic = str(variables.get('topic') or template.template_name).strip()
        company = str(variables.get('company_name') or '').strip() or topic
        period = str(variables.get('reporting_period') or '').strip() or 'joriy hisobot davri'
        industry = str(variables.get('optional_industry') or '').strip() or 'asosiy bozor'
        region = str(variables.get('optional_region') or '').strip() or 'ustuvor hududlar'
        context = str(variables.get('optional_company_context') or '').strip() or industry

        raw = {
            'slides_content': [
                self._slide_item('cover', {
                    'title': company,
                    'subtitle': f"{period} bo'yicha natijalar va keyingi ustuvorliklar",
                }),
                self._slide_item('table_of_contents', {
                    'toc_title': 'Mundarija',
                    'section_1_title': 'Bozor',
                    'section_2_title': 'Natijalar',
                    'section_3_title': 'Drayverlar',
                    'section_4_title': 'Hisobotlar',
                    'section_5_title': 'KPI tarkibi',
                    'section_6_title': 'Keyingi qadam',
                }),
                self._slide_item('opening_summary', {
                    'title': 'Yil yakuni',
                    'body': f"{company} uchun {period} davri {context} bozorida o'sish sifati va foyda intizomini birga tiklash bosqichi bo'ldi. Ushbu hisobot daromad tarkibi, operatsion qarorlar va keyingi ustuvorliklarni jamlaydi.",
                }),
                self._slide_item('market_context', {
                    'title': 'Bozor va ichki fon',
                    'paragraph_1': f"{industry} segmentida talab saqlanib qoldi, biroq xaridor qarori ancha ehtiyotkor bo'ldi. Tez qiymat, ishonchli servis va aniq iqtisodiy samara bera olgan takliflar kuchliroq ishladi.",
                    'paragraph_2': f"Ichkarida {company} segment tanlovi, xarajat intizomi va jamoalararo bajarish ritmiga tayanib ishladi. Shu yondashuv {region} bo'yicha ko'rinishni yaxshilab, yil oxirida sifatliroq tushum yaratdi.",
                }),
                self._slide_item('key_highlights', {
                    'title': 'Yilning asosiy natijalari',
                    'intro_text': f"{company} yili quyidagi tayanch natijalar bilan yopdi.",
                    'bullet_1': 'Takroriy tushum ulushi oshdi',
                    'bullet_2': 'Marjali mahsulotlar salmogʻi kengaydi',
                    'bullet_3': 'Churn bosimi pasaydi',
                    'bullet_4': "Prognoz aniqligi sezilarli yaxshilandi",
                }),
                self._slide_item('comparison_drivers', {
                    'title': 'Ikki asosiy drayver',
                    'left_heading': 'Daromad sifati',
                    'left_text': "O'sish sof hajm bilan emas, yuqori qaytuvli segmentlar va takroriy paketlarga tayanish orqali tezlashdi.",
                    'right_heading': 'Marja sifati',
                    'right_text': "Marja narx intizomi, servis tannarxini boshqarish va past qaytuvli tashabbuslarni yopish hisobiga tiklandi.",
                }),
                self._slide_item('strategic_themes', {
                    'title': 'Strategik mavzular',
                    'theme_1_heading': 'Segment tanlovi',
                    'theme_1_text': "Past qaytuvli hajmdan ko'ra foydaliroq segmentlarga ustuvorlik berildi.",
                    'theme_2_heading': 'Mijoz qadri',
                    'theme_2_text': "Takroriy xarid va xizmat chuqurligi barqaror ko'tarildi.",
                    'theme_3_heading': 'Bajarish ritmi',
                    'theme_3_text': "Sotuv, servis va moliya umumiy boshqaruv ritmiga o'tdi.",
                }),
                self._slide_item('operating_pillars', {
                    'title': "Biznesni ushlab turgan to'rtta tayanch",
                    'pillar_1_heading': 'Taklif',
                    'pillar_1_text': "Portfel tez qaytadigan toifalar atrofida qayta yig'ildi.",
                    'pillar_2_heading': 'Kanal',
                    'pillar_2_text': "Konversiyasi yuqori kanallarga resurs ko'proq yo'naltirildi.",
                    'pillar_3_heading': 'Narx',
                    'pillar_3_text': "Chegirma tartibi qayta ko'rilib, foyda sifati himoya qilindi.",
                    'pillar_4_heading': 'Xizmat',
                    'pillar_4_text': 'Yetkazish sifati va javob tezligi yagona standartga keltirildi.',
                }),
                self._slide_item('review_grid', {
                    'title': 'Natijani izohlovchi omillar',
                    'item_1_heading': 'Takroriy tushum',
                    'item_1_text': 'Barqaror baza yil yakunini mustahkamladi.',
                    'item_2_heading': 'Narx intizomi',
                    'item_2_text': 'Marja uchun zararli chegirmalar qisqartirildi.',
                    'item_3_heading': 'Onboarding',
                    'item_3_text': 'Yangi mijozlar qiymatni tezroq ola boshladi.',
                    'item_4_heading': 'Xarajat nazorati',
                    'item_4_text': 'Past samarali faoliyatdan resurs qayta taqsimlandi.',
                    'item_5_heading': 'Jamoa uygʻunligi',
                    'item_5_text': "Jamoalar bir xil boshqaruv ritmiga o'tdi.",
                    'item_6_heading': 'Prognoz aniqligi',
                    'item_6_text': "Reja va fakt o'rtasidagi tafovut qisqardi.",
                }),
                self._slide_item('quote_highlight', {
                    'quote': 'Barqaror natija hajm bilan emas, intizomli tanlov bilan quriladi.',
                    'author': 'Rahbariyat xulosasi',
                }),
                self._slide_item('dashboard_overview', {
                    'title': "Boshqaruv paneli ko'rinishi",
                    'body': f"Ichki panel {company} uchun tushum sifati, segmentlar dinamikasi, churn xavfi va choraklik bajarilish holatini bir ekranda jamlaydi.",
                }),
                self._slide_item('yearly_timeline', {
                    'title': "Yil davomida bajarilgan yo'l",
                    'step_1_heading': 'Yanvar',
                    'step_1_text': 'Fokus segmentlar qayta tanlandi',
                    'step_2_heading': 'Fevral',
                    'step_2_text': 'Narx paketlari soddalashtirildi',
                    'step_3_heading': 'Aprel',
                    'step_3_text': 'Onboarding tezligi oshirildi',
                    'step_4_heading': 'May',
                    'step_4_text': 'Sotuv playbooki yagona qilindi',
                    'step_5_heading': 'Iyul',
                    'step_5_text': 'Past qaytuvli tashabbuslar yopildi',
                    'step_6_heading': 'Sentabr',
                    'step_6_text': 'Hududiy kengayish kalibrlashtirildi',
                    'step_7_heading': 'Noyabr',
                    'step_7_text': 'Qayta xarid dasturi kuchaydi',
                    'step_8_heading': 'Dekabr',
                    'step_8_text': 'Yil yakuni prioritetlari jamlandi',
                }),
                self._slide_item('value_chain', {
                    'title': 'Qiymat zanjiri qanday ishladi',
                    'block_1_heading': 'Talab',
                    'block_1_text': 'Talab signallari aniqroq oʻqildi va prioritetlar toraytirildi.',
                    'block_2_heading': 'Taklif',
                    'block_2_text': 'Taklif kuchli konversiya beradigan paketlarga yigʻildi.',
                    'block_3_heading': 'Yetkazish',
                    'block_3_text': 'Servis sifati va tezligi bir xil standartga keltirildi.',
                    'block_4_heading': 'Natija',
                    'block_4_text': "Marja, prognoz va churn bo'yicha yaxshilanish ko'rindi.",
                }),
                self._slide_item('reporting_catalog', {
                    'title': 'Rahbariyat kuzatadigan hisobotlar',
                    'catalog_kicker': 'Asosiy turlar',
                    'item_1_heading': 'Daromad',
                    'item_1_text': "kanal va segment kesimi",
                    'item_2_heading': 'Marja',
                    'item_2_text': "paket va mahsulot talqini",
                    'item_3_heading': 'Ushlab qolish',
                    'item_3_text': 'kohort va churn signali',
                    'item_4_heading': 'Sotuv oqimi',
                    'item_4_text': 'bosqich va konversiya holati',
                    'item_5_heading': 'Xarajat',
                    'item_5_text': 'funksiya va tashabbus nazorati',
                    'item_6_heading': 'Likvidlik',
                    'item_6_text': 'naqd oqim va bosimlar',
                    'item_7_heading': 'Prognoz',
                    'item_7_text': 'reja va fakt aniqligi',
                    'item_8_heading': 'Audit',
                    'item_8_text': 'ichki nazorat va ishonch',
                }),
                self._slide_item('initiatives_milestones', {
                    'title': 'Tashabbuslar va tayanch nuqtalar',
                    'left_heading': 'Asosiy tashabbuslar',
                    'left_point_1': "Yuqori qaytuvli segmentlar ulushi kengaytirildi.",
                    'left_point_2': "Xizmat va yetkazish ritmi soddalashtirildi.",
                    'left_point_3': "Prognoz intizomi kuchaytirilib, qarorlar tezlashtirildi.",
                    'right_heading': 'Muvofiqlashtirilgan marralar',
                    'milestone_1_label': 'Q1',
                    'milestone_1_text': 'Segment va paketlar qayta yigʻildi',
                    'milestone_2_label': 'Q2',
                    'milestone_2_text': 'Narx va xizmat standarti birlashtirildi',
                    'milestone_3_label': 'Q4',
                    'milestone_3_text': 'Kelasi yil uchun tayanch reja qulflandi',
                }),
                self._slide_item('segment_distribution', {
                    'title': 'Daromad tarkibi',
                    'segment_1_name': 'Takroriy paketlar',
                    'segment_1_percent': 44,
                    'segment_1_text': "Barqaror tushum va yuqori ko'rinish beradi.",
                    'segment_2_name': 'Kengayish savdosi',
                    'segment_2_percent': 31,
                    'segment_2_text': "Mavjud bazada marjali o'sish manbai bo'ldi.",
                    'segment_3_name': 'Yangi mijozlar',
                    'segment_3_percent': 25,
                    'segment_3_text': 'Selektiv jalb qilinish tufayli sifat saqlandi.',
                }),
                self._slide_item('kpi_summary', {
                    'title': "Asosiy KPI ko'rsatkichlar",
                    'kpi_1_value': '18%',
                    'kpi_1_text': 'takroriy tushum oʻsishi',
                    'kpi_2_value': '92%',
                    'kpi_2_text': 'asosiy bazani ushlab qolish',
                    'kpi_3_value': '11 kun',
                    'kpi_3_text': "birinchi qiymatgacha bo'lgan vaqt",
                    'kpi_4_value': '23%',
                    'kpi_4_text': 'yuqori marjali mahsulot ulushi',
                    'kpi_5_value': '1.4x',
                    'kpi_5_text': 'reja-fakt aniqligi yaxshilanishi',
                }),
                self._slide_item('period_comparison', {
                    'title': 'Operatsion marja talqini',
                    'period_1_name': 'Yil boshida',
                    'period_1_text': "Marja hajmni ushlab qolish va qayta narxlash o'rtasida siqildi. Past qaytuvli faoliyat daromadning sifatsiz qismini kattalashtirdi.",
                    'period_2_name': 'Yil oxirida',
                    'period_2_text': "Segment tanlovi va xarajat intizomi hisobiga operatsion foyda sifati tiklandi. Foyda kamroq shovqin va yaxshiroq ko'rinish bilan shakllandi.",
                }),
            ]
        }
        return self._sanitize_result(template, payload, raw)

    def _build_science_newsletter_fallback(
        self,
        template: MagicTemplateSpec,
        payload: dict[str, Any],
    ) -> list[dict[str, Any]]:
        variables = dict(payload.get('variables') or {})
        language = str(variables.get('language') or 'uz').strip().lower()
        if language != 'uz':
            return self._sanitize_result(template, payload, {'slides_content': []})

        topic = str(variables.get('topic') or template.template_name).strip()
        focus_area = str(variables.get('optional_focus_area') or '').strip() or 'maktab ilmiy loyihalari'
        event_name = str(variables.get('optional_event_name') or '').strip() or "ilmiy ko'rgazma soni"

        raw = {
            'slides_content': [
                self._slide_item('cover', {
                    'title': topic,
                    'subtitle': f"{event_name} uchun qisqa ilmiy yangiliklar soni",
                }),
                self._slide_item('table_of_contents', {
                    'toc_title': 'Mundarija',
                    'section_1_title': "E'lonlar",
                    'section_1_desc': 'dolzarb yangilanishlar',
                    'section_2_title': 'Tahlil',
                    'section_2_desc': "chuqurroq qarash",
                    'section_3_title': 'Qisqa',
                    'section_3_desc': 'tezkor xabarlar',
                    'section_4_title': 'Jamoa',
                    'section_4_desc': 'ishtirokchi va fikrlar',
                    'section_5_title': 'Tadbirlar',
                    'section_5_desc': 'raqamlar va xulosa',
                }),
                self._slide_item('section_divider_announcements', {
                    'section_number': '01',
                    'section_title': "E'lonlar",
                    'section_subtitle': 'asosiy yangiliklar va tayanch signal',
                }),
                self._slide_item('lead_announcement', {
                    'title': 'Asosiy eʼlon',
                    'body': f"{event_name} doirasida {topic} bo‘yicha eng faol loyihalar alohida saralandi va {focus_area} yo‘nalishidagi ishlar uchun kengroq namoyish maydoni ajratildi.",
                }),
                self._slide_item('highlights_of_the_month', {
                    'title': 'Oy yakunlari',
                    'highlight_1_heading': 'Loyiha',
                    'highlight_1_text': "Yangi tajriba ishlari ko'proq amaliy yo'naldi.",
                    'highlight_2_heading': 'Mentor',
                    'highlight_2_text': "Maslahat sessiyalari ishtirokni jonlantirdi.",
                    'highlight_3_heading': 'Namoyish',
                    'highlight_3_text': "Ko'rgazma burchaklari ancha tartibli ko'rindi.",
                }),
                self._slide_item('visual_photo_break', {
                    'title': 'Tajriba jarayoni markazda',
                    'caption': "Ko'rgazmadagi eng kuchli taassurot amaliy namoyishlar orqali berildi.",
                }),
                self._slide_item('section_divider_in_depth', {
                    'section_number': '02',
                    'section_title': 'Tahlil',
                    'section_subtitle': 'loyiha ichiga chuqurroq nazar',
                }),
                self._slide_item('in_depth_cards', {
                    'title': 'Uch asosiy yoʻnalish',
                    'topic_1_heading': 'Metod',
                    'topic_1_text': "Tajriba usuli aniq bo'lsa, natija ishonchliroq ko'rinadi.",
                    'topic_2_heading': 'Dalil',
                    'topic_2_text': "Jadval va kuzatuv natijani sodda, lekin taʼsirli ko'rsatadi.",
                    'topic_3_heading': 'Talqin',
                    'topic_3_text': "Natijani tushuntirish loyihaning haqiqiy qiymatini ochadi.",
                }),
                self._slide_item('in_depth_timeline', {
                    'title': 'Rivojlanish bosqichi',
                    'year_1': '2023',
                    'year_1_heading': 'Boshlanish',
                    'year_1_text': "Qiziqish uyg'otgan ilk kichik sinovlar o'tkazildi.",
                    'year_2': '2024',
                    'year_2_heading': 'Kengayish',
                    'year_2_text': "Loyihalar soni va mentorlar ishtiroki ko'paydi.",
                    'year_3': '2025',
                    'year_3_heading': 'Sifat',
                    'year_3_text': "Natijalarni himoya qilish usuli yaxshilandi.",
                    'focus_heading': 'Bugun',
                    'focus_text': "Amaliy taʼsir va taqdimot sifati asosiy mezonga aylandi.",
                }),
                self._slide_item('sector_news', {
                    'title': 'Soha yangiliklari',
                    'item_1_heading': 'Kimyo',
                    'item_1_text': "Oddiy vositalar bilan kuchli tajriba ko'rsatildi.",
                    'item_2_heading': 'Biologiya',
                    'item_2_text': "Mahalliy muammolarga mos kuzatuv ishlari ko'paydi.",
                    'item_3_heading': 'Fizika',
                    'item_3_text': "Model va prototiplar ancha tushunarli bo'ldi.",
                    'item_4_heading': 'Texno',
                    'item_4_text': "Qurilma va dastur uyg'unligi yaxshiroq namoyon bo'ldi.",
                    'item_5_heading': 'Ekologiya',
                    'item_5_text': "Atrof-muhitga oid loyihalar ko'proq eʼtibor oldi.",
                    'item_6_heading': 'Matematika',
                    'item_6_text': "Tahlil va hisoblash loyihaning dalilini kuchaytirdi.",
                }),
                self._slide_item('in_brief_list', {
                    'title': 'Qisqacha',
                    'item_1': 'Mentor uchrashuvlari oshdi',
                    'item_2': "Amaliy maketlar ko'paydi",
                    'item_3': "Savol-javoblar jonlandi",
                    'item_4': 'Sinf ichki sinovi kuchaydi',
                    'item_5': 'Poster sifati yaxshilandi',
                    'item_6': "Natija yozuvi soddalashdi",
                    'item_7': "Ko'rgazma tartibi yaxshilandi",
                    'item_8': 'Hamkorlik kuchaydi',
                }),
                self._slide_item('participant_reviews', {
                    'title': 'Ishtirokchi fikrlari',
                    'review_1_text': "Tajriba qilish bizga mavzuni yaqinlashtirdi.",
                    'review_1_author': 'Aziza J.',
                    'review_2_text': "Savollar ko'p bo'lsa ham, himoya ancha foydali bo'ldi.",
                    'review_2_author': 'Kamron T.',
                    'review_3_text': "Mentor bilan ishlash natijani tartiblab berdi.",
                    'review_3_author': 'Dilshod R.',
                    'review_4_text': "Ko'rgazmada tushuntirish qilish ishonchni oshirdi.",
                    'review_4_author': 'Malika S.',
                }),
                self._slide_item('news_columns', {
                    'title': 'Yangiliklar',
                    'column_1_heading': 'Saralash',
                    'column_1_text': "Dastlabki tanlovda amaliy taʼsir ko'rsatgan loyihalar oldinga chiqdi.",
                    'column_2_heading': 'Taqdimot',
                    'column_2_text': "Qisqa va aniq tushuntirish tashrif buyuruvchilar eʼtiborini ushlab turdi.",
                    'column_3_heading': 'Aloqa',
                    'column_3_text': "Ota-ona va ustoz bilan muloqot loyihalarga qo'shimcha kuch berdi.",
                }),
                self._slide_item('our_numbers', {
                    'title': 'Bizning raqamlar',
                    'header_project': 'Loyiha',
                    'header_metric_a': 'Ball',
                    'header_metric_b': 'Tashrif',
                    'header_metric_c': 'Savol',
                    'row_1_project': 'Ekologiya',
                    'row_1_metric_a': '92',
                    'row_1_metric_b': '180',
                    'row_1_metric_c': '34',
                    'row_2_project': 'Robotika',
                    'row_2_metric_a': '95',
                    'row_2_metric_b': '210',
                    'row_2_metric_c': '41',
                    'row_3_project': 'Biologiya',
                    'row_3_metric_a': '89',
                    'row_3_metric_b': '165',
                    'row_3_metric_c': '29',
                }),
                self._slide_item('welcome_team', {
                    'title': 'Xush kelibsiz',
                    'person_1_name': 'Mohira Xasan',
                    'person_1_role': "Eksperimentlar koordinatori va mentorlar bilan aloqa mas'uli.",
                    'person_2_name': 'Javohir Orip',
                    'person_2_role': "Ko'rgazma tartibi va ishtirokchilar yo'nalishini boshqaradi.",
                }),
                self._slide_item('past_events_review', {
                    'title': "Oʻtgan tadbirlar sharhi",
                    'event_1_heading': 'Start',
                    'event_1_text': "Birlamchi loyiha g'oyalari saralandi.",
                    'event_2_heading': 'Sinov',
                    'event_2_text': "Ilk tajribalar va kichik namoyishlar o'tdi.",
                    'event_3_heading': 'Himoya',
                    'event_3_text': "Qisqa taqdimot formati sinovdan o'tdi.",
                    'event_4_heading': "Ko'rgazma",
                    'event_4_text': "Tashrif buyuruvchilar bilan faol muloqot bo'ldi.",
                    'event_5_heading': 'Yakun',
                    'event_5_text': "Eng kuchli ishlar keyingi bosqichga yo'llandi.",
                }),
            ]
        }
        return self._sanitize_result(template, payload, raw)

    def _build_science_epidemiology_fallback(
        self,
        template: MagicTemplateSpec,
        payload: dict[str, Any],
    ) -> list[dict[str, Any]]:
        variables = dict(payload.get('variables') or {})
        language = str(variables.get('language') or 'uz').strip().lower()
        if language != 'uz':
            return self._sanitize_result(template, payload, {'slides_content': []})

        topic = str(variables.get('topic') or template.template_name).strip()
        focus_area = str(variables.get('optional_focus_area') or '').strip() or 'epidemiologik kuzatuv'
        region = str(variables.get('optional_region') or '').strip() or 'tanlangan kuzatuv hududi'

        raw = {
            'slides_content': [
                self._slide_item('cover', {
                    'title': topic,
                    'subtitle': f"{focus_area} va {region} kesimidagi ilmiy himoya",
                }),
                self._slide_item('table_of_contents', {
                    'toc_title': 'Mundarija',
                    'section_1_title': 'Muammo',
                    'section_1_desc': 'dolzarblik va fon',
                    'section_2_title': 'Faraz',
                    'section_2_desc': 'tekshiriladigan bogʻlanishlar',
                    'section_3_title': 'Maqsad',
                    'section_3_desc': 'vazifa va indikatorlar',
                    'section_4_title': 'Usul',
                    'section_4_desc': 'maʼlumot va tanlanma',
                    'section_5_title': 'Tahlil',
                    'section_5_desc': 'dinamika va talqin',
                    'section_6_title': 'Xulosa',
                    'section_6_desc': 'amaliy natijalar',
                }),
                self._slide_item('section_divider_statement', {
                    'section_number': '01',
                    'section_title': 'Tadqiqot bayoni',
                    'section_subtitle': 'muammo va dolzarblik',
                }),
                self._slide_item('introduction', {
                    'title': 'Kirish',
                    'body': f"{topic} mavzusi {region} sharoitida xatar omillari, kuzatuv sifati va javob choralarining bir-biriga qanday taʼsir qilishini ko'rsatadi. Himoya ushbu bogʻlanishlarni tizimli epidemiologik nuqtai nazardan tahlil qiladi.",
                }),
                self._slide_item('purpose_statement', {
                    'title': "Tadqiqotni yo'naltiruvchi savollar",
                    'question_1': 'Qaysi omillar asosiy?',
                    'text_1': 'Xatar kuchayishiga eng ko‘p taʼsir qilayotgan muhit va xulq omillarini ajratish maqsad qilinadi.',
                    'question_2': 'Qaysi signal erta ko‘rinadi?',
                    'text_2': 'Kuzatuv tizimida xavf oshishini eng erta ko‘rsatadigan indikatorlar aniqlanadi.',
                    'question_3': 'Qaysi javob samaraliroq?',
                    'text_3': 'Profilaktika, monitoring va tezkor choralar qaysi kombinatsiyada yaxshiroq natija berishi baholanadi.',
                }),
                self._slide_item('quote_slide', {
                    'quote': "Xatarni boshqarish uchun avval uni qayerda to'planayotganini ko'rish kerak.",
                    'author': 'Tadqiqot tezisi',
                }),
                self._slide_item('current_situation_problems', {
                    'title': 'Hozirgi holat va muammolar',
                    'current_situation_heading': 'Hozirgi holat',
                    'current_situation_text': f"{region} doirasida kuzatuv signallari bir xil emas: ayrim nuqtalarda xavf tez ko'tariladi, ayrimida esa nazorat kechikadi.",
                    'problem_1_heading': 'Kechikkan signal',
                    'problem_1_text': 'Erta ogohlantirish belgilari yetarlicha ushlanmaydi.',
                    'problem_2_heading': 'Notekis nazorat',
                    'problem_2_text': "Hududlar bo'yicha monitoring sifati bir xil emas.",
                    'problem_3_heading': 'Sust javob',
                    'problem_3_text': 'Choralar yuqori xatar nuqtalariga kech moslashtiriladi.',
                }),
                self._slide_item('section_divider_hypothesis', {
                    'section_number': '02',
                    'section_title': 'Farazlar',
                    'section_subtitle': 'tekshiriladigan bogʻlanishlar',
                }),
                self._slide_item('hypotheses', {
                    'title': 'Asosiy farazlar',
                    'hypothesis_1_heading': 'Aloqa zichligi',
                    'hypothesis_1_text': 'Zich aloqa yuqori nuqtalarda xatar yuklamasi tezroq oshadi.',
                    'hypothesis_2_heading': 'Erta kuzatuv',
                    'hypothesis_2_text': 'Erta signal resursni aniqroq taqsimlashga yordam beradi.',
                    'hypothesis_3_heading': 'Qatlamli javob',
                    'hypothesis_3_text': 'Qatlamli javob yakka choradan barqarorroq samara beradi.',
                }),
                self._slide_item('study_objectives', {
                    'title': 'Tadqiqot maqsadlari',
                    'objective_1_heading': 'Omillarni ajratish',
                    'objective_1_text': "Asosiy xatar drayverlarini ustuvorlik bo'yicha saralash.",
                    'objective_2_heading': 'Signalni topish',
                    'objective_2_text': 'Erta ogohlantirish indikatorlarini aniqlash va solishtirish.',
                    'objective_3_heading': 'Javobni baholash',
                    'objective_3_text': "Chora kombinatsiyalarining nisbiy taʼsirini ko'rsatish.",
                    'objective_4_heading': 'Tavsiya berish',
                    'objective_4_text': 'Amaliy nazorat va monitoring tavsiyalarini shakllantirish.',
                }),
                self._slide_item('concept_review', {
                    'title': 'Tayanch epidemiologik tushunchalar',
                    'concept_1_heading': 'Xatar',
                    'concept_1_text': 'Xavfni oshiruvchi omillar majmui.',
                    'concept_2_heading': 'Ekspozitsiya',
                    'concept_2_text': 'Taʼsirga uchrash darajasi.',
                    'concept_3_heading': 'Klaster',
                    'concept_3_text': 'Holatlarning bir joyda jamlanishi.',
                    'concept_4_heading': 'Vaqt oynasi',
                    'concept_4_text': 'Signal ko‘rinadigan kuzatuv oralig‘i.',
                    'concept_5_heading': 'Profilaktika',
                    'concept_5_text': 'Oldini oluvchi choralar majmui.',
                }),
                self._slide_item('literature_review', {
                    'title': 'Adabiyotlar sharhi',
                    'entry_1_ref': 'WHO yoʻriqnomalari',
                    'entry_1_note': 'Erta kuzatuv va qatlamli javobning ahamiyatini ko‘rsatadi.',
                    'entry_2_ref': 'Hududiy hisobotlar',
                    'entry_2_note': 'Xatar notekis taqsimlanishini tasdiqlaydi.',
                    'entry_3_ref': 'Maktab kuzatuvlari',
                    'entry_3_note': 'Zich aloqa muhitida signal tezroq paydo bo‘lishini ko‘rsatadi.',
                    'entry_4_ref': 'Profilaktika tadqiqotlari',
                    'entry_4_note': 'Barqaror choralar takroriy to‘lqinni pasaytirishini bildiradi.',
                    'entry_5_ref': 'Javob strategiyalari',
                    'entry_5_note': 'Mahalliy sharoitga moslashuv natijani yaxshilashini ko‘rsatadi.',
                }),
                self._slide_item('theoretical_framework', {
                    'title': 'Nazariy ramka',
                    'key_terms_heading': 'Tayanch atamalar',
                    'key_term_1': 'xatar zanjiri',
                    'key_term_2': 'ekspozitsiya',
                    'key_term_3': 'signal kechikishi',
                    'key_term_4': 'profilaktika kuchi',
                    'key_term_5': 'aloqa zichligi',
                    'key_term_6': 'klasterlashuv',
                    'key_term_7': 'adaptiv javob',
                    'theories_heading': 'Asosiy yondashuvlar',
                    'theory_1_heading': 'Zanjir modeli',
                    'theory_1_text': 'Xatar signal, muhit va aloqa orqali kuchayadi.',
                    'theory_2_heading': 'Moslashuvchan nazorat',
                    'theory_2_text': 'Joy va vaqtga mos nazorat samarani oshiradi.',
                    'framework_heading': 'Bizning talqin',
                    'framework_text': 'Ramka xatar, signal va javob kuchi o‘rtasidagi bog‘lanishni bir tizimda o‘qiydi.',
                }),
                self._slide_item('methodology', {
                    'title': 'Metodologiya',
                    'method_label': 'Yondashuv',
                    'data_type_label': "Ma'lumot turi",
                    'data_type_text': 'Kuzatuv yozuvlari, holat kesimlari va jarayon qaydlari ishlatildi.',
                    'motives_label': 'Tadqiqot sababi',
                    'motives_text': 'Xatar va javob o‘rtasidagi uzilishni aniqlash ko‘zlandi.',
                    'collection_label': "Yig'ish usuli",
                    'collection_text': 'Nuqtalar vaqt, joy va aloqa zichligi bo‘yicha saralandi.',
                    'sampling_label': 'Tanlanma',
                    'sampling_text': 'Yuqori xatar va sust nazorat ko‘ringan holatlar tanlandi.',
                }),
                self._slide_item('analysis_development_phases', {
                    'title': 'Tahlilning ikki bosqichi',
                    'phase_1_heading': 'Xatarni ajratish',
                    'phase_1_point_1': 'Zich aloqa nuqtalari belgilandi.',
                    'phase_1_point_2': 'Muhit omillari alohida baholandi.',
                    'phase_1_point_3': 'Xulqiy naqshlar qiyoslandi.',
                    'phase_2_heading': 'Javobni tekshirish',
                    'phase_2_point_1': 'Monitoring uzilishlari topildi.',
                    'phase_2_point_2': 'Profilaktika kuchi o‘lchandi.',
                    'phase_2_point_3': 'Tezkor choralar taʼsiri solishtirildi.',
                }),
                self._slide_item('analysis_graph', {
                    'title': 'Kuzatuv natijalari dinamikasi',
                    'graph_caption': 'Toʻrtta holatdagi nisbiy yuklama va xavf bosimi qiyoslanadi.',
                    'note_1_heading': 'Zich aloqa',
                    'note_1_text': 'Eng katta bosimni saqlab qoldi.',
                    'note_2_heading': 'Nazorat uzilishi',
                    'note_2_text': 'Xavfni ancha kuchaytirdi.',
                }),
                self._slide_item('analysis_results', {
                    'title': 'Natijalar tahlili',
                    'highlight_percent': '37%',
                    'result_1_heading': 'Asosiy xatar',
                    'result_1_text': 'Zich aloqa va sust nazorat birga xavfni oshirdi.',
                    'result_2_heading': 'Himoya omili',
                    'result_2_text': 'Erta monitoring choʻqqi bosimini kamaytirdi.',
                    'result_3_heading': 'Amaliy xulosa',
                    'result_3_text': 'Qatlamli javob yakka choradan samaraliroq bo‘ldi.',
                }),
                self._slide_item('discussion', {
                    'title': 'Muhokama',
                    'intro_heading': 'Tayanch dalil',
                    'intro_text': "Dalillar xatarni biologik emas, tashkiliy va xulqiy muhit bilan birga o'qish kerakligini ko'rsatadi.",
                    'bullet_1': 'Xatar bir nechta maydonda yigʻiladi',
                    'bullet_2': 'Erta signal kechikishni kamaytiradi',
                    'bullet_3': 'Profilaktika doimiy bo‘lsa kuchliroq',
                    'bullet_4': 'Moslashgan javob natijani yaxshilaydi',
                    'discussion_1_heading': 'Talqin',
                    'discussion_1_text': 'Kasallanish ko‘rsatkichini xatar mantiqidan ajratib o‘qish noto‘g‘ri xulosaga olib keladi.',
                    'discussion_2_heading': 'Cheklov',
                    'discussion_2_text': 'Bir xil chora turli maydonlarda bir xil natija bermaydi.',
                }),
                self._slide_item('study_steps', {
                    'title': 'Tadqiqot jarayonining besh qadami',
                    'step_1_heading': 'Kuzatuv',
                    'step_1_text': 'Yuqori xatar nuqtalari belgilandi.',
                    'step_2_heading': 'Saralash',
                    'step_2_text': 'Asosiy omillar ustuvorlashtirildi.',
                    'step_3_heading': 'Qiyoslash',
                    'step_3_text': 'Hudud va vaziyatlar solishtirildi.',
                    'step_4_heading': 'Sinov',
                    'step_4_text': 'Javob choralari taʼsiri tekshirildi.',
                    'step_5_heading': 'Tavsiya',
                    'step_5_text': 'Amaliy va yengil choralar tanlandi.',
                }),
                self._slide_item('conclusions', {
                    'title': 'Asosiy xulosalar',
                    'point_1_heading': 'Xatar',
                    'point_1_text': 'Xatar qatlamli birikma sifatida namoyon bo‘ladi.',
                    'point_2_heading': 'Monitoring',
                    'point_2_text': 'Erta signal nazorat sifatini keskin oshiradi.',
                    'point_3_heading': 'Javob',
                    'point_3_text': 'Maqsadli chora umumiy cheklovdan kuchliroq ishlaydi.',
                    'point_4_heading': 'Amaliyot',
                    'point_4_text': 'Qatlamli yondashuv o‘lchanadigan foyda beradi.',
                }),
            ]
        }
        return self._sanitize_result(template, payload, raw)


    def _build_history_thesis_defense_fallback(
        self,
        template: MagicTemplateSpec,
        payload: dict[str, Any],
    ) -> list[dict[str, Any]]:
        variables = dict(payload.get('variables') or {})
        language = str(variables.get('language') or 'uz').strip().lower()
        if language != 'uz':
            return self._sanitize_result(template, payload, {'slides_content': []})

        topic = str(variables.get('topic') or template.template_name).strip()
        focus_period = str(variables.get('optional_focus_period') or '').strip() or 'asosiy tadqiqot davri'
        region = str(variables.get('optional_region') or '').strip() or 'asosiy tadqiqot hududi'

        raw = {
            'slides_content': [
                self._slide_item('cover', {
                    'title': topic,
                    'subtitle': f"{focus_period} va {region} kesimidagi tarixiy himoya",
                }),
                self._slide_item('table_of_contents', {
                    'toc_title': 'Mundarija',
                    'part_1_title': 'Masala',
                    'part_1_desc': 'mavzu va tarixiy doira',
                    'part_2_title': 'Farazlar',
                    'part_2_desc': 'tayanch tezislar',
                    'part_3_title': 'Talqin',
                    'part_3_desc': 'mavjud yondashuvlar',
                    'part_4_title': 'Usul',
                    'part_4_desc': 'ramka va jadval',
                    'part_5_title': 'Tahlil',
                    'part_5_desc': 'bosqich va natijalar',
                    'part_6_title': 'Xulosa',
                    'part_6_desc': 'meros va yakun',
                }),
                self._slide_item('section_divider_intro', {
                    'section_number': '01',
                    'section_title': 'Tadqiqot masalasi',
                    'section_subtitle': 'doira va tezis',
                }),
                self._slide_item('core_questions', {
                    'title': "Tadqiqotning tayanch savollari",
                    'point_1_heading': 'Tarixiy kontekst',
                    'point_1_text': f"{topic} {focus_period} davrida qaysi bosimlar ichida shakllangani ochiladi.",
                    'point_2_heading': 'Asosiy aktorlar',
                    'point_2_text': f"Jarayonni oldinga surgan guruhlar va institutlar {region} doirasida bog'lanadi.",
                    'point_3_heading': 'Meros',
                    'point_3_text': f"{topic}ning uzoq davom etgan ijtimoiy va madaniy ta'siri talqin qilinadi.",
                }),
                self._slide_item('image_break', {
                    'title': 'Tarixiy manzara',
                }),
                self._slide_item('visual_support', {
                    'title_line_1': 'Tarixiy',
                    'title_line_2': 'fon',
                    'body': f"{focus_period} davrining siyosiy, ma'rifiy va jamoaviy muhiti {topic}ning ichki mantiqini ochadi.",
                }),
                self._slide_item('quote_slide', {
                    'quote': f"{topic}ni tushunish uchun voqea bilan birga uni ko'targan ijtimoiy kayfiyatni ham o'qish kerak.",
                    'author': 'Tadqiqot tezisi',
                }),
                self._slide_item('hypotheses', {
                    'title': 'Asosiy farazlar',
                    'hypothesis_1_heading': 'Institutsional omil',
                    'hypothesis_1_text': f"Maktab, matbuot va jamoa kabi institutlar {topic} rivojiga tayanch bo'lgan.",
                    'hypothesis_2_heading': "Ijtimoiy ko'prik",
                    'hypothesis_2_text': "Jarayon faqat elita ichida qolmay, mahalliy qatlamlar bilan ko'prik hosil qilgan.",
                    'hypothesis_3_heading': 'Uzoq oqibat',
                    'hypothesis_3_text': f"{topic}ning fikriy va madaniy ta'siri keyingi davrlarga cho'zilgan.",
                }),
                self._slide_item('main_value', {
                    'main_value': '3 qatlam',
                    'caption': "bosim, aktor va meros kesimi",
                }),
                self._slide_item('stat_triplet', {
                    'stat_1_value': '2 yoʻnalish',
                    'stat_1_label': 'islohot va fikr',
                    'stat_2_value': '3 makon',
                    'stat_2_label': 'maktab, matbuot, jamoa',
                    'stat_3_value': '1 meros',
                    'stat_3_label': 'uzoq davom etgan taʼsir',
                }),
                self._slide_item('section_pause', {
                    'title': 'Tahlil bosqichi',
                }),
                self._slide_item('current_situation', {
                    'title': 'Mavjud talqin va muammolar',
                    'situation_heading': 'Hozirgi talqin',
                    'situation_text': f"{topic} bo'yicha mavjud talqinlar ko'pincha voqea markazida to'xtab qoladi, ammo uning ijtimoiy bazasi va uzoq taʼsiri yetarli darajada ochilmaydi.",
                    'problems_heading': "Asosiy bo'shliqlar",
                    'problem_1_heading': 'Darholik',
                    'problem_1_text': 'Jarayonning ichki tayanchlari yetarli ko‘rsatilmaydi.',
                    'problem_2_heading': 'Hudud',
                    'problem_2_text': 'Hududiy farqlar umumiy bayon ichida yoʻqolib ketadi.',
                    'problem_3_heading': 'Meros',
                    'problem_3_text': 'Uzoq oqibatlar voqeadan alohida tahlil qilinmaydi.',
                }),
                self._slide_item('study_objectives', {
                    'title': 'Tadqiqot vazifalari',
                    'objective_1_heading': 'Kontekstni tiklash',
                    'objective_1_text': 'Siyosiy va ijtimoiy bosimlarni tizimlashtirish.',
                    'objective_2_heading': "Aktorlarni bog'lash",
                    'objective_2_text': "Institutlar va ziyolilar rolini bir makonda ko'rsatish.",
                    'objective_3_heading': 'Merosni ochish',
                    'objective_3_text': "Qisqa hodisadan uzoq davom etgan taʼsirga o'tish.",
                }),
                self._slide_item('insights_grid', {
                    'title': "Tadqiqotni tutib turgan olti g'oya",
                    'insight_1_heading': 'Matbuot',
                    'insight_1_text': 'G‘oyani ommaga ko‘targan asosiy kanal bo‘ldi.',
                    'insight_2_heading': 'Maktab',
                    'insight_2_text': 'Yangi ijtimoiy tayanchni shakllantirdi.',
                    'insight_3_heading': 'Ziyolilar',
                    'insight_3_text': 'Talqinni siyosiy mazmunga bogʻladi.',
                    'insight_4_heading': 'Hudud',
                    'insight_4_text': 'Farqlar harakat yo‘nalishini o‘zgartirdi.',
                    'insight_5_heading': 'Meros',
                    'insight_5_text': 'Natija keyingi davrlarga uzatildi.',
                    'insight_6_heading': 'Talqin',
                    'insight_6_text': 'Voqea va gʻoya birgalikda o‘qiladi.',
                }),
                self._slide_item('literature_review', {
                    'title': 'Adabiyotlar sharhi',
                    'entry_1_ref': 'Tarixiy monografiyalar',
                    'entry_1_note': 'Siyosiy fonni izohlashda tayanch bo‘ladi.',
                    'entry_2_ref': 'Matbuot tahlillari',
                    'entry_2_note': 'Jamoatchilik tilini va tezislarni ko‘rsatadi.',
                    'entry_3_ref': "Ta'lim tarixchilari",
                    'entry_3_note': 'Institutlarning ijtimoiy rolini ochib beradi.',
                    'entry_4_ref': 'Hududiy tadqiqotlar',
                    'entry_4_note': 'Mahalliy farqlarni ko‘rsatadi.',
                    'entry_5_ref': 'Madaniy xotira ishlari',
                    'entry_5_note': 'Merosning uzoq davom etishini tahlil qiladi.',
                }),
                self._slide_item('theoretical_framework', {
                    'title': 'Nazariy ramka',
                    'framework_kicker': 'Tarixiy o‘qish',
                    'key_terms_heading': 'Tayanch tushunchalar',
                    'key_term_1': 'modernlashuv',
                    'key_term_2': 'jamoaviy xotira',
                    'key_term_3': "ma'rifiy harakat",
                    'theories_heading': 'Yondashuvlar',
                    'framework_heading': 'Bizning yondashuv',
                    'framework_text': f"{topic} voqea sifatida emas, balki institutlar, fikr ishlab chiqarish va ijtimoiy qabul o'rtasidagi tarixiy aloqa sifatida talqin qilinadi.",
                    'theory_1_heading': 'Modernlashuv nazariyasi',
                    'theory_1_text': "Islohot va ta'limning tizimli rolini tushuntiradi.",
                    'theory_2_heading': 'Madaniy xotira',
                    'theory_2_text': 'Merosning keyingi davrlarda qanday yashashini ko‘rsatadi.',
                }),
                self._slide_item('schedule_table', {
                    'title': 'Tadqiqot jadvali',
                    'header_task': 'Vazifa',
                    'header_description': 'Tavsif',
                    'header_date': 'Muddat',
                    'header_status': 'Holat',
                    'row_1_task': "Manba yig'ish",
                    'row_1_description': 'birlamchi va ikkilamchi manbalarni saralash',
                    'row_1_date': '1-oy',
                    'row_1_status': 'yakun',
                    'row_2_task': 'Kontekst xaritasi',
                    'row_2_description': "davr va hudud bo'yicha tahlil karkasi",
                    'row_2_date': '2-oy',
                    'row_2_status': 'yakun',
                    'row_3_task': "Faraz tekshiruvi",
                    'row_3_description': 'asosiy aktorlar va kuchlarni qiyoslash',
                    'row_3_date': '3-oy',
                    'row_3_status': 'jarayon',
                    'row_4_task': 'Talqin yozuvi',
                    'row_4_description': "boblararo dalillarni bog'lash",
                    'row_4_date': '4-oy',
                    'row_4_status': 'reja',
                    'row_5_task': 'Himoya tayyorgarligi',
                    'row_5_description': "yakuniy xulosa va ko'rgazma",
                    'row_5_date': '5-oy',
                    'row_5_status': 'navbat',
                }),
                self._slide_item('methodology', {
                    'title': 'Metodologiya',
                    'pillar_1_heading': 'Usul',
                    'pillar_1_text': 'Tarixiy qiyos va manba tahlili birlashtirildi.',
                    'pillar_2_heading': "Ma'lumot",
                    'pillar_2_text': 'Matbuot, arxiv va ikkilamchi ishlardan foydalanildi.',
                    'pillar_3_heading': 'Motiv',
                    'pillar_3_text': 'Jarayon ichki tayanchlarini ajratish maqsad qilindi.',
                    'pillar_4_heading': "Yig'ish",
                    'pillar_4_text': "Manbalar davr, hudud va aktor bo'yicha saralandi.",
                    'pillar_5_heading': 'Tanlanma',
                }),
                self._slide_item('analysis_phases', {
                    'title': 'Tahlilning ikki bosqichi',
                    'phase_1_heading': 'Kontekstni qurish',
                    'phase_1_point_1': 'Davr bosimlari ajratildi',
                    'phase_1_point_2': "Aktorlar bog'lanishi chizildi",
                    'phase_1_point_3': 'Hududiy farqlar solishtirildi',
                    'phase_2_heading': 'Merosni talqin qilish',
                    'phase_2_point_1': "Qisqa natijalar qayta o'qildi",
                    'phase_2_point_2': "Uzoq ta'sir yo'nalishlari belgilandi",
                    'phase_2_point_3': 'Talqin va fakt birlashtirildi',
                }),
                self._slide_item('analysis_graph', {
                    'title': 'Tahliliy taqsimot',
                    'chart_caption': "Asosiy kuchlarning nisbiy og'irligi va o'zaro bogʻlanishi ko'rsatiladi.",
                    'note_1_heading': 'Matbuot',
                    'note_1_text': "G'oyani tez yoygan.",
                    'note_2_heading': 'Meros',
                    'note_2_text': 'Keyingi davrni bog‘lagan.',
                }),
                self._slide_item('results_summary', {
                    'title': 'Natijalar tahlili',
                    'result_1_heading': 'Kontekst',
                    'result_1_text': 'Voqea murakkab tarixiy bosimlar ichida yuz bergan.',
                    'result_2_heading': 'Aktorlar',
                    'result_2_text': "Institut va ziyolilar birgalikda yo'nalish bergan.",
                    'result_3_heading': 'Meros',
                    'result_3_text': "Ta'sir keyingi ijtimoiy va madaniy qatlamlarga o'tgan.",
                    'highlight_value': '65%',
                }),
                self._slide_item('regional_map', {
                    'title': 'Hududiy tayanch nuqtalar',
                    'location_1_name': 'Markaz',
                    'location_1_text': "g'oya va matbuot markazi",
                    'location_2_name': 'Hudud',
                    'location_2_text': "ta'lim va ijtimoiy ta'sir nuqtasi",
                    'location_3_name': 'Chekka',
                    'location_3_text': "an'ana va islohot kesishgan makon",
                }),
                self._slide_item('discussion', {
                    'title': 'Muhokama',
                    'intro_heading': 'Asosiy talqin',
                    'intro_text': f"{topic} bo'yicha bahs markazida voqeaning o'zi emas, uni tayyorlagan ijtimoiy zamin va undan keyin yashab qolgan meros turadi.",
                    'discussion_1_heading': 'Kuchli tomon',
                    'discussion_1_text': "Tadqiqot voqea, aktor va merosni bitta mantiqda bog'laydi.",
                    'discussion_2_heading': 'Bahsli nuqta',
                    'discussion_2_text': "Hududiy farqlarni chuqurroq ochish keyingi bosqich uchun muhim.",
                }),
                self._slide_item('conclusions', {
                    'title': 'Asosiy xulosalar',
                    'point_1_heading': 'Kontekst',
                    'point_1_text': 'Mavzu siyosiy va ijtimoiy bosimlar ichida shakllangan.',
                    'point_2_heading': 'Aktorlar',
                    'point_2_text': 'Institutlar va ziyolilar jarayonni birga harakatlantirgan.',
                    'point_3_heading': 'Hudud',
                    'point_3_text': 'Hududiy farqlar umumiy talqinni chuqurlashtiradi.',
                    'point_4_heading': 'Meros',
                    'point_4_text': "Mavzuning ta'siri keyingi davrlarda ham davom etgan.",
                }),
                self._slide_item('bibliography', {
                    'title': 'Qisqa bibliografiya',
                    'entry_1': "Tarixiy monografiyalar to'plami",
                    'entry_2': 'Davr matbuoti va sharhlar',
                    'entry_3': "Ta'lim va islohot tarixiga oid tadqiqotlar",
                    'entry_4': 'Hududiy tarix ishlari',
                    'entry_5': 'Madaniy xotira va meros adabiyoti',
                    'entry_6': 'Biografik manbalar va maktublar',
                    'entry_7': 'Ijtimoiy tarix bo‘yicha maqolalar',
                    'entry_8': 'Taqriz va sharhlar',
                    'entry_9': 'Arxivga tayanuvchi umumiy ishlar',
                    'entry_10': 'Merosga oid tahliliy nashrlar',
                }),
            ]
        }
        return self._sanitize_result(template, payload, raw)


class MagicPptxRenderer:
    def __init__(self, *, output_dir: str | Path | None = None) -> None:
        base_dir = Path(output_dir) if output_dir else Path(tempfile.gettempdir()) / 'slide_generator_magic_outputs'
        base_dir.mkdir(parents=True, exist_ok=True)
        self.output_dir = base_dir
        self._active_language = 'uz'

    def render(
        self,
        *,
        template: MagicTemplateSpec,
        order_payload: dict[str, Any],
        content_plan: list[dict[str, Any]],
    ) -> str:
        presentation = Presentation(template.template_path)
        content_by_slide_index = {int(item['slide_index']): item for item in content_plan}
        self._active_language = str((order_payload.get('variables') or {}).get('language') or 'uz').strip().lower()

        for slide_spec in template.generated_slides:
            slide_index = int(slide_spec['index']) - 1
            slide = presentation.slides[slide_index]
            slide_content = content_by_slide_index.get(int(slide_spec['index']), {})
            content = dict(slide_content.get('content') or {})
            self._fill_slide(
                slide=slide,
                slide_spec=slide_spec,
                content=content,
            )
            self._cleanup_placeholder_shapes(slide, content=content)

        delete_indexes = [
            int(slide['index']) - 1
            for slide in template.slides
            if not slide.get('generate') and (slide.get('skip_in_final') or slide.get('optional_variant'))
        ]
        for index in sorted(delete_indexes, reverse=True):
            self._delete_slide(presentation, index)
        for index in range(len(presentation.slides) - 1, -1, -1):
            if self._is_placeholder_only_slide(presentation.slides[index]):
                self._delete_slide(presentation, index)

        variables = dict(order_payload.get('variables') or {})
        template_name = str(order_payload.get('template_name') or template.template_name)
        topic = str(variables.get('topic') or template_name).strip()
        safe_name = self._safe_filename(f"{template_name}_{topic[:80]}")
        file_path = self.output_dir / f'{safe_name}.pptx'
        presentation.save(file_path)
        return str(file_path)

    def _fill_slide(self, *, slide, slide_spec: dict[str, Any], content: dict[str, Any]) -> None:
        render_strategy = str(slide_spec.get('render_strategy') or 'token_replace').strip().lower()
        if render_strategy != 'token_replace':
            slide_id = str(slide_spec.get('id') or '?')
            raise ValueError(f"Qo'llab-quvvatlanmaydigan render strategy: {slide_id} -> {render_strategy}")
        self._fill_tokenized_slide(slide=slide, slide_spec=slide_spec, content=content)
        self._cleanup_placeholder_shapes(slide, content=content)

    def _fill_tokenized_slide(self, *, slide, slide_spec: dict[str, Any], content: dict[str, Any]) -> None:
        slot_constraints = dict(slide_spec.get('slot_constraints') or {})
        token_values = {
            self._slot_token(slot_name): self._stringify_value(content.get(slot_name))
            for slot_name in (slide_spec.get('content_slots') or {}).keys()
        }
        self._replace_tokens_in_shapes(slide.shapes, token_values=token_values, slot_constraints=slot_constraints)

    @staticmethod
    def _slot_token(slot_name: str) -> str:
        return f'[[{slot_name}]]'

    def _replace_tokens_in_shapes(self, shapes, *, token_values: dict[str, str], slot_constraints: dict[str, Any]) -> None:
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self._replace_tokens_in_shapes(shape.shapes, token_values=token_values, slot_constraints=slot_constraints)
                continue
            if getattr(shape, 'has_table', False):
                self._replace_tokens_in_table(shape.table, token_values=token_values, slot_constraints=slot_constraints)
            if getattr(shape, 'has_text_frame', False):
                matched_slots = self._replace_tokens_in_text_frame(shape.text_frame, token_values)
                if matched_slots:
                    self._apply_text_fit(
                        shape.text_frame,
                        matched_slots=matched_slots,
                        slot_constraints=slot_constraints,
                        container_shape=shape,
                    )

    def _replace_tokens_in_table(self, table, *, token_values: dict[str, str], slot_constraints: dict[str, Any]) -> None:
        for row in table.rows:
            for cell in row.cells:
                matched_slots = self._replace_tokens_in_text_frame(cell.text_frame, token_values)
                if matched_slots:
                    self._apply_text_fit(
                        cell.text_frame,
                        matched_slots=matched_slots,
                        slot_constraints=slot_constraints,
                        container_shape=None,
                    )

    def _replace_tokens_in_text_frame(self, text_frame, token_values: dict[str, str]) -> list[str]:
        matched_slots: list[str] = []
        for paragraph in text_frame.paragraphs:
            token_name = str(paragraph.text or '').strip()
            if not token_name or token_name not in token_values:
                continue
            self._write_paragraph(paragraph, token_values[token_name])
            matched_slots.append(token_name[2:-2])
        return matched_slots

    def _apply_text_fit(
        self,
        text_frame,
        *,
        matched_slots: Sequence[str],
        slot_constraints: dict[str, Any],
        container_shape: Any | None,
    ) -> None:
        strategy = 'shrink'
        for slot_name in matched_slots:
            config = slot_constraints.get(slot_name) or {}
            slot_strategy = str(config.get('fit') or '').strip().lower()
            if slot_strategy:
                strategy = slot_strategy
                break

        text_frame.word_wrap = True
        self._compact_text_frame_spacing(text_frame, matched_slots=matched_slots, container_shape=container_shape)
        if strategy == 'expand':
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            return
        if strategy == 'none':
            text_frame.auto_size = MSO_AUTO_SIZE.NONE
            return
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        self._fit_text_frame(
            text_frame,
            matched_slots=matched_slots,
            slot_constraints=slot_constraints,
            container_shape=container_shape,
        )

    def _fit_text_frame(
        self,
        text_frame,
        *,
        matched_slots: Sequence[str],
        slot_constraints: dict[str, Any],
        container_shape: Any | None,
    ) -> None:
        detected_size = self._detect_text_frame_max_size(text_frame)
        suggested_size = self._suggest_text_frame_max_size(
            text_frame,
            matched_slots=matched_slots,
            slot_constraints=slot_constraints,
            container_shape=container_shape,
        )
        max_size = min(detected_size, suggested_size)
        font_family = self._detect_text_frame_font_family(text_frame)
        if max_size <= 0:
            return
        try:
            if font_family:
                text_frame.fit_text(max_size=max_size, font_family=font_family)
            else:
                text_frame.fit_text(max_size=max_size)
        except Exception:
            pass

    @staticmethod
    def _detect_text_frame_max_size(text_frame) -> int:
        sizes: list[int] = []
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                size = getattr(run.font, 'size', None)
                if size is not None and getattr(size, 'pt', None):
                    sizes.append(int(round(size.pt)))
        return max(sizes) if sizes else 28

    @staticmethod
    def _detect_text_frame_font_family(text_frame) -> str | None:
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                name = str(getattr(run.font, 'name', '') or '').strip()
                if name:
                    return name
        return None

    def _suggest_text_frame_max_size(
        self,
        text_frame,
        *,
        matched_slots: Sequence[str],
        slot_constraints: dict[str, Any],
        container_shape: Any | None,
    ) -> int:
        caps = [self._slot_font_cap(slot_name) for slot_name in matched_slots if self._slot_font_cap(slot_name) > 0]
        suggested = min(caps) if caps else self._detect_text_frame_max_size(text_frame)

        non_empty = [str(paragraph.text or '').strip() for paragraph in text_frame.paragraphs if str(paragraph.text or '').strip()]
        char_count = sum(len(item) for item in non_empty)
        paragraph_count = max(1, len(non_empty))

        if paragraph_count >= 4:
            suggested = min(suggested, 14)
        elif paragraph_count == 3:
            suggested = min(suggested, 16)
        elif paragraph_count == 2:
            suggested = min(suggested, 18)

        if char_count > 180:
            suggested = min(suggested, 13)
        elif char_count > 130:
            suggested = min(suggested, 14)
        elif char_count > 90:
            suggested = min(suggested, 16)
        elif char_count > 60:
            suggested = min(suggested, 18)

        if container_shape is not None:
            width_pt = int(getattr(container_shape, 'width', 0) / 12_700)
            height_pt = int(getattr(container_shape, 'height', 0) / 12_700)
            if width_pt and width_pt < 220:
                suggested = min(suggested, 18)
            if width_pt and width_pt < 170:
                suggested = min(suggested, 16)
            if width_pt and width_pt < 140:
                suggested = min(suggested, 14)
            if height_pt and height_pt < 100:
                suggested = min(suggested, 17)
            if height_pt and height_pt < 72:
                suggested = min(suggested, 15)
            if width_pt and height_pt and width_pt < 180 and height_pt < 90 and paragraph_count >= 2:
                suggested = min(suggested, 13)

        for slot_name in matched_slots:
            config = slot_constraints.get(slot_name) or {}
            if config.get('fit') == 'expand':
                suggested = max(suggested, 18)

        return max(11, suggested)

    @staticmethod
    def _slot_font_cap(slot_name: str) -> int:
        slot = str(slot_name or '').strip().lower()
        if not slot:
            return 28
        if slot == 'title':
            return 28
        if slot.endswith('_title'):
            return 24
        if slot.endswith(('_heading', '_label', '_name')):
            return 18
        if slot.endswith(('_subtitle', '_desc', '_caption')):
            return 16
        if slot.startswith('bullet_') or slot in {'intro_text', 'caption'}:
            return 15
        if slot.endswith(('_text', '_body')):
            return 15
        if slot.endswith(('_value', '_percent', '_number', '_year')):
            return 20
        return 18

    def _compact_text_frame_spacing(
        self,
        text_frame,
        *,
        matched_slots: Sequence[str],
        container_shape: Any | None,
    ) -> None:
        non_empty = [str(paragraph.text or '').strip() for paragraph in text_frame.paragraphs if str(paragraph.text or '').strip()]
        char_count = sum(len(item) for item in non_empty)
        paragraph_count = max(1, len(non_empty))
        width_pt = int(getattr(container_shape, 'width', 0) / 12_700) if container_shape is not None else 0
        height_pt = int(getattr(container_shape, 'height', 0) / 12_700) if container_shape is not None else 0
        crowded = (
            paragraph_count >= 3
            or char_count > 80
            or (width_pt and width_pt < 180)
            or (height_pt and height_pt < 90)
            or any(slot.startswith('bullet_') for slot in matched_slots)
        )

        if crowded:
            try:
                text_frame.margin_left = Pt(3)
                text_frame.margin_right = Pt(3)
                text_frame.margin_top = Pt(1)
                text_frame.margin_bottom = Pt(1)
            except Exception:
                pass

        for paragraph in text_frame.paragraphs:
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0 if crowded else 1)
            if crowded:
                paragraph.line_spacing = 0.92 if paragraph_count >= 3 else 0.96

    def _fill_toc_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        shapes = self._collect_text_shapes(slide.shapes)
        title_candidates: list[ShapeInfo] = []
        desc_candidates: list[ShapeInfo] = []
        toc_title_shapes: list[ShapeInfo] = []

        for info in shapes:
            normalized = self._normalize_text(info.text)
            if not normalized:
                continue
            if normalized == 'table of contents':
                toc_title_shapes.append(info)
                continue
            if re.fullmatch(r'\d+', normalized):
                continue
            if normalized in DESC_PLACEHOLDER_TEXTS:
                desc_candidates.append(info)
            else:
                title_candidates.append(info)

        if not title_candidates:
            return False

        title_candidates.sort(key=lambda item: (item.top, item.left))
        desc_candidates.sort(key=lambda item: (item.top, item.left))

        title_values = [content.get(key, '') for key in slot_keys if key.endswith('_title') or key.startswith('part_')]
        desc_values = [content.get(key, '') for key in slot_keys if key.endswith('_desc')]

        for info in toc_title_shapes:
            self._write_shape_from_value(info.shape, self._localized_toc_title())
        for info, value in zip(title_candidates, title_values):
            self._write_shape_paragraphs(info.shape, [self._stringify_value(value)])
        for info, value in zip(desc_candidates, desc_values):
            self._write_shape_paragraphs(info.shape, [self._stringify_value(value)])
        return True

    def _fill_table_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        table_shape = self._first_table_shape(slide.shapes)
        if table_shape is None:
            return False

        text_slot_keys = [key for key in slot_keys if key not in {'rows', 'columns', 'table_headers'}]
        self._fill_generic_slide(slide, text_slot_keys, content)

        headers = list(content.get('columns') or content.get('table_headers') or [])
        rows = list(content.get('rows') or [])
        table = table_shape.table

        if headers:
            for column_index, value in enumerate(headers[: len(table.columns)]):
                self._write_cell(table.cell(0, column_index), self._stringify_value(value))

        body_start_row = 1 if len(table.rows) > 1 else 0
        for row_index in range(body_start_row, len(table.rows)):
            values = rows[row_index - body_start_row] if row_index - body_start_row < len(rows) else []
            if not isinstance(values, list):
                values = [values]
            for column_index in range(len(table.columns)):
                cell_value = values[column_index] if column_index < len(values) else ''
                self._write_cell(table.cell(row_index, column_index), self._stringify_value(cell_value))
        return True

    def _fill_specialized_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        handlers = (
            self._fill_title_only_slide,
            self._fill_quote_slide,
            self._fill_large_numeric_callout_slide,
            self._fill_three_stat_cards_slide,
            self._fill_four_point_conclusion_slide,
            self._fill_three_hypotheses_slide,
            self._fill_three_point_statement_slide,
            self._fill_split_milestones_slide,
            self._fill_percentage_distribution_slide,
            self._fill_kpi_slide,
            self._fill_platform_guide_slide,
            self._fill_dos_and_donts_slide,
            self._fill_subject_task_grid_slide,
            self._fill_materials_checklist_slide,
            self._fill_review_checklist_slide,
            self._fill_five_step_process_slide,
            self._fill_three_step_planning_slide,
            self._fill_people_slide,
            self._fill_phase_points_slide,
            self._fill_framework_slide,
            self._fill_literature_review_five_slide,
            self._fill_methodology_matrix_slide,
            self._fill_discussion_two_columns_slide,
            self._fill_five_column_methodology_slide,
            self._fill_chart_plus_notes_slide,
            self._fill_three_way_result_analysis_slide,
            self._fill_geo_context_map_slide,
            self._fill_heading_text_grid_slide,
            self._fill_label_text_matrix_slide,
            self._fill_rows_without_table_slide,
        )
        for handler in handlers:
            if handler(slide, slot_keys, content):
                return True
        return False

    def _fill_title_only_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        non_meta_keys = [key for key in slot_keys if key not in {'title', 'subtitle', 'author'}]
        if non_meta_keys or 'title' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        candidates.sort(key=self._shape_sort_key)
        self._write_shape_from_value(candidates[0].shape, content.get('title'))
        if len(candidates) > 1 and 'subtitle' in content:
            self._write_shape_from_value(candidates[1].shape, content.get('subtitle'))
        if len(candidates) > 2 and 'author' in content:
            self._write_shape_from_value(candidates[2].shape, content.get('author'))
        for info in candidates[1 + int('subtitle' in content) + int('author' in content):]:
            self._clear_shape(info.shape)
        return True

    def _fill_quote_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'quote' not in slot_keys or 'author' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if len(candidates) < 2:
            return False

        quote_shape = max(candidates, key=lambda item: len(item.text))
        author_candidates = [info for info in candidates if info is not quote_shape]
        author_shape = min(author_candidates, key=lambda item: (len(item.text), item.top, item.left))

        self._write_shape_from_value(quote_shape.shape, content.get('quote'))
        self._write_shape_from_value(author_shape.shape, content.get('author'))
        for info in author_candidates:
            if info is not author_shape and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_large_numeric_callout_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'main_value' not in slot_keys or 'caption' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        if 'title' in slot_keys:
            self._consume_title_shape(candidates, content)
        if len(candidates) < 2:
            return False

        value_shape = max(candidates, key=lambda info: (getattr(info.shape, 'height', 0), len(info.text)))
        caption_candidates = [info for info in candidates if info is not value_shape]
        if not caption_candidates:
            return False
        caption_shape = max(caption_candidates, key=lambda info: (len(info.text), getattr(info.shape, 'width', 0)))

        self._write_shape_from_value(value_shape.shape, content.get('main_value'))
        self._write_shape_from_value(caption_shape.shape, content.get('caption'))

        for info in candidates:
            if info not in [value_shape, caption_shape] and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_three_stat_cards_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        required = {'stat_1_value', 'stat_1_label', 'stat_2_value', 'stat_2_label', 'stat_3_value', 'stat_3_label'}
        if not required.issubset(set(slot_keys)):
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        if 'title' in slot_keys:
            self._consume_title_shape(candidates, content)
        non_empty = [info for info in candidates if self._normalize_text(info.text)]
        if len(non_empty) < 6:
            return False

        single_column = max(info.left for info in non_empty) - min(info.left for info in non_empty) < 1_000_000
        if single_column:
            ordered = sorted(non_empty, key=self._shape_sort_key)
            value_shapes = ordered[::2][:3]
            label_shapes = ordered[1::2][:3]
        else:
            value_shapes = sorted(
                [info for info in non_empty if self._looks_like_value_shape(info) or len(info.text) <= 14],
                key=self._shape_sort_key,
            )[:3]
            label_shapes = sorted([info for info in non_empty if info not in value_shapes], key=self._shape_sort_key)[:3]

        if len(value_shapes) < 3 or len(label_shapes) < 3:
            return False

        for index, shape in enumerate(value_shapes, start=1):
            self._write_shape_from_value(shape.shape, content.get(f'stat_{index}_value'))
        for index, shape in enumerate(label_shapes, start=1):
            self._write_shape_from_value(shape.shape, content.get(f'stat_{index}_label'))

        for info in non_empty:
            if info not in value_shapes and info not in label_shapes and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_three_point_statement_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        required = {
            'point_1_heading',
            'point_1_text',
            'point_2_heading',
            'point_2_text',
            'point_3_heading',
            'point_3_text',
        }
        if not required.issubset(set(slot_keys)):
            return False
        return self._fill_card_pairs_slide(
            slide=slide,
            content=content,
            heading_keys=[f'point_{index}_heading' for index in range(1, 4)],
            text_keys=[f'point_{index}_text' for index in range(1, 4)],
        )

    def _fill_three_hypotheses_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        required = {
            'hypothesis_1_heading',
            'hypothesis_1_text',
            'hypothesis_2_heading',
            'hypothesis_2_text',
            'hypothesis_3_heading',
            'hypothesis_3_text',
        }
        if not required.issubset(set(slot_keys)):
            return False
        return self._fill_card_pairs_slide(
            slide=slide,
            content=content,
            heading_keys=[f'hypothesis_{index}_heading' for index in range(1, 4)],
            text_keys=[f'hypothesis_{index}_text' for index in range(1, 4)],
            clear_numbers=True,
        )

    def _fill_platform_guide_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'part_1_heading' not in slot_keys or 'part_1_steps' not in slot_keys or 'part_2_steps' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        title_shape = self._consume_title_shape(candidates, content)
        if title_shape is None and candidates:
            title_shape = min(candidates, key=self._shape_sort_key)
            self._write_shape_from_value(title_shape.shape, content.get('title'))
            candidates.remove(title_shape)
        if not candidates:
            return False

        body_shape = max(
            candidates,
            key=lambda info: (getattr(info.shape, 'height', 0), info.paragraph_count, len(info.text)),
        )
        paragraphs: list[str] = []
        part_1_heading = self._stringify_value(content.get('part_1_heading'))
        if part_1_heading:
            paragraphs.append(part_1_heading)
        paragraphs.extend(self._coerce_string_list(content.get('part_1_steps')))
        part_2_heading = self._stringify_value(content.get('part_2_heading'))
        if part_2_heading:
            paragraphs.append(part_2_heading)
        paragraphs.extend(self._coerce_string_list(content.get('part_2_steps')))
        self._write_shape_paragraphs(body_shape.shape, paragraphs or [''])

        for info in candidates:
            if info is not body_shape and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_dos_and_donts_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'ok_items' not in slot_keys or 'not_ok_items' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        title_shape = self._consume_title_shape(candidates, content)
        if title_shape is None:
            return False

        text_candidates = [info for info in candidates if self._normalize_text(info.text)]
        if len(text_candidates) < 4:
            return False

        label_shapes = sorted(text_candidates, key=lambda info: (info.top, info.left))[:2]
        body_shapes = sorted(
            [info for info in text_candidates if info not in label_shapes],
            key=lambda info: (info.left, info.top),
        )[:2]
        label_shapes.sort(key=lambda info: info.left)
        body_shapes.sort(key=lambda info: info.left)
        if len(label_shapes) < 2 or len(body_shapes) < 2:
            return False

        self._write_shape_from_value(label_shapes[0].shape, content.get('ok_label'))
        self._write_shape_from_value(label_shapes[1].shape, content.get('not_ok_label'))
        self._write_shape_paragraphs(body_shapes[0].shape, self._coerce_string_list(content.get('ok_items')))
        self._write_shape_paragraphs(body_shapes[1].shape, self._coerce_string_list(content.get('not_ok_items')))

        for info in text_candidates:
            if info not in label_shapes and info not in body_shapes and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_subject_task_grid_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'tasks_heading' not in slot_keys or 'subjects_heading' not in slot_keys or 'items' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        title_shape = self._consume_title_shape(candidates, content)
        if title_shape is None:
            return False

        text_candidates = [info for info in candidates if self._normalize_text(info.text)]
        if len(text_candidates) < 4:
            return False

        heading_shapes = sorted(text_candidates, key=lambda info: (info.top, info.left))[:2]
        body_shapes = sorted(
            [info for info in text_candidates if info not in heading_shapes],
            key=lambda info: (info.top, info.left),
        )[-2:]
        heading_shapes.sort(key=lambda info: info.left)
        body_shapes.sort(key=lambda info: info.left)
        if len(heading_shapes) < 2 or len(body_shapes) < 2:
            return False

        items = [item for item in content.get('items') or [] if isinstance(item, dict)]
        task_lines = [self._stringify_value(item.get('task')) for item in items if self._stringify_value(item.get('task'))]
        subject_lines = [
            self._stringify_value(item.get('subject')) for item in items if self._stringify_value(item.get('subject'))
        ]

        self._write_shape_from_value(heading_shapes[0].shape, content.get('tasks_heading'))
        self._write_shape_from_value(heading_shapes[1].shape, content.get('subjects_heading'))
        self._write_shape_paragraphs(body_shapes[0].shape, task_lines or [''])
        self._write_shape_paragraphs(body_shapes[1].shape, subject_lines or [''])

        for info in text_candidates:
            if (
                info not in heading_shapes
                and info not in body_shapes
                and (
                    self._looks_like_source_placeholder(info.text)
                    or self._normalize_text(info.text) in {'english', 'art', 'french', 'math'}
                )
            ):
                self._clear_shape(info.shape)
        return True

    def _fill_materials_checklist_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'materials' not in slot_keys or 'note_label' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        title_shape = self._consume_title_shape(candidates, content)
        if title_shape is None:
            return False

        text_candidates = [info for info in candidates if self._normalize_text(info.text)]
        if len(text_candidates) < 2:
            return False

        heading_shape = min(text_candidates, key=self._shape_sort_key)
        item_shapes = sorted([info for info in text_candidates if info is not heading_shape], key=self._shape_sort_key)
        materials = self._coerce_string_list(content.get('materials'))
        if not item_shapes:
            return False

        self._write_shape_from_value(heading_shape.shape, content.get('note_label'))
        for shape, value in zip(item_shapes, materials):
            self._write_shape_from_value(shape.shape, value)
        for shape in item_shapes[len(materials):]:
            self._clear_shape(shape.shape)
        return True

    def _fill_review_checklist_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'items' not in slot_keys or 'title' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if len(candidates) < 2:
            return False

        candidates.sort(key=self._shape_sort_key)
        title_shape = candidates[0]
        body_shape = max(candidates[1:], key=lambda info: (getattr(info.shape, 'height', 0), len(info.text)))
        self._write_shape_from_value(title_shape.shape, content.get('title'))
        self._write_shape_paragraphs(body_shape.shape, self._coerce_string_list(content.get('items')) or [''])

        for info in candidates[1:]:
            if info is not body_shape and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_three_step_planning_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        required = {
            'step_1_heading',
            'step_1_text',
            'step_2_heading',
            'step_2_text',
            'step_3_heading',
            'step_3_text',
        }
        if not required.issubset(set(slot_keys)):
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        title_shape = self._consume_title_shape(candidates, content)
        if title_shape is None:
            return False

        text_candidates = [info for info in candidates if self._normalize_text(info.text)]
        heading_shapes = sorted(
            [info for info in text_candidates if getattr(info.shape, 'height', 0) <= 500_000],
            key=lambda info: info.left,
        )
        body_shapes = sorted(
            [info for info in text_candidates if getattr(info.shape, 'height', 0) > 500_000],
            key=lambda info: info.left,
        )
        if len(heading_shapes) < 3 or len(body_shapes) < 3:
            return False

        for index, shape in enumerate(heading_shapes[:3], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'step_{index}_heading'))
        for index, shape in enumerate(body_shapes[:3], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'step_{index}_text'))
        return True

    def _fill_split_milestones_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'left_body' not in slot_keys or 'right_month_1' not in slot_keys or 'right_text_1' not in slot_keys:
            return False

        candidates = self._collect_non_static_shapes(slide, slot_keys)
        title_shape = self._consume_title_shape(candidates, content)
        if title_shape is None:
            return False

        left_candidates = sorted(
            [info for info in candidates if info.left < title_shape.left + 3_000_000],
            key=self._shape_sort_key,
        )
        right_candidates = sorted(
            [info for info in candidates if info.left >= title_shape.left + 3_000_000],
            key=self._shape_sort_key,
        )
        if not left_candidates or len(right_candidates) < 6:
            return False

        heading_like = [info for info in left_candidates if not self._looks_like_body_shape(info)]
        body_like = [info for info in left_candidates if self._looks_like_body_shape(info)]
        if 'left_intro' in content and heading_like:
            self._write_shape_from_value(heading_like[0].shape, content.get('left_intro'))
        if body_like:
            self._write_shape_from_value(body_like[0].shape, content.get('left_body'))

        month_shapes = [info for info in right_candidates if not self._looks_like_body_shape(info)]
        text_shapes = [info for info in right_candidates if self._looks_like_body_shape(info)]
        if len(month_shapes) < 3 or len(text_shapes) < 3:
            return False

        month_shapes.sort(key=self._shape_sort_key)
        text_shapes.sort(key=self._shape_sort_key)
        used_shapes: list[ShapeInfo] = []
        for index, shape in enumerate(month_shapes[:3], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'right_month_{index}'))
            used_shapes.append(shape)
        for index, shape in enumerate(text_shapes[:3], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'right_text_{index}'))
            used_shapes.append(shape)
        if heading_like:
            used_shapes.append(heading_like[0])
        if body_like:
            used_shapes.append(body_like[0])
        for info in candidates:
            if info not in used_shapes and (info in right_candidates or self._looks_like_source_placeholder(info.text)):
                self._clear_shape(info.shape)
        return True

    def _fill_percentage_distribution_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'segment_1_name' not in slot_keys or 'segment_1_percent' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        self._consume_title_shape(candidates, content)
        if not candidates:
            return False

        value_shapes = sorted([info for info in candidates if self._looks_like_value_shape(info)], key=self._shape_sort_key)
        text_shapes = sorted([info for info in candidates if self._looks_like_body_shape(info)], key=self._shape_sort_key)
        name_shapes = sorted(
            [info for info in candidates if info not in value_shapes and info not in text_shapes],
            key=self._shape_sort_key,
        )
        if len(value_shapes) < 3 or len(name_shapes) < 3 or len(text_shapes) < 3:
            return False

        used_shapes: list[ShapeInfo] = []
        for index, shape in enumerate(name_shapes[:3], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'segment_{index}_name'))
            used_shapes.append(shape)
        for index, shape in enumerate(value_shapes[:3], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'segment_{index}_percent'))
            used_shapes.append(shape)
        for index, shape in enumerate(text_shapes[:3], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'segment_{index}_text'))
            used_shapes.append(shape)
        for info in candidates:
            if info not in used_shapes and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_kpi_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'kpi_1_value' not in slot_keys or 'kpi_1_text' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        self._consume_title_shape(candidates, content)
        if not candidates:
            return False

        value_shapes = sorted([info for info in candidates if self._looks_like_value_shape(info)], key=self._shape_sort_key)
        text_shapes = sorted([info for info in candidates if self._looks_like_body_shape(info)], key=self._shape_sort_key)
        if len(value_shapes) < 5 or len(text_shapes) < 5:
            return False

        used_shapes: list[ShapeInfo] = []
        for index, shape in enumerate(value_shapes[:5], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'kpi_{index}_value'))
            used_shapes.append(shape)
        for index, shape in enumerate(text_shapes[:5], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'kpi_{index}_text'))
            used_shapes.append(shape)
        for info in candidates:
            if info not in used_shapes and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_people_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        name_slots = [key for key in slot_keys if re.fullmatch(r'person_\d+_name', key)]
        if not name_slots:
            return False

        candidates = [
            info
            for info in self._collect_text_shapes(slide.shapes)
            if self._normalize_text(info.text) not in STATIC_TEXT_EXACT
        ]
        self._consume_title_shape(candidates, content)
        if not candidates:
            return False

        clusters = self._cluster_shapes_by_left(candidates)
        if len(clusters) < len(name_slots):
            return False

        for index, cluster in enumerate(clusters[: len(name_slots)], start=1):
            cluster.sort(key=self._shape_sort_key)
            if cluster:
                self._write_shape_from_value(cluster[0].shape, content.get(f'person_{index}_name'))
            if len(cluster) > 1:
                self._write_shape_from_value(cluster[1].shape, content.get(f'person_{index}_role'))
        return True

    def _fill_phase_points_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'phase_1_points' not in slot_keys or 'phase_2_points' not in slot_keys:
            return False

        candidates = self._collect_non_static_shapes(slide, slot_keys)
        self._consume_title_shape(candidates, content)
        if not candidates:
            return False

        heading_shapes = sorted([info for info in candidates if not self._looks_like_body_shape(info)], key=self._shape_sort_key)
        body_shapes = sorted([info for info in candidates if self._looks_like_body_shape(info)], key=self._shape_sort_key)
        if len(heading_shapes) < 2 or len(body_shapes) < 2:
            return False

        self._write_shape_from_value(heading_shapes[0].shape, content.get('phase_1_heading'))
        self._write_shape_from_value(heading_shapes[1].shape, content.get('phase_2_heading'))
        self._write_shape_from_value(body_shapes[0].shape, content.get('phase_1_points'))
        self._write_shape_from_value(body_shapes[1].shape, content.get('phase_2_points'))
        for info in candidates:
            if info not in heading_shapes[:2] and info not in body_shapes[:2] and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_framework_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'framework_text' not in slot_keys or not any(key.startswith('theory_') for key in slot_keys):
            return False

        candidates = self._collect_non_static_shapes(slide, slot_keys)
        title_shape = self._consume_title_shape(candidates, content)
        if title_shape is None or not candidates:
            return False

        heading_shapes = sorted([info for info in candidates if not self._looks_like_body_shape(info)], key=self._shape_sort_key)
        body_shapes = sorted([info for info in candidates if self._looks_like_body_shape(info)], key=self._shape_sort_key)
        if not body_shapes:
            return False

        left_body = min(body_shapes, key=lambda item: (item.left, item.top))
        right_body = max(body_shapes, key=lambda item: (item.left, -item.top))
        middle_body = [info for info in body_shapes if info is not left_body and info is not right_body]
        middle_headings = [info for info in heading_shapes if info.left > left_body.left + 900_000 and info.left < right_body.left - 900_000]

        if 'key_terms_heading' in content:
            left_headings = sorted([info for info in heading_shapes if info.left <= left_body.left + 900_000], key=self._shape_sort_key)
            if left_headings:
                self._write_shape_from_value(left_headings[0].shape, content.get('key_terms_heading'))

        if 'key_terms_items' in content:
            self._write_shape_from_value(left_body.shape, content.get('key_terms_items'))
        else:
            key_term_values = self._collect_indexed_slot_values(content, prefix='key_term_')
            if key_term_values:
                self._write_shape_paragraphs(left_body.shape, key_term_values)

        if 'framework_heading' in content:
            right_headings = sorted([info for info in heading_shapes if info.left >= right_body.left - 900_000], key=self._shape_sort_key)
            if right_headings:
                self._write_shape_from_value(right_headings[0].shape, content.get('framework_heading'))
        self._write_shape_from_value(right_body.shape, content.get('framework_text'))

        theory_heading_keys = [key for key in slot_keys if re.fullmatch(r'theory_\d+_heading', key)]
        theory_text_keys = [key for key in slot_keys if re.fullmatch(r'theory_\d+_text', key)]
        middle_headings.sort(key=self._shape_sort_key)
        middle_body.sort(key=self._shape_sort_key)
        if len(middle_headings) > len(theory_heading_keys):
            middle_headings = middle_headings[-len(theory_heading_keys):]
        used_shapes = [left_body, right_body]
        for shape, key in zip(middle_headings, theory_heading_keys):
            self._write_shape_from_value(shape.shape, content.get(key))
            used_shapes.append(shape)
        for shape, key in zip(middle_body, theory_text_keys):
            self._write_shape_from_value(shape.shape, content.get(key))
            used_shapes.append(shape)
        for info in candidates:
            if info not in used_shapes and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_literature_review_five_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'entry_1_ref' not in slot_keys or 'entry_1_note' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        self._consume_title_shape(candidates, content)
        if not candidates:
            return False

        body_shape = max(candidates, key=lambda info: (getattr(info.shape, 'height', 0), len(info.text)))
        paragraphs: list[str] = []
        for index in range(1, 6):
            ref = self._stringify_value(content.get(f'entry_{index}_ref'))
            note = self._stringify_value(content.get(f'entry_{index}_note'))
            if ref and note:
                paragraphs.append(f'{ref} — {note}')
            elif ref:
                paragraphs.append(ref)
            elif note:
                paragraphs.append(note)
        if not paragraphs:
            return False

        self._write_shape_paragraphs(body_shape.shape, paragraphs)
        for info in candidates:
            if info is not body_shape and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_methodology_matrix_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        required = {
            'method_label',
            'method_text',
            'data_type_label',
            'data_type_text',
            'motives_label',
            'motives_text',
            'collection_label',
            'collection_text',
            'sampling_label',
            'sampling_text',
        }
        if not required.issubset(set(slot_keys)):
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        self._consume_title_shape(candidates, content)
        if len(candidates) < 9:
            return False

        slide_center = 4_914_000
        short_shapes = [info for info in candidates if not self._looks_like_body_shape(info)]
        body_shapes = [info for info in candidates if self._looks_like_body_shape(info)]
        if len(short_shapes) < 5 or len(body_shapes) < 4:
            return False

        center_shape = min(short_shapes, key=lambda info: abs(info.left - slide_center) + abs(info.top - 2_200_000))
        remaining_labels = [info for info in short_shapes if info is not center_shape]
        remaining_labels.sort(key=lambda info: (info.left, info.top))
        body_shapes.sort(key=lambda info: (info.left, info.top))

        if len(remaining_labels) < 4 or len(body_shapes) < 4:
            return False

        self._write_shape_from_value(
            center_shape.shape,
            f"{self._stringify_value(content.get('method_label'))}\n{self._stringify_value(content.get('method_text'))}",
        )

        corner_pairs = [
            ('data_type_label', 'data_type_text'),
            ('motives_label', 'motives_text'),
            ('collection_label', 'collection_text'),
            ('sampling_label', 'sampling_text'),
        ]
        for shape, (label_key, text_key) in zip(remaining_labels[:4], corner_pairs):
            self._write_shape_from_value(shape.shape, content.get(label_key))
        for shape, (_, text_key) in zip(body_shapes[:4], corner_pairs):
            self._write_shape_from_value(shape.shape, content.get(text_key))

        used = [center_shape, *remaining_labels[:4], *body_shapes[:4]]
        for info in candidates:
            if info not in used and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_discussion_two_columns_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'discussion_1_heading' not in slot_keys or 'discussion_2_heading' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        self._consume_title_shape(candidates, content)
        if not candidates:
            return False

        left_candidates = [info for info in candidates if info.left < 4_500_000]
        right_candidates = [info for info in candidates if info.left >= 4_500_000]
        if not left_candidates or len(right_candidates) < 4:
            return False

        left_body = max(left_candidates, key=lambda info: (getattr(info.shape, 'height', 0), len(info.text)))
        left_lines: list[str] = []
        intro = self._stringify_value(content.get('intro_text'))
        if intro:
            left_lines.append(intro)
        list_value = content.get('bullet_points') or content.get('supporting_points')
        left_lines.extend(self._coerce_string_list(list_value))
        if left_lines:
            self._write_shape_paragraphs(left_body.shape, left_lines)

        right_text_candidates = [info for info in right_candidates if self._normalize_text(info.text) and not re.fullmatch(r'\d+', self._normalize_text(info.text))]
        heading_shapes = sorted([info for info in right_text_candidates if not self._looks_like_body_shape(info)], key=self._shape_sort_key)
        body_shapes = sorted([info for info in right_text_candidates if self._looks_like_body_shape(info)], key=self._shape_sort_key)
        if len(heading_shapes) < 2 or len(body_shapes) < 2:
            return False

        self._write_shape_from_value(heading_shapes[0].shape, content.get('discussion_1_heading'))
        self._write_shape_from_value(body_shapes[0].shape, content.get('discussion_1_text'))
        self._write_shape_from_value(heading_shapes[1].shape, content.get('discussion_2_heading'))
        self._write_shape_from_value(body_shapes[1].shape, content.get('discussion_2_text'))

        used = [left_body, heading_shapes[0], heading_shapes[1], body_shapes[0], body_shapes[1]]
        for info in candidates:
            if info not in used and (
                info in left_candidates
                or self._looks_like_source_placeholder(info.text)
                or re.fullmatch(r'\d+', self._normalize_text(info.text))
            ):
                self._clear_shape(info.shape)
        return True

    def _fill_five_step_process_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        required = {f'step_{index}_{suffix}' for index in range(1, 6) for suffix in ('heading', 'text')}
        if not required.issubset(set(slot_keys)):
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        self._consume_title_shape(candidates, content)
        usable = [info for info in candidates if not re.fullmatch(r'\d+', self._normalize_text(info.text))]
        clusters = self._cluster_shapes_by_left(sorted(usable, key=lambda info: (info.left, info.top)), tolerance=900_000)
        if len(clusters) < 5:
            return False

        step_clusters = clusters[:5]
        used: list[ShapeInfo] = []
        for index, cluster in enumerate(step_clusters, start=1):
            cluster = sorted(cluster, key=self._shape_sort_key)
            if len(cluster) < 2:
                continue
            heading_shape = min(cluster, key=lambda info: (len(info.text), getattr(info.shape, 'height', 0)))
            text_shape = max(cluster, key=lambda info: (len(info.text), getattr(info.shape, 'height', 0)))
            self._write_shape_from_value(heading_shape.shape, content.get(f'step_{index}_heading'))
            self._write_shape_from_value(text_shape.shape, content.get(f'step_{index}_text'))
            used.extend([heading_shape, text_shape])

        for info in candidates:
            if info not in used and (self._looks_like_source_placeholder(info.text) or re.fullmatch(r'\d+', self._normalize_text(info.text))):
                self._clear_shape(info.shape)
        return True

    def _fill_five_column_methodology_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        required = {f'pillar_{index}_{suffix}' for index in range(1, 6) for suffix in ('heading', 'text')}
        if not required.issubset(set(slot_keys)):
            return False
        return self._fill_card_pairs_slide(
            slide=slide,
            content=content,
            heading_keys=[f'pillar_{index}_heading' for index in range(1, 6)],
            text_keys=[f'pillar_{index}_text' for index in range(1, 5 + 1)],
        )

    def _fill_chart_plus_notes_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'chart_caption' not in slot_keys or 'supporting_notes' not in slot_keys:
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        self._consume_title_shape(candidates, content)
        if not candidates:
            return False

        left_body = min(candidates, key=lambda info: (info.left, -getattr(info.shape, 'width', 0)))
        self._write_shape_paragraphs(
            left_body.shape,
            [self._stringify_value(content.get('chart_caption')), *self._coerce_string_list(content.get('supporting_notes'))],
        )

        right_candidates = [info for info in candidates if info is not left_body and info.left > left_body.left + 2_000_000]
        heading_shapes = sorted([info for info in right_candidates if not self._looks_like_body_shape(info)], key=self._shape_sort_key)
        body_shapes = sorted([info for info in right_candidates if self._looks_like_body_shape(info)], key=self._shape_sort_key)

        data_points = [item for item in (content.get('data_points') or []) if isinstance(item, dict)]
        for shape, item in zip(heading_shapes[:2], data_points[:2]):
            self._write_shape_from_value(shape.shape, item.get('label'))
        for shape, item in zip(body_shapes[:2], data_points[:2]):
            value = item.get('value')
            label = item.get('label')
            self._write_shape_from_value(shape.shape, f'{label}: {value}' if label and value is not None else value)

        used = [left_body, *heading_shapes[:2], *body_shapes[:2]]
        for info in candidates:
            if info not in used and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_three_way_result_analysis_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        required = {
            'result_1_heading',
            'result_1_text',
            'result_2_heading',
            'result_2_text',
            'result_3_heading',
            'result_3_text',
            'highlight_value',
        }
        if not required.issubset(set(slot_keys)):
            return False

        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        self._consume_title_shape(candidates, content)
        numeric_shapes = [info for info in candidates if self._looks_like_value_shape(info)]
        non_numeric = [info for info in candidates if info not in numeric_shapes and not self._looks_like_source_placeholder(info.text)]
        heading_shapes = sorted([info for info in non_numeric if not self._looks_like_body_shape(info)], key=self._shape_sort_key)
        body_shapes = sorted([info for info in non_numeric if self._looks_like_body_shape(info)], key=self._shape_sort_key)
        if len(heading_shapes) < 3 or len(body_shapes) < 3:
            return False

        for index, shape in enumerate(heading_shapes[:3], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'result_{index}_heading'))
        for index, shape in enumerate(body_shapes[:3], start=1):
            self._write_shape_from_value(shape.shape, content.get(f'result_{index}_text'))
        if numeric_shapes:
            highlight_shape = max(numeric_shapes, key=lambda info: getattr(info.shape, 'height', 0))
            self._write_shape_from_value(highlight_shape.shape, content.get('highlight_value'))

        used = list(heading_shapes[:3] + body_shapes[:3] + (numeric_shapes[:1] if numeric_shapes else []))
        for info in candidates:
            if info not in used and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_geo_context_map_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        required = {
            'location_1_name',
            'location_1_text',
            'location_2_name',
            'location_2_text',
            'location_3_name',
            'location_3_text',
        }
        if not required.issubset(set(slot_keys)):
            return False
        return self._fill_card_pairs_slide(
            slide=slide,
            content=content,
            heading_keys=[f'location_{index}_name' for index in range(1, 4)],
            text_keys=[f'location_{index}_text' for index in range(1, 4)],
            clear_numbers=True,
        )

    def _fill_four_point_conclusion_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        required = {
            'point_1_heading',
            'point_1_text',
            'point_2_heading',
            'point_2_text',
            'point_3_heading',
            'point_3_text',
            'point_4_heading',
            'point_4_text',
        }
        if not required.issubset(set(slot_keys)):
            return False
        return self._fill_card_pairs_slide(
            slide=slide,
            content=content,
            heading_keys=[f'point_{index}_heading' for index in range(1, 5)],
            text_keys=[f'point_{index}_text' for index in range(1, 5)],
            clear_numbers=True,
        )

    def _fill_heading_text_grid_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        heading_keys = self._sorted_indexed_keys(slot_keys, suffixes=('heading', 'name', 'label', 'author'))
        text_keys = self._sorted_indexed_keys(slot_keys, suffixes=('text', 'note'))
        value_keys = self._sorted_indexed_keys(slot_keys, suffixes=('value', 'percent', 'year'))
        if len(heading_keys) < 2 and len(text_keys) < 2:
            return False

        candidates = self._collect_non_static_shapes(slide, slot_keys)
        self._consume_title_shape(candidates, content)
        if not candidates:
            return False

        value_shapes = sorted([info for info in candidates if self._looks_like_value_shape(info)], key=self._shape_sort_key)
        body_shapes = sorted([info for info in candidates if self._looks_like_body_shape(info)], key=self._shape_sort_key)
        heading_shapes = sorted(
            [info for info in candidates if info not in value_shapes and info not in body_shapes],
            key=self._shape_sort_key,
        )

        if len(heading_shapes) < len(heading_keys) and len(body_shapes) < len(text_keys):
            return False

        used_shapes: list[ShapeInfo] = []
        for shape, key in zip(heading_shapes, heading_keys):
            self._write_shape_from_value(shape.shape, content.get(key))
            used_shapes.append(shape)

        text_shapes = list(body_shapes)
        if len(text_shapes) < len(text_keys):
            spillover_shapes = [
                info
                for info in sorted(candidates, key=self._shape_sort_key)
                if info not in used_shapes and info not in value_shapes and info not in text_shapes
            ]
            text_shapes.extend(spillover_shapes[: max(0, len(text_keys) - len(text_shapes))])

        for shape, key in zip(text_shapes, text_keys):
            self._write_shape_from_value(shape.shape, content.get(key))
            used_shapes.append(shape)
        for shape, key in zip(value_shapes, value_keys):
            self._write_shape_from_value(shape.shape, content.get(key))
            used_shapes.append(shape)
        for info in candidates:
            if info not in used_shapes and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_label_text_matrix_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        label_keys = [key for key in slot_keys if key.endswith('_label')]
        text_keys = [key for key in slot_keys if key.endswith('_text')]
        if len(label_keys) < 3 or len(text_keys) < 3:
            return False

        candidates = self._collect_non_static_shapes(slide, slot_keys)
        self._consume_title_shape(candidates, content)
        if not candidates:
            return False

        label_shapes = sorted([info for info in candidates if not self._looks_like_body_shape(info)], key=self._shape_sort_key)
        text_shapes = sorted([info for info in candidates if self._looks_like_body_shape(info)], key=self._shape_sort_key)
        if len(label_shapes) < len(label_keys) or len(text_shapes) < len(text_keys):
            return False

        for shape, key in zip(label_shapes, label_keys):
            self._write_shape_from_value(shape.shape, content.get(key))
        for shape, key in zip(text_shapes, text_keys):
            self._write_shape_from_value(shape.shape, content.get(key))
        used_shapes = list(label_shapes[: len(label_keys)] + text_shapes[: len(text_keys)])
        for info in candidates:
            if info not in used_shapes and self._looks_like_source_placeholder(info.text):
                self._clear_shape(info.shape)
        return True

    def _fill_rows_without_table_slide(self, slide, slot_keys: list[str], content: dict[str, Any]) -> bool:
        if 'rows' not in slot_keys or self._first_table_shape(slide.shapes) is not None:
            return False

        candidates = self._collect_non_static_shapes(slide, slot_keys)
        title_shape = self._consume_title_shape(candidates, content)
        rows_text = self._format_rows_as_text(content.get('rows'))
        if not rows_text:
            return False

        if candidates:
            self._write_shape_paragraphs(candidates[0].shape, rows_text.splitlines())
            for info in candidates[1:]:
                if self._looks_like_source_placeholder(info.text):
                    self._clear_shape(info.shape)
            return True
        if title_shape is not None:
            self._write_shape_paragraphs(
                title_shape.shape,
                [self._stringify_value(content.get('title')), *rows_text.splitlines()],
            )
            return True
        return False

    def _fill_generic_slide(self, slide, slot_keys: Sequence[str], content: dict[str, Any]) -> None:
        candidates = [
            info
            for info in self._collect_text_shapes(slide.shapes)
            if not self._is_static_shape_text(info.text, slot_keys)
        ]
        if not candidates:
            return

        candidates.sort(key=lambda item: (item.top, item.left))
        items = deque(self._build_slot_items(slot_keys, content))
        shape_index = 0

        while items and shape_index < len(candidates):
            info = candidates[shape_index]
            current = items[0]
            remaining_shapes = len(candidates) - shape_index

            if current['group']:
                paragraphs = list(current['paragraphs'])
                if info.paragraph_count > 1:
                    self._write_shape_paragraphs(info.shape, paragraphs[: max(info.paragraph_count, len(paragraphs))])
                    items.popleft()
                    shape_index += 1
                    continue

                if len(paragraphs) <= remaining_shapes and all(candidate.paragraph_count == 1 for candidate in candidates[shape_index:shape_index + len(paragraphs)]):
                    items.popleft()
                    for paragraph in paragraphs:
                        self._write_shape_paragraphs(candidates[shape_index].shape, [paragraph])
                        shape_index += 1
                    continue

                self._write_shape_paragraphs(info.shape, paragraphs)
                items.popleft()
                shape_index += 1
                continue

            paragraph_target: list[str] = []
            while items and not items[0]['group'] and len(paragraph_target) < info.paragraph_count:
                paragraph_target.extend(items[0]['paragraphs'])
                items.popleft()
            if not paragraph_target:
                paragraph_target = ['']
            self._write_shape_paragraphs(info.shape, paragraph_target)
            shape_index += 1

    def _fill_card_pairs_slide(
        self,
        *,
        slide,
        content: dict[str, Any],
        heading_keys: Sequence[str],
        text_keys: Sequence[str],
        clear_numbers: bool = False,
    ) -> bool:
        candidates = self._collect_text_shapes(slide.shapes)
        if not candidates:
            return False

        self._consume_title_shape(candidates, content)
        usable = []
        for info in candidates:
            normalized = self._normalize_text(info.text)
            if clear_numbers and re.fullmatch(r'\d+%?', normalized):
                continue
            usable.append(info)

        heading_shapes = sorted([info for info in usable if not self._looks_like_body_shape(info)], key=self._shape_sort_key)
        body_shapes = sorted([info for info in usable if self._looks_like_body_shape(info)], key=self._shape_sort_key)
        if len(heading_shapes) < len(heading_keys) or len(body_shapes) < len(text_keys):
            return False

        used_shapes: list[ShapeInfo] = []
        for shape, key in zip(heading_shapes[: len(heading_keys)], heading_keys):
            self._write_shape_from_value(shape.shape, content.get(key))
            used_shapes.append(shape)
        for shape, key in zip(body_shapes[: len(text_keys)], text_keys):
            self._write_shape_from_value(shape.shape, content.get(key))
            used_shapes.append(shape)

        for info in candidates:
            normalized = self._normalize_text(info.text)
            if info not in used_shapes and (
                self._looks_like_source_placeholder(info.text) or (clear_numbers and re.fullmatch(r'\d+%?', normalized))
            ):
                self._clear_shape(info.shape)
        return True

    def _build_slot_items(self, slot_keys: Sequence[str], content: dict[str, Any]) -> list[dict[str, Any]]:
        items: list[dict[str, Any]] = []
        for key in slot_keys:
            if key in {'rows', 'columns', 'table_headers'}:
                continue
            value = content.get(key)
            if isinstance(value, list):
                paragraphs = [self._stringify_value(item) for item in value if self._stringify_value(item)]
                items.append({'key': key, 'group': True, 'paragraphs': paragraphs or ['']})
                continue
            items.append({'key': key, 'group': False, 'paragraphs': [self._stringify_value(value)]})
        return items

    def _collect_text_shapes(self, shapes) -> list[ShapeInfo]:
        result: list[ShapeInfo] = []
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                result.extend(self._collect_text_shapes(shape.shapes))
                continue
            if getattr(shape, 'has_text_frame', False):
                text = '\n'.join(paragraph.text for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
                if not text.strip():
                    continue
                paragraph_count = sum(1 for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()) or 1
                result.append(
                    ShapeInfo(
                        shape=shape,
                        text=text.strip(),
                        top=int(getattr(shape, 'top', 0)),
                        left=int(getattr(shape, 'left', 0)),
                        paragraph_count=paragraph_count,
                    )
                )
        return result

    def _collect_non_static_shapes(self, slide, slot_keys: Sequence[str]) -> list[ShapeInfo]:
        return [
            info
            for info in self._collect_text_shapes(slide.shapes)
            if not self._is_static_shape_text(info.text, slot_keys)
        ]

    def _first_table_shape(self, shapes):
        for shape in shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                table = self._first_table_shape(shape.shapes)
                if table is not None:
                    return table
                continue
            if getattr(shape, 'has_table', False):
                return shape
        return None

    @staticmethod
    def _normalize_text(text: str) -> str:
        return ' '.join(str(text or '').strip().lower().split())

    def _is_static_shape_text(self, text: str, slot_keys: Sequence[str]) -> bool:
        normalized = self._normalize_text(text)
        if not normalized:
            return True
        if normalized in STATIC_TEXT_EXACT:
            return True
        if re.fullmatch(r'\d+', normalized):
            return True
        if (
            ('follow the link in the graph' in normalized or 'to modify this graph' in normalized)
            and not any(slot in slot_keys for slot in {'caption', 'chart_caption', 'body', 'supporting_notes'})
        ):
            return True
        if (
            'replace the image on the screen with your own' in normalized
            and not any(slot in slot_keys for slot in {'caption', 'body'})
        ):
            return True
        if (
            'replace the images on the screen with your own' in normalized
            and not any(slot in slot_keys for slot in {'caption', 'body'})
        ):
            return True
        return any(fragment in normalized for fragment in STATIC_TEXT_CONTAINS)

    def _cleanup_placeholder_shapes(self, slide, *, content: dict[str, Any]) -> None:
        content_blob = self._normalize_text(' '.join(self._flatten_content_strings(content)))
        for info in self._collect_text_shapes(slide.shapes):
            if self._looks_like_source_placeholder(info.text, content_blob=content_blob):
                self._clear_shape(info.shape)

    def _is_placeholder_only_slide(self, slide) -> bool:
        texts = [info.text for info in self._collect_text_shapes(slide.shapes)]
        if not texts:
            return False
        normalized = [self._normalize_text(text) for text in texts if self._normalize_text(text)]
        if not normalized:
            return False
        placeholder_count = sum(1 for text in normalized if self._looks_like_source_placeholder(text))
        return placeholder_count == len(normalized)

    def _consume_title_shape(self, candidates: list[ShapeInfo], content: dict[str, Any]) -> ShapeInfo | None:
        title_value = self._stringify_value(content.get('title'))
        if not title_value or not candidates:
            return None

        title_shape = min(candidates, key=self._shape_sort_key)
        self._write_shape_from_value(title_shape.shape, title_value)
        candidates.remove(title_shape)

        duplicate_text = self._normalize_text(title_shape.text)
        if duplicate_text:
            duplicates = [
                info
                for info in list(candidates)
                if self._normalize_text(info.text) == duplicate_text and abs(info.top - title_shape.top) < 1_200_000
            ]
            for info in duplicates:
                candidates.remove(info)
        return title_shape

    @staticmethod
    def _shape_sort_key(info: ShapeInfo) -> tuple[int, int]:
        return (info.top, info.left)

    @staticmethod
    def _looks_like_value_shape(info: ShapeInfo) -> bool:
        normalized = ' '.join(str(info.text or '').split())
        return bool(re.fullmatch(r'[\+\-]?\d+(?:[.,]\d+)?%?', normalized))

    @staticmethod
    def _looks_like_body_shape(info: ShapeInfo) -> bool:
        normalized = ' '.join(str(info.text or '').split())
        return info.paragraph_count > 1 or len(normalized) > 42

    def _write_shape_from_value(self, shape, value: Any) -> None:
        if isinstance(value, list):
            paragraphs = [self._stringify_value(item) for item in value if self._stringify_value(item)]
            self._write_shape_paragraphs(shape, paragraphs or [''])
            return
        self._write_shape_paragraphs(shape, [self._stringify_value(value)])

    def _clear_shape(self, shape) -> None:
        self._write_shape_paragraphs(shape, [''])

    @staticmethod
    def _cluster_shapes_by_left(candidates: Sequence[ShapeInfo], tolerance: int = 1_500_000) -> list[list[ShapeInfo]]:
        clusters: list[list[ShapeInfo]] = []
        for info in sorted(candidates, key=lambda item: (item.left, item.top)):
            if not clusters or abs(info.left - clusters[-1][0].left) > tolerance:
                clusters.append([info])
            else:
                clusters[-1].append(info)
        return clusters

    @staticmethod
    def _collect_indexed_slot_values(content: dict[str, Any], *, prefix: str) -> list[str]:
        values: list[tuple[int, str]] = []
        pattern = re.compile(rf'{re.escape(prefix)}(\d+)$')
        for key, value in content.items():
            match = pattern.fullmatch(key)
            if not match:
                continue
            text = str(value or '').strip()
            if text:
                values.append((int(match.group(1)), text))
        values.sort(key=lambda item: item[0])
        return [text for _, text in values]

    def _format_rows_as_text(self, rows: Any) -> str:
        if not isinstance(rows, list):
            return ''
        lines: list[str] = []
        for row in rows:
            if isinstance(row, list):
                parts = [self._stringify_value(item) for item in row if self._stringify_value(item)]
            elif isinstance(row, dict):
                parts = [self._stringify_value(item) for item in row.values() if self._stringify_value(item)]
            else:
                parts = [self._stringify_value(row)] if self._stringify_value(row) else []
            if parts:
                lines.append(' - '.join(parts))
        return '\n'.join(lines)

    def _coerce_string_list(self, value: Any) -> list[str]:
        if value is None:
            return []
        if isinstance(value, list):
            return [self._stringify_value(item) for item in value if self._stringify_value(item)]
        text = self._stringify_value(value)
        if not text:
            return []
        return [part.strip(' -•\t') for part in re.split(r'[\r\n]+', text) if part.strip(' -•\t')]

    @staticmethod
    def _sorted_indexed_keys(slot_keys: Sequence[str], *, suffixes: tuple[str, ...]) -> list[str]:
        pattern = re.compile(rf'(.+?)_(\d+)_({"|".join(re.escape(suffix) for suffix in suffixes)})$')
        matched: list[tuple[int, str, str]] = []
        for key in slot_keys:
            match = pattern.fullmatch(key)
            if match:
                matched.append((int(match.group(2)), match.group(3), key))
        matched.sort(key=lambda item: (item[0], item[1]))
        return [key for _, _, key in matched]

    def _looks_like_source_placeholder(self, text: str, *, content_blob: str = '') -> bool:
        normalized = self._normalize_text(text)
        if not normalized:
            return False
        if normalized in SOURCE_PLACEHOLDER_EXACT:
            return True
        if normalized in {'clay boards', 'papyrus', 'paper', 'printing', 'typewriter'} and normalized not in content_blob:
            return True
        if re.fullmatch(r'(left|right|tasks?|subjects?) \d+', normalized):
            return True
        if normalized == 'table of contents' and normalized not in content_blob:
            return True
        if normalized in {'mercury', 'venus', 'mars', 'jupiter', 'saturn', 'neptune', 'pluto'} and normalized not in content_blob:
            return True
        return any(fragment in normalized for fragment in SOURCE_PLACEHOLDER_CONTAINS)

    def _localized_toc_title(self) -> str:
        language = self._active_language or 'uz'
        if language.startswith('en'):
            return 'Table of contents'
        if language.startswith('ru'):
            return 'Содержание'
        return 'Mundarija'

    def _flatten_content_strings(self, value: Any) -> list[str]:
        if value is None:
            return []
        if isinstance(value, dict):
            result: list[str] = []
            for item in value.values():
                result.extend(self._flatten_content_strings(item))
            return result
        if isinstance(value, list):
            result: list[str] = []
            for item in value:
                result.extend(self._flatten_content_strings(item))
            return result
        text = str(value).strip()
        return [text] if text else []

    def _write_shape_paragraphs(self, shape, paragraphs: Sequence[str]) -> None:
        text_frame = shape.text_frame
        target_values = list(paragraphs) if paragraphs else ['']
        existing_paragraphs = list(text_frame.paragraphs)
        while len(existing_paragraphs) < len(target_values):
            existing_paragraphs.append(text_frame.add_paragraph())

        for index, paragraph in enumerate(existing_paragraphs):
            value = target_values[index] if index < len(target_values) else ''
            self._write_paragraph(paragraph, value)

    @staticmethod
    def _write_paragraph(paragraph: _Paragraph, value: str) -> None:
        cleaned = str(value or '')
        if paragraph.runs:
            paragraph.runs[0].text = cleaned
            for run in paragraph.runs[1:]:
                run.text = ''
        else:
            paragraph.text = cleaned
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)

    def _write_cell(self, cell, value: str) -> None:
        text_frame = cell.text_frame
        paragraphs = list(text_frame.paragraphs)
        if not paragraphs:
            paragraphs = [text_frame.paragraphs[0]]
        self._write_paragraph(paragraphs[0], value)
        for paragraph in paragraphs[1:]:
            self._write_paragraph(paragraph, '')

    @staticmethod
    def _stringify_value(value: Any) -> str:
        if value is None:
            return ''
        if isinstance(value, list):
            return '\n'.join(str(item).strip() for item in value if str(item).strip())
        if isinstance(value, dict):
            return '\n'.join(f'{key}: {item}' for key, item in value.items() if str(item).strip())
        return str(value).strip()

    @staticmethod
    def _safe_filename(value: str) -> str:
        cleaned = re.sub(r'[^A-Za-z0-9А-Яа-яЎўҚқҒғҲҳ_\- ]+', '', value).strip()
        cleaned = re.sub(r'\s+', '_', cleaned)
        return cleaned or 'magic_slide'

    @staticmethod
    def _delete_slide(presentation: Presentation, slide_index: int) -> None:
        slide_id_list = presentation.slides._sldIdLst
        slide_ids = list(slide_id_list)
        slide_id = slide_ids[slide_index]
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        del slide_id_list[slide_index]


class MagicGenerationService:
    def __init__(
        self,
        *,
        registry: MagicTemplateRegistry,
        content_generator: MagicContentGenerator,
        renderer: MagicPptxRenderer,
    ) -> None:
        self.registry = registry
        self.content_generator = content_generator
        self.renderer = renderer

    def has_template(self, template_id: str) -> bool:
        return self.registry.has(template_id)

    def generate(self, order_payload: dict[str, Any]) -> tuple[str, MagicTemplateSpec]:
        template_id = str(order_payload.get('template_id') or '')
        template = self.registry.get(template_id)
        content_plan = self.content_generator.build_content(template, order_payload)
        file_path = self.renderer.render(
            template=template,
            order_payload=order_payload,
            content_plan=content_plan,
        )
        return file_path, template
