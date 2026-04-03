from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "templates" / "education" / "education-2.pptx"
TARGET = ROOT / "templates" / "education" / "education-2-tokenized.pptx"


def find_shape(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"Shape topilmadi: slide={slide.slide_id} shape_id={shape_id}")


def set_shape_text(slide, shape_id: int, text: str) -> None:
    find_shape(slide, shape_id).text = text


def set_shape_paragraphs(slide, shape_id: int, values: list[str]) -> None:
    shape = find_shape(slide, shape_id)
    paragraphs = list(shape.text_frame.paragraphs)
    while len(paragraphs) < len(values):
        paragraphs.append(shape.text_frame.add_paragraph())
    for index, paragraph in enumerate(paragraphs):
        paragraph.text = values[index] if index < len(values) else ""


def set_table_tokens(slide, shape_id: int, values: list[list[str]]) -> None:
    table = find_shape(slide, shape_id).table
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            token = values[row_index][col_index]
            if token is not None:
                cell.text = token


def main() -> None:
    prs = Presentation(SOURCE)

    # Slide 1
    slide = prs.slides[0]
    set_shape_text(slide, 717, "[[title]]")
    set_shape_text(slide, 718, "[[author]]")

    # Slide 3
    slide = prs.slides[2]
    set_shape_text(slide, 750, "[[toc_title]]")
    set_shape_text(slide, 744, "[[section_1_title]]")
    set_shape_text(slide, 732, "[[section_1_desc]]")
    set_shape_text(slide, 745, "[[section_2_title]]")
    set_shape_text(slide, 733, "[[section_2_desc]]")
    set_shape_text(slide, 746, "[[section_3_title]]")
    set_shape_text(slide, 734, "[[section_3_desc]]")
    set_shape_text(slide, 747, "[[section_4_title]]")
    set_shape_text(slide, 735, "[[section_4_desc]]")
    set_shape_text(slide, 748, "[[section_5_title]]")
    set_shape_text(slide, 736, "[[section_5_desc]]")
    set_shape_text(slide, 749, "[[section_6_title]]")
    set_shape_text(slide, 737, "[[section_6_desc]]")

    # Slide 4
    slide = prs.slides[3]
    set_shape_text(slide, 755, "[[title]]")
    set_shape_text(slide, 756, "[[body]]")

    # Slide 5
    slide = prs.slides[4]
    set_shape_text(slide, 768, "[[title]]")
    set_shape_text(slide, 770, "[[objective_text]]")
    set_shape_text(slide, 769, "[[supporting_context]]")

    # Slide 6
    slide = prs.slides[5]
    set_shape_text(slide, 775, "[[title]]")
    set_shape_paragraphs(
        slide,
        776,
        [
            "[[intro_line]]",
            "[[bullet_1]]",
            "[[bullet_2]]",
            "[[bullet_3]]",
            "[[bullet_4]]",
        ],
    )

    # Slide 7
    slide = prs.slides[6]
    set_shape_text(slide, 782, "[[title]]")
    set_shape_text(slide, 785, "[[card_1_heading]]")
    set_shape_text(slide, 784, "[[card_1_text]]")
    set_shape_text(slide, 786, "[[card_2_heading]]")
    set_shape_text(slide, 783, "[[card_2_text]]")

    # Slide 8
    slide = prs.slides[7]
    set_shape_text(slide, 807, "[[title]]")
    set_shape_text(slide, 811, "[[type_1_heading]]")
    set_shape_text(slide, 808, "[[type_1_text]]")
    set_shape_text(slide, 812, "[[type_2_heading]]")
    set_shape_text(slide, 809, "[[type_2_text]]")
    set_shape_text(slide, 813, "[[type_3_heading]]")
    set_shape_text(slide, 810, "[[type_3_text]]")

    # Slide 9
    slide = prs.slides[8]
    set_shape_text(slide, 840, "[[title]]")
    set_shape_text(slide, 845, "[[factor_1_heading]]")
    set_shape_text(slide, 841, "[[factor_1_text]]")
    set_shape_text(slide, 847, "[[factor_2_heading]]")
    set_shape_text(slide, 842, "[[factor_2_text]]")
    set_shape_text(slide, 846, "[[factor_3_heading]]")
    set_shape_text(slide, 843, "[[factor_3_text]]")
    set_shape_text(slide, 848, "[[factor_4_heading]]")
    set_shape_text(slide, 844, "[[factor_4_text]]")

    # Slide 10
    slide = prs.slides[9]
    set_shape_text(slide, 890, "[[title]]")
    set_shape_text(slide, 895, "[[aspect_1_heading]]")
    set_shape_text(slide, 891, "[[aspect_1_text]]")
    set_shape_text(slide, 896, "[[aspect_2_heading]]")
    set_shape_text(slide, 888, "[[aspect_2_text]]")
    set_shape_text(slide, 897, "[[aspect_3_heading]]")
    set_shape_text(slide, 889, "[[aspect_3_text]]")
    set_shape_text(slide, 898, "[[aspect_4_heading]]")
    set_shape_text(slide, 892, "[[aspect_4_text]]")
    set_shape_text(slide, 899, "[[aspect_5_heading]]")
    set_shape_text(slide, 893, "[[aspect_5_text]]")
    set_shape_text(slide, 900, "[[aspect_6_heading]]")
    set_shape_text(slide, 894, "[[aspect_6_text]]")

    # Slide 11
    slide = prs.slides[10]
    set_shape_text(slide, 946, "[[quote]]")
    set_shape_text(slide, 947, "[[author]]")

    # Slide 12
    slide = prs.slides[11]
    set_shape_text(slide, 953, "[[title]]")

    # Slide 13
    slide = prs.slides[12]
    set_shape_paragraphs(slide, 976, ["[[headline_line_1]]", "[[headline_line_2]]"])
    set_shape_text(slide, 977, "[[body]]")

    # Slide 14
    slide = prs.slides[13]
    set_shape_text(slide, 1022, "[[title]]")
    set_shape_text(slide, 1023, "[[caption]]")

    # Slide 15
    slide = prs.slides[14]
    set_shape_text(slide, 1029, "[[title]]")
    set_shape_text(slide, 1030, "[[caption]]")

    # Slide 16
    slide = prs.slides[15]
    set_shape_text(slide, 1037, "[[title]]")
    set_shape_text(slide, 1038, "[[caption]]")

    # Slide 17
    slide = prs.slides[16]
    set_shape_text(slide, 1048, "[[title]]")
    set_shape_text(slide, 1045, "[[location_1_name]]")
    set_shape_text(slide, 1049, "[[location_1_text]]")
    set_shape_text(slide, 1047, "[[location_2_name]]")
    set_shape_text(slide, 1051, "[[location_2_text]]")
    set_shape_text(slide, 1046, "[[location_3_name]]")
    set_shape_text(slide, 1050, "[[location_3_text]]")

    # Slide 18
    slide = prs.slides[17]
    set_shape_text(slide, 1126, "[[title]]")
    set_shape_text(slide, 1134, "[[event_1_heading]]")
    set_shape_text(slide, 1135, "[[event_1_text]]")
    set_shape_text(slide, 1132, "[[event_2_heading]]")
    set_shape_text(slide, 1133, "[[event_2_text]]")
    set_shape_text(slide, 1137, "[[event_3_heading]]")
    set_shape_text(slide, 1138, "[[event_3_text]]")
    set_shape_text(slide, 1139, "[[event_4_heading]]")
    set_shape_text(slide, 1140, "[[event_4_text]]")
    set_shape_text(slide, 1144, "[[milestone_heading]]")
    set_shape_text(slide, 1145, "[[milestone_text]]")

    # Slide 19
    slide = prs.slides[18]
    set_shape_text(slide, 1152, "[[title]]")
    set_shape_text(slide, 1151, "[[subtitle]]")
    set_shape_text(slide, 1153, "[[node_1_heading]]")
    set_shape_text(slide, 1157, "[[node_1_text]]")
    set_shape_text(slide, 1154, "[[node_2_heading]]")
    set_shape_text(slide, 1158, "[[node_2_text]]")
    set_shape_text(slide, 1161, "[[node_3_heading]]")
    set_shape_text(slide, 1162, "[[node_3_text]]")
    set_shape_text(slide, 1155, "[[node_4_heading]]")
    set_shape_text(slide, 1159, "[[node_4_text]]")
    set_shape_text(slide, 1156, "[[node_5_heading]]")
    set_shape_text(slide, 1160, "[[node_5_text]]")

    # Slide 20
    slide = prs.slides[19]
    set_shape_text(slide, 1172, "[[title]]")
    set_table_tokens(
        slide,
        1173,
        [
            ["[[header_aspect]]", "[[header_subject_a]]", "[[header_subject_b]]", "[[header_subject_c]]"],
            ["[[row_1_aspect]]", "[[row_1_subject_a]]", "[[row_1_subject_b]]", "[[row_1_subject_c]]"],
            ["[[row_2_aspect]]", "[[row_2_subject_a]]", "[[row_2_subject_b]]", "[[row_2_subject_c]]"],
            ["[[row_3_aspect]]", "[[row_3_subject_a]]", "[[row_3_subject_b]]", "[[row_3_subject_c]]"],
            ["[[row_4_aspect]]", "[[row_4_subject_a]]", "[[row_4_subject_b]]", "[[row_4_subject_c]]"]
        ],
    )

    # Slide 21
    slide = prs.slides[20]
    set_shape_text(slide, 1178, "[[title]]")
    set_shape_text(slide, 1179, "[[chart_caption]]")
    set_shape_text(slide, 1180, "[[series_1_label]]")
    set_shape_text(slide, 1181, "[[series_1_text]]")
    set_shape_text(slide, 1184, "[[series_2_label]]")
    set_shape_text(slide, 1185, "[[series_2_text]]")
    set_shape_text(slide, 1182, "[[series_3_label]]")
    set_shape_text(slide, 1183, "[[series_3_text]]")
    set_shape_text(slide, 1186, "[[series_4_label]]")
    set_shape_text(slide, 1187, "[[series_4_text]]")

    # Slide 22
    slide = prs.slides[21]
    set_shape_text(slide, 1216, "[[title]]")
    set_table_tokens(
        slide,
        1217,
        [
            ["[[header_data]]", "[[header_description]]"],
            ["[[row_1_data]]", "[[row_1_description]]"],
            ["[[row_2_data]]", "[[row_2_description]]"],
            ["[[row_3_data]]", "[[row_3_description]]"],
            ["[[row_4_data]]", "[[row_4_description]]"],
            ["[[row_5_data]]", "[[row_5_description]]"]
        ],
    )

    prs.save(TARGET)
    print(TARGET)


if __name__ == "__main__":
    main()
