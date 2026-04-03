from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "templates" / "education" / "education-1.pptx"
TARGET = ROOT / "templates" / "education" / "education-1-tokenized.pptx"


def find_shape(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"Shape topilmadi: slide={slide.slide_id} shape_id={shape_id}")


def set_shape_text(slide, shape_id: int, text: str) -> None:
    find_shape(slide, shape_id).text = text


def set_shape_paragraphs(slide, shape_id: int, values: list[str]) -> None:
    shape = find_shape(slide, shape_id)
    paragraphs = list(shape.text_frame.paragraphs)
    while len(paragraphs) < len(values):
        paragraphs.append(shape.text_frame.add_paragraph())
    for index, paragraph in enumerate(paragraphs):
        paragraph.text = values[index] if index < len(values) else ""


def main() -> None:
    prs = Presentation(SOURCE)

    slide = prs.slides[0]
    set_shape_paragraphs(slide, 2006, ["[[title_line_1]]", "[[title_line_2]]"])
    set_shape_text(slide, 2007, "[[subtitle]]")

    slide = prs.slides[2]
    set_shape_text(slide, 2020, "[[toc_title]]")
    set_shape_text(slide, 2021, "[[section_1_title]]")
    set_shape_text(slide, 2018, "[[section_1_desc]]")
    set_shape_text(slide, 2022, "[[section_2_title]]")
    set_shape_text(slide, 2019, "[[section_2_desc]]")
    set_shape_text(slide, 2023, "[[section_3_title]]")
    set_shape_text(slide, 2025, "[[section_3_desc]]")
    set_shape_text(slide, 2024, "[[section_4_title]]")
    set_shape_text(slide, 2026, "[[section_4_desc]]")
    set_shape_text(slide, 2032, "[[section_5_title]]")
    set_shape_text(slide, 2031, "[[section_5_desc]]")
    set_shape_text(slide, 2033, "[[section_6_title]]")
    set_shape_text(slide, 2034, "[[section_6_desc]]")

    slide = prs.slides[3]
    set_shape_text(slide, 2043, "[[toc_title]]")
    set_shape_text(slide, 2044, "[[section_1_title]]")
    set_shape_text(slide, 2041, "[[section_1_desc]]")
    set_shape_text(slide, 2045, "[[section_2_title]]")
    set_shape_text(slide, 2042, "[[section_2_desc]]")
    set_shape_text(slide, 2046, "[[section_3_title]]")
    set_shape_text(slide, 2048, "[[section_3_desc]]")
    set_shape_text(slide, 2047, "[[section_4_title]]")
    set_shape_text(slide, 2049, "[[section_4_desc]]")

    slide = prs.slides[4]
    set_shape_text(slide, 2058, "[[title]]")
    set_shape_text(slide, 2059, "[[body]]")

    slide = prs.slides[5]
    set_shape_text(slide, 2065, "[[title]]")
    set_shape_paragraphs(
        slide,
        2066,
        [
            "[[intro_note]]",
            "",
            "[[step_1]]",
            "[[step_2]]",
            "[[step_3]]",
            "[[step_4]]",
            "[[step_5]]",
            "[[step_6]]",
            "[[step_7]]",
        ],
    )

    slide = prs.slides[6]
    set_shape_text(slide, 2072, "[[title]]")
    set_shape_paragraphs(
        slide,
        2073,
        [
            "[[step_1]]",
            "[[step_2]]",
            "[[step_3]]",
            "[[step_4]]",
            "[[step_5]]",
            "[[step_6]]",
            "[[tip]]",
            "[[warning]]",
        ],
    )

    slide = prs.slides[7]
    set_shape_text(slide, 2084, "[[title]]")
    set_shape_paragraphs(
        slide,
        2085,
        [
            "[[part_1_heading]]",
            "",
            "[[part_1_step_1]]",
            "[[part_1_step_2]]",
            "[[part_1_step_3]]",
            "[[part_1_step_4]]",
            "",
            "[[part_2_heading]]",
            "[[part_2_step_1]]",
            "[[part_2_step_2]]",
            "[[part_2_step_3]]",
        ],
    )

    slide = prs.slides[8]
    set_shape_text(slide, 2607, "[[title]]")
    set_shape_paragraphs(slide, 2610, ["[[ok_1]]", "[[ok_2]]", "[[ok_3]]"])
    set_shape_paragraphs(slide, 2611, ["[[not_ok_1]]", "[[not_ok_2]]", "[[not_ok_3]]"])
    set_shape_text(slide, 2612, "[[ok_label]]")
    set_shape_text(slide, 2613, "[[not_ok_label]]")

    slide = prs.slides[9]
    set_shape_text(slide, 2618, "[[title]]")
    set_shape_paragraphs(slide, 2621, ["[[task_1]]", "[[task_2]]", "[[task_3]]", "[[task_4]]"])
    set_shape_paragraphs(slide, 2622, ["[[subject_1]]", "[[subject_2]]", "[[subject_3]]", "[[subject_4]]"])
    set_shape_text(slide, 2623, "[[tasks_heading]]")
    set_shape_text(slide, 2624, "[[subjects_heading]]")

    slide = prs.slides[10]
    set_shape_text(slide, 2630, "[[title]]")
    set_shape_text(slide, 2632, "[[note_label]]")
    set_shape_text(slide, 2633, "[[material_1]]")
    set_shape_text(slide, 2634, "[[material_2]]")
    set_shape_text(slide, 2635, "[[material_3]]")
    set_shape_text(slide, 2636, "[[material_4]]")
    set_shape_text(slide, 2637, "[[material_5]]")
    set_shape_text(slide, 2638, "[[material_6]]")

    slide = prs.slides[11]
    set_shape_text(slide, 2682, "[[title]]")
    set_shape_paragraphs(slide, 2683, ["[[item_1]]", "[[item_2]]", "[[item_3]]"])

    slide = prs.slides[12]
    set_shape_text(slide, 2696, "[[title]]")
    set_shape_text(slide, 2697, "[[step_1_heading]]")
    set_shape_text(slide, 2698, "[[step_1_text]]")
    set_shape_text(slide, 2699, "[[step_2_heading]]")
    set_shape_text(slide, 2700, "[[step_2_text]]")
    set_shape_text(slide, 2701, "[[step_3_heading]]")
    set_shape_text(slide, 2702, "[[step_3_text]]")

    slide = prs.slides[13]
    set_shape_text(slide, 2992, "[[title]]")
    set_shape_text(slide, 2994, "[[card_1_heading]]")
    set_shape_text(slide, 2995, "[[card_1_text]]")
    set_shape_text(slide, 2997, "[[card_2_heading]]")
    set_shape_text(slide, 2998, "[[card_2_text]]")
    set_shape_text(slide, 3001, "[[card_3_heading]]")
    set_shape_text(slide, 3000, "[[card_3_text]]")
    set_shape_text(slide, 3004, "[[card_4_heading]]")
    set_shape_text(slide, 3005, "[[card_4_text]]")

    slide = prs.slides[14]
    set_shape_text(slide, 3022, "[[title]]")
    set_shape_text(slide, 3026, "[[item_1_heading]]")
    set_shape_text(slide, 3031, "[[item_1_text]]")
    set_shape_text(slide, 3027, "[[item_2_heading]]")
    set_shape_text(slide, 3030, "[[item_2_text]]")
    set_shape_text(slide, 3028, "[[item_3_heading]]")
    set_shape_text(slide, 3029, "[[item_3_text]]")

    prs.save(TARGET)
    print(TARGET)


if __name__ == "__main__":
    main()
