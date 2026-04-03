from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "templates" / "business" / "business-2.pptx"
TARGET = ROOT / "templates" / "business" / "business-2-tokenized.pptx"


def set_shape_text(slide, shape_id: int, text: str) -> None:
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            shape.text = text
            return
    raise ValueError(f"Shape topilmadi: slide={slide.slide_id} shape_id={shape_id}")


def set_shape_paragraphs(slide, shape_id: int, values: list[str]) -> None:
    for shape in slide.shapes:
        if shape.shape_id != shape_id:
            continue
        paragraphs = list(shape.text_frame.paragraphs)
        while len(paragraphs) < len(values):
            paragraphs.append(shape.text_frame.add_paragraph())
        for index, paragraph in enumerate(paragraphs):
            paragraph.text = values[index] if index < len(values) else ""
        return
    raise ValueError(f"Shape topilmadi: slide={slide.slide_id} shape_id={shape_id}")


def duplicate_slide(prs: Presentation, source_index: int):
    source = prs.slides[source_index]
    blank_layout = prs.slide_layouts[6]
    destination = prs.slides.add_slide(blank_layout)

    for shape in list(destination.shapes):
        destination.shapes._spTree.remove(shape._element)

    for shape in source.shapes:
        destination.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")

    return destination


def reorder_slides(prs: Presentation, order: list[int]) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_nodes = list(slide_id_list)
    reordered = [slide_nodes[index - 1] for index in order]
    for node in list(slide_id_list):
        slide_id_list.remove(node)
    for node in reordered:
        slide_id_list.append(node)


def main() -> None:
    prs = Presentation(SOURCE)

    duplicate_slide(prs, 2)  # slide 10 -> section divider
    duplicate_slide(prs, 4)  # slide 11 -> two-concept comparison
    duplicate_slide(prs, 3)  # slide 12 -> compact list
    duplicate_slide(prs, 2)  # slide 13 -> section divider

    reorder_slides(prs, [1, 2, 3, 4, 5, 6, 10, 11, 12, 7, 13, 8, 9])

    # Slide 1
    slide = prs.slides[0]
    set_shape_text(slide, 2056, "[[title]]")
    set_shape_text(slide, 2057, "[[subtitle]]")

    # Slide 2
    slide = prs.slides[1]
    set_shape_text(slide, 2097, "[[toc_title]]")
    set_shape_text(slide, 2089, "[[section_1_title]]")
    set_shape_text(slide, 2090, "[[section_1_desc]]")
    set_shape_text(slide, 2093, "[[section_2_title]]")
    set_shape_text(slide, 2091, "[[section_2_desc]]")
    set_shape_text(slide, 2092, "[[section_3_title]]")
    set_shape_text(slide, 2094, "[[section_3_desc]]")
    set_shape_text(slide, 2095, "[[section_4_title]]")
    set_shape_text(slide, 2096, "[[section_4_desc]]")

    # Slide 3
    slide = prs.slides[2]
    set_shape_text(slide, 2106, "[[section_title]]")
    set_shape_text(slide, 2107, "[[section_number]]")
    set_shape_text(slide, 2108, "[[section_subtitle]]")

    # Slide 4
    slide = prs.slides[3]
    set_shape_text(slide, 2114, "[[title]]")
    set_shape_paragraphs(
        slide,
        2113,
        [
            "[[intro_text]]",
            "[[bullet_1]]",
            "[[bullet_2]]",
            "",
            "",
            "",
            "",
        ],
    )

    # Slide 5
    slide = prs.slides[4]
    set_shape_text(slide, 2242, "[[title]]")
    set_shape_text(slide, 2238, "[[concept_1_heading]]")
    set_shape_text(slide, 2240, "[[concept_1_text]]")
    set_shape_text(slide, 2239, "[[concept_2_heading]]")
    set_shape_text(slide, 2241, "[[concept_2_text]]")

    # Slide 6
    slide = prs.slides[5]
    set_shape_text(slide, 2263, "[[title]]")
    set_shape_text(slide, 2260, "[[column_1_heading]]")
    set_shape_text(slide, 2293, "[[column_1_text]]")
    set_shape_text(slide, 2262, "[[column_2_heading]]")
    set_shape_text(slide, 2294, "[[column_2_text]]")
    set_shape_text(slide, 2261, "[[column_3_heading]]")
    set_shape_text(slide, 2295, "[[column_3_text]]")

    # Slide 7
    slide = prs.slides[6]
    set_shape_text(slide, 2106, "[[section_title]]")
    set_shape_text(slide, 2107, "[[section_number]]")
    set_shape_text(slide, 2108, "[[section_subtitle]]")

    # Slide 8
    slide = prs.slides[7]
    set_shape_text(slide, 2242, "[[title]]")
    set_shape_text(slide, 2238, "[[concept_1_heading]]")
    set_shape_text(slide, 2240, "[[concept_1_text]]")
    set_shape_text(slide, 2239, "[[concept_2_heading]]")
    set_shape_text(slide, 2241, "[[concept_2_text]]")

    # Slide 9
    slide = prs.slides[8]
    set_shape_text(slide, 2114, "[[title]]")
    set_shape_paragraphs(
        slide,
        2113,
        [
            "[[intro_text]]",
            "[[bullet_1]]",
            "[[bullet_2]]",
            "",
            "",
            "",
            "",
        ],
    )

    # Slide 10
    slide = prs.slides[9]
    set_shape_text(slide, 2300, "[[title]]")
    set_shape_text(slide, 2301, "[[caption]]")

    # Slide 11
    slide = prs.slides[10]
    set_shape_text(slide, 2106, "[[section_title]]")
    set_shape_text(slide, 2107, "[[section_number]]")
    set_shape_text(slide, 2108, "[[section_subtitle]]")

    # Slide 12
    slide = prs.slides[11]
    set_shape_text(slide, 2399, "[[title]]")
    set_shape_text(slide, 2392, "[[item_1_heading]]")
    set_shape_text(slide, 2401, "[[item_1_text]]")
    set_shape_text(slide, 2396, "[[item_2_heading]]")
    set_shape_text(slide, 2400, "[[item_2_text]]")
    set_shape_text(slide, 2394, "[[item_3_heading]]")
    set_shape_text(slide, 2402, "[[item_3_text]]")
    set_shape_text(slide, 2398, "[[item_4_heading]]")
    set_shape_text(slide, 2403, "[[item_4_text]]")

    # Slide 13
    slide = prs.slides[12]
    set_shape_text(slide, 2454, "[[title]]")
    table = next(shape.table for shape in slide.shapes if getattr(shape, "has_table", False))
    token_grid = [
        ["[[table_stub_header]]", "[[table_column_1]]", "[[table_column_2]]", "[[table_column_3]]"],
        ["[[row_1_label]]", "[[row_1_value_1]]", "[[row_1_value_2]]", "[[row_1_value_3]]"],
        ["[[row_2_label]]", "[[row_2_value_1]]", "[[row_2_value_2]]", "[[row_2_value_3]]"],
        ["[[row_3_label]]", "[[row_3_value_1]]", "[[row_3_value_2]]", "[[row_3_value_3]]"],
    ]
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            cell.text = token_grid[row_index][col_index]

    prs.save(TARGET)
    print(TARGET)


if __name__ == "__main__":
    main()
