from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "templates" / "science" / "science-2.pptx"
TARGET = ROOT / "templates" / "science" / "science-2-tokenized.pptx"


def find_shape(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"Shape topilmadi: slide={slide.slide_id} shape_id={shape_id}")


def set_shape_text(slide, shape_id: int, text: str) -> None:
    find_shape(slide, shape_id).text = text


def set_shape_paragraphs(slide, shape_id: int, values: list[str]) -> None:
    shape = find_shape(slide, shape_id)
    paragraphs = list(shape.text_frame.paragraphs)
    while len(paragraphs) < len(values):
        paragraphs.append(shape.text_frame.add_paragraph())
    for index, paragraph in enumerate(paragraphs):
        paragraph.text = values[index] if index < len(values) else ""


def main() -> None:
    prs = Presentation(SOURCE)

    # Slide 1
    slide = prs.slides[0]
    set_shape_text(slide, 206, "[[title]]")
    set_shape_text(slide, 207, "[[subtitle]]")

    # Slide 2
    slide = prs.slides[1]
    set_shape_text(slide, 224, "[[toc_title]]")
    set_shape_text(slide, 227, "[[section_1_title]]")
    set_shape_text(slide, 225, "[[section_1_desc]]")
    set_shape_text(slide, 230, "[[section_2_title]]")
    set_shape_text(slide, 228, "[[section_2_desc]]")
    set_shape_text(slide, 233, "[[section_3_title]]")
    set_shape_text(slide, 231, "[[section_3_desc]]")
    set_shape_text(slide, 236, "[[section_4_title]]")
    set_shape_text(slide, 234, "[[section_4_desc]]")
    set_shape_text(slide, 239, "[[section_5_title]]")
    set_shape_text(slide, 237, "[[section_5_desc]]")
    set_shape_text(slide, 242, "[[section_6_title]]")
    set_shape_text(slide, 240, "[[section_6_desc]]")

    # Slide 3
    slide = prs.slides[2]
    set_shape_text(slide, 254, "[[section_title]]")
    set_shape_text(slide, 255, "[[section_subtitle]]")
    set_shape_text(slide, 256, "[[section_number]]")

    # Slide 4
    slide = prs.slides[3]
    set_shape_text(slide, 264, "[[title]]")
    set_shape_text(slide, 265, "[[body]]")

    # Slide 5
    slide = prs.slides[4]
    set_shape_text(slide, 272, "[[title]]")
    set_shape_text(slide, 276, "[[question_1]]")
    set_shape_text(slide, 273, "[[text_1]]")
    set_shape_text(slide, 277, "[[question_2]]")
    set_shape_text(slide, 274, "[[text_2]]")
    set_shape_text(slide, 278, "[[question_3]]")
    set_shape_text(slide, 275, "[[text_3]]")

    # Slide 6
    slide = prs.slides[5]
    set_shape_text(slide, 284, "[[quote]]")
    set_shape_text(slide, 285, "[[author]]")

    # Slide 7
    slide = prs.slides[6]
    set_shape_text(slide, 336, "[[title]]")
    set_shape_text(slide, 337, "[[current_situation_heading]]")
    set_shape_text(slide, 338, "[[current_situation_text]]")
    set_shape_text(slide, 341, "[[problem_1_heading]]")
    set_shape_text(slide, 342, "[[problem_1_text]]")
    set_shape_text(slide, 345, "[[problem_2_heading]]")
    set_shape_text(slide, 346, "[[problem_2_text]]")
    set_shape_text(slide, 349, "[[problem_3_heading]]")
    set_shape_text(slide, 350, "[[problem_3_text]]")

    # Slide 8
    slide = prs.slides[7]
    set_shape_text(slide, 358, "[[section_title]]")
    set_shape_text(slide, 359, "[[section_subtitle]]")
    set_shape_text(slide, 360, "[[section_number]]")

    # Slide 9
    slide = prs.slides[8]
    set_shape_text(slide, 368, "[[title]]")
    set_shape_text(slide, 372, "[[hypothesis_1_heading]]")
    set_shape_text(slide, 369, "[[hypothesis_1_text]]")
    set_shape_text(slide, 373, "[[hypothesis_2_heading]]")
    set_shape_text(slide, 370, "[[hypothesis_2_text]]")
    set_shape_text(slide, 374, "[[hypothesis_3_heading]]")
    set_shape_text(slide, 371, "[[hypothesis_3_text]]")

    # Slide 10
    slide = prs.slides[9]
    set_shape_text(slide, 386, "[[title]]")
    set_shape_text(slide, 390, "[[objective_1_heading]]")
    set_shape_text(slide, 387, "[[objective_1_text]]")
    set_shape_text(slide, 391, "[[objective_2_heading]]")
    set_shape_text(slide, 388, "[[objective_2_text]]")
    set_shape_text(slide, 392, "[[objective_3_heading]]")
    set_shape_text(slide, 389, "[[objective_3_text]]")
    set_shape_text(slide, 395, "[[objective_4_heading]]")
    set_shape_text(slide, 394, "[[objective_4_text]]")

    # Slide 11
    slide = prs.slides[10]
    set_shape_text(slide, 400, "[[title]]")
    set_shape_text(slide, 405, "[[concept_1_heading]]")
    set_shape_text(slide, 401, "[[concept_1_text]]")
    set_shape_text(slide, 404, "[[concept_2_heading]]")
    set_shape_text(slide, 402, "[[concept_2_text]]")
    set_shape_text(slide, 403, "[[concept_3_heading]]")
    set_shape_text(slide, 406, "[[concept_3_text]]")
    set_shape_text(slide, 409, "[[concept_4_heading]]")
    set_shape_text(slide, 407, "[[concept_4_text]]")
    set_shape_text(slide, 410, "[[concept_5_heading]]")
    set_shape_text(slide, 408, "[[concept_5_text]]")

    # Slide 12
    slide = prs.slides[11]
    set_shape_text(slide, 417, "[[title]]")
    set_shape_paragraphs(
        slide,
        418,
        [
            "[[entry_1_ref]]",
            "[[entry_1_note]]",
            "[[entry_2_ref]]",
            "[[entry_2_note]]",
            "[[entry_3_ref]]",
            "[[entry_3_note]]",
            "[[entry_4_ref]]",
            "[[entry_4_note]]",
            "[[entry_5_ref]]",
            "[[entry_5_note]]",
        ],
    )

    # Slide 13
    slide = prs.slides[12]
    set_shape_text(slide, 424, "[[title]]")
    set_shape_text(slide, 425, "[[key_terms_heading]]")
    set_shape_paragraphs(
        slide,
        426,
        [
            "[[key_term_1]]",
            "[[key_term_2]]",
            "[[key_term_3]]",
            "[[key_term_4]]",
            "[[key_term_5]]",
            "[[key_term_6]]",
            "[[key_term_7]]",
        ],
    )
    set_shape_text(slide, 431, "[[theories_heading]]")
    set_shape_text(slide, 427, "[[theory_1_heading]]")
    set_shape_text(slide, 428, "[[theory_1_text]]")
    set_shape_text(slide, 429, "[[theory_2_heading]]")
    set_shape_text(slide, 430, "[[theory_2_text]]")
    set_shape_text(slide, 432, "[[framework_heading]]")
    set_shape_text(slide, 435, "[[framework_text]]")

    # Slide 14
    slide = prs.slides[13]
    set_shape_text(slide, 458, "[[title]]")
    set_shape_text(slide, 461, "[[method_label]]")
    set_shape_text(slide, 462, "[[data_type_label]]")
    set_shape_text(slide, 459, "[[data_type_text]]")
    set_shape_text(slide, 463, "[[motives_label]]")
    set_shape_text(slide, 460, "[[motives_text]]")
    set_shape_text(slide, 466, "[[collection_label]]")
    set_shape_text(slide, 464, "[[collection_text]]")
    set_shape_text(slide, 467, "[[sampling_label]]")
    set_shape_text(slide, 465, "[[sampling_text]]")

    # Slide 15
    slide = prs.slides[14]
    set_shape_text(slide, 477, "[[title]]")
    set_shape_text(slide, 480, "[[phase_1_heading]]")
    set_shape_paragraphs(slide, 479, ["[[phase_1_point_1]]", "[[phase_1_point_2]]", "[[phase_1_point_3]]"])
    set_shape_text(slide, 481, "[[phase_2_heading]]")
    set_shape_paragraphs(slide, 478, ["[[phase_2_point_1]]", "[[phase_2_point_2]]", "[[phase_2_point_3]]"])

    # Slide 16
    slide = prs.slides[15]
    set_shape_text(slide, 486, "[[title]]")
    set_shape_text(slide, 492, "[[graph_caption]]")
    set_shape_text(slide, 488, "[[note_1_heading]]")
    set_shape_text(slide, 489, "[[note_1_text]]")
    set_shape_text(slide, 490, "[[note_2_heading]]")
    set_shape_text(slide, 491, "[[note_2_text]]")
    set_shape_text(slide, 487, "")

    # Slide 17
    slide = prs.slides[16]
    set_shape_text(slide, 503, "[[title]]")
    set_shape_text(slide, 502, "[[highlight_percent]]")
    set_shape_text(slide, 509, "[[result_1_heading]]")
    set_shape_text(slide, 510, "[[result_1_text]]")
    set_shape_text(slide, 511, "[[result_2_heading]]")
    set_shape_text(slide, 512, "[[result_2_text]]")
    set_shape_text(slide, 504, "[[result_3_heading]]")
    set_shape_text(slide, 505, "[[result_3_text]]")
    set_shape_text(slide, 507, "")

    # Slide 18
    slide = prs.slides[17]
    set_shape_text(slide, 599, "[[title]]")
    set_shape_text(slide, 601, "[[intro_heading]]")
    set_shape_paragraphs(
        slide,
        600,
        [
            "[[intro_text]]",
            "",
            "[[bullet_1]]",
            "[[bullet_2]]",
            "[[bullet_3]]",
            "[[bullet_4]]",
        ],
    )
    set_shape_text(slide, 602, "[[discussion_1_heading]]")
    set_shape_text(slide, 603, "[[discussion_1_text]]")
    set_shape_text(slide, 604, "[[discussion_2_heading]]")
    set_shape_text(slide, 605, "[[discussion_2_text]]")

    # Slide 19
    slide = prs.slides[18]
    set_shape_text(slide, 614, "[[title]]")
    set_shape_text(slide, 615, "[[step_1_heading]]")
    set_shape_text(slide, 616, "[[step_1_text]]")
    set_shape_text(slide, 617, "[[step_2_heading]]")
    set_shape_text(slide, 618, "[[step_2_text]]")
    set_shape_text(slide, 619, "[[step_3_heading]]")
    set_shape_text(slide, 620, "[[step_3_text]]")
    set_shape_text(slide, 621, "[[step_4_heading]]")
    set_shape_text(slide, 622, "[[step_4_text]]")
    set_shape_text(slide, 623, "[[step_5_heading]]")
    set_shape_text(slide, 624, "[[step_5_text]]")

    # Slide 20
    slide = prs.slides[19]
    set_shape_text(slide, 643, "[[title]]")
    set_shape_text(slide, 648, "[[point_1_heading]]")
    set_shape_text(slide, 644, "[[point_1_text]]")
    set_shape_text(slide, 649, "[[point_2_heading]]")
    set_shape_text(slide, 645, "[[point_2_text]]")
    set_shape_text(slide, 650, "[[point_3_heading]]")
    set_shape_text(slide, 646, "[[point_3_text]]")
    set_shape_text(slide, 651, "[[point_4_heading]]")
    set_shape_text(slide, 647, "[[point_4_text]]")

    prs.save(TARGET)
    print(TARGET)


if __name__ == "__main__":
    main()
