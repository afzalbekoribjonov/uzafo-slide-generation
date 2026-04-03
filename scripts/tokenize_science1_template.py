from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "templates" / "science" / "science-1.pptx"
TARGET = ROOT / "templates" / "science" / "science-1-tokenized.pptx"


def find_shape(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"Shape topilmadi: slide={slide.slide_id} shape_id={shape_id}")


def set_shape_text(slide, shape_id: int, text: str) -> None:
    find_shape(slide, shape_id).text = text


def set_shape_paragraphs(slide, shape_id: int, values: list[str]) -> None:
    shape = find_shape(slide, shape_id)
    paragraphs = list(shape.text_frame.paragraphs)
    while len(paragraphs) < len(values):
        paragraphs.append(shape.text_frame.add_paragraph())
    for index, paragraph in enumerate(paragraphs):
        paragraph.text = values[index] if index < len(values) else ""


def set_table_tokens(slide, shape_id: int, values: list[list[str]]) -> None:
    table = find_shape(slide, shape_id).table
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            token = values[row_index][col_index]
            if token is not None:
                cell.text = token


def main() -> None:
    prs = Presentation(SOURCE)

    # Slide 1
    slide = prs.slides[0]
    set_shape_text(slide, 128, "[[title]]")
    set_shape_text(slide, 129, "[[subtitle]]")

    # Slide 3
    slide = prs.slides[2]
    set_shape_text(slide, 140, "[[toc_title]]")
    set_shape_text(slide, 141, "[[section_1_title]]")
    set_shape_text(slide, 142, "[[section_1_desc]]")
    set_shape_text(slide, 143, "[[section_2_title]]")
    set_shape_text(slide, 144, "[[section_2_desc]]")
    set_shape_text(slide, 145, "[[section_3_title]]")
    set_shape_text(slide, 146, "[[section_3_desc]]")
    set_shape_text(slide, 147, "[[section_4_title]]")
    set_shape_text(slide, 148, "[[section_4_desc]]")
    set_shape_text(slide, 149, "[[section_5_title]]")
    set_shape_text(slide, 150, "[[section_5_desc]]")

    # Slide 4
    slide = prs.slides[3]
    set_shape_text(slide, 160, "[[section_number]]")
    set_shape_text(slide, 161, "[[section_title]]")
    set_shape_text(slide, 162, "[[section_subtitle]]")

    # Slide 5
    slide = prs.slides[4]
    set_shape_text(slide, 167, "[[title]]")
    set_shape_text(slide, 169, "[[body]]")

    # Slide 6
    slide = prs.slides[5]
    set_shape_text(slide, 216, "[[title]]")
    set_shape_text(slide, 222, "[[highlight_1_heading]]")
    set_shape_text(slide, 223, "[[highlight_1_text]]")
    set_shape_text(slide, 224, "[[highlight_2_heading]]")
    set_shape_text(slide, 225, "[[highlight_2_text]]")
    set_shape_text(slide, 229, "[[highlight_3_heading]]")
    set_shape_text(slide, 230, "[[highlight_3_text]]")

    # Slide 7
    slide = prs.slides[6]
    set_shape_text(slide, 236, "[[title]]")
    set_shape_text(slide, 237, "[[caption]]")

    # Slide 8
    slide = prs.slides[7]
    set_shape_text(slide, 252, "[[section_number]]")
    set_shape_text(slide, 250, "[[section_title]]")
    set_shape_text(slide, 251, "[[section_subtitle]]")

    # Slide 9
    slide = prs.slides[8]
    set_shape_text(slide, 257, "[[title]]")
    set_shape_text(slide, 264, "[[topic_1_heading]]")
    set_shape_text(slide, 267, "[[topic_1_text]]")
    set_shape_text(slide, 268, "[[topic_2_heading]]")
    set_shape_text(slide, 269, "[[topic_2_text]]")
    set_shape_text(slide, 265, "[[topic_3_heading]]")
    set_shape_text(slide, 266, "[[topic_3_text]]")

    # Slide 10
    slide = prs.slides[9]
    set_shape_text(slide, 277, "[[title]]")
    set_shape_text(slide, 290, "[[year_1]]")
    set_shape_text(slide, 278, "[[year_1_heading]]")
    set_shape_text(slide, 279, "[[year_1_text]]")
    set_shape_text(slide, 291, "[[year_2]]")
    set_shape_text(slide, 282, "[[year_2_heading]]")
    set_shape_text(slide, 283, "[[year_2_text]]")
    set_shape_text(slide, 292, "[[year_3]]")
    set_shape_text(slide, 280, "[[year_3_heading]]")
    set_shape_text(slide, 281, "[[year_3_text]]")
    set_shape_text(slide, 284, "[[focus_heading]]")
    set_shape_text(slide, 285, "[[focus_text]]")

    # Slide 11
    slide = prs.slides[10]
    set_shape_text(slide, 298, "[[title]]")
    set_shape_text(slide, 309, "[[item_1_heading]]")
    set_shape_text(slide, 310, "[[item_1_text]]")
    set_shape_text(slide, 301, "[[item_2_heading]]")
    set_shape_text(slide, 299, "[[item_2_text]]")
    set_shape_text(slide, 304, "[[item_3_heading]]")
    set_shape_text(slide, 300, "[[item_3_text]]")
    set_shape_text(slide, 305, "[[item_4_heading]]")
    set_shape_text(slide, 306, "[[item_4_text]]")
    set_shape_text(slide, 302, "[[item_5_heading]]")
    set_shape_text(slide, 303, "[[item_5_text]]")
    set_shape_text(slide, 307, "[[item_6_heading]]")
    set_shape_text(slide, 308, "[[item_6_text]]")

    # Slide 12
    slide = prs.slides[11]
    set_shape_text(slide, 323, "[[title]]")
    set_shape_paragraphs(
        slide,
        324,
        [
            "",
            "[[item_1]]",
            "[[item_2]]",
            "[[item_3]]",
            "[[item_4]]",
            "[[item_5]]",
            "[[item_6]]",
            "[[item_7]]",
            "[[item_8]]",
        ],
    )

    # Slide 13
    slide = prs.slides[12]
    set_shape_text(slide, 352, "[[title]]")
    set_shape_text(slide, 363, "[[review_1_text]]")
    set_shape_text(slide, 366, "[[review_1_author]]")
    set_shape_text(slide, 362, "[[review_2_text]]")
    set_shape_text(slide, 365, "[[review_2_author]]")
    set_shape_text(slide, 368, "[[review_3_text]]")
    set_shape_text(slide, 367, "[[review_3_author]]")
    set_shape_text(slide, 361, "[[review_4_text]]")
    set_shape_text(slide, 364, "[[review_4_author]]")

    # Slide 14
    slide = prs.slides[13]
    set_shape_text(slide, 410, "[[title]]")
    set_shape_text(slide, 411, "[[column_1_heading]]")
    set_shape_text(slide, 412, "[[column_1_text]]")
    set_shape_text(slide, 413, "[[column_2_heading]]")
    set_shape_text(slide, 414, "[[column_2_text]]")
    set_shape_text(slide, 415, "[[column_3_heading]]")
    set_shape_text(slide, 416, "[[column_3_text]]")

    # Slide 15
    slide = prs.slides[14]
    set_shape_text(slide, 449, "[[title]]")
    set_table_tokens(
        slide,
        448,
        [
            ["[[header_project]]", "[[header_metric_a]]", "[[header_metric_b]]", "[[header_metric_c]]"],
            ["[[row_1_project]]", "[[row_1_metric_a]]", "[[row_1_metric_b]]", "[[row_1_metric_c]]"],
            ["[[row_2_project]]", "[[row_2_metric_a]]", "[[row_2_metric_b]]", "[[row_2_metric_c]]"],
            ["[[row_3_project]]", "[[row_3_metric_a]]", "[[row_3_metric_b]]", "[[row_3_metric_c]]"],
        ],
    )

    # Slide 16
    slide = prs.slides[15]
    set_shape_text(slide, 733, "[[title]]")
    set_shape_text(slide, 741, "[[person_1_name]]")
    set_shape_text(slide, 740, "[[person_1_role]]")
    set_shape_text(slide, 738, "[[person_2_name]]")
    set_shape_text(slide, 739, "[[person_2_role]]")

    # Slide 17
    slide = prs.slides[16]
    set_shape_text(slide, 755, "[[title]]")
    set_shape_text(slide, 765, "[[event_1_heading]]")
    set_shape_text(slide, 764, "[[event_1_text]]")
    set_shape_text(slide, 763, "[[event_2_heading]]")
    set_shape_text(slide, 762, "[[event_2_text]]")
    set_shape_text(slide, 757, "[[event_3_heading]]")
    set_shape_text(slide, 756, "[[event_3_text]]")
    set_shape_text(slide, 759, "[[event_4_heading]]")
    set_shape_text(slide, 758, "[[event_4_text]]")
    set_shape_text(slide, 761, "[[event_5_heading]]")
    set_shape_text(slide, 760, "[[event_5_text]]")

    prs.save(TARGET)
    print(TARGET)


if __name__ == "__main__":
    main()
