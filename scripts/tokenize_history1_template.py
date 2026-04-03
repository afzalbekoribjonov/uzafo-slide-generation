from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "templates" / "history" / "history-1.pptx"
TARGET = ROOT / "templates" / "history" / "history-1-tokenized.pptx"


def find_shape(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"Shape topilmadi: slide={slide.slide_id} shape_id={shape_id}")


def set_shape_text(slide, shape_id: int, text: str) -> None:
    find_shape(slide, shape_id).text = text


def set_shape_paragraphs(slide, shape_id: int, values: list[str]) -> None:
    shape = find_shape(slide, shape_id)
    paragraphs = list(shape.text_frame.paragraphs)
    while len(paragraphs) < len(values):
        paragraphs.append(shape.text_frame.add_paragraph())
    for index, paragraph in enumerate(paragraphs):
        paragraph.text = values[index] if index < len(values) else ""


def set_table_tokens(slide, shape_id: int, values: list[list[str]]) -> None:
    table = find_shape(slide, shape_id).table
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            token = values[row_index][col_index]
            if token is not None:
                cell.text = token


def main() -> None:
    prs = Presentation(SOURCE)

    slide = prs.slides[0]
    set_shape_paragraphs(slide, 220, ["[[title_line_1]]", "[[title_line_2]]"])
    set_shape_text(slide, 221, "[[subtitle]]")

    slide = prs.slides[2]
    set_shape_text(slide, 470, "[[toc_title]]")
    set_shape_text(slide, 471, "[[section_1_title]]")
    set_shape_text(slide, 473, "[[section_1_desc]]")
    set_shape_text(slide, 474, "[[section_2_title]]")
    set_shape_text(slide, 476, "[[section_2_desc]]")
    set_shape_text(slide, 477, "[[section_3_title]]")
    set_shape_text(slide, 479, "[[section_3_desc]]")
    set_shape_text(slide, 480, "[[section_4_title]]")
    set_shape_text(slide, 482, "[[section_4_desc]]")

    slide = prs.slides[3]
    set_shape_text(slide, 489, "[[section_title]]")
    set_shape_text(slide, 490, "[[section_number]]")
    set_shape_paragraphs(slide, 491, ["[[section_subtitle_line_1]]", "[[section_subtitle_line_2]]"])

    slide = prs.slides[4]
    set_shape_text(slide, 4589, "[[title]]")
    set_shape_text(slide, 4590, "[[body]]")

    slide = prs.slides[5]
    set_shape_text(slide, 8688, "[[quote]]")
    set_shape_text(slide, 8687, "[[author]]")

    slide = prs.slides[6]
    set_shape_text(slide, 8702, "[[title]]")

    slide = prs.slides[7]
    set_shape_text(slide, 8932, "[[title]]")
    set_shape_text(slide, 8933, "[[card_1_heading]]")
    set_shape_text(slide, 8934, "[[card_1_text]]")
    set_shape_text(slide, 8935, "[[card_2_heading]]")
    set_shape_text(slide, 8936, "[[card_2_text]]")
    set_shape_text(slide, 8937, "[[card_3_heading]]")
    set_shape_text(slide, 8938, "[[card_3_text]]")

    slide = prs.slides[8]
    set_shape_text(slide, 9133, "[[title]]")
    set_shape_text(slide, 9134, "[[step_1_heading]]")
    set_shape_text(slide, 9135, "[[step_1_text]]")
    set_shape_text(slide, 9136, "[[step_2_heading_line_1]]")
    set_shape_text(slide, 9137, "[[step_2_text]]")
    set_shape_text(slide, 9154, "[[step_3_heading_line_1]]")
    set_shape_text(slide, 9155, "[[step_3_text]]")

    slide = prs.slides[9]
    set_shape_text(slide, 9200, "[[title]]")
    set_shape_text(slide, 9201, "[[step_4_heading_line_1]]")
    set_shape_text(slide, 9202, "[[step_4_text]]")
    set_shape_text(slide, 9203, "[[step_5_heading_line_1]]")
    set_shape_text(slide, 9204, "[[step_5_text]]")

    slide = prs.slides[10]
    set_shape_text(slide, 9260, "[[title]]")
    set_shape_text(slide, 9261, "[[left_heading]]")
    set_shape_text(slide, 9262, "[[right_heading]]")
    set_shape_text(slide, 9263, "[[left_text]]")
    set_shape_text(slide, 9264, "[[right_text]]")

    slide = prs.slides[11]
    set_shape_text(slide, 9882, "[[title]]")
    set_shape_text(slide, 10063, "[[label_1]]")
    set_shape_text(slide, 10064, "[[label_2]]")
    set_shape_text(slide, 10065, "[[label_3]]")
    set_shape_text(slide, 10067, "[[label_4]]")
    set_shape_text(slide, 10071, "[[label_5]]")
    set_shape_text(slide, 10072, "[[label_6]]")

    slide = prs.slides[12]
    set_shape_text(slide, 10369, "[[section_title]]")
    set_shape_text(slide, 10370, "[[section_number]]")
    set_shape_paragraphs(slide, 10371, ["[[section_subtitle_line_1]]", "[[section_subtitle_line_2]]"])

    slide = prs.slides[13]
    set_shape_text(slide, 11884, "[[title]]")
    set_shape_text(slide, 11885, "[[card_1_heading]]")
    set_shape_text(slide, 11886, "[[card_1_text]]")
    set_shape_text(slide, 11887, "[[card_2_heading]]")
    set_shape_text(slide, 11888, "[[card_2_text]]")
    set_shape_text(slide, 11889, "[[card_3_heading]]")
    set_shape_text(slide, 11890, "[[card_3_text]]")
    set_shape_text(slide, 11891, "[[card_4_heading]]")
    set_shape_text(slide, 11892, "[[card_4_text]]")

    slide = prs.slides[14]
    set_shape_text(slide, 13119, "[[title]]")
    set_table_tokens(
        slide,
        13120,
        [
            ["[[row_1_label]]", "[[row_1_text]]"],
            ["[[row_2_label]]", "[[row_2_text]]"],
            ["[[row_3_label]]", "[[row_3_text]]"],
            ["[[row_4_label]]", "[[row_4_text]]"],
        ],
    )

    slide = prs.slides[15]
    set_shape_text(slide, 13303, "[[title]]")
    set_shape_text(slide, 13304, "")
    set_shape_text(slide, 13305, "[[point_1_label]]")
    set_shape_text(slide, 13306, "[[point_2_label]]")
    set_shape_text(slide, 13307, "[[point_3_label]]")
    set_shape_text(slide, 13308, "[[point_4_label]]")

    slide = prs.slides[16]
    set_shape_text(slide, 13494, "[[title]]")

    slide = prs.slides[17]
    set_shape_text(slide, 13501, "[[title]]")
    set_shape_text(slide, 13502, "[[era_label]]")
    set_shape_text(slide, 13503, "[[intro_text]]")
    set_shape_text(slide, 13690, "[[detail_1_heading]]")
    set_shape_text(slide, 13691, "[[detail_1_text]]")
    set_shape_text(slide, 13693, "[[detail_2_heading]]")
    set_shape_text(slide, 13694, "[[detail_2_text]]")

    slide = prs.slides[18]
    set_shape_text(slide, 13704, "[[title]]")
    set_shape_text(slide, 13715, "[[stage_1]]")
    set_shape_text(slide, 13716, "[[stage_2]]")
    set_shape_text(slide, 13717, "[[stage_3]]")

    slide = prs.slides[19]
    set_shape_text(slide, 14106, "[[event_1_year]]")
    set_shape_text(slide, 14109, "[[event_1_text]]")
    set_shape_text(slide, 14108, "[[event_2_year]]")
    set_shape_text(slide, 14111, "[[event_2_text]]")
    set_shape_text(slide, 14107, "[[event_3_year]]")
    set_shape_text(slide, 14110, "[[event_3_text]]")

    slide = prs.slides[20]
    set_shape_text(slide, 14284, "[[title]]")
    set_shape_text(slide, 14285, "[[region_1_name]]")
    set_shape_text(slide, 14286, "[[region_1_text]]")
    set_shape_text(slide, 14287, "[[region_2_name]]")
    set_shape_text(slide, 14288, "[[region_2_text]]")

    slide = prs.slides[21]
    set_shape_text(slide, 14316, "[[highlight_value]]")
    set_shape_text(slide, 14317, "[[highlight_text]]")

    slide = prs.slides[22]
    set_shape_text(slide, 14373, "[[title]]")
    set_shape_text(slide, 14374, "[[point_1_heading]]")
    set_shape_text(slide, 14375, "[[point_1_text]]")
    set_shape_text(slide, 14376, "[[point_2_heading]]")
    set_shape_text(slide, 14377, "[[point_2_text]]")
    set_shape_text(slide, 14378, "[[point_3_heading]]")
    set_shape_text(slide, 14379, "[[point_3_text]]")
    set_shape_text(slide, 14380, "[[point_4_heading]]")
    set_shape_text(slide, 14381, "[[point_4_text]]")
    set_shape_text(slide, 14382, "[[point_5_heading]]")
    set_shape_text(slide, 14383, "[[point_5_text]]")
    set_shape_text(slide, 14384, "[[point_6_heading]]")
    set_shape_text(slide, 14385, "[[point_6_text]]")

    slide = prs.slides[23]
    set_shape_text(slide, 15793, "[[title]]")
    set_shape_text(slide, 15795, "[[event_1_year]]")
    set_shape_text(slide, 15796, "[[event_1_heading]]")
    set_shape_text(slide, 15797, "[[event_1_text]]")
    set_shape_text(slide, 15798, "[[event_2_year]]")
    set_shape_text(slide, 15800, "[[event_2_heading]]")
    set_shape_text(slide, 15801, "[[event_2_text]]")
    set_shape_text(slide, 15803, "[[event_3_year]]")
    set_shape_text(slide, 15804, "[[event_3_heading]]")
    set_shape_text(slide, 15805, "[[event_3_text]]")

    slide = prs.slides[24]
    set_shape_text(slide, 17039, "[[title]]")
    set_shape_text(slide, 17041, "[[event_4_year]]")
    set_shape_text(slide, 17042, "[[event_4_heading]]")
    set_shape_text(slide, 17043, "[[event_4_text]]")
    set_shape_text(slide, 17044, "[[event_5_year]]")
    set_shape_text(slide, 17046, "[[event_5_heading]]")
    set_shape_text(slide, 17047, "[[event_5_text]]")
    set_shape_text(slide, 17049, "[[event_6_year]]")
    set_shape_text(slide, 17050, "[[event_6_heading]]")
    set_shape_text(slide, 17051, "[[event_6_text]]")

    slide = prs.slides[25]
    set_shape_text(slide, 17086, "[[title]]")
    set_shape_text(slide, 17087, "[[body]]")

    slide = prs.slides[26]
    set_shape_text(slide, 17141, "[[title]]")
    set_shape_paragraphs(slide, 17142, ["[[entry_1]]", "[[entry_2]]", "[[entry_3]]", "[[entry_4]]"])

    prs.save(TARGET)
    print(TARGET)


if __name__ == "__main__":
    main()
