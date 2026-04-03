from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "templates" / "business" / "business-1.pptx"
TARGET = ROOT / "templates" / "business" / "business-1-tokenized.pptx"


def find_shape(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"Shape topilmadi: slide={slide.slide_id} shape_id={shape_id}")


def set_shape_text(slide, shape_id: int, text: str) -> None:
    find_shape(slide, shape_id).text = text


def set_shape_paragraphs(slide, shape_id: int, values: list[str]) -> None:
    shape = find_shape(slide, shape_id)
    paragraphs = list(shape.text_frame.paragraphs)
    while len(paragraphs) < len(values):
        paragraphs.append(shape.text_frame.add_paragraph())
    for index, paragraph in enumerate(paragraphs):
        paragraph.text = values[index] if index < len(values) else ""


def main() -> None:
    prs = Presentation(SOURCE)

    # Slide 1
    slide = prs.slides[0]
    set_shape_text(slide, 932, "[[title]]")
    set_shape_text(slide, 933, "[[subtitle]]")

    # Slide 2
    slide = prs.slides[1]
    set_shape_text(slide, 992, "[[toc_title]]")
    set_shape_text(slide, 997, "[[section_1_title]]")
    set_shape_text(slide, 999, "[[section_2_title]]")
    set_shape_text(slide, 1000, "[[section_3_title]]")
    set_shape_text(slide, 998, "[[section_4_title]]")
    set_shape_text(slide, 1003, "[[section_5_title]]")
    set_shape_text(slide, 1004, "[[section_6_title]]")

    # Slide 3
    slide = prs.slides[2]
    set_shape_text(slide, 1078, "[[title]]")
    set_shape_text(slide, 1079, "[[body]]")

    # Slide 4
    slide = prs.slides[3]
    set_shape_text(slide, 1120, "[[title]]")
    set_shape_text(slide, 1121, "[[paragraph_1]]")
    set_shape_text(slide, 1122, "[[paragraph_2]]")

    # Slide 5
    slide = prs.slides[4]
    set_shape_text(slide, 1127, "[[title]]")
    set_shape_paragraphs(
        slide,
        1128,
        [
            "[[intro_text]]",
            "[[bullet_1]]",
            "[[bullet_2]]",
            "[[bullet_3]]",
            "[[bullet_4]]",
        ],
    )

    # Slide 6
    slide = prs.slides[5]
    set_shape_text(slide, 1151, "[[title]]")
    set_shape_text(slide, 1154, "[[left_heading]]")
    set_shape_text(slide, 1152, "[[left_text]]")
    set_shape_text(slide, 1155, "[[right_heading]]")
    set_shape_text(slide, 1153, "[[right_text]]")

    # Slide 7
    slide = prs.slides[6]
    set_shape_text(slide, 1166, "[[title]]")
    set_shape_text(slide, 1169, "[[theme_1_heading]]")
    set_shape_text(slide, 1167, "[[theme_1_text]]")
    set_shape_text(slide, 1170, "[[theme_2_heading]]")
    set_shape_text(slide, 1168, "[[theme_2_text]]")
    set_shape_text(slide, 1172, "[[theme_3_heading]]")
    set_shape_text(slide, 1171, "[[theme_3_text]]")

    # Slide 8
    slide = prs.slides[7]
    set_shape_text(slide, 1194, "[[title]]")
    set_shape_text(slide, 1197, "[[pillar_1_heading]]")
    set_shape_text(slide, 1195, "[[pillar_1_text]]")
    set_shape_text(slide, 1198, "[[pillar_2_heading]]")
    set_shape_text(slide, 1196, "[[pillar_2_text]]")
    set_shape_text(slide, 1201, "[[pillar_3_heading]]")
    set_shape_text(slide, 1199, "[[pillar_3_text]]")
    set_shape_text(slide, 1202, "[[pillar_4_heading]]")
    set_shape_text(slide, 1200, "[[pillar_4_text]]")

    # Slide 9
    slide = prs.slides[8]
    set_shape_text(slide, 1241, "[[title]]")
    set_shape_text(slide, 1244, "[[item_1_heading]]")
    set_shape_text(slide, 1242, "[[item_1_text]]")
    set_shape_text(slide, 1245, "[[item_2_heading]]")
    set_shape_text(slide, 1243, "[[item_2_text]]")
    set_shape_text(slide, 1247, "[[item_3_heading]]")
    set_shape_text(slide, 1246, "[[item_3_text]]")
    set_shape_text(slide, 1250, "[[item_4_heading]]")
    set_shape_text(slide, 1248, "[[item_4_text]]")
    set_shape_text(slide, 1251, "[[item_5_heading]]")
    set_shape_text(slide, 1249, "[[item_5_text]]")
    set_shape_text(slide, 1253, "[[item_6_heading]]")
    set_shape_text(slide, 1252, "[[item_6_text]]")

    # Slide 10
    slide = prs.slides[9]
    set_shape_text(slide, 1283, "[[quote]]")
    set_shape_text(slide, 1282, "[[author]]")

    # Slide 11
    slide = prs.slides[10]
    set_shape_text(slide, 1399, "[[title]]")
    set_shape_text(slide, 1400, "[[body]]")

    # Slide 12
    slide = prs.slides[11]
    set_shape_text(slide, 1464, "[[title]]")
    set_shape_text(slide, 1467, "[[step_1_heading]]")
    set_shape_text(slide, 1465, "[[step_1_text]]")
    set_shape_text(slide, 1468, "[[step_2_heading]]")
    set_shape_text(slide, 1466, "[[step_2_text]]")
    set_shape_text(slide, 1470, "[[step_3_heading]]")
    set_shape_text(slide, 1469, "[[step_3_text]]")
    set_shape_text(slide, 1478, "[[step_4_heading]]")
    set_shape_text(slide, 1477, "[[step_4_text]]")
    set_shape_text(slide, 1473, "[[step_5_heading]]")
    set_shape_text(slide, 1471, "[[step_5_text]]")
    set_shape_text(slide, 1474, "[[step_6_heading]]")
    set_shape_text(slide, 1472, "[[step_6_text]]")
    set_shape_text(slide, 1476, "[[step_7_heading]]")
    set_shape_text(slide, 1475, "[[step_7_text]]")
    set_shape_text(slide, 1480, "[[step_8_heading]]")
    set_shape_text(slide, 1479, "[[step_8_text]]")

    # Slide 13
    slide = prs.slides[12]
    set_shape_text(slide, 1500, "[[title]]")
    set_shape_text(slide, 1503, "[[block_1_heading]]")
    set_shape_text(slide, 1501, "[[block_1_text]]")
    set_shape_text(slide, 1504, "[[block_2_heading]]")
    set_shape_text(slide, 1502, "[[block_2_text]]")
    set_shape_text(slide, 1507, "[[block_3_heading]]")
    set_shape_text(slide, 1505, "[[block_3_text]]")
    set_shape_text(slide, 1508, "[[block_4_heading]]")
    set_shape_text(slide, 1506, "[[block_4_text]]")

    # Slide 14
    slide = prs.slides[13]
    set_shape_text(slide, 1564, "[[title]]")
    set_shape_text(slide, 1569, "[[catalog_kicker]]")
    set_shape_text(slide, 1566, "[[item_1_heading]]")
    set_shape_text(slide, 1565, "[[item_1_text]]")
    set_shape_text(slide, 1568, "[[item_2_heading]]")
    set_shape_text(slide, 1567, "[[item_2_text]]")
    set_shape_text(slide, 1571, "[[item_3_heading]]")
    set_shape_text(slide, 1570, "[[item_3_text]]")
    set_shape_text(slide, 1573, "[[item_4_heading]]")
    set_shape_text(slide, 1572, "[[item_4_text]]")
    set_shape_text(slide, 1575, "[[item_5_heading]]")
    set_shape_text(slide, 1574, "[[item_5_text]]")
    set_shape_text(slide, 1577, "[[item_6_heading]]")
    set_shape_text(slide, 1576, "[[item_6_text]]")
    set_shape_text(slide, 1579, "[[item_7_heading]]")
    set_shape_text(slide, 1578, "[[item_7_text]]")
    set_shape_text(slide, 1581, "[[item_8_heading]]")
    set_shape_text(slide, 1580, "[[item_8_text]]")

    # Slide 15
    slide = prs.slides[14]
    set_shape_text(slide, 1607, "[[title]]")
    set_shape_text(slide, 1609, "[[left_heading]]")
    set_shape_paragraphs(
        slide,
        1608,
        [
            "[[left_point_1]]",
            "[[left_point_2]]",
            "[[left_point_3]]",
        ],
    )
    set_shape_text(slide, 1610, "[[right_heading]]")
    set_shape_text(slide, 1613, "[[milestone_1_label]]")
    set_shape_text(slide, 1612, "[[milestone_1_text]]")
    set_shape_text(slide, 1616, "[[milestone_2_label]]")
    set_shape_text(slide, 1614, "[[milestone_2_text]]")
    set_shape_text(slide, 1617, "[[milestone_3_label]]")
    set_shape_text(slide, 1615, "[[milestone_3_text]]")

    # Slide 16
    slide = prs.slides[15]
    set_shape_text(slide, 1633, "[[title]]")
    set_shape_text(slide, 1636, "[[segment_1_name]]")
    set_shape_text(slide, 1629, "[[segment_1_percent]]")
    set_shape_text(slide, 1635, "[[segment_1_text]]")
    set_shape_text(slide, 1639, "[[segment_2_name]]")
    set_shape_text(slide, 1630, "[[segment_2_percent]]")
    set_shape_text(slide, 1637, "[[segment_2_text]]")
    set_shape_text(slide, 1640, "[[segment_3_name]]")
    set_shape_text(slide, 1631, "[[segment_3_percent]]")
    set_shape_text(slide, 1638, "[[segment_3_text]]")

    # Slide 17
    slide = prs.slides[16]
    set_shape_text(slide, 1650, "[[title]]")
    set_shape_text(slide, 1652, "[[kpi_1_value]]")
    set_shape_text(slide, 1651, "[[kpi_1_text]]")
    set_shape_text(slide, 1654, "[[kpi_2_value]]")
    set_shape_text(slide, 1653, "[[kpi_2_text]]")
    set_shape_text(slide, 1656, "[[kpi_3_value]]")
    set_shape_text(slide, 1655, "[[kpi_3_text]]")
    set_shape_text(slide, 1658, "[[kpi_4_value]]")
    set_shape_text(slide, 1657, "[[kpi_4_text]]")
    set_shape_text(slide, 1660, "[[kpi_5_value]]")
    set_shape_text(slide, 1659, "[[kpi_5_text]]")

    # Slide 18
    slide = prs.slides[17]
    set_shape_text(slide, 1665, "[[title]]")
    set_shape_text(slide, 1671, "[[period_1_name]]")
    set_shape_text(slide, 1668, "[[period_1_text]]")
    set_shape_text(slide, 1670, "[[period_2_name]]")
    set_shape_text(slide, 1669, "[[period_2_text]]")

    prs.save(TARGET)
    print(TARGET)


if __name__ == "__main__":
    main()
