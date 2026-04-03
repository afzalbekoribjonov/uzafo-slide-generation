from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SOURCE = ROOT / "templates" / "history" / "history-2.pptx"
TARGET = ROOT / "templates" / "history" / "history-2-tokenized.pptx"


def find_shape(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise ValueError(f"Shape topilmadi: slide={slide.slide_id} shape_id={shape_id}")


def set_shape_text(slide, shape_id: int, text: str) -> None:
    find_shape(slide, shape_id).text = text


def set_shape_paragraphs(slide, shape_id: int, values: list[str]) -> None:
    shape = find_shape(slide, shape_id)
    paragraphs = list(shape.text_frame.paragraphs)
    while len(paragraphs) < len(values):
        paragraphs.append(shape.text_frame.add_paragraph())
    for index, paragraph in enumerate(paragraphs):
        paragraph.text = values[index] if index < len(values) else ""


def set_table_tokens(slide, shape_id: int, values: list[list[str]]) -> None:
    table = find_shape(slide, shape_id).table
    for row_index, row in enumerate(table.rows):
        for col_index, cell in enumerate(row.cells):
            token = values[row_index][col_index]
            if token is not None:
                cell.text = token


def main() -> None:
    prs = Presentation(SOURCE)

    # Slide 1
    slide = prs.slides[0]
    set_shape_text(slide, 7805, "[[title]]")
    set_shape_text(slide, 7806, "[[subtitle]]")

    # Slide 3
    slide = prs.slides[2]
    set_shape_text(slide, 7820, "[[toc_title]]")
    set_shape_text(slide, 7833, "[[part_1_title]]")
    set_shape_text(slide, 7822, "[[part_1_desc]]")
    set_shape_text(slide, 7836, "[[part_2_title]]")
    set_shape_text(slide, 7824, "[[part_2_desc]]")
    set_shape_text(slide, 7838, "[[part_3_title]]")
    set_shape_text(slide, 7830, "[[part_3_desc]]")
    set_shape_text(slide, 7834, "[[part_4_title]]")
    set_shape_text(slide, 7823, "[[part_4_desc]]")
    set_shape_text(slide, 7835, "[[part_5_title]]")
    set_shape_text(slide, 7821, "[[part_5_desc]]")
    set_shape_text(slide, 7837, "[[part_6_title]]")
    set_shape_text(slide, 7829, "[[part_6_desc]]")

    # Slide 4
    slide = prs.slides[3]
    set_shape_text(slide, 7843, "[[section_title]]")
    set_shape_text(slide, 7844, "[[section_subtitle]]")
    set_shape_text(slide, 7845, "[[section_number]]")

    # Slide 5
    slide = prs.slides[4]
    set_shape_text(slide, 7894, "[[title]]")
    set_shape_text(slide, 7898, "[[point_1_heading]]")
    set_shape_text(slide, 7895, "[[point_1_text]]")
    set_shape_text(slide, 7899, "[[point_2_heading]]")
    set_shape_text(slide, 7896, "[[point_2_text]]")
    set_shape_text(slide, 7900, "[[point_3_heading]]")
    set_shape_text(slide, 7897, "[[point_3_text]]")

    # Slide 6
    slide = prs.slides[5]
    set_shape_text(slide, 7906, "[[title]]")

    # Slide 7
    slide = prs.slides[6]
    set_shape_paragraphs(slide, 8045, ["[[title_line_1]]", "[[title_line_2]]"])
    set_shape_text(slide, 8046, "[[body]]")

    # Slide 8
    slide = prs.slides[7]
    set_shape_text(slide, 8139, "[[quote]]")
    set_shape_text(slide, 8140, "[[author]]")

    # Slide 9
    slide = prs.slides[8]
    set_shape_text(slide, 8145, "[[title]]")
    set_shape_text(slide, 8152, "[[hypothesis_1_heading]]")
    set_shape_text(slide, 8146, "[[hypothesis_1_text]]")
    set_shape_text(slide, 8153, "[[hypothesis_2_heading]]")
    set_shape_text(slide, 8147, "[[hypothesis_2_text]]")
    set_shape_text(slide, 8154, "[[hypothesis_3_heading]]")
    set_shape_text(slide, 8148, "[[hypothesis_3_text]]")

    # Slide 10
    slide = prs.slides[9]
    set_shape_text(slide, 8159, "[[main_value]]")
    set_shape_text(slide, 8160, "[[caption]]")

    # Slide 11
    slide = prs.slides[10]
    set_shape_text(slide, 8166, "[[stat_1_value]]")
    set_shape_text(slide, 8165, "[[stat_1_label]]")
    set_shape_text(slide, 8167, "[[stat_2_value]]")
    set_shape_text(slide, 8168, "[[stat_2_label]]")
    set_shape_text(slide, 8169, "[[stat_3_value]]")
    set_shape_text(slide, 8170, "[[stat_3_label]]")

    # Slide 12
    slide = prs.slides[11]
    set_shape_text(slide, 8175, "[[title]]")

    # Slide 13
    slide = prs.slides[12]
    set_shape_text(slide, 8180, "[[title]]")
    set_shape_text(slide, 8183, "[[situation_heading]]")
    set_shape_text(slide, 8184, "[[situation_text]]")
    set_shape_text(slide, 8189, "[[problem_1_heading]]")
    set_shape_text(slide, 8190, "[[problem_1_text]]")
    set_shape_text(slide, 8194, "[[problem_2_heading]]")
    set_shape_text(slide, 8195, "[[problem_2_text]]")
    set_shape_text(slide, 8199, "[[problem_3_heading]]")
    set_shape_text(slide, 8200, "[[problem_3_text]]")
    set_shape_text(slide, 8205, "[[problems_heading]]")

    # Slide 14
    slide = prs.slides[13]
    set_shape_text(slide, 8325, "[[title]]")
    set_shape_text(slide, 8329, "[[objective_1_heading]]")
    set_shape_text(slide, 8326, "[[objective_1_text]]")
    set_shape_text(slide, 8330, "[[objective_2_heading]]")
    set_shape_text(slide, 8327, "[[objective_2_text]]")
    set_shape_text(slide, 8331, "[[objective_3_heading]]")
    set_shape_text(slide, 8328, "[[objective_3_text]]")

    # Slide 15
    slide = prs.slides[14]
    set_shape_text(slide, 8390, "[[title]]")
    set_shape_text(slide, 8395, "[[insight_1_heading]]")
    set_shape_text(slide, 8391, "[[insight_1_text]]")
    set_shape_text(slide, 8394, "[[insight_2_heading]]")
    set_shape_text(slide, 8392, "[[insight_2_text]]")
    set_shape_text(slide, 8400, "[[insight_3_heading]]")
    set_shape_text(slide, 8398, "[[insight_3_text]]")
    set_shape_text(slide, 8393, "[[insight_4_heading]]")
    set_shape_text(slide, 8396, "[[insight_4_text]]")
    set_shape_text(slide, 8401, "[[insight_5_heading]]")
    set_shape_text(slide, 8397, "[[insight_5_text]]")
    set_shape_text(slide, 8402, "[[insight_6_heading]]")
    set_shape_text(slide, 8399, "[[insight_6_text]]")

    # Slide 16
    slide = prs.slides[15]
    set_shape_text(slide, 8510, "[[title]]")
    set_shape_paragraphs(
        slide,
        8511,
        [
            "[[entry_1_ref]]",
            "[[entry_1_note]]",
            "[[entry_2_ref]]",
            "[[entry_2_note]]",
            "[[entry_3_ref]]",
            "[[entry_3_note]]",
            "[[entry_4_ref]]",
            "[[entry_4_note]]",
            "[[entry_5_ref]]",
            "[[entry_5_note]]"
        ],
    )

    # Slide 17
    slide = prs.slides[16]
    set_shape_text(slide, 8516, "[[title]]")
    set_shape_text(slide, 8517, "[[framework_kicker]]")
    set_shape_text(slide, 8518, "[[key_terms_heading]]")
    set_shape_text(slide, 8519, "[[theories_heading]]")
    set_shape_paragraphs(slide, 8520, ["[[key_term_1]]", "[[key_term_2]]", "[[key_term_3]]"])
    set_shape_text(slide, 8525, "[[theory_1_heading]]")
    set_shape_text(slide, 8526, "[[theory_1_text]]")
    set_shape_text(slide, 8527, "[[theory_2_heading]]")
    set_shape_text(slide, 8528, "[[theory_2_text]]")
    set_shape_text(slide, 8523, "[[framework_heading]]")
    set_shape_text(slide, 8524, "[[framework_text]]")

    # Slide 18
    slide = prs.slides[17]
    set_shape_text(slide, 8537, "[[title]]")
    set_table_tokens(
        slide,
        8538,
        [
            ["[[header_task]]", "[[header_description]]", "[[header_date]]", None, None, None, None, None, None, "[[header_status]]"],
            ["[[row_1_task]]", "[[row_1_description]]", "[[row_1_date]]", None, None, None, None, None, None, "[[row_1_status]]"],
            ["[[row_2_task]]", "[[row_2_description]]", "[[row_2_date]]", None, None, None, None, None, None, "[[row_2_status]]"],
            ["[[row_3_task]]", "[[row_3_description]]", "[[row_3_date]]", None, None, None, None, None, None, "[[row_3_status]]"],
            ["[[row_4_task]]", "[[row_4_description]]", "[[row_4_date]]", None, None, None, None, None, None, "[[row_4_status]]"],
            ["[[row_5_task]]", "[[row_5_description]]", "[[row_5_date]]", None, None, None, None, None, None, "[[row_5_status]]"],
        ],
    )

    # Slide 19
    slide = prs.slides[18]
    set_shape_text(slide, 8543, "[[title]]")
    set_shape_text(slide, 8553, "[[pillar_1_heading]]")
    set_shape_text(slide, 8549, "[[pillar_1_text]]")
    set_shape_text(slide, 8558, "[[pillar_2_heading]]")
    set_shape_text(slide, 8550, "[[pillar_2_text]]")
    set_shape_text(slide, 8559, "[[pillar_3_heading]]")
    set_shape_text(slide, 8551, "[[pillar_3_text]]")
    set_shape_text(slide, 8560, "[[pillar_4_heading]]")
    set_shape_text(slide, 8552, "[[pillar_4_text]]")
    set_shape_text(slide, 8561, "[[pillar_5_heading]]")

    # Slide 20
    slide = prs.slides[19]
    set_shape_text(slide, 8611, "[[title]]")
    set_shape_text(slide, 8614, "[[phase_1_heading]]")
    set_shape_paragraphs(slide, 8613, ["[[phase_1_point_1]]", "[[phase_1_point_2]]", "[[phase_1_point_3]]"])
    set_shape_text(slide, 8615, "[[phase_2_heading]]")
    set_shape_paragraphs(slide, 8612, ["[[phase_2_point_1]]", "[[phase_2_point_2]]", "[[phase_2_point_3]]"])

    # Slide 21
    slide = prs.slides[20]
    set_shape_text(slide, 8620, "[[title]]")
    set_shape_text(slide, 8628, "[[chart_caption]]")
    set_shape_text(slide, 8622, "[[note_1_heading]]")
    set_shape_text(slide, 8623, "[[note_1_text]]")
    set_shape_text(slide, 8624, "[[note_2_heading]]")
    set_shape_text(slide, 8625, "[[note_2_text]]")
    set_shape_text(slide, 8621, "")

    # Slide 22
    slide = prs.slides[21]
    set_shape_text(slide, 8634, "[[title]]")
    set_shape_text(slide, 8635, "[[result_1_heading]]")
    set_shape_text(slide, 8636, "[[result_1_text]]")
    set_shape_text(slide, 8637, "[[result_2_heading]]")
    set_shape_text(slide, 8638, "[[result_2_text]]")
    set_shape_text(slide, 8642, "[[result_3_heading]]")
    set_shape_text(slide, 8643, "[[result_3_text]]")
    set_shape_text(slide, 8646, "[[highlight_value]]")
    set_shape_text(slide, 8647, "")

    # Slide 23
    slide = prs.slides[22]
    set_shape_text(slide, 8657, "[[title]]")
    set_shape_text(slide, 8654, "[[location_1_name]]")
    set_shape_text(slide, 8658, "[[location_1_text]]")
    set_shape_text(slide, 8656, "[[location_2_name]]")
    set_shape_text(slide, 8660, "[[location_2_text]]")
    set_shape_text(slide, 8655, "[[location_3_name]]")
    set_shape_text(slide, 8659, "[[location_3_text]]")

    # Slide 24
    slide = prs.slides[23]
    set_shape_text(slide, 8723, "[[title]]")
    set_shape_text(slide, 8725, "[[intro_heading]]")
    set_shape_text(slide, 8724, "[[intro_text]]")
    set_shape_text(slide, 8728, "[[discussion_1_heading]]")
    set_shape_text(slide, 8729, "[[discussion_1_text]]")
    set_shape_text(slide, 8730, "[[discussion_2_heading]]")
    set_shape_text(slide, 8731, "[[discussion_2_text]]")

    # Slide 25
    slide = prs.slides[24]
    set_shape_text(slide, 8736, "[[title]]")
    set_shape_text(slide, 8741, "[[point_1_heading]]")
    set_shape_text(slide, 8737, "[[point_1_text]]")
    set_shape_text(slide, 8742, "[[point_2_heading]]")
    set_shape_text(slide, 8738, "[[point_2_text]]")
    set_shape_text(slide, 8743, "[[point_3_heading]]")
    set_shape_text(slide, 8739, "[[point_3_text]]")
    set_shape_text(slide, 8744, "[[point_4_heading]]")
    set_shape_text(slide, 8740, "[[point_4_text]]")

    # Slide 26
    slide = prs.slides[25]
    set_shape_text(slide, 8749, "[[title]]")
    set_shape_paragraphs(
        slide,
        8750,
        [
            "[[entry_1]]",
            "[[entry_2]]",
            "[[entry_3]]",
            "[[entry_4]]",
            "[[entry_5]]",
            "[[entry_6]]",
            "[[entry_7]]",
            "[[entry_8]]",
            "[[entry_9]]",
            "[[entry_10]]",
        ],
    )

    prs.save(TARGET)
    print(TARGET)


if __name__ == "__main__":
    main()
